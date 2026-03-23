import base64 as b64
import calendar
import datetime
import io
import json
import mimetypes
import re
import sqlite3
from dataclasses import dataclass
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple
import pandas as pd
import streamlit as st
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from pathlib import Path
from modules.cnp_utils import decode_cnp, cnp_clean, cnp_validate, cnp_parse, cnp_birthdate
from docx import Document
from graphviz import Digraph
from PIL import Image, ImageDraw, ImageFont
import math
from docx.shared import Inches
import os
import bcrypt 
import tempfile

BASE_DIR = Path(__file__).resolve().parent

# CSS global pentru a forța dark mode și controale consistente, indiferent de tema OS/browser.
FORCE_DARK_CSS = """
<style>
/* A) Force dark controls, indiferent de OS */
:root { color-scheme: dark !important; }
html, body { background: #070B12 !important; color: #E5E7EB !important; }

/* B) Overlay global dezactivat – tema globală se ocupă de contrast */

/* C) Container-e Streamlit */
[data-testid="stAppViewContainer"],
[data-testid="stMain"],
[data-testid="stSidebar"] {
  background: transparent !important;
}

/* D) TextInput/PasswordInput: vopsește WRAPPER-ul BaseWeb (nu doar input-ul) */
div[data-testid="stTextInput"] div[data-baseweb="input"] > div,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] > div,
div[data-testid="stTextInput"] div[data-baseweb="input"] > div > div,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] > div > div {
  background-color: #0F172A !important;
  border: 1px solid #263244 !important;
  box-shadow: none !important;
}

/* input efectiv */
div[data-testid="stTextInput"] input {
  background-color: #0F172A !important;
  color: #E5E7EB !important;
}

/* label + help */
div[data-testid="stTextInput"] label,
div[data-testid="stTextInput"] p {
  color: #E5E7EB !important;
}

/* placeholder */
div[data-testid="stTextInput"] input::placeholder {
  color: #9CA3AF !important;
}

/* icon / buton eye (password) să nu fie alb */
div[data-testid="stTextInput"] [data-baseweb="icon"],
div[data-testid="stTextInput"] button,
div[data-testid="stTextInput"] [role="button"] {
  background-color: #0F172A !important;
  color: #E5E7EB !important;
}

/* H) Fix pentru "eye" (password reveal) - endEnhancer BaseWeb */
div[data-testid="stTextInput"] div[data-baseweb="input"] > div > div:last-child,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] > div > div:last-child {
  background-color: #0F172A !important;
  border-left: 1px solid #263244 !important;
}

/* Butonul din endEnhancer (ochiul) */
div[data-testid="stTextInput"] div[data-baseweb="input"] > div > div:last-child button,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] > div > div:last-child button {
  background-color: #0F172A !important;
  border: 0 !important;
  box-shadow: none !important;
}

/* Icon-ul să fie vizibil pe fundal închis */
div[data-testid="stTextInput"] div[data-baseweb="input"] > div > div:last-child svg,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] > div > div:last-child svg {
  color: #E5E7EB !important;
  fill: #E5E7EB !important;
}

/* Dacă BaseWeb mai pune background pe elemente interne */
div[data-testid="stTextInput"] div[data-baseweb="input"] > div > div:last-child *,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] > div > div:last-child * {
  background-color: #0F172A !important;
}

/* --- FIX "cursiv" pentru password eye (BaseWeb endEnhancer) --- */

/* 1) Wrapper-ul principal al input-ului: el devine "capsula" unică */
div[data-testid="stTextInput"] div[data-baseweb="input"] > div,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] > div {
  background: #0F172A !important;
  border: 1px solid #263244 !important;
  border-radius: 14px !important;
  overflow: hidden !important;      /* cheie: face eye-ul să fie tăiat după radius */
  box-shadow: none !important;
}

/* 2) Input-ul efectiv: transparent ca să se vadă background-ul wrapper-ului */
div[data-testid="stTextInput"] input {
  background: transparent !important;
  color: #E5E7EB !important;
  border: 0 !important;
  box-shadow: none !important;
}

/* 3) End enhancer container (zona eye): forțat transparent ca să nu bage alb */
div[data-testid="stTextInput"] [data-baseweb="end-enhancer"],
div[data-testid="stTextInput"] [data-baseweb="endEnhancer"],
div[data-testid="stTextInput"] [data-baseweb="end-enhancer"] *,
div[data-testid="stTextInput"] [data-baseweb="endEnhancer"] * {
  background: transparent !important;
  box-shadow: none !important;
}

/* 4) Butonul eye: complet transparent + fără borduri */
div[data-testid="stTextInput"] button {
  background: transparent !important;
  border: 0 !important;
  box-shadow: none !important;
}

/* 5) Icon-ul eye vizibil pe dark */
div[data-testid="stTextInput"] button svg,
div[data-testid="stTextInput"] [data-baseweb="icon"] svg {
  color: #E5E7EB !important;
  fill: #E5E7EB !important;
}

/* 6) Focus pe întreg wrapper-ul (nu ring alb pe bucăți) */
div[data-testid="stTextInput"] div[data-baseweb="input"] > div:focus-within,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] > div:focus-within {
  border-color: #334155 !important;
  box-shadow: 0 0 0 2px rgba(51, 65, 85, 0.35) !important;
}

/* Ciocan final: dacă BaseWeb pune inline background alb pe eye */
div[data-testid="stTextInput"] [data-baseweb="end-enhancer"] [style*="background"],
div[data-testid="stTextInput"] [data-baseweb="endEnhancer"] [style*="background"] {
  background: transparent !important;
}

/* ====== LOGIN INPUTS: target după aria-label (robust) ====== */

/* Paletă (schimbăm culoarea ca să fie clar premium și uniform) */
:root{
  --login-bg: #0B1220;        /* mai închis, “premium” */
  --login-border: #2B3A55;    /* bordură calmă */
  --login-border-focus: #3B4B6A;
  --login-text: #E5E7EB;
  --login-muted: #9CA3AF;
}

/* 1) Control unitar (capsulă) pentru UTILIZATOR și PAROLĂ */
div[data-testid="stTextInput"]:has(input[aria-label="Utilizator"]) div[data-baseweb="input"] > div,
div[data-testid="stTextInput"]:has(input[aria-label="Parolă"])     div[data-baseweb="input"] > div{
  background: var(--login-bg) !important;
  border: 1px solid var(--login-border) !important;
  border-radius: 14px !important;
  overflow: hidden !important;              /* cheie: eye nu mai “sparge” capsula */
  box-shadow: none !important;
}

/* 2) Tot ce e în interior devine transparent ca să se vadă doar background-ul capsulei */
div[data-testid="stTextInput"]:has(input[aria-label="Utilizator"]) div[data-baseweb="input"] * ,
div[data-testid="stTextInput"]:has(input[aria-label="Parolă"])     div[data-baseweb="input"] * {
  background-color: transparent !important;
  box-shadow: none !important;
}

/* 3) Input-ul efectiv */
div[data-testid="stTextInput"]:has(input[aria-label="Utilizator"]) input,
div[data-testid="stTextInput"]:has(input[aria-label="Parolă"])     input{
  background: transparent !important;
  color: var(--login-text) !important;
  border: 0 !important;
  outline: none !important;
}

/* placeholder */
div[data-testid="stTextInput"]:has(input[aria-label="Utilizator"]) input::placeholder,
div[data-testid="stTextInput"]:has(input[aria-label="Parolă"])     input::placeholder{
  color: var(--login-muted) !important;
}

/* 4) Eye button: forțat să fie parte din capsulă (fără alb) */
div[data-testid="stTextInput"]:has(input[aria-label="Parolă"]) button,
div[data-testid="stTextInput"]:has(input[aria-label="Parolă"]) [role="button"]{
  background: transparent !important;
  border: 0 !important;
  outline: none !important;
  box-shadow: none !important;
  appearance: none !important;
  -webkit-appearance: none !important;
}

/* Icon eye vizibil */
div[data-testid="stTextInput"]:has(input[aria-label="Parolă"]) svg{
  color: var(--login-text) !important;
  fill: var(--login-text) !important;
  opacity: 0.9 !important;
}

/* 5) Focus pe capsulă (nu ring alb pe bucăți) */
div[data-testid="stTextInput"]:has(input[aria-label="Utilizator"]) div[data-baseweb="input"] > div:focus-within,
div[data-testid="stTextInput"]:has(input[aria-label="Parolă"])     div[data-baseweb="input"] > div:focus-within{
  border-color: var(--login-border-focus) !important;
  box-shadow: 0 0 0 2px rgba(59, 75, 106, 0.35) !important;
}

/* Dacă totuși BaseWeb forțează stiluri albe pe endEnhancer, ultima soluție:
   ascundem complet eye-ul ca să păstrăm controlul curat și uniform. */
div[data-testid="stTextInput"] [data-baseweb="end-enhancer"],
div[data-testid="stTextInput"] [data-baseweb="endEnhancer"]{
  display: none !important;
}

/* focus fără ring alb */
div[data-testid="stTextInput"] input:focus {
  outline: none !important;
}
div[data-testid="stTextInput"] div[data-baseweb="input"] > div:focus-within,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] > div:focus-within {
  border-color: #334155 !important;
  box-shadow: 0 0 0 2px rgba(51, 65, 85, 0.35) !important;
}

/* E) Chrome autofill (alb/galben) */
input:-webkit-autofill,
input:-webkit-autofill:hover,
input:-webkit-autofill:focus {
  -webkit-text-fill-color: #E5E7EB !important;
  box-shadow: 0 0 0px 1000px #0F172A inset !important;
  transition: background-color 9999s ease-out 0s;
}

/* (fost F) Butoane coerente – mutat pe rail-uri specifice, nu global */

/* G) Dacă există “glass” (blur/transparență), forțează-l să fie aproape opac */
.glass, .glass-card, .glass-panel, .login-card {
  background: rgba(15,23,42,0.92) !important;
  backdrop-filter: none !important;
  -webkit-backdrop-filter: none !important;
  border: 1px solid rgba(38,50,68,1) !important;
}

/* Ultimate override pentru zona albă de la eye: tot ce este în controlul parolei
   și are background inline devine aceeași culoare ca input-ul. */
div[data-testid="stTextInput"]:has(input[aria-label="Parolă"]) *[style*="background"] {
  background-color: var(--login-bg) !important;
  border-color: var(--login-bg) !important;
  box-shadow: none !important;
}

/* === FINAL GENERAL FIX: tratează BaseWeb base-input ca o singură capsulă pentru TOATE stTextInput === */
div[data-testid="stTextInput"] div[data-baseweb="base-input"] {
  background-color: var(--login-bg) !important;
  border: 1px solid var(--login-border) !important;
  border-radius: 14px !important;
  overflow: hidden !important;
  box-shadow: none !important;
}

div[data-testid="stTextInput"] div[data-baseweb="base-input"] input,
div[data-testid="stTextInput"] div[data-baseweb="base-input"] button {
  background-color: transparent !important;
  border: 0 !important;
  box-shadow: none !important;
}

div[data-testid="stTextInput"] div[data-baseweb="base-input"] svg {
  color: var(--login-text) !important;
  fill: var(--login-text) !important;
}

/* Orice background inline din interiorul base-input devine transparent, ca să nu mai apară alb */
div[data-testid="stTextInput"] div[data-baseweb="base-input"] *[style*="background"] {
  background-color: transparent !important;
  box-shadow: none !important;
}

/* ===== LOGIN V2: bară AUTENTIFICARE + câmpuri consistente ===== */

.login-card{
  background: transparent !important;
  border-radius: 0 !important;
  border: none !important;
  box-shadow: none !important;
}

.login-header-bar{
  width: 100%;
  text-align: center;
  padding: 11px 20px;
  margin-bottom: 16px;
  border-radius: 999px;
  background: linear-gradient(135deg, rgba(15,23,42,0.98), rgba(15,23,42,0.90));
  border: 1px solid rgba(148,163,184,0.70);
  color: #FFFFFF;
  font-size: 0.98rem;
  font-weight: 820;
  letter-spacing: 0.20em;
  text-transform: uppercase;
  box-shadow: 0 18px 42px rgba(0,0,0,0.80);
}

.login-card .login-field{
  margin-bottom: 12px;
}

.login-card .login-field-label{
  font-size: 0.86rem;
  font-weight: 600;
  color: rgba(226,232,240,0.92);
  margin-bottom: 4px;
}

.login-card .login-field div[data-testid="stTextInputRootElement"]{
  background: radial-gradient(circle at top left, rgba(15,23,42,0.98), rgba(15,23,42,0.88)) !important;
  border-radius: 999px !important;
  border: 1px solid rgba(148,163,184,0.70) !important;
  box-shadow: 0 18px 40px rgba(0,0,0,0.85) !important;
  overflow: hidden !important;
}

.login-card .login-field input{
  height: 46px !important;
  padding: 0 18px !important;
  background: transparent !important;
  border: none !important;
  outline: none !important;
  color: #F9FAFB !important;
  font-weight: 500 !important;
}

.login-card .login-field input::placeholder{
  color: rgba(148,163,184,0.95) !important;
}

.login-card .login-field div[data-testid="stTextInputRootElement"]:focus-within{
  border-color: rgba(56,189,248,0.90) !important;
  box-shadow: 0 0 0 1px rgba(56,189,248,0.85), 0 0 0 6px rgba(8,47,73,0.85) !important;
}
</style>
"""

APP_THEME_CSS = """
<style>
/* ============ SOCRATES DARK THEME PACK (BASE) ============ */
:root{
  color-scheme: dark !important;          /* ignoră tema OS */
  --bg: #070B12;
  --surface: #0B1220;
  --surface2: #0F172A;
  --border: #263244;
  --text: #F8FAFC;                        /* alb curat */
  --muted: #CBD5E1;                       /* text secundar lizibil */
  --muted2: #94A3B8;                      /* captions/help */
  --accent: rgba(34,197,94,0.55);

  --tab-bg: #0B1220;
  --tab-active: #0F172A;
  --tab-border: #263244;
}

html, body { background: var(--bg) !important; color: var(--text) !important; }
[data-testid="stAppViewContainer"] { background: var(--bg) !important; color: var(--text) !important; }
/* Sidebar lasă background-ul de imagine/gradient; doar culoarea textului rămâne controlată */
[data-testid="stSidebar"] { background: transparent !important; color: var(--text) !important; }

/* Tabs de bază – restul e rafinat în patch-ul post-render */
div[data-testid="stTabs"] [role="tablist"]{
  border-bottom: 0 !important;
  gap: 10px !important;
}
div[data-testid="stTabs"] button[role="tab"]{
  background: var(--tab-bg) !important;
  color: var(--muted2) !important;
  border: 1px solid var(--tab-border) !important;
  border-radius: 999px !important;
  padding: 8px 14px !important;
  margin: 0 !important;
}
div[data-testid="stTabs"] button[role="tab"][aria-selected="true"]{
  background: var(--tab-active) !important;
  color: var(--text) !important;
  border-color: rgba(34,197,94,0.6) !important;
  box-shadow: 0 0 0 2px rgba(34,197,94,0.15) !important;
}
div[data-testid="stTabs"] [role="tabpanel"]{
  padding-top: 12px !important;
}
</style>
"""

APP_POST_RENDER_CSS = """
<style>
/* =========================
   PHASE 1: RESET (neutralize)
   ========================= */
:root { color-scheme: dark !important; }

div[data-testid="stTextInputRootElement"],
div[data-testid="stNumberInputRootElement"],
div[data-testid="stTextAreaRootElement"],
div[data-testid="stSelectbox"],
div[data-testid="stMultiSelect"],
div[data-testid="stDateInput"],
div[data-testid="stTimeInput"],
div[data-testid="stFileUploader"],
div[data-testid="stTabs"],
div[data-testid="stForm"],
div[data-testid="stExpander"]{
  background: transparent !important;
  border: none !important;
  box-shadow: none !important;
}

/* BaseWeb wrappers: scoatem “skin”-ul */
div[data-baseweb],
div[data-baseweb] > div,
div[data-baseweb] > div > div{
  background: transparent !important;
  box-shadow: none !important;
}

/* =========================
   PHASE 2: THEME (re-color)
   ========================= */
:root{
  --bg: #070B12;
  --surface: #0B1220;
  --surface2: #0F172A;
  --border: #263244;
  --text: #F8FAFC;
  --muted: #CBD5E1;
  --muted2: #94A3B8;
  --accent: rgba(34,197,94,0.55);
}

html, body { background: var(--bg) !important; color: var(--text) !important; }
[data-testid="stAppViewContainer"]{ background: transparent !important; color: var(--text) !important; }
[data-testid="stMain"]{ color: var(--text) !important; }
/* Sidebar păstrează albastrul din background (imagine/gradient); doar textul e controlat */
[data-testid="stSidebar"]{ background: transparent !important; color: var(--text) !important; }

/* Headings + markdown lizibile (fără gri spălăcit) */
h1,h2,h3,h4,h5,h6 { color: var(--text) !important; opacity: 1 !important; }
[data-testid="stMarkdownContainer"] * { color: var(--text) !important; opacity: 1 !important; }

/* Captions/help: nu le lăsăm invizibile */
[data-testid="stCaptionContainer"] *{
  color: var(--muted) !important;
  opacity: 1 !important;
}

/* INPUT ROOT (stabil) – capsulă unitară ca în login */
div[data-testid="stTextInputRootElement"],
div[data-testid="stNumberInputRootElement"],
div[data-testid="stTextAreaRootElement"]{
  background: var(--surface) !important;
  border: 1px solid var(--border) !important;
  border-radius: 14px !important;
  overflow: hidden !important;
}

/* input/textarea efectiv */
div[data-testid="stTextInputRootElement"] input,
div[data-testid="stNumberInputRootElement"] input,
div[data-testid="stTextAreaRootElement"] textarea{
  background: transparent !important;
  color: var(--text) !important;
  caret-color: var(--text) !important;
  border: 0 !important;
  box-shadow: none !important;
}

/* placeholder */
div[data-testid="stTextInputRootElement"] input::placeholder,
div[data-testid="stTextAreaRootElement"] textarea::placeholder{
  color: var(--muted2) !important;
  opacity: 1 !important;
}

/* end-enhancer/eye etc: transparent (să nu bage alb) */
div[data-testid="stTextInputRootElement"] button,
div[data-testid="stTextInputRootElement"] [role="button"],
div[data-testid="stTextInputRootElement"] svg{
  background: transparent !important;
  border: 0 !important;
  box-shadow: none !important;
  color: var(--text) !important;
  fill: var(--text) !important;
}

/* autofill */
input:-webkit-autofill,
input:-webkit-autofill:hover,
input:-webkit-autofill:focus,
textarea:-webkit-autofill{
  -webkit-text-fill-color: var(--text) !important;
  box-shadow: 0 0 0px 1000px var(--surface) inset !important;
  transition: background-color 9999s ease-out 0s;
}

/* Focus coerent */
div[data-testid="stTextInputRootElement"]:focus-within,
div[data-testid="stNumberInputRootElement"]:focus-within,
div[data-testid="stTextAreaRootElement"]:focus-within{
  border-color: var(--accent) !important;
  box-shadow: 0 0 0 2px rgba(34,197,94,0.15) !important;
}

/* Tabs uniforme */
div[data-testid="stTabs"] button[role="tab"]{
  background: var(--surface) !important;
  color: var(--muted) !important;
  border: 1px solid var(--border) !important;
  border-radius: 999px !important;
  padding: 8px 14px !important;
}
div[data-testid="stTabs"] button[role="tab"][aria-selected="true"]{
  background: var(--surface2) !important;
  color: var(--text) !important;
  border-color: var(--accent) !important;
  box-shadow: 0 0 0 2px rgba(34,197,94,0.12) !important;
}

/* ===== Sidebar buttons: Logout + Help (target by key) ===== */
[data-testid="stSidebar"] .st-key-sidebar_btn_logout button,
[data-testid="stSidebar"] .st-key-sidebar_btn_help button {
  background: transparent !important;         /* transparent, cum ai cerut */
  color: #F8FAFC !important;                  /* text alb */
  border: 1px solid #263244 !important;       /* border subtil */
  border-radius: 999px !important;            /* pill */
  box-shadow: none !important;
  width: 100% !important;                     /* să arate “uniform” în coloane */
  padding: 10px 14px !important;
}

/* hover/focus: fără alburi, doar accent discret */
[data-testid="stSidebar"] .st-key-sidebar_btn_logout button:hover,
[data-testid="stSidebar"] .st-key-sidebar_btn_help button:hover {
  border-color: rgba(34, 197, 94, 0.55) !important;
  background: rgba(15, 23, 42, 0.35) !important;
}

[data-testid="stSidebar"] .st-key-sidebar_btn_logout button:focus,
[data-testid="stSidebar"] .st-key-sidebar_btn_help button:focus,
[data-testid="stSidebar"] .st-key-sidebar_btn_logout button:focus-visible,
[data-testid="stSidebar"] .st-key-sidebar_btn_help button:focus-visible {
  outline: none !important;
  box-shadow: 0 0 0 2px rgba(34, 197, 94, 0.18) !important;
}

/* dacă Streamlit pune background pe wrapper intern, îl “spălăm” */
[data-testid="stSidebar"] .st-key-sidebar_btn_logout * ,
[data-testid="stSidebar"] .st-key-sidebar_btn_help * {
  background-color: transparent !important;
}

/* variantă de rezervă: dacă key-ul e pe div[data-testid="stButton"] direct */
[data-testid="stSidebar"] div[data-testid="stButton"].st-key-sidebar_btn_logout button,
[data-testid="stSidebar"] div[data-testid="stButton"].st-key-sidebar_btn_help button { 
  background: transparent !important;
  color: #F8FAFC !important;
  border: 1px solid #263244 !important;
  border-radius: 999px !important;
}

/* ===== TOP BAR: match sidebar palette + glass menus ===== */

:root{
  /* ajustează fin dacă vrei mai albastru/mai verde */
  --topA: rgba(10, 30, 45, 0.72);     /* petrol dark */
  --topB: rgba(14, 55, 70, 0.62);     /* petrol/teal */
  --glass: rgba(7, 11, 18, 0.35);
  --glass2: rgba(15, 23, 42, 0.35);
  --border: rgba(38, 50, 68, 0.85);
  --text: #F8FAFC;
  --accent: rgba(34,197,94,0.45);
}

/* 1) Header: gradient în ton cu sidebar */
header[data-testid="stHeader"]{
  background: linear-gradient(90deg, var(--topA), var(--topB)) !important;
  backdrop-filter: blur(12px) !important;
  -webkit-backdrop-filter: blur(12px) !important;
  border-bottom: 1px solid var(--border) !important;
}

/* 2) Decorative strip: transparent */
div[data-testid="stDecoration"]{
  background: transparent !important;
  box-shadow: none !important;
}

/* 3) Toolbar: păstrat, dar glass */
div[data-testid="stToolbar"]{
  background: transparent !important;
}
div[data-testid="stToolbar"] *{
  color: var(--text) !important;
  fill: var(--text) !important;
  opacity: 1 !important;
}

/* 4) Butoane din toolbar (Deploy + …) glass – va fi rafinat de blocul final kebaB */
div[data-testid="stToolbar"] button,
div[data-testid="stToolbar"] [role="button"]{
  background: var(--glass2) !important;
  border: 1px solid var(--border) !important;
  border-radius: 12px !important;
  box-shadow: none !important;
}
div[data-testid="stToolbar"] button:hover,
div[data-testid="stToolbar"] [role="button"]:hover{
  background: rgba(15,23,42,0.55) !important;
  border-color: var(--accent) !important;
}
div[data-testid="stToolbar"] button:focus,
div[data-testid="stToolbar"] button:focus-visible{
  outline: none !important;
  box-shadow: 0 0 0 2px rgba(34,197,94,0.18) !important;
}

/* 5) Kebab menu dropdown (popup/portal): glass dark */
div[data-baseweb="popover"],
div[data-baseweb="popover"] > div{
  background: rgba(10, 30, 45, 0.72) !important;
  backdrop-filter: blur(14px) !important;
  -webkit-backdrop-filter: blur(14px) !important;
  border: 1px solid var(--border) !important;
  border-radius: 14px !important;
  box-shadow: none !important;
}

/* meniul în sine */
[role="menu"],
[role="menu"] > *{
  background: transparent !important;
  color: var(--text) !important;
}

/* item-urile din meniu */
[role="menuitem"]{
  color: var(--text) !important;
}
[role="menuitem"]:hover{
  background: rgba(15,23,42,0.55) !important;
}

/* text/icon în popover */
div[data-baseweb="popover"] *{
  color: var(--text) !important;
  fill: var(--text) !important;
  stroke: var(--text) !important;
}

/* fallback pentru popup-uri */
[role="dialog"],
[role="dialog"] > div{
  background: rgba(10, 30, 45, 0.72) !important;
  backdrop-filter: blur(14px) !important;
  -webkit-backdrop-filter: blur(14px) !important;
  border: 1px solid rgba(38,50,68,0.85) !important;
  border-radius: 14px !important;
  box-shadow: none !important;
}

/* ===== TOP BAR: match SIDEBAR palette (same gradient) ===== */
header[data-testid="stHeader"]{
  /* copiat ca stil din sidebar: radial-gradient(900px 600px at 0% 0%, rgba(15,23,42,0.90), rgba(15,23,42,0.80)) */
  background: radial-gradient(900px 600px at 0% 0%,
              rgba(15,23,42,0.86),
              rgba(15,23,42,0.72)) !important;
  backdrop-filter: blur(12px) !important;
  -webkit-backdrop-filter: blur(12px) !important;
  border-bottom: 1px solid rgba(15,23,42,0.85) !important;
}

div[data-testid="stDecoration"]{
  background: transparent !important;
  box-shadow: none !important;
}

/* ===== Toolbar buttons (Deploy + …) -> glass ===== */
div[data-testid="stToolbar"]{
  background: transparent !important;
}
div[data-testid="stToolbar"] button,
div[data-testid="stToolbar"] [role="button"]{
  background: rgba(15,23,42,0.35) !important;
  border: 1px solid rgba(148,163,184,0.18) !important;
  border-radius: 12px !important;
  box-shadow: none !important;
}
div[data-testid="stToolbar"] button:hover,
div[data-testid="stToolbar"] [role="button"]:hover{
  background: rgba(15,23,42,0.55) !important;
  border-color: rgba(34,197,94,0.45) !important;
}
div[data-testid="stToolbar"] *{
  color: rgba(248,250,252,0.95) !important;
  fill: rgba(248,250,252,0.95) !important;
  stroke: rgba(248,250,252,0.95) !important;
}

/* ===== “…” dropdown / menus (Streamlit portals) -> GLASS ===== */
/* Target general pentru popovers/menus (Streamlit le randă în portal în afara app container) */
body > div[role="dialog"],
body > div[role="dialog"] > div,
body > div[data-baseweb="popover"],
body > div[data-baseweb="popover"] > div,
div[data-baseweb="popover"],
div[data-baseweb="popover"] > div{
  background: rgba(15,23,42,0.82) !important;
  backdrop-filter: blur(14px) !important;
  -webkit-backdrop-filter: blur(14px) !important;
  border: 1px solid rgba(148,163,184,0.18) !important;
  border-radius: 14px !important;
  box-shadow: 0 12px 40px rgba(0,0,0,0.35) !important;
}

/* Meniul și item-urile: fără “white tiles” */
body > div[role="dialog"] [role="menu"],
body > div[role="dialog"] [role="menu"] *{
  background: transparent !important;
  color: rgba(248,250,252,0.95) !important;
  fill: rgba(248,250,252,0.95) !important;
}

/* Butoanele din meniu (acolo unde se vede alb) */
body > div[role="dialog"] button,
body > div[role="dialog"] a,
body > div[role="dialog"] [role="menuitem"]{
  background: transparent !important;
  color: rgba(248,250,252,0.95) !important;
  border: 0 !important;
  box-shadow: none !important;
}

/* Hover pe meniu */
body > div[role="dialog"] [role="menuitem"]:hover,
body > div[role="dialog"] button:hover{
  background: rgba(15,23,42,0.55) !important;
}

/* Separatoare din dropdown */
body > div[role="dialog"] hr{
  border-color: rgba(148,163,184,0.18) !important;
}

/* ===== STREAMLIT TOOLBAR MENU (kebab "...") - make it readable ALWAYS ===== */

/* 1) Containerul popup-ului: glass dark */
body > div[role="dialog"],
body > div[data-baseweb="popover"],
div[data-baseweb="popover"] {
  background: rgba(15,23,42,0.82) !important;
  backdrop-filter: blur(14px) !important;
  -webkit-backdrop-filter: blur(14px) !important;
  border: 1px solid rgba(148,163,184,0.18) !important;
  border-radius: 14px !important;
  box-shadow: 0 12px 40px rgba(0,0,0,0.35) !important;
  color: rgba(248,250,252,0.95) !important;
}

/* 2) HAMMER: elimină *orice* background alb din interior */
body > div[role="dialog"] * ,
body > div[data-baseweb="popover"] * ,
div[data-baseweb="popover"] * {
  background: transparent !important;
  background-color: transparent !important;
  color: rgba(248,250,252,0.95) !important;
  opacity: 1 !important;
  box-shadow: none !important;
}

/* 3) Item-uri meniu: "chip" dark, lizibil */
body > div[role="dialog"] [role="menuitem"],
body > div[role="dialog"] button,
body > div[role="dialog"] a,
body > div[data-baseweb="popover"] [role="menuitem"],
body > div[data-baseweb="popover"] button,
body > div[data-baseweb="popover"] a {
  border-radius: 10px !important;
  padding: 10px 12px !important;
}

/* 4) Hover/focus (discret, premium) */
body > div[role="dialog"] [role="menuitem"]:hover,
body > div[role="dialog"] button:hover,
body > div[data-baseweb="popover"] [role="menuitem"]:hover,
body > div[data-baseweb="popover"] button:hover {
  background: rgba(15,23,42,0.55) !important;
  border: 1px solid rgba(34,197,94,0.35) !important;
}

/* 5) Separatoare (liniile dintre secțiuni) */
body > div[role="dialog"] hr,
body > div[data-baseweb="popover"] hr {
  border-color: rgba(148,163,184,0.18) !important;
}

/* ===== STREAMLIT KEBAB ICON: remove white tile by replacing icon ===== */

/* prindem butonul icon-only din toolbar (ăla cu SVG) */
header[data-testid="stHeader"] div[data-testid="stToolbar"] button:has(svg),
header[data-testid="stHeader"] div[data-testid="stToolbar"] [role="button"]:has(svg){
  position: relative !important;
  width: 40px !important;
  min-width: 40px !important;
  height: 40px !important;
  padding: 0 !important;

  background: rgba(15,23,42,0.45) !important;     /* glass */
  border: 1px solid rgba(148,163,184,0.18) !important;
  border-radius: 12px !important;
  box-shadow: none !important;
  overflow: hidden !important;
}

/* ascundem TOT conținutul intern (inclusiv pătratul alb) */
header[data-testid="stHeader"] div[data-testid="stToolbar"] button:has(svg) > *,
header[data-testid="stHeader"] div[data-testid="stToolbar"] [role="button"]:has(svg) > *{
  opacity: 0 !important;
}

/* desenăm noi kebab-ul (3 puncte verticale) */
header[data-testid="stHeader"] div[data-testid="stToolbar"] button:has(svg)::after,
header[data-testid="stHeader"] div[data-testid="stToolbar"] [role="button"]:has(svg)::after{
  content: "" !important;
  position: absolute !important;
  left: 50% !important;
  top: 50% !important;
  width: 4px !important;
  height: 4px !important;
  border-radius: 999px !important;
  background: rgba(248,250,252,0.95) !important;
  transform: translate(-50%, -50%) !important;

  /* 3 puncte */
  box-shadow:
    0 -8px 0 rgba(248,250,252,0.95),
    0  8px 0 rgba(248,250,252,0.95) !important;
}

/* hover premium */
header[data-testid="stHeader"] div[data-testid="stToolbar"] button:has(svg):hover,
header[data-testid="stHeader"] div[data-testid="stToolbar"] [role="button"]:has(svg):hover{
  background: rgba(15,23,42,0.65) !important;
  border-color: rgba(34,197,94,0.45) !important;
}
</style>
"""

GLASS_ALWAYS_CSS = """
<style>
/* ===== GLASS ALWAYS (no more white bars, no hover-only) ===== */
:root{
  --glass-bg: rgba(15,23,42,0.42);
  --glass-bg2: rgba(15,23,42,0.60);
  --glass-border: rgba(148,163,184,0.18);
  --glass-text: rgba(248,250,252,0.95);
  --glass-muted: rgba(203,213,225,0.80);
  --glass-accent: rgba(34,197,94,0.45);
}

/* Scope: aplică pe conținut + sidebar (nu pe portal popovers) */
[data-testid="stAppViewContainer"],
[data-testid="stSidebar"]{
  color: var(--glass-text) !important;
}

/* ---------- PHASE 1: RESET (neutralize white skins) ---------- */
/* Anulează “white tiles” puse de module (emp-wrap etc) */
[data-testid="stAppViewContainer"] .stTextInput input,
[data-testid="stAppViewContainer"] .stSelectbox div,
[data-testid="stAppViewContainer"] .stMultiSelect div,
[data-testid="stSidebar"] .stTextInput input,
[data-testid="stSidebar"] .stSelectbox div,
[data-testid="stSidebar"] .stMultiSelect div{
  background: transparent !important;
  background-color: transparent !important;
  box-shadow: none !important;
}

/* ---------- PHASE 2: GLASS (always visible, not only hover) ---------- */
/* Inputs / selects / textareas (BaseWeb wrappers) */
[data-testid="stAppViewContainer"] div[data-baseweb="input"] > div,
[data-testid="stAppViewContainer"] div[data-baseweb="base-input"] > div,
[data-testid="stAppViewContainer"] div[data-baseweb="textarea"] > div,
[data-testid="stAppViewContainer"] div[data-baseweb="select"] > div,
[data-testid="stSidebar"] div[data-baseweb="input"] > div,
[data-testid="stSidebar"] div[data-baseweb="base-input"] > div,
[data-testid="stSidebar"] div[data-baseweb="textarea"] > div,
[data-testid="stSidebar"] div[data-baseweb="select"] > div{
  background: var(--glass-bg) !important;
  border: 1px solid var(--glass-border) !important;
  border-radius: 14px !important;
  box-shadow: none !important;
  overflow: hidden !important;
  backdrop-filter: blur(10px) !important;
  -webkit-backdrop-filter: blur(10px) !important;
}

/* Nu vrem “devine transparent doar pe hover” — hover rămâne glass, puțin mai solid */
[data-testid="stAppViewContainer"] div[data-baseweb="input"] > div:hover,
[data-testid="stAppViewContainer"] div[data-baseweb="base-input"] > div:hover,
[data-testid="stAppViewContainer"] div[data-baseweb="textarea"] > div:hover,
[data-testid="stAppViewContainer"] div[data-baseweb="select"] > div:hover,
[data-testid="stSidebar"] div[data-baseweb="input"] > div:hover,
[data-testid="stSidebar"] div[data-baseweb="base-input"] > div:hover,
[data-testid="stSidebar"] div[data-baseweb="textarea"] > div:hover,
[data-testid="stSidebar"] div[data-baseweb="select"] > div:hover{
  background: var(--glass-bg2) !important;
}

/* input/textarea text */
[data-testid="stAppViewContainer"] input,
[data-testid="stAppViewContainer"] textarea,
[data-testid="stSidebar"] input,
[data-testid="stSidebar"] textarea{
  background: transparent !important;
  color: var(--glass-text) !important;
  caret-color: var(--glass-text) !important;
}

/* placeholder lizibil */
[data-testid="stAppViewContainer"] input::placeholder,
[data-testid="stAppViewContainer"] textarea::placeholder,
[data-testid="stSidebar"] input::placeholder,
[data-testid="stSidebar"] textarea::placeholder{
  color: var(--glass-muted) !important;
  opacity: 1 !important;
}

/* focus */
[data-testid="stAppViewContainer"] div[data-baseweb="input"] > div:focus-within,
[data-testid="stAppViewContainer"] div[data-baseweb="base-input"] > div:focus-within,
[data-testid="stAppViewContainer"] div[data-baseweb="textarea"] > div:focus-within,
[data-testid="stAppViewContainer"] div[data-baseweb="select"] > div:focus-within{
  border-color: var(--glass-accent) !important;
  box-shadow: 0 0 0 2px rgba(34,197,94,0.14) !important;
}

/* Form submit buttons + normal buttons: glass always (nu doar hover) */
[data-testid="stAppViewContainer"] div[data-testid="stFormSubmitButton"] button,
[data-testid="stAppViewContainer"] .stButton > button,
[data-testid="stSidebar"] .stButton > button{
  background: var(--glass-bg) !important;
  border: 1px solid var(--glass-border) !important;
  color: var(--glass-text) !important;
  border-radius: 14px !important;
  box-shadow: none !important;
  backdrop-filter: blur(10px) !important;
  -webkit-backdrop-filter: blur(10px) !important;
}
[data-testid="stAppViewContainer"] div[data-testid="stFormSubmitButton"] button:hover,
[data-testid="stAppViewContainer"] .stButton > button:hover,
[data-testid="stSidebar"] .stButton > button:hover{
  background: var(--glass-bg2) !important;
  border-color: var(--glass-accent) !important;
}

/* Download buttons (Exportă) – glass, fără fundal alb */
[data-testid="stAppViewContainer"] .stDownloadButton > button,
[data-testid="stSidebar"] .stDownloadButton > button,
[data-testid="stAppViewContainer"] div[data-testid="stDownloadButton"] button,
[data-testid="stSidebar"] div[data-testid="stDownloadButton"] button{
  background: var(--glass-bg) !important;
  border: 1px solid var(--glass-border) !important;
  color: var(--glass-text) !important;
  border-radius: 14px !important;
  box-shadow: none !important;
  backdrop-filter: blur(10px) !important;
  -webkit-backdrop-filter: blur(10px) !important;
}
[data-testid="stAppViewContainer"] .stDownloadButton > button:hover,
[data-testid="stSidebar"] .stDownloadButton > button:hover,
[data-testid="stAppViewContainer"] div[data-testid="stDownloadButton"] button:hover,
[data-testid="stSidebar"] div[data-testid="stDownloadButton"] button:hover{
  background: var(--glass-bg2) !important;
  border-color: var(--glass-accent) !important;
}

/* Link buttons (ex. 🚀 Deschide Pontaj) – glass, nu alb */
[data-testid="stAppViewContainer"] .stLinkButton > a,
[data-testid="stSidebar"] .stLinkButton > a,
[data-testid="stAppViewContainer"] div[data-testid="stLinkButton"] > a,
[data-testid="stSidebar"] div[data-testid="stLinkButton"] > a{
  display: inline-flex !important;
  align-items: center !important;
  justify-content: center !important;
  background: var(--glass-bg) !important;
  border: 1px solid var(--glass-border) !important;
  color: var(--glass-text) !important;
  border-radius: 999px !important;
  padding: 0.5rem 1rem !important;
  text-decoration: none !important;
  box-shadow: none !important;
  backdrop-filter: blur(10px) !important;
  -webkit-backdrop-filter: blur(10px) !important;
}
[data-testid="stAppViewContainer"] .stLinkButton > a:hover,
[data-testid="stSidebar"] .stLinkButton > a:hover,
[data-testid="stAppViewContainer"] div[data-testid="stLinkButton"] > a:hover,
[data-testid="stSidebar"] div[data-testid="stLinkButton"] > a:hover{
  background: var(--glass-bg2) !important;
  border-color: var(--glass-accent) !important;
}

/* Data editor (tabele de editare, ex. Stat de funcții) – fundal mai deschis, font închis */
[data-testid="stAppViewContainer"] .stDataFrame table,
[data-testid="stAppViewContainer"] .stDataFrame tbody tr,
[data-testid="stAppViewContainer"] .stDataFrame tbody td,
[data-testid="stAppViewContainer"] .stDataFrame thead th{
  background: #E5E7EB !important;          /* fundal deschis, tip “foaie de calcul” */
  color: #0B1220 !important;               /* text închis pentru lizibilitate */
}
[data-testid="stAppViewContainer"] .stDataFrame tbody tr:nth-child(even) td{
  background: #CBD5E1 !important;          /* bandare discretă pe rânduri */
}

/* Expanders: summary bar să fie glass (nu alb) */
[data-testid="stAppViewContainer"] div[data-testid="stExpander"] summary{
  background: var(--glass-bg) !important;
  border: 1px solid var(--glass-border) !important;
  border-radius: 14px !important;
  color: var(--glass-text) !important;
  backdrop-filter: blur(10px) !important;
  -webkit-backdrop-filter: blur(10px) !important;
}

/* Card-uri / benzi făcute din div-uri custom în theme.css */
.emp-card, .home-hero, .home-block, .dosar-card, .org-card, .stat-card{
  background: var(--glass-bg) !important;
  border: 1px solid var(--glass-border) !important;
  backdrop-filter: blur(10px) !important;
  -webkit-backdrop-filter: blur(10px) !important;
}

/* ===== Fix FINAL: kebab / Main menu button să NU mai fie alb ===== */

/* prinde butonul "..." indiferent cum e etichetat */
div[data-testid="stToolbar"] button[aria-label="Main menu"],
div[data-testid="stToolbar"] button[title="Main menu"],
div[data-testid="stToolbar"] button[aria-haspopup="menu"],
div[data-testid="stToolbar"] [role="button"][aria-label="Main menu"],
div[data-testid="stToolbar"] [role="button"][aria-haspopup="menu"]{
  background: rgba(15,23,42,0.45) !important;
  border: 1px solid rgba(148,163,184,0.18) !important;
  border-radius: 12px !important;
  box-shadow: none !important;
  outline: none !important;
}

/* "ciocan": orice copil intern să nu mai aibă background alb */
div[data-testid="stToolbar"] button[aria-label="Main menu"] *,
div[data-testid="stToolbar"] button[title="Main menu"] *,
div[data-testid="stToolbar"] button[aria-haspopup="menu"] *,
div[data-testid="stToolbar"] [role="button"][aria-label="Main menu"] *,
div[data-testid="stToolbar"] [role="button"][aria-haspopup="menu"] *{
  background: transparent !important;
  background-color: transparent !important;
  box-shadow: none !important;
}

/* dacă pătratul alb vine din SVG rect / shape */
div[data-testid="stToolbar"] button[aria-label="Main menu"] svg rect,
div[data-testid="stToolbar"] button[title="Main menu"] svg rect,
div[data-testid="stToolbar"] button[aria-haspopup="menu"] svg rect{
  fill: transparent !important;
}

/* icon-ul (cele 3 puncte) să fie alb */
div[data-testid="stToolbar"] button[aria-label="Main menu"] svg,
div[data-testid="stToolbar"] button[title="Main menu"] svg,
div[data-testid="stToolbar"] button[aria-haspopup="menu"] svg,
div[data-testid="stToolbar"] button[aria-label="Main menu"] svg *,
div[data-testid="stToolbar"] button[title="Main menu"] svg *,
div[data-testid="stToolbar"] button[aria-haspopup="menu"] svg *{
  fill: rgba(248,250,252,0.95) !important;
  stroke: rgba(248,250,252,0.95) !important;
  color: rgba(248,250,252,0.95) !important;
}

/* hover premium */
div[data-testid="stToolbar"] button[aria-label="Main menu"]:hover,
div[data-testid="stToolbar"] button[title="Main menu"]:hover,
div[data-testid="stToolbar"] button[aria-haspopup="menu"]:hover{
  background: rgba(15,23,42,0.65) !important;
  border-color: rgba(34,197,94,0.45) !important;
}

/* fallback: dacă există inline style="background:#fff" pe vreun copil intern */
div[data-testid="stToolbar"] button[aria-haspopup="menu"] [style*="background"]{
  background: transparent !important;
  background-color: transparent !important;
}
</style>
"""

LOGIN_FORCE_INPUTS_CSS = """
<style>
/* ---- LOGIN: FORCE DARK INPUTS (super robust) ---- */
:root { color-scheme: dark !important; }

/* prinde orice variantă BaseWeb: input, base-input, etc */
div[data-testid="stTextInput"] [data-baseweb*="input"],
div[data-testid="stTextInput"] [data-baseweb*="input"] > div,
div[data-testid="stTextInput"] [data-baseweb*="input"] > div > div {
  background: #0B1220 !important;
  border: 1px solid #263244 !important;
  box-shadow: none !important;
  border-radius: 14px !important;
}

/* input efectiv – text clar, alb */
div[data-testid="stTextInput"] input {
  background: #0B1220 !important;
  color: #F8FAFC !important;          /* text alb în casetă */
  caret-color: #F8FAFC !important;
  border: 0 !important;
  box-shadow: none !important;
}

/* placeholder – tot deschis, nu gri prăfuit */
div[data-testid="stTextInput"] input::placeholder {
  color: #F8FAFC !important;
  opacity: 0.7 !important;
}

/* autofill */
input:-webkit-autofill,
input:-webkit-autofill:hover,
input:-webkit-autofill:focus {
  -webkit-text-fill-color: #F8FAFC !important;
  box-shadow: 0 0 0px 1000px #0B1220 inset !important;
  transition: background-color 9999s ease-out 0s;
}

/* dacă există eye/end-enhancer, îl facem transparent ca să nu bage alb */
div[data-testid="stTextInput"] [data-baseweb*="end"] ,
div[data-testid="stTextInput"] [data-baseweb*="end"] * {
  background: transparent !important;
  box-shadow: none !important;
}
div[data-testid="stTextInput"] button {
  background: transparent !important;
  border: 0 !important;
  box-shadow: none !important;
}
div[data-testid="stTextInput"] svg {
  color: #F8FAFC !important;
  fill: #F8FAFC !important;
}

/* override explicit pentru login (în interiorul .login-wrap) – culoare foarte vizibilă */
.login-wrap .stTextInput input,
.login-wrap input[type="text"],
.login-wrap input[type="password"]{
  color: #22C55E !important;              /* verde accent, foarte vizibil */
  caret-color: #22C55E !important;
  font-weight: 600 !important;
}
.login-wrap .stTextInput input::placeholder,
.login-wrap input[type="text"]::placeholder,
.login-wrap input[type="password"]::placeholder{
  color: #4ADE80 !important;              /* verde mai deschis pentru placeholder */
  opacity: 1 !important;
}
</style>
"""

DASHBOARD_PATCH_CSS = """
<style>
:root{
  color-scheme: dark !important;
  --bg: #070B12;
  --surface: #0B1220;
  --surface2: #0F172A;
  --border: #263244;
  --text: #F8FAFC;
  --muted: #CBD5E1;
  --muted2: #94A3B8;
  --accent: rgba(34,197,94,0.55);
}

/* 1) Headings clare (Scurtături inclus) */
h1,h2,h3,h4 { color: var(--text) !important; opacity: 1 !important; }
[data-testid="stMarkdownContainer"] * { opacity: 1 !important; }

/* 2) Text secundar lizibil */
[data-testid="stCaptionContainer"] *,
small, .stMarkdown p, .stMarkdown span {
  color: var(--muted) !important;
  opacity: 1 !important;
}

/* 3) Dacă există paneluri “glass” cu opacity, forțăm să NU afecteze textul */
.glass, .glass-card, .glass-panel, .card, .panel {
  opacity: 1 !important;
}

/* 4) QUICK SEARCH: capsulă unitară, o singură culoare, fără borduri vizibile */
.home-section div[data-testid="stTextInput"]:has(input[aria-label*="Căutare"]) [data-baseweb*="input"],
.home-section div[data-testid="stTextInput"]:has(input[aria-label*="Căutare"]) [data-baseweb*="base-input"] {
  background: var(--surface) !important;
  border-radius: 999px !important;
  border: 0 !important;
  box-shadow: none !important;
  overflow: hidden !important;
}
.home-section div[data-testid="stTextInput"]:has(input[aria-label*="Căutare"]) [data-baseweb*="input"] > div,
.home-section div[data-testid="stTextInput"]:has(input[aria-label*="Căutare"]) [data-baseweb*="input"] * {
  background: transparent !important;
  border: 0 !important;
  box-shadow: none !important;
}
.home-section div[data-testid="stTextInput"] input[aria-label*="Căutare"] {
  background: transparent !important;
  color: var(--text) !important;
  caret-color: var(--text) !important;
  border: 0 !important;
  box-shadow: none !important;
}
.home-section div[data-testid="stTextInput"] input[aria-label*="Căutare"]::placeholder{
  color: var(--muted2) !important;
  opacity: 1 !important;
}

/* autofill */
input:-webkit-autofill,
input:-webkit-autofill:hover,
input:-webkit-autofill:focus {
  -webkit-text-fill-color: var(--text) !important;
  box-shadow: 0 0 0px 1000px var(--surface) inset !important;
  transition: background-color 9999s ease-out 0s;
}

/* focus */
div[data-testid="stTextInput"] [data-baseweb*="input"] > div:focus-within{
  border-color: var(--accent) !important;
  box-shadow: 0 0 0 2px rgba(34,197,94,0.15) !important;
}

</style>
"""

DEBUG_CSS = "body{outline:6px solid red !important;}"
# inject_css(DEBUG_CSS)  # activează 1 run, apoi comentează


def apply_app_theme():
    import streamlit as st
    st.markdown(APP_THEME_CSS, unsafe_allow_html=True)


def apply_premium_theme():
    """Premium dashboard theming: glass cards + left rail."""
    import streamlit as st
    st.markdown(
        """
        <style>
        /* Workspace centrat + padding */
        .block-container{
          width: 100% !important;
          max-width: 100% !important;
          padding-top: 1.15rem !important;
          padding-bottom: 2rem !important;
        }

        /* Card glass reutilizabil */
        .s-card{
          background: rgba(7,11,18,0.70);
          border: 1px solid rgba(148,163,184,0.18);
          border-radius: 16px;
          padding: 16px 18px;
          margin-bottom: 14px;
          backdrop-filter: blur(14px) saturate(140%);
          -webkit-backdrop-filter: blur(14px) saturate(140%);
        }
        .s-card:hover{
          background: rgba(7,11,18,0.82);
          border-color: rgba(59,130,246,0.35);
        }

        .s-title{ font-size: 16px; font-weight: 800; margin: 0 0 10px 0; }
        .s-muted{ opacity: .78; font-size: 12px; margin-top: 4px; }

        /* Left rail buttons: full width în coloană îngustă */
        .action-rail .stButton>button{
          width: 100%;
          height: 38px;
          border-radius: 12px;
          text-align: left;
          padding: 0 14px;
          background: rgba(15,23,42,0.70);
          border: 1px solid rgba(148,163,184,0.22);
          color: #F8FAFC;
          box-shadow: none;
        }
        .action-rail .stButton>button:hover{
          background: rgba(37,99,235,0.26);
          border-color: rgba(59,130,246,0.40);
        }
        .action-rail .stButton>button:focus{
          box-shadow: 0 0 0 2px rgba(59,130,246,0.35);
          outline: none;
        }

        /* Search row input: colțuri mai rotunjite */
        .search-row .stTextInput>div>div input{
          border-radius: 12px !important;
        }

        /* --- RAIL (scurtături): coloană îngustă + butoane egale, aliniate stânga --- */
        .rail{
          max-width: 320px;
        }

        .rail .stButton,
        .rail .stButton > button{
          width: 100% !important;
        }

        .rail .stButton > button{
          height: 42px !important;
          border-radius: 12px !important;
          padding: 0 14px !important;

          display: flex !important;
          align-items: center !important;
          justify-content: flex-start !important;   /* lipit stânga */
          gap: 10px !important;

          text-align: left !important;
          background: rgba(255,255,255,0.06) !important;
          border: 1px solid rgba(255,255,255,0.10) !important;
          color: rgba(255,255,255,0.92) !important;
        }

        .rail .stButton > button:hover{
          background: rgba(59,130,246,0.16) !important;
          border-color: rgba(59,130,246,0.28) !important;
        }

        /* ===== SECTION LAYOUT + BUTTON STACK (Dashboard/Home) ===== */

        /* Carduri moderne pentru secțiunile de pe „Acasă” */
        .section-card{
          /* delimitare ultra-fină: fundal abia mai deschis + border foarte discret */
          background: rgba(15,23,42,0.14);
          border: 1px solid rgba(148,163,184,0.12);
          border-radius: 12px;
          padding: 18px 22px;
          margin-bottom: 18px;
        }

        .section-title{
          font-size: 20px;
          font-weight: 600;
          margin-bottom: 16px;
          letter-spacing: 0.3px;
          color: #e6edf5;
        }

        .section-subtitle{
          color: #94a3b8;
          margin-bottom: 20px;
        }

        /* Header de pagină (titlu + acțiuni) – folosit pe Acasă, Angajați, Organigramă etc. */
        .page-header{
          display: flex;
          align-items: center;
          justify-content: space-between;
          margin-bottom: 36px;
        }

        .page-title,
        h1{
          font-size: 36px;
          font-weight: 700;
          letter-spacing: 0.3px;
          margin: 0;
          color: #e6edf5;
        }

        .page-actions{
          display: flex;
          gap: 12px;
        }

        .btn-primary{
          background: #1d4ed8;
          border: none;
          padding: 10px 18px;
          border-radius: 10px;
          font-weight: 500;
          cursor: pointer;
          color: #e5e7eb;
        }

        .btn-primary:hover{
          background: #1e40af;
        }

        /* Stil aplicat pe butoanele Streamlit din header (în loc de <button>) */
        .page-actions div[data-testid="stButton"] > button{
          background: #1d4ed8 !important;
          border: none !important;
          padding: 10px 18px !important;
          border-radius: 10px !important;
          font-weight: 500 !important;
          cursor: pointer !important;
        }

        .page-actions div[data-testid="stButton"] > button:hover{
          background: #1e40af !important;
        }

        /* Wrapper: toate sectiunile sunt aliniate la stanga */
        .section-wrap{
          display: block;
        }

        /* Butoane “stacked”: narrow, identice, stanga */
        .btnstack{
          max-width: 320px;
        }

        .btnstack .stButton,
        .btnstack .stButton > button{
          width: 100% !important;
        }

        .btnstack .stButton > button{
          height: 42px !important;
          border-radius: 12px !important;
          padding: 0 14px !important;

          display: flex !important;
          align-items: center !important;
          justify-content: flex-start !important; /* stanga */
          gap: 10px !important;

          text-align: left !important;
          background: rgba(255,255,255,0.06) !important;
          border: 1px solid rgba(255,255,255,0.10) !important;
          color: rgba(255,255,255,0.92) !important;
        }

        .btnstack .stButton > button:hover{
          background: rgba(59,130,246,0.16) !important;
          border-color: rgba(59,130,246,0.28) !important;
        }

        /* ===== EMPLOYEES: FILTER CARD (page_angajati) ===== */
        .filter-card{
          background: rgba(10, 18, 35, 0.45);
          border: 1px solid rgba(255,255,255,0.10);
          border-radius: 18px;
          padding: 18px 18px;
          backdrop-filter: blur(14px) saturate(140%);
          -webkit-backdrop-filter: blur(14px) saturate(140%);
        }

        .filter-h2{ font-size: 26px; font-weight: 900; margin: 0 0 14px 0; }
        .filter-h3{ font-size: 14px; font-weight: 800; opacity: .9; margin: 0 0 8px 0; }

        .vgap-12{ margin-top: 12px; }
        .vgap-16{ margin-top: 16px; }

        .filter-card .stTextInput>div>div input{
          border-radius: 14px !important;
          height: 44px !important;
        }

        .filter-card .stButton>button{
          height: 44px !important;
          border-radius: 14px !important;
        }

        .filter-card details{
          border-radius: 14px !important;
          overflow: hidden;
        }
        .filter-card details > summary{
          padding: 12px 14px !important;
          border: 1px solid rgba(255,255,255,0.10) !important;
          border-radius: 14px !important;
          background: rgba(255,255,255,0.04) !important;
        }
        .filter-card details[open] > summary{
          border-bottom-left-radius: 0 !important;
          border-bottom-right-radius: 0 !important;
        }
        .filter-card details > div{
          border: 1px solid rgba(255,255,255,0.10) !important;
          border-top: none !important;
          padding: 14px 14px 16px 14px !important;
          border-bottom-left-radius: 14px !important;
          border-bottom-right-radius: 14px !important;
          background: rgba(255,255,255,0.03) !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


SIDEBAR_BLUE = "#163A4C"  # ajustează dacă sidebar-ul are altă nuanță

# Păstrăm o constantă pentru compatibilitate cu apelurile existente;
# fundalul efectiv al Home-ului este controlat în apply_page_background.
HOME_BG_CSS = ""


def apply_centered_layout(max_width_px: int = 1280):
    """Centrează conținutul (block-container) pe toate paginile."""
    import streamlit as st
    st.markdown(
        f"""
        <style>
        section.main .block-container {{
            width: 100% !important;
            max-width: 100% !important;
            margin-left: 0 !important;
            margin-right: 0 !important;
            padding-left: 1.5rem !important;
            padding-right: 1.5rem !important;
            padding-top: 0.8rem !important;
            padding-bottom: 2.2rem !important;
        }}

        [data-testid="stAppViewContainer"],
        main[data-testid="stMain"],
        div[data-testid="stMainBlockContainer"] {{
            width: 100% !important;
            max-width: 100% !important;
        }}

        div[data-testid="stAppViewContainer"] {{
            overflow-x: hidden;
        }}

        @media (max-width: 900px) {{
          section.main .block-container {{
            padding-left: 1.1rem !important;
            padding-right: 1.1rem !important;
          }}
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )


def apply_page_background(page_name: str, home_bg_css: str | None = None):
    """
    Controlează fundalul pentru Home vs restul paginilor.
    - Home: poate avea imagine (home_bg_css)
    - Restul: fundal solid SIDEBAR_BLUE, fără imagine
    """
    import streamlit as st

    if page_name == "Acasă":
        # Fundal cu statuie.jpeg (ca înainte), doar pe Home
        img_path = BASE_DIR / "assets" / "statuie.jpeg"
        if not img_path.exists():
            # fallback foarte safe
            img_path = BASE_DIR / "assets" / "fundal.jpeg"

        if img_path.exists():
            try:
                img_bytes = img_path.read_bytes()
                img_b64 = b64.b64encode(img_bytes).decode("utf-8")
                mime = mimetypes.guess_type(str(img_path))[0] or "image/png"
                st.markdown(
                    f"""
                    <style>
                    div[data-testid="stAppViewContainer"] {{
                      background:
                        linear-gradient(180deg, rgba(10,18,35,0.25), rgba(10,18,35,0.55)),
                        url("data:{mime};base64,{img_b64}") center/cover no-repeat fixed !important;
                    }}
                    body {{
                      background: {SIDEBAR_BLUE} !important;
                    }}
                    </style>
                    """,
                    unsafe_allow_html=True,
                )
            except Exception:
                # dacă ceva nu merge, cădem pe fundal simplu albastru
                st.markdown(
                    f"""
                    <style>
                    div[data-testid="stAppViewContainer"] {{
                      background: {SIDEBAR_BLUE} !important;
                    }}
                    body {{
                      background: {SIDEBAR_BLUE} !important;
                    }}
                    </style>
                    """,
                    unsafe_allow_html=True,
                )
        else:
            st.markdown(
                f"""
                <style>
                div[data-testid="stAppViewContainer"] {{
                  background: {SIDEBAR_BLUE} !important;
                }}
                body {{
                  background: {SIDEBAR_BLUE} !important;
                }}
                </style>
                """,
                unsafe_allow_html=True,
            )
    else:
        st.markdown(
            f"""
            <style>
            div[data-testid="stAppViewContainer"] {{
              background: {SIDEBAR_BLUE} !important;
              background-image: none !important;
            }}
            body {{
              background: {SIDEBAR_BLUE} !important;
            }}
            .stApp {{
              background: transparent !important;
            }}
            </style>
            """,
            unsafe_allow_html=True,
        )


def apply_login_fix():
    import streamlit as st
    st.markdown(LOGIN_FORCE_INPUTS_CSS, unsafe_allow_html=True)


def apply_dashboard_patch() -> None:
    inject_css(DASHBOARD_PATCH_CSS)


def apply_force_dark_ui() -> None:
    st.markdown(FORCE_DARK_CSS, unsafe_allow_html=True)


def inject_css(css: str | None = None):
    import streamlit as st
    # helper pentru injectare directă (debug/patch-uri)
    if css is not None:
        st.markdown(f"<style>{css}</style>", unsafe_allow_html=True)
        return

    css_path = BASE_DIR / "assets" / "theme.css"
    if css_path.exists():
        st.markdown(f"<style>{css_path.read_text(encoding='utf-8')}</style>", unsafe_allow_html=True)

        # Login background: doar când user-ul NU este autentificat.
        # După login, fundalul este controlat de apply_page_background().
        try:
            logged_in = bool(st.session_state.get("logged_in", False))
        except Exception:
            logged_in = False

        if not logged_in:
            # Fundal pentru pagina de login: statuie.jpeg, cu fallback pe grecia.png apoi fundal.jpeg
            img_path = BASE_DIR / "assets" / "statuie.jpeg"
            if not img_path.exists():
                alt_path = BASE_DIR / "assets" / "grecia.png"
                if alt_path.exists():
                    img_path = alt_path
                else:
                    alt_path2 = BASE_DIR / "assets" / "fundal.jpeg"
                    if alt_path2.exists():
                        img_path = alt_path2

            if img_path.exists():
                try:
                    img_bytes = img_path.read_bytes()
                    img_b64 = b64.b64encode(img_bytes).decode("utf-8")
                    mime = mimetypes.guess_type(str(img_path))[0] or "image/png"
                    st.markdown(
                        f"""
                        <style>
                        .stApp {{
                          background-color:#05060a;
                          background-image: url("data:{mime};base64,{img_b64}");
                          background-size: cover;
                          background-position: center center;
                          background-attachment: fixed;
                          background-repeat: no-repeat;
                        }}
                        </style>
                        """,
                        unsafe_allow_html=True,
                    )
                except Exception:
                    pass
    else:
        st.error(f"CSS file not found: {css_path}")


def apply_app_post_render_fix() -> None:
    import streamlit as st
    st.markdown(APP_POST_RENDER_CSS, unsafe_allow_html=True)


def apply_glass_always():
    import streamlit as st
    st.markdown(GLASS_ALWAYS_CSS, unsafe_allow_html=True)


def apply_toolbar_runtime_fix():
    """
    Fix pentru butonul alb din toolbar (kebab / main menu).
    CSS-ul nu îl poate bate mereu (inline !important / rerender).
    Aici îl forțăm din runtime: glass + remove white backgrounds.
    """
    import streamlit.components.v1 as components

    components.html(
        """
        <script>
        (function () {
          const GLASS_BG = 'rgba(15,23,42,0.45)';
          const BORDER   = '1px solid rgba(148,163,184,0.18)';
          const TEXT     = 'rgba(248,250,252,0.95)';

          function isWhite(bg) {
            return bg === 'rgb(255, 255, 255)' || bg === 'rgba(255, 255, 255, 1)';
          }

          function patchToolbar() {
            const doc = window.parent?.document || document;
            const header = doc.querySelector('header[data-testid="stHeader"]');
            const toolbar = doc.querySelector('div[data-testid="stToolbar"]');

            if (!header || !toolbar) return;

            // 1) Glass pe toate butoanele din toolbar
            const buttons = toolbar.querySelectorAll('button, [role="button"]');
            buttons.forEach(btn => {
              btn.style.setProperty('background', GLASS_BG, 'important');
              btn.style.setProperty('border', BORDER, 'important');
              btn.style.setProperty('border-radius', '12px', 'important');
              btn.style.setProperty('box-shadow', 'none', 'important');
              btn.style.setProperty('outline', 'none', 'important');
              btn.style.setProperty('color', TEXT, 'important');

              // 2) Orice copil cu background alb (sau inline background) -> transparent
              btn.querySelectorAll('*').forEach(el => {
                try {
                  const cs = getComputedStyle(el);
                  const bg = cs.backgroundColor;

                  if (isWhite(bg)) {
                    el.style.setProperty('background', 'transparent', 'important');
                    el.style.setProperty('background-color', 'transparent', 'important');
                  }

                  const st = el.getAttribute('style') || '';
                  if (st.toLowerCase().includes('background')) {
                    el.style.setProperty('background', 'transparent', 'important');
                    el.style.setProperty('background-color', 'transparent', 'important');
                  }

                  if (el.tagName && el.tagName.toLowerCase() === 'rect') {
                    el.style.setProperty('fill', 'transparent', 'important');
                  }
                } catch (e) {}
              });

              // 3) dacă există SVG: îl facem alb
              btn.querySelectorAll('svg, svg *').forEach(svgEl => {
                svgEl.style.setProperty('fill', TEXT, 'important');
                svgEl.style.setProperty('stroke', TEXT, 'important');
                svgEl.style.setProperty('color', TEXT, 'important');
              });
            });

            // 4) Extra: caută în toolbar orice element mic (<= 44px) cu background alb și îl “stinge”
            toolbar.querySelectorAll('*').forEach(el => {
              try {
                const cs = getComputedStyle(el);
                if (isWhite(cs.backgroundColor) && el.offsetWidth <= 44 && el.offsetHeight <= 44) {
                  el.style.setProperty('background', 'transparent', 'important');
                  el.style.setProperty('background-color', 'transparent', 'important');
                }
              } catch (e) {}
            });
          }

          // Patch imediat
          patchToolbar();

          // Repatch la rerender (Streamlit mai rescrie header-ul)
          const doc = window.parent?.document || document;
          const obs = new MutationObserver(() => patchToolbar());
          obs.observe(doc.documentElement, {subtree:true, childList:true, attributes:true});

          // (opțional) oprește observer după 20s
          setTimeout(() => { try { obs.disconnect(); } catch(e) {} }, 20000);
        })();
        </script>
        """,
        height=0,
        scrolling=False,
    )


DB_PATH_DEFAULT = str(Path(__file__).resolve().parent.parent / "data" / "ANCPI.db")

ROOT = Path(__file__).resolve().parent.parent 
DATA = ROOT / "data"

ANCP_DB = DATA / "ANCPI.db"
HR_DB   = DATA / "socrates_hr.db"

conn_ancp = sqlite3.connect(str(ANCP_DB), check_same_thread=False)
conn_hr   = sqlite3.connect(str(HR_DB), check_same_thread=False)


# ============================================================
# CNP – validare + parsare (modul unic: modules.cnp_utils)
# ============================================================
# cnp_clean, cnp_validate, cnp_parse, cnp_birthdate, decode_cnp importate din modules.cnp_utils

def dates_match(d1: Optional[datetime.date], d2: Optional[datetime.date]) -> bool:
    if d1 is None or d2 is None:
        return False
    return d1.year == d2.year and d1.month == d2.month and d1.day == d2.day

# -------------------------------------------------------------
# Helper: export DataFrame to XLSX bytes without hard dependency
# -------------------------------------------------------------
def df_to_xlsx_bytes(df, sheet_name: str = 'Sheet1'):
    """Return XLSX bytes or None if no engine available.

    Tries xlsxwriter first (preferred on servers), then openpyxl.
    """
    import io
    buf = io.BytesIO()
    for engine in ('xlsxwriter', 'openpyxl'):
        try:
            with pd.ExcelWriter(buf, engine=engine) as writer:
                df.to_excel(writer, index=False, sheet_name=sheet_name)
            buf.seek(0)
            return buf.getvalue()
        except Exception:
            buf = io.BytesIO()
            continue
    return None

# -------------------------------------------------------------
# Home – istoric navigare recentă (în session_state)
# -------------------------------------------------------------
def _init_home_state():
    st.session_state.setdefault("home_recent", [])  # list[dict]
    st.session_state.setdefault("home_last_query", "")

def add_recent(kind: str, label: str, target: str, payload: dict | None = None):
    """kind: 'employee'/'doc'/'nav' etc. target: numele exact din radio (ex: 'Angajați')."""
    _init_home_state()
    item = {
        "ts": datetime.datetime.now().isoformat(timespec="seconds"),
        "kind": kind,
        "label": label,
        "target": target,
        "payload": payload or {},
    }
    # dedup by label+target
    rec = [x for x in st.session_state.home_recent if not (x["label"] == label and x["target"] == target)]
    st.session_state.home_recent = [item] + rec
    st.session_state.home_recent = st.session_state.home_recent[:8]

def goto(page_label: str):
    """Navighează către un modul existent fără să schimbăm backend-ul."""
    st.session_state["main_choice"] = page_label
    st.rerun()

# -------------------------------------------------------------
# Compat helpers (alias) – folosite in pagini
# -------------------------------------------------------------

def _df_to_xlsx_bytes(df, sheet_name: str = 'Sheet1'):
    """Wrapper peste df_to_xlsx_bytes.

    Returnează bytes XLSX sau None dacă nu există engine (xlsxwriter/openpyxl).
    Caller-ul decide eventualul fallback (ex: CSV).
    """
    return df_to_xlsx_bytes(df, sheet_name=sheet_name)


def _df_to_csv_bytes(df) -> bytes:
    """Fallback CSV (UTF-8)."""
    try:
        return df.to_csv(index=False).encode('utf-8')
    except Exception:
        return (str(df)).encode('utf-8')
    

st.set_page_config(page_title="SOCRATES@HR", layout="wide")
apply_app_theme()
apply_force_dark_ui()
inject_css()

# IMPORTANT: asta trebuie să fie ULTIMUL CSS global înainte de bg-glow,
# ca să nu fie calculat de theme.css sau alte injectări.
apply_glass_always()

st.markdown('<div class="bg-glow"></div>', unsafe_allow_html=True)

# -------------------------------------------------------------
# CNP -> Data nașterii (folosim cnp_birthdate din modules.cnp_utils)
# -------------------------------------------------------------
import datetime as _dt  # alias sigur (evită conflicte cu variabile locale numite datetime)

def _format_ro_date(d):
    try:
        if isinstance(d, _dt.date):
            return d.strftime("%d.%m.%Y")
    except Exception:
        pass
    return ""
# -------------------------------------------------------------
# DOCX: înlocuire placeholder-e (CIM / șabloane)
# -------------------------------------------------------------
def _docx_replace_placeholders(doc: Document, mapping: dict) -> None:
    """Înlocuiește placeholder-ele {{CHEIE}} într-un DOCX.

    Suport:
    - paragrafe (run-uri sparte de Word)
    - tabele (celule + nested tables)
    - header/footer

    Notă: acceptă mapping cu chei fie "CHEIE" fie "{{CHEIE}}".
    """
    # Normalizează cheile: CHEIE (fără acolade)
    norm = {}
    for k, v in (mapping or {}).items():
        kk = str(k)
        kk = kk.strip()
        if kk.startswith("{{") and kk.endswith("}}"):
            kk = kk[2:-2].strip()
        norm[kk] = "" if v is None else str(v)

    def _replace_in_paragraph(p):
        # Concatenează toate run-urile ca să prindă placeholder-ele sparte
        if not getattr(p, "runs", None):
            return
        full = "".join(r.text for r in p.runs)
        if not full:
            return

        new_text = full
        for kk, vv in norm.items():
            new_text = new_text.replace(f"{{{{{kk}}}}}", vv)

        if new_text != full:
            # păstrăm formatarea aproximativ: textul în primul run, restul gol
            p.runs[0].text = new_text
            for r in p.runs[1:]:
                r.text = ""

    def _replace_in_table(t):
        for row in t.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _replace_in_paragraph(p)
                for nt in cell.tables:
                    _replace_in_table(nt)

    def _replace_in_part(part):
        # paragrafe
        for p in part.paragraphs:
            _replace_in_paragraph(p)
        # tabele
        for t in part.tables:
            _replace_in_table(t)

    # Document body
    _replace_in_part(doc)

    # Headers / Footers (pentru fiecare secțiune)
    try:
        for sec in doc.sections:
            _replace_in_part(sec.header)
            _replace_in_part(sec.footer)
    except Exception:
        pass

# -------------------------------------------------------------
# Documente generate din cod (fără șabloane DOCX)
# -------------------------------------------------------------
def _docx_add_paragraph(doc: Document, text: str, *, bold: bool = False, size: int = 11, align: str | None = None):
    """Adaugă paragraf cu setări de bază (folosit la CIM / Act adițional generate din cod)."""
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    p = doc.add_paragraph()
    run = p.add_run(text or "")
    run.bold = bool(bold)
    run.font.size = Pt(int(size))
    if align:
        a = align.lower().strip()
        if a == "center":
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        elif a == "right":
            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        elif a == "justify":
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    return p


# -------------------------------------------------------------
# STANDARDIZARE: ANTET + SEMNĂTURI (pentru toate documentele generate din cod)
# -------------------------------------------------------------
def _docx_apply_antet_si_semnaturi(doc: "Document") -> None:
    """Aplică antet (sus) și semnături (jos) pe document, pe baza Configurării.
    - Antet: denumire, CUI, adresă, tel/fax/email (+ opțional IBAN/bancă) și siglă (stânga, opțional).
    - Semnături: Conducător unitate + Responsabil HR (în footer).
    """
    try:
        cfg = load_config()
    except Exception:
        cfg = {}

    den = (cfg.get("denumire_unitate") or "").strip()
    cui = (cfg.get("cui") or "").strip()
    adr = (cfg.get("adresa") or "").strip()
    tel = (cfg.get("telefon") or "").strip()
    fax = (cfg.get("fax") or "").strip()
    email = (cfg.get("email") or "").strip()
    iban = (cfg.get("cont_bancar") or "").strip()
    banca = (cfg.get("banca") or "").strip()
    show_iban = bool(cfg.get("afiseaza_iban_in_antet", False))

    sigla_path = (cfg.get("sigla_path") or "").strip()

    # --- ANTET (header) ---
    try:
        section = doc.sections[0]
        header = section.header

        # Curățare minimă (evită antet dublu dacă funcția e chemată de 2 ori)
        if len(header.paragraphs) == 1 and (header.paragraphs[0].text or "").strip() == "":
            pass

        # tabel 2 coloane: siglă | text
        table = header.add_table(rows=1, cols=2)
        table.autofit = True
        c0 = table.cell(0, 0)
        c1 = table.cell(0, 1)

        # siglă stânga (opțional)
        if sigla_path and os.path.exists(sigla_path):
            try:
                p = c0.paragraphs[0]
                run = p.add_run()
                run.add_picture(sigla_path, width=Inches(1.0))
            except Exception:
                pass

        # text antet (dreapta)
        p = c1.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT

        if den:
            r = p.add_run(den + "\n")
            r.bold = True
            r.font.size = Pt(11)

        lines = []
        if cui:
            lines.append(f"CUI/CIF: {cui}")
        if adr:
            lines.append(adr)

        contact_parts = []
        if tel:
            contact_parts.append(f"Tel: {tel}")
        if fax:
            contact_parts.append(f"Fax: {fax}")
        if email:
            contact_parts.append(f"Email: {email}")
        if contact_parts:
            lines.append(" | ".join(contact_parts))

        if show_iban:
            iban_line = []
            if iban:
                iban_line.append(f"IBAN: {iban}")
            if banca:
                iban_line.append(f"Banca: {banca}")
            if iban_line:
                lines.append(" | ".join(iban_line))

        if lines:
            p.add_run("\n".join(lines))

    except Exception:
        # nu blocăm generarea dacă header-ul nu poate fi setat
        pass

    # --- SEMNĂTURI (footer) ---
    try:
        section = doc.sections[0]
        footer = section.footer

        # tabel 2 coloane: conducător | HR
        table_f = footer.add_table(rows=1, cols=2)
        table_f.autofit = True
        lc = table_f.cell(0, 0)
        rc = table_f.cell(0, 1)

        cond_nume = (cfg.get("conducator_nume") or "").strip()
        cond_func = (cfg.get("conducator_functie") or "").strip()
        hr_nume = (cfg.get("responsabil_hr_nume") or "").strip()
        hr_func = (cfg.get("responsabil_hr_functie") or "").strip()

        # stânga: conducător
        p1 = lc.paragraphs[0]
        if cond_func:
            r = p1.add_run(cond_func + "\n")
            r.bold = True
        if cond_nume:
            p1.add_run(cond_nume)
        else:
            p1.add_run("")

        # dreapta: HR
        p2 = rc.paragraphs[0]
        if hr_func:
            r = p2.add_run(hr_func + "\n")
            r.bold = True
        if hr_nume:
            p2.add_run(hr_nume)
        else:
            p2.add_run("")

    except Exception:
        pass

def build_act_aditional_cim_docx_bytes(data: dict) -> bytes:
    """Generează DOCX (bytes) pentru Act adițional la CIM, după modelul 2021."""
    import io

    doc = Document()

    _docx_apply_antet_si_semnaturi(doc)

    _docx_add_paragraph(doc, "ACT ADIȚIONAL LA C.I.M.", bold=True, size=14, align="center")
    _docx_add_paragraph(doc, "")
    _docx_add_paragraph(doc, f"S.C. {data.get('ANGAJATOR_DEN','')}")
    _docx_add_paragraph(doc, f"Sediul {data.get('ANGAJATOR_SEDIU','')}")
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "ACT ADIȚIONAL", bold=True, size=13, align="center")
    _docx_add_paragraph(doc, f"NR. {data.get('AA_NR','')} / {data.get('AA_DATA','')}", bold=True)
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(
        doc,
        "La contractul de muncă înregistrat în Registrul electronic de evidență a salariaților "
        f"sub nr. {data.get('REVISAL_NR','')} / {data.get('REVISAL_DATA','')}, "
        f"Subsemnatul {data.get('REPREZENTANT_LEGAL','')}, reprezentant legal al "
        f"S.C. {data.get('ANGAJATOR_DEN','')} și salariatul(a) {data.get('SALARIAT_NUME_COMPLET','')} "
        "am convenit modificarea următoarelor elemente ale contractului de muncă sus-menționat, după cum urmează:"
    )
    _docx_add_paragraph(doc, "")
    _docx_add_paragraph(doc, data.get("AA_CONTINUT_MODIFICARI",""))
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(
        doc,
        "Restul clauzelor contractului individual de muncă sunt conform prevederilor Legii nr. 53/2003 "
        "și rămân neschimbate."
    )
    _docx_add_paragraph(doc, "")
    _docx_add_paragraph(
        doc,
        "Prezentul act adițional s-a încheiat în 2 (două) exemplare, câte unul pentru fiecare parte contractată, "
        f"urmând să își producă efectele începând cu data de {data.get('AA_DATA_EFECT','')}."
    )

    _docx_add_paragraph(doc, "")
    _docx_add_paragraph(doc, "ANGAJATOR                                                                       ANGAJAT", bold=True)
    _docx_add_paragraph(doc, "")
    _docx_add_paragraph(doc, "_________________________                                               _________________________")

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def build_contract_cim_docx_bytes(data: dict) -> bytes:
    """Generează DOCX (bytes) pentru Contract individual de muncă, după modelul 2021."""
    import io

    doc = Document()

    _docx_apply_antet_si_semnaturi(doc)

    _docx_add_paragraph(doc, "CONTRACT INDIVIDUAL DE MUNCĂ", bold=True, size=14, align="center")
    _docx_add_paragraph(doc, f"încheiat și înregistrat sub nr. {data.get('CIM_NR','')}/{data.get('CIM_DATA','')} în Revisal", size=11)
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "A. Părțile contractului", bold=True)
    _docx_add_paragraph(
        doc,
        f"Angajator - {data.get('ANGAJATOR_DEN','')}, cu sediul în {data.get('ANGAJATOR_SEDIU','')}, "
        f"cod fiscal {data.get('ANGAJATOR_CUI','')}, reprezentată legal prin {data.get('REPREZENTANT_LEGAL','')} "
        f"în calitate de {data.get('REPREZENTANT_FUNCTIE','')}."
    )
    _docx_add_paragraph(
        doc,
        f"și salariatul/salariata - {data.get('SALARIAT_NUME_COMPLET','')}, domiciliat(ă) în {data.get('SALARIAT_DOMICILIU','')}, "
        f"CNP {data.get('SALARIAT_CNP','')}."
    )
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "B. Obiectul contractului:", bold=True)
    _docx_add_paragraph(doc, data.get("OBIECT",""))
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "C. Durata contractului:", bold=True)
    _docx_add_paragraph(doc, data.get("DURATA_TEXT",""))
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "D. Locul de muncă", bold=True)
    _docx_add_paragraph(doc, data.get("LOC_MUNCA",""))
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "E. Felul muncii", bold=True)
    _docx_add_paragraph(doc, f"Funcția/meseria: {data.get('FUNCTIE','')} (COR: {data.get('COD_COR','')})")
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "F. Atribuțiile postului", bold=True)
    _docx_add_paragraph(doc, data.get("ATRIBUTII","Atribuțiile postului sunt prevăzute în fișa postului, anexă la CIM."))
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "H. Durata muncii", bold=True)
    _docx_add_paragraph(doc, data.get("PROGRAM_TEXT",""))
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "I. Concediul", bold=True)
    _docx_add_paragraph(doc, data.get("CONCEDIU_TEXT",""))
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(doc, "J. Salariul", bold=True)
    _docx_add_paragraph(doc, data.get("SALARIU_TEXT",""))
    _docx_add_paragraph(doc, "")

    _docx_add_paragraph(
        doc,
        "Prezentul contract individual de muncă s-a încheiat în două exemplare, câte unul pentru fiecare parte."
    )
    _docx_add_paragraph(doc, "")
    _docx_add_paragraph(doc, "Angajator,                                     Salariat", bold=True)
    _docx_add_paragraph(doc, "")
    _docx_add_paragraph(doc, "_________________________                      _________________________")
    _docx_add_paragraph(doc, "")

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# -------------------------------------------------------------
# L153 (anexe salarizare) – grilă din XLSX (Normalized)
# -------------------------------------------------------------

# -------------------------------------------------------------
# DECIZII (DOCX) generate din cod (fără șabloane)
# -------------------------------------------------------------
def build_decizie_generic_docx_bytes(data: dict) -> bytes:
    """Generează DOCX (bytes) pentru o decizie (template generic)."""
    import io
    import re as _re
    from docx import Document

    doc = Document()

    _docx_apply_antet_si_semnaturi(doc)

    def p(text="", **kw):
        return _docx_add_paragraph(doc, text or "", **kw)

    # antet
    p(f"SC {data.get('COMPANY_NAME','')}".strip())
    if data.get('COMPANY_SEDIU'):
        p(f"Sediul {data.get('COMPANY_SEDIU')}")
    if data.get('COMPANY_REGCOM'):
        p(f"Reg.Com. {data.get('COMPANY_REGCOM')}")
    p("")

    # titlu
    titlu = (data.get("DEC_TITLU") or "DECIZIE").strip()
    p(titlu, bold=True, size=14, align="center")
    p("")
    p(f"DECIZIA NR. {data.get('DEC_NR','')}".strip(), bold=True, align="center")
    p(f"Data {data.get('DEC_DATA','')}".strip(), align="center")
    p("")

    # emitent
    admin = (data.get("ADMIN_NAME") or "").strip()
    if admin:
        p(f"Administratorul {data.get('COMPANY_NAME','')} {admin}".strip())
    else:
        p(f"Administratorul SC {data.get('COMPANY_NAME','')}".strip())
    p("")
    emp_name = (data.get("EMP_NAME") or "").strip()
    emp_fun = (data.get("EMP_FUNCTIE") or "").strip()
    if emp_name:
        p(f"Având în vedere situația salariatului/salariatei {emp_name}" + (f", având funcția de {emp_fun}," if emp_fun else ","))
    p("")
    considerente = (data.get("CONSIDERENTE") or "").strip()
    if considerente:
        for line in considerente.splitlines():
            p(line)

    p("")
    p("D E C I D E:", bold=True, align="center")
    p("")
    arts = data.get("ARTICLES") or []
    if isinstance(arts, str):
        arts = [arts]

    for i, art in enumerate(arts, start=1):
        art_txt = (art or "").strip()
        if not art_txt:
            continue
        if _re.match(r"^\s*Art\.?\s*\d+", art_txt, flags=_re.IGNORECASE):
            p(art_txt)
        else:
            p(f"Art.{i}. {art_txt}")

    tribunal = (data.get("TRIBUNAL") or "").strip()
    if tribunal:
        p("")
        p(f"Prezenta poate fi contestată conform legii la {tribunal}.")

    p("")
    executor = (data.get("EXECUTOR") or "").strip()
    if executor:
        p(f"Cu ducerea la îndeplinire a prezentei se însărcinează {executor} și se comunică celui în cauză.")

    p("")
    p("ADMINISTRATOR,", bold=True)
    p("")
    p("L.S.")
    p("")
    p("_________________________")

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _l153_norm_df(df: pd.DataFrame) -> pd.DataFrame:
    """Normalizează grila L153 din XLSX: coloane standard + filtrează rândurile valide."""
    if df is None or getattr(df, "empty", True):
        return pd.DataFrame(columns=["Anexa","Tabel","Funcție","Studii","Grad","Treaptă/Gradație","Salariu","Coeficient"])

    # normalize col names
    cols_map = {str(c).strip(): str(c).strip() for c in df.columns}
    df2 = df.copy()

    # accept variante de denumiri
    rename = {}
    for c in df2.columns:
        cu = str(c).strip().lower()
        if cu in ("functie","funcție","functia"):
            rename[c] = "Funcție"
        elif cu == "anexa":
            rename[c] = "Anexa"
        elif cu == "tabel":
            rename[c] = "Tabel"
        elif cu in ("studii",):
            rename[c] = "Studii"
        elif cu in ("grad",):
            rename[c] = "Grad"
        elif "treapt" in cu or "gradatie" in cu or "gradație" in cu:
            rename[c] = "Treaptă/Gradație"
        elif "salariu" in cu:
            rename[c] = "Salariu"
        elif "coef" in cu:
            rename[c] = "Coeficient"
    if rename:
        df2 = df2.rename(columns=rename)

    # asigură coloane
    for c in ["Anexa","Tabel","Funcție","Studii","Grad","Treaptă/Gradație","Salariu","Coeficient"]:
        if c not in df2.columns:
            df2[c] = ""

    # curățare
    for c in ["Anexa","Tabel","Funcție","Studii","Grad","Treaptă/Gradație"]:
        df2[c] = df2[c].fillna("").astype(str).str.strip()

    # Salariu numeric
    try:
        df2["Salariu"] = pd.to_numeric(df2["Salariu"], errors="coerce")
    except Exception:
        df2["Salariu"] = pd.NA

    # Coeficient numeric (poate lipsi)
    try:
        df2["Coeficient"] = pd.to_numeric(df2["Coeficient"], errors="coerce")
    except Exception:
        df2["Coeficient"] = pd.NA

    # păstrăm doar rândurile cu salariu (opțiuni finale)
    df2 = df2[df2["Salariu"].notna()].copy()

    return df2


@st.cache_data(show_spinner=False)
def _load_l153_grid_from_path(xlsx_path: str) -> pd.DataFrame:
    """Încarcă grila L153 dintr-un fișier XLSX (sheet: Normalized)."""
    df = pd.read_excel(xlsx_path, sheet_name="Normalized")
    return _l153_norm_df(df)


@st.cache_data(show_spinner=False)
def _load_l153_grid_from_bytes(xlsx_bytes: bytes) -> pd.DataFrame:
    """Încarcă grila L153 din bytes (uploader)."""
    bio = BytesIO(xlsx_bytes)
    df = pd.read_excel(bio, sheet_name="Normalized")
    return _l153_norm_df(df)




# ------------------------------
# L153 grilă salarizare – DB (nomenclator)
# ------------------------------
def _lege153_grid_ensure_table(conn: sqlite3.Connection) -> None:
    # ensure_schema() o creează, dar apelăm și aici ca safety
    try:
        ensure_schema(conn)
    except Exception:
        pass

def _lege153_grid_replace(conn: sqlite3.Connection, df: pd.DataFrame) -> int:
    """Înlocuiește complet nomenclatorul L153 (lege153_grid) cu df normalizat."""
    _lege153_grid_ensure_table(conn)
    cur = conn.cursor()
    cur.execute("DELETE FROM lege153_grid")
    conn.commit()

    if df is None or df.empty:
        return 0

    cols = {c.lower(): c for c in df.columns}
    # așteptăm coloanele exact ca în 'Normalized'
    anexa_c = cols.get("anexa", "Anexa")
    tabel_c = cols.get("tabel", "Tabel")
    functie_c = cols.get("funcție") or cols.get("functie") or "Funcție"
    studii_c = cols.get("studii", "Studii")
    grad_c = cols.get("grad", "Grad")
    treapta_c = cols.get("treaptă/gradație") or cols.get("treapta/gradație") or cols.get("treapta") or "Treaptă/Gradație"
    salariu_c = cols.get("salariu", "Salariu")
    coef_c = cols.get("coeficient", "Coeficient")

    rows = []
    for _, r in df.iterrows():
        an = str(r.get(anexa_c, "")).strip()
        tb = str(r.get(tabel_c, "")).strip()
        fn = str(r.get(functie_c, "")).strip()
        if not an or not tb or not fn:
            continue
        sal = r.get(salariu_c, None)
        coef = r.get(coef_c, None)
        # păstrăm doar rândurile finale (cu salariu sau coeficient)
        if (sal is None or str(sal).strip() == "" or str(sal).strip().lower() == "nan") and (
            coef is None or str(coef).strip() == "" or str(coef).strip().lower() == "nan"
        ):
            continue
        rows.append((
            an, tb, fn,
            str(r.get(studii_c, "") or "").strip(),
            str(r.get(grad_c, "") or "").strip(),
            str(r.get(treapta_c, "") or "").strip(),
            float(sal) if str(sal).strip() not in ("", "nan", "None") else None,
            float(coef) if str(coef).strip() not in ("", "nan", "None") else None,
        ))

    cur.executemany(
        """
        INSERT INTO lege153_grid(anexa, tabel, functie, studii, grad, treapta, salariu, coeficient)
        VALUES(?,?,?,?,?,?,?,?)
        """,
        rows,
    )
    conn.commit()
    return len(rows)

def _lege153_grid_is_empty(conn: sqlite3.Connection) -> bool:
    _lege153_grid_ensure_table(conn)
    cur = conn.cursor()
    cur.execute("SELECT COUNT(1) FROM lege153_grid")
    return (cur.fetchone() or [0])[0] == 0

@st.cache_data(show_spinner=False)
def _lege153_grid_df(conn_path: str) -> pd.DataFrame:
    """Cache pe path DB (pentru viteză în selectoare)."""
    conn = sqlite3.connect(conn_path, check_same_thread=False)
    try:
        cur = conn.cursor()
        cur.execute("""
            SELECT anexa, tabel, functie, studii, grad, treapta, salariu, coeficient
            FROM lege153_grid
        """)
        rows = cur.fetchall()
        df = pd.DataFrame(rows, columns=["Anexa","Tabel","Funcție","Studii","Grad","Treaptă/Gradație","Salariu","Coeficient"])
        return _l153_norm_df(df)
    finally:
        conn.close()

def _lege153_grid_load(conn: sqlite3.Connection) -> pd.DataFrame:
    """Încărcare grilă L153 din DB (cu cache)."""
    try:
        db_path = conn.execute("PRAGMA database_list").fetchone()[2]
    except Exception:
        db_path = ""
    if not db_path:
        # fallback: fără cache
        cur = conn.cursor()
        cur.execute("""
            SELECT anexa, tabel, functie, studii, grad, treapta, salariu, coeficient
            FROM lege153_grid
        """)
        rows = cur.fetchall()
        df = pd.DataFrame(rows, columns=["Anexa","Tabel","Funcție","Studii","Grad","Treaptă/Gradație","Salariu","Coeficient"])
        return _l153_norm_df(df)
    return _lege153_grid_df(db_path)

def _l153_salary_ensure_table(conn: sqlite3.Connection) -> None:
    """Tabel istoric selecție anexă/tabel/funcție -> salariu/coeficient per angajat."""
    cur = conn.cursor()
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_l153 (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER NOT NULL,
            anexa TEXT,
            tabel TEXT,
            functie TEXT,
            studii TEXT,
            grad TEXT,
            treapta TEXT,
            salariu REAL,
            coeficient REAL,
            data_aplicare TEXT,
            is_active INTEGER DEFAULT 1,
            created_at TEXT
        );
        """
    )
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_l153_emp ON employee_l153(employee_id)")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_l153_active ON employee_l153(employee_id, is_active)")

    # ---------------------------------------------------------
    # L153 – grilă salarizare (nomenclator intern)
    # ---------------------------------------------------------
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS lege153_grid (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            anexa TEXT NOT NULL,
            tabel TEXT NOT NULL,
            functie TEXT NOT NULL,
            studii TEXT,
            grad TEXT,
            treapta TEXT,
            salariu REAL,
            coeficient REAL
        );
        """
    )
    cur.execute("CREATE INDEX IF NOT EXISTS idx_lege153_grid_a_t ON lege153_grid(anexa, tabel)")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_lege153_grid_a_t_f ON lege153_grid(anexa, tabel, functie)")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_lege153_grid_a_t_f_s ON lege153_grid(anexa, tabel, functie, studii)")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_lege153_grid_a_t_f_s_g ON lege153_grid(anexa, tabel, functie, studii, grad)")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_lege153_grid_all ON lege153_grid(anexa, tabel, functie, studii, grad, treapta)")

    conn.commit()


def _employee_l153_set_active(conn: sqlite3.Connection, employee_id: int, row_id: int) -> None:
    cur = conn.cursor()
    cur.execute("UPDATE employee_l153 SET is_active=0 WHERE employee_id=?", (int(employee_id),))
    cur.execute("UPDATE employee_l153 SET is_active=1 WHERE id=? AND employee_id=?", (int(row_id), int(employee_id)))
    conn.commit()


def _employee_l153_insert(conn: sqlite3.Connection, employee_id: int, payload: dict) -> int:
    _l153_salary_ensure_table(conn)
    cur = conn.cursor()
    # dezactivează vechiul activ
    cur.execute("UPDATE employee_l153 SET is_active=0 WHERE employee_id=?", (int(employee_id),))
    cur.execute(
        """
        INSERT INTO employee_l153(
            employee_id, anexa, tabel, functie, studii, grad, treapta, salariu, coeficient,
            data_aplicare, is_active, created_at
        ) VALUES (?,?,?,?,?,?,?,?,?,?,1,?)
        """,
        (
            int(employee_id),
            str(payload.get("anexa","") or "").strip(),
            str(payload.get("tabel","") or "").strip(),
            str(payload.get("functie","") or "").strip(),
            str(payload.get("studii","") or "").strip(),
            str(payload.get("grad","") or "").strip(),
            str(payload.get("treapta","") or "").strip(),
            float(payload.get("salariu") or 0.0) if payload.get("salariu") is not None else None,
            float(payload.get("coeficient") or 0.0) if payload.get("coeficient") is not None else None,
            str(payload.get("data_aplicare","") or "").strip(),
            datetime.datetime.now().isoformat(timespec="seconds"),
        ),
    )
    conn.commit()
    return int(cur.lastrowid or 0)


def _employee_l153_list(conn: sqlite3.Connection, employee_id: int, limit: int | None = None) -> list:
    _l153_salary_ensure_table(conn)
    cur = conn.cursor()
    cur.execute(
        """
        SELECT id, anexa, tabel, functie, studii, grad, treapta, salariu, coeficient, data_aplicare, is_active, created_at
        FROM employee_l153
        WHERE employee_id=?
        ORDER BY created_at DESC, id DESC
        """,
        (int(employee_id),),
    )
    rows = cur.fetchall() or []
    return rows[:int(limit)] if limit else rows

def _employee_l153_get_active(conn: sqlite3.Connection, employee_id: int) -> dict:
    _l153_salary_ensure_table(conn)
    cur = conn.cursor()
    cur.execute(
        """
        SELECT anexa, tabel, functie, studii, grad, treapta, salariu, coeficient, data_aplicare
        FROM employee_l153
        WHERE employee_id=? AND is_active=1
        ORDER BY created_at DESC, id DESC
        LIMIT 1
        """,
        (int(employee_id),),
    )
    r = cur.fetchone()
    if not r:
        return {}
    return {
        "anexa": r[0] or "",
        "tabel": r[1] or "",
        "functie": r[2] or "",
        "studii": r[3] or "",
        "grad": r[4] or "",
        "treapta": r[5] or "",
        "salariu": r[6],
        "coeficient": r[7],
        "data_aplicare": r[8] or "",
    }

def _color_for_unit(name_upper: str):
    """Determină rolul și nivelul + culorile pentru un nume de unitate."""
    # Culori pe nivel
    level_colors = {
        "directie":      ("#457B9D", "white"),
        "serviciu":      ("#F4A261", "black"),
        "compartiment":  ("#9D4EDD", "white"),
        "birou":         ("#F7DC6F", "black"),
        "unitate":       ("#A8DADC", "black"),
        "subunitate":    ("#CED4DA", "black"),
        "default":       ("#D3E3F3", "black"),
    }

    # Culori pe rol
    role_colors = {
        "director_general": ("#0B3C5D", "white"),
        "director_adjunct": ("#145A32", "white"),
        "director":         ("#C0392B", "white"),
    }

    # 1) Rol
    role = None
    if "DIRECTOR GENERAL" in name_upper:
        role = "director_general"
    elif "DIRECTOR ADJ" in name_upper or "ADJUNCT" in name_upper:
        role = "director_adjunct"
    elif "DIRECTOR" in name_upper:
        role = "director"

    # 2) Nivel
    level = "default"
    if ("DIREC" in name_upper or "DIRECȚ" in name_upper or "DIRECŢ" in name_upper) and role is None:
        level = "directie"
    elif "SERVICIU" in name_upper or "SERV." in name_upper:
        level = "serviciu"
    elif "COMPARTIMENT" in name_upper or "COMP." in name_upper:
        level = "compartiment"
    elif "BIROU" in name_upper or "BIR." in name_upper:
        level = "birou"
    elif "UNITATE" in name_upper:
        level = "unitate"
    elif "SUBUNITATE" in name_upper:
        level = "subunitate"

    if role is not None:
        fill, font = role_colors[role]
    else:
        fill, font = level_colors.get(level, level_colors["default"])

    return role, level, fill, font


# ============================================
# 1) RADIALĂ – Director General în centru
# ============================================

def build_org_digraph_radial(
    df_units,
    *,
    fontname: str = "Calibri",
    font_size: int = 14,
    box_scale: float = 1.0,
    paper: str = "A4",
    wrap: bool = True,
    vertical_stretch: float = 1.4,
):
    """Organigramă radială (DG în centru) cu controale pentru font/casete/pagină.

    Optimizări (cerute):
    - casete mai mari (padding + lățime minimă)
    - wrap mai agresiv pentru a păstra layout-ul compact
    """
    dot = Digraph(format="png", engine="twopi")

    try:
        fs = max(8, int(font_size))
    except Exception:
        fs = 14
    try:
        bs = max(0.6, float(box_scale))
    except Exception:
        bs = 1.0

    # ↕️ aerisire verticală (folosită la canvas/page height)
    try:
        vs = max(0.8, float(vertical_stretch))
    except Exception:
        vs = 1.4

    paper_u = (paper or "A4").strip().upper()
    dot.graph_attr.update(
        overlap="false",
        splines="ortho",
        outputorder="edgesfirst",
        dpi=_org_paper_dpi(paper_u) if "_org_paper_dpi" in globals() else ("280" if paper_u == "A3" else "300"),
    )
    if "_org_paper_size" in globals():
        dot.graph_attr["size"] = _org_paper_size(paper_u)
        dot.graph_attr["ratio"] = "fill"

    min_width = 4.0 * bs * (1.0 if paper_u == "A4" else 1.1)
    pad_x = 0.75 * bs
    pad_y = 0.50 * bs

    dot.attr(
        "node",
        shape="box",
        style="filled,rounded",
        fontsize=str(fs),
        fontname=str(fontname or "Calibri"),
        margin=f"{pad_x:.2f},{pad_y:.2f}",
        width=f"{min_width:.2f}",
        fixedsize="false",
        penwidth="2.2",
    )

    # Root = DG sau prima rădăcină
    root_row = None
    try:
        m = df_units["name"].astype(str).str.upper().str.contains("DIRECTOR GENERAL")
        if m.any():
            root_row = df_units[m].iloc[0]
    except Exception:
        root_row = None

    if root_row is None:
        try:
            roots = df_units[df_units["parent_id"].isna()]
            if not roots.empty:
                root_row = roots.iloc[0]
        except Exception:
            root_row = None

    root_id = str(root_row["id"]) if root_row is not None else None
    if root_id:
        dot.graph_attr["root"] = root_id

    max_chars = 20 if paper_u == "A4" else 24
    try:
        max_chars_eff = int(max_chars * (0.95 + 0.15 * bs))
    except Exception:
        max_chars_eff = max_chars

    for _, row in df_units.iterrows():
        node_id = str(row["id"])
        raw_label = str(row.get("name") or "").strip() or f"ID {row.get('id')}"
        label = _wrap_org_label(raw_label, max_chars=max_chars_eff) if wrap and "_wrap_org_label" in globals() else raw_label

        role, _, fill, fontc = _color_for_unit(raw_label.upper())

        node_kwargs = {}
        if role == "director_general":
            node_kwargs.update(
                peripheries="2",
                penwidth="2.8",
                fontsize=str(fs + 6),
                fontname=f"{fontname} Bold",
                width=f"{max(min_width, 5.0 * bs):.2f}",
            )

        dot.node(node_id, label=label, fillcolor=fill, fontcolor=fontc, **node_kwargs)

    for _, row in df_units.iterrows():
        if pd.notna(row.get("parent_id")):
            try:
                dot.edge(str(int(row["parent_id"])), str(int(row["id"])))
            except Exception:
                continue

    return dot
def _wrap_org_label(text: str, max_chars: int = 24) -> str:
    """Împarte textul pe rânduri (Graphviz) pentru lizibilitate."""
    try:
        s = (text or "").strip()
    except Exception:
        return str(text)

    if not s:
        return ""

    # Normalize whitespace
    s = re.sub(r"\s+", " ", s).strip()

    # Dacă are deja rânduri, îl lăsăm
    if "\n" in s:
        return s

    # Split pe delimitatori naturali
    parts = re.split(r"\s*(?:[-–—/|]|\(|\)|,|;|:)\s*", s)
    words = []
    # dacă split-ul a produs prea multe bucăți mici, revenim pe words simple
    if len(parts) <= 1:
        words = s.split(" ")
    else:
        # reconstruim cu spații (păstrăm totuși ordinea)
        words = s.split(" ")

    lines = []
    cur = ""
    for w in words:
        if not w:
            continue
        if not cur:
            cur = w
        elif len(cur) + 1 + len(w) <= max_chars:
            cur = cur + " " + w
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)

    # Limităm la maxim 4 rânduri ca să nu devină prea înalt
    if len(lines) > 4:
        lines = lines[:3] + [" ".join(lines[3:])]

    return "\n".join(lines)


def _org_paper_scale(paper: str) -> float:
    p = (paper or "screen").strip().upper()
    if p == "A3":
        return 1.25
    if p == "A0":
        return 2.1
    if p == "A4":
        return 1.0
    return 1.0



def _org_paper_size(paper: str) -> str:
    p = (paper or "A4").strip().upper()
    if p == "A3":
        return "11.69,16.54"  # inches
    if p == "A0":
        return "33.11,46.81"
    # default A4
    return "8.27,11.69"

def _org_paper_dpi(paper: str) -> str:
    p = (paper or "A4").strip().upper()
    if p == "A3":
        return "280"
    if p == "A0":
        return "220"
    return "300"

def build_org_digraph_vertical(
    df_units,
    *,
    fontname: str = "Calibri",
    font_size: int = 14,
    box_scale: float = 1.0,
    paper: str = "A4",
    wrap: bool = True,
    vertical_stretch: float = 1.4,
):
    """Organigramă verticală (DG SUS) – compactă, lizibilă, fără shrink."""

    # --- siguranță parametri ---
    try:
        fs = max(8, int(font_size))
    except Exception:
        fs = 14
    try:
        bs = max(0.6, float(box_scale))
    except Exception:
        bs = 1.0

    
    # ↕️ aerisire pe înălțime (niveluri) – influențează înălțimea canvas-ului/exportului
    try:
        vs = max(0.8, float(vertical_stretch))
    except Exception:
        vs = 1.4
    paper_u = (paper or "A4").strip().upper()

    # --- SPAȚIERE MINIMĂ (NU depinde de zoom) ---
    if paper_u == "A4":
        ranksep = 0.55
        nodesep = 0.18
        dpi_fallback = "300"
        max_chars = 18
    else:  # A3 etc.
        ranksep = 0.65
        nodesep = 0.22
        dpi_fallback = "280"
        max_chars = 22

    # --- Graphviz ---
    dot = Digraph(format="png", engine="dot")
    dot.graph_attr.update(
        rankdir="TB",
        dpi=_org_paper_dpi(paper_u) if "_org_paper_dpi" in globals() else dpi_fallback,
        ranksep=f"{ranksep:.2f}",
        nodesep=f"{nodesep:.2f}",
        overlap="false",
        splines="ortho",
        outputorder="edgesfirst",
        newrank="true",
    )

    # ↕️ Forțează „pagina/canvas”-ul să crească pe verticală (fără shrink).
    # IMPORTANT: nu seta ratio="fill" aici.
    if "_org_paper_size" in globals():
        base_size = _org_paper_size(paper_u)
    else:
        base_size = "8.27,11.69"  # A4 fallback (inches)

    w_str, h_str = [x.strip() for x in base_size.split(",")]
    try:
        w = float(w_str)
        h = float(h_str) * vs
        dot.graph_attr["size"] = f"{w:.2f},{h:.2f}!"
    except Exception:
        pass

    # --- CASATE (mari + padding real) ---
    min_width = 4.2 * bs * (1.0 if paper_u == "A4" else 1.1)
    pad_x = 0.85 * bs
    pad_y = 0.55 * bs

    dot.attr(
        "node",
        shape="box",
        style="filled,rounded",
        fontsize=str(fs),
        fontname=str(fontname or "Calibri"),
        margin=f"{pad_x:.2f},{pad_y:.2f}",
        width=f"{min_width:.2f}",
        fixedsize="false",
        penwidth="2.4",
    )

    # --- ROOT (Director General) ---
    root_row = None
    try:
        m = df_units["name"].astype(str).str.upper().str.contains("DIRECTOR GENERAL")
        if m.any():
            root_row = df_units[m].iloc[0]
    except Exception:
        root_row = None

    if root_row is None:
        roots = df_units[df_units["parent_id"].isna()]
        if not roots.empty:
            root_row = roots.iloc[0]

    root_id = str(root_row["id"]) if root_row is not None else None
    if root_id:
        dot.graph_attr["root"] = root_id

    # --- WRAP adaptat ---
    try:
        max_chars_eff = int(max_chars * (0.95 + 0.15 * bs))
    except Exception:
        max_chars_eff = max_chars

    # --- NODURI ---
    for _, row in df_units.iterrows():
        try:
            node_id = str(int(row["id"]))
        except Exception:
            continue

        raw_label = str(row.get("name") or "").strip() or f"ID {row.get('id')}"
        label = (
            _wrap_org_label(raw_label, max_chars=max_chars_eff)
            if wrap and "_wrap_org_label" in globals()
            else raw_label
        )

        role, _, fill, fontc = _color_for_unit(raw_label.upper())

        node_kwargs = {}
        if role == "director_general":
            node_kwargs.update(
                peripheries="2",
                penwidth="2.8",
                fontsize=str(fs + 6),
                fontname=f"{fontname} Bold",
                width=f"{max(min_width, 5.2 * bs):.2f}",
            )

        dot.node(node_id, label=label, fillcolor=fill, fontcolor=fontc, **node_kwargs)

    # --- MUCHII ---
    for _, row in df_units.iterrows():
        if pd.notna(row.get("parent_id")):
            try:
                dot.edge(str(int(row["parent_id"])), str(int(row["id"])))
            except Exception:
                pass

    return dot

def build_director_top_graph(
    df_units,
    *,
    fontname: str = "Calibri",
    font_size: int = 14,
    box_scale: float = 1.0,
    paper: str = "A4",
    wrap: bool = True,
):
    """
    Director General sus, sub el adjuncți / directori, apoi structurile lor.
    Cu control pe font + mărime casete + format pagină.
    """
    scale = _org_paper_scale(paper) * max(0.6, float(box_scale))

    dot = Digraph(format="png", engine="dot")
    dot.graph_attr.update(
        rankdir="TB",
        ranksep=f"{2.8 * (0.85 + 0.15 * scale):.2f}",
        nodesep=f"{0.15 * (0.85 + 0.15 * scale):.2f}",
        overlap="false",
        splines="ortho",
        newrank="true",
    )

    margin_x = 0.9 * (0.85 + 0.15 * scale)
    margin_y = 0.6 * (0.85 + 0.15 * scale)

    # Casete late, pentru prezentare
    base_w = 3.5
    base_h = 1.2
    w = base_w * scale
    h = base_h * scale

    dot.attr(
        "node",
        shape="box",
        style="filled,rounded",
        fontsize=str(int(max(8, int(font_size)))),
        fontname=str(fontname or "Calibri"),
        margin=f"{margin_x:.2f},{margin_y:.2f}",
        penwidth="2.2",
        width=f"{w:.2f}",
        height=f"{h:.2f}",
        fixedsize="true",
    )

    # find root
    root_row = None
    try:
        m = df_units["name"].astype(str).str.upper().str.contains("DIRECTOR GENERAL")
        if m.any():
            root_row = df_units[m].iloc[0]
    except Exception:
        root_row = None
    if root_row is None:
        roots = df_units[df_units["parent_id"].isna()]
        if not roots.empty:
            root_row = roots.iloc[0]

    root_id = str(root_row["id"]) if root_row is not None else None
    if root_id:
        dot.graph_attr["root"] = root_id

    adj_ids = []
    max_chars = int(26 * (0.9 + 0.25 * scale))

    for _, row in df_units.iterrows():
        node_id = str(row["id"])
        raw_label = str(row["name"]) if pd.notna(row.get("name")) else f"ID {row.get('id')}"
        label = _wrap_org_label(raw_label, max_chars=max_chars) if wrap else raw_label
        role, _, fill, fontc = _color_for_unit(raw_label.upper())

        node_kwargs = {}
        if role == "director_general":
            node_kwargs.update(
                peripheries="2",
                penwidth="2.6",
                fontsize=str(int(max(font_size + 8, font_size * 1.5))),
                fontname=f"{fontname} Bold",
            )

        dot.node(node_id, label=label, fillcolor=fill, fontcolor=fontc, **node_kwargs)

        if root_row is not None and pd.notna(row.get("parent_id")):
            try:
                if int(row["parent_id"]) == int(root_row["id"]) and role in ("director_adjunct", "director"):
                    adj_ids.append(node_id)
            except Exception:
                pass

    # Rank root
    if root_id:
        with dot.subgraph() as s0:
            s0.attr(rank="same")
            s0.node(root_id)

    # Rank adjuncți sub DG
    if adj_ids:
        with dot.subgraph() as s1:
            s1.attr(rank="same")
            for nid in adj_ids:
                s1.node(nid)

    # Muchii
    for _, row in df_units.iterrows():
        if pd.notna(row.get("parent_id")):
            try:
                dot.edge(str(int(row["parent_id"])), str(row["id"]))
            except Exception:
                pass

    return dot


def build_org_digraph_manual_positions(
    df_units,
    positions: dict,
    *,
    fontname: str = "Calibri",
    font_size: int = 14,
    box_scale: float = 1.0,
    paper: str = "A4",
    wrap: bool = True,
    vertical_stretch: float = 1.4,
):
    """Graphviz 'neato' cu poziții manuale (pos=x,y!) + muchii după parent_id.
    positions: {id(str/int): (x,y)} în coordonate editor (px).
    """
    from graphviz import Digraph
    import pandas as pd

    df = df_units.copy()
    cols = {c.upper(): c for c in df.columns}

    # normalize
    rename = {}
    if "ID_UNITATE" in cols and "id" not in df.columns:
        rename[cols["ID_UNITATE"]] = "id"
    if "DENUMIRE" in cols and "name" not in df.columns:
        rename[cols["DENUMIRE"]] = "name"
    if "PARENT_ID" in cols and "parent_id" not in df.columns:
        rename[cols["PARENT_ID"]] = "parent_id"
    if "TIP" in cols and "type" not in df.columns:
        rename[cols["TIP"]] = "type"
    if rename:
        df = df.rename(columns=rename)

    for need in ["id", "name", "parent_id"]:
        if need not in df.columns:
            if need == "id":
                df["id"] = range(1, len(df) + 1)
            elif need == "name":
                df["name"] = ""
            else:
                df["parent_id"] = pd.NA

    df["id"] = pd.to_numeric(df["id"], errors="coerce")
    df["parent_id"] = pd.to_numeric(df["parent_id"], errors="coerce")

    try:
        vs = float(vertical_stretch)
    except Exception:
        vs = 1.4
    if vs <= 0:
        vs = 1.0

    # DPI per paper (folosim funcția deja existentă dacă există)
    paper_u = (paper or "A4").strip().upper()
    try:
        dpi = _org_paper_dpi(paper_u)  # type: ignore
    except Exception:
        dpi = 220 if paper_u == "A3" else 180

    dot = Digraph(format="png", engine="neato")
    dot.graph_attr.update(
        overlap="false",
        splines="ortho",
        outputorder="edgesfirst",
        dpi=str(int(dpi)),
    )
    dot.node_attr.update(
        shape="box",
        style="rounded,filled",
        fontname=fontname,
        fontsize=str(int(max(8, int(font_size)))),
    )

    # dimensiuni casetă (în inch la Graphviz). Ajustăm empiric.
    try:
        bs = float(box_scale)
    except Exception:
        bs = 1.0
    bs = max(0.6, bs)

    width_in = 1.8 * bs
    height_in = 0.65 * bs

    def _wrap_label(s: str, max_len: int = 22) -> str:
        if not wrap:
            return s
        s = (s or "").strip()
        if len(s) <= max_len:
            return s
        words = s.split()
        out, line = [], ""
        for w in words:
            if len(line) + len(w) + 1 <= max_len:
                line = (line + " " + w).strip()
            else:
                out.append(line)
                line = w
        if line:
            out.append(line)
        return "\n".join(out) if out else s

    # nodes
    for _, r in df.iterrows():
        if pd.isna(r["id"]):
            continue
        nid = str(int(r["id"]))
        label = _wrap_label(str(r.get("name", "")))
        tval = r.get("type")
        if tval is not None and not pd.isna(tval) and str(tval).strip() != "":
            label = f"{label}\n({str(tval)})"

        # coordonate
        p = positions.get(nid) or positions.get(int(nid)) or positions.get(str(nid))
        if not p:
            # fallback: nu avem coordonate -> le lăsăm la (0,0)
            x, y = 0.0, 0.0
        else:
            x, y = float(p[0]), float(p[1])

        # inversăm Y ca să fie "sus -> jos" și aplicăm vertical_stretch
        y = -y * vs

        dot.node(
            nid,
            label=label,
            width=str(width_in),
            height=str(height_in),
            fixedsize="true",
            pos=f"{x},{y}!",
        )

    # edges
    for _, r in df.iterrows():
        try:
            cid = int(r["id"])
        except Exception:
            continue
        pid = r.get("parent_id", None)
        if pid is None or pd.isna(pid):
            continue
        try:
            dot.edge(str(int(pid)), str(int(cid)))
        except Exception:
            pass

    return dot


def generate_org_chart_files(
    df_units,
    layout_mode: str = "Ierarhică",
    *,
    paper: str = "A4",
    style_preset: str = "blue",
    strict_columns: bool = True,
    fontname: str = "Calibri",
    font_size: int = 14,
    box_scale: float = 1.0,
    wrap: bool = True,
    vertical_stretch: float = 1.4,
    use_manual_layout: bool = False,
    layout_positions: dict | None = None,
):
    """
    Generează fișierele pentru organigrama vizuală (Graphviz) și returnează un dict:
      {png, jpg, pdf, docx, dot}

    - Dacă use_manual_layout=True și layout_positions există, exportul folosește coordonatele manuale (neato).
    - Altfel, folosește layout automat (Radial/Ierarhic).
    """
    if df_units is None or getattr(df_units, "empty", True):
        return None

    import pandas as pd
    import io
    from PIL import Image
    from docx import Document
    from docx.shared import Inches

    df = df_units.copy()
    cols = {c.upper(): c for c in df.columns}

    # normalize columns
    rename = {}
    if "ID_UNITATE" in cols and "id" not in df.columns:
        rename[cols["ID_UNITATE"]] = "id"
    if "DENUMIRE" in cols and "name" not in df.columns:
        rename[cols["DENUMIRE"]] = "name"
    if "PARENT_ID" in cols and "parent_id" not in df.columns:
        rename[cols["PARENT_ID"]] = "parent_id"
    if "TIP" in cols and "type" not in df.columns:
        rename[cols["TIP"]] = "type"
    if rename:
        df = df.rename(columns=rename)

    for need in ["id", "name", "parent_id"]:
        if need not in df.columns:
            if need == "id":
                df["id"] = range(1, len(df) + 1)
            elif need == "name":
                df["name"] = ""
            else:
                df["parent_id"] = pd.NA

    df["id"] = pd.to_numeric(df["id"], errors="coerce")
    df["parent_id"] = pd.to_numeric(df["parent_id"], errors="coerce")

    lm = (layout_mode or "").strip()

    # --- build DOT ---
    dot = None
    if bool(use_manual_layout) and isinstance(layout_positions, dict) and len(layout_positions) > 0:
        dot = build_org_digraph_manual_positions(
            df,
            layout_positions,
            fontname=fontname,
            font_size=int(font_size),
            box_scale=float(box_scale),
            paper=paper,
            wrap=wrap,
            vertical_stretch=float(vertical_stretch),
        )
        lm = "Manual (Drag&Drop)"
    else:
        if "Radial" in lm or "Centru" in lm:
            dot = build_org_digraph_radial(
                df,
                fontname=fontname,
                font_size=int(font_size),
                box_scale=float(box_scale),
                paper=paper,
                wrap=wrap,
            )
        else:
            dot = build_org_digraph_vertical(
                df,
                fontname=fontname,
                font_size=int(font_size),
                box_scale=float(box_scale),
                paper=paper,
                wrap=wrap,
                vertical_stretch=float(vertical_stretch),
            )

    dot_src = dot.source

    # --- render bytes from DOT ---
    try:
        png_bytes, pdf_bytes, _pptx = generate_org_files_from_dot(dot_src)
    except Exception as e:
        try:
            st.error(f"Eroare la randarea Graphviz: {e}")
        except Exception:
            pass
        return None

    # JPG (din PNG)
    jpg_buf = io.BytesIO()
    Image.open(io.BytesIO(png_bytes)).convert("RGB").save(jpg_buf, format="JPEG", quality=92)
    jpg_bytes = jpg_buf.getvalue()

    # DOCX (o pagină cu titlu + imagine)
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    doc.add_heading("Organigramă", level=1)
    doc.add_paragraph(
        f"Layout: {lm} | Format: {paper} | Font: {fontname} | Font size: {int(font_size)} | Casete: {float(box_scale):.2f}"
    )

    img_stream = io.BytesIO(png_bytes)
    try:
        doc.add_picture(img_stream, width=Inches(6.5))
    except Exception:
        img_stream.seek(0)
        doc.add_picture(img_stream)

    out = io.BytesIO()
    doc.save(out)
    docx_bytes = out.getvalue()

    return {
        "png": png_bytes,
        "jpg": jpg_bytes,
        "pdf": pdf_bytes,
        "docx": docx_bytes,
        "dot": dot_src,
    }

# ============================================================
# ORGANIGRAMĂ INTERACTIVĂ (Cytoscape.js) – expand/collapse + export
# ============================================================

def _org_df_normalize_for_cyto(df_units: pd.DataFrame) -> pd.DataFrame:
    """Normalizează df_units la coloane: id, name, parent_id, type (opțional)."""
    import pandas as pd
    if df_units is None or getattr(df_units, "empty", True):
        return pd.DataFrame(columns=["id", "name", "parent_id", "type"])

    df = df_units.copy()
    cols = {c.upper(): c for c in df.columns}

    rename = {}
    if "ID_UNITATE" in cols and "id" not in df.columns:
        rename[cols["ID_UNITATE"]] = "id"
    if "DENUMIRE" in cols and "name" not in df.columns:
        rename[cols["DENUMIRE"]] = "name"
    if "PARENT_ID" in cols and "parent_id" not in df.columns:
        rename[cols["PARENT_ID"]] = "parent_id"
    if "TIP" in cols and "type" not in df.columns:
        rename[cols["TIP"]] = "type"
    if rename:
        df = df.rename(columns=rename)

    for need in ["id", "name", "parent_id"]:
        if need not in df.columns:
            if need == "id":
                df["id"] = range(1, len(df) + 1)
            elif need == "name":
                df["name"] = ""
            else:
                df["parent_id"] = pd.NA
    if "type" not in df.columns:
        df["type"] = ""

    df["id"] = pd.to_numeric(df["id"], errors="coerce")
    df["parent_id"] = pd.to_numeric(df["parent_id"], errors="coerce")
    df = df.dropna(subset=["id"]).copy()
    df["id"] = df["id"].astype(int)
    # păstrăm parent_id ca int sau None
    df["parent_id"] = df["parent_id"].apply(lambda x: int(x) if pd.notna(x) else None)
    df["name"] = df["name"].astype(str).fillna("")
    df["type"] = df["type"].astype(str).fillna("")
    return df

# ============================================================
# PAGINA PRINCIPALĂ – Home Dashboard\
# ============================================================


# Funcție de navigare centralizată (pentru butoane, istoric recent etc.)
def goto(page_label: str):
    st.session_state["main_choice"] = page_label
    st.rerun()

def page_home(conn, cfg):
    _init_home_state()
    apply_premium_theme()

    # Marker pentru a delimita clar scope-ul Dashboard-ului (Home)
    st.markdown('<span id="dash-scope"></span>', unsafe_allow_html=True)

    user = st.session_state.get("username") or "administrator"

    # === HERO / WELCOME SECTION ===
    st.markdown(
        """
        <div class="home-hero">
          <div class="title">Bun venit în platformă</div>
          <div class="subtitle">
            Platformă instituțională pentru administrarea resurselor umane, concepută pentru acces rapid,
            claritate operațională și administrare eficientă.
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.markdown("")

    # === MODULE PRINCIPALE: GRID DE CARDURI PREMIUM ===
    st.markdown('<div class="home-grid-title">◆ Module principale</div>', unsafe_allow_html=True)
    col1, col2, col3 = st.columns(3)

    with col1:
        st.markdown(
            """
            <div class="home-card">
              <div class="home-card-title">Organigramă</div>
              <div class="home-card-text">
                Vizualizează structura organizațională și relațiile dintre departamente, funcții și angajați.
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.markdown(
            """
            <div class="home-card">
              <div class="home-card-title">Angajați</div>
              <div class="home-card-text">
                Accesează rapid evidența personalului și gestionează informațiile esențiale ale angajaților.
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )

    with col2:
        st.markdown(
            """
            <div class="home-card">
              <div class="home-card-title">Stat de funcții</div>
              <div class="home-card-text">
                Consultă structura posturilor și distribuția funcțiilor din cadrul organizației.
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.markdown(
            """
            <div class="home-card">
              <div class="home-card-title">Pontaj</div>
              <div class="home-card-text">
                Monitorizează prezența, timpul de lucru și situațiile administrative asociate.
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )

    with col3:
        st.markdown(
            """
            <div class="home-card">
              <div class="home-card-title">Dosar profesional</div>
              <div class="home-card-text">
                Administrează documentele, istoricul profesional și informațiile relevante pentru fiecare angajat.
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.markdown(
            """
            <div class="home-card">
              <div class="home-card-title">Configurări</div>
              <div class="home-card-text">
                Gestionează setările aplicației și personalizează administrarea platformei.
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )

    # Pagina Acasă rămâne deliberat simplă: hero + overview de module.

# ============================================================
# ORGANIGRAMĂ INTERACTIVĂ (Cytoscape.js) – expand/collapse + export
# ============================================================ 

def render_org_cytoscape_collapsible(
    df_units: pd.DataFrame | None = None,
    *,
    elements_payload: dict | None = None,   # ✅ NEW
    height: int = 780,
    key: str = "org_cyto",
    box_scale: float = 1.0,
    font_size: int = 14,
    spacing: float = 1.0,
    start_collapsed: bool = True,
    enable_drag: bool = True,
    enable_save_layout: bool = True,
    use_saved_layout: bool = True,
    color_overrides: dict | None = None,
):
    """
    Organigramă interactivă (Cytoscape.js):
    - click pe nod: expand/collapse subarbore
    - pan/zoom
    - drag noduri (opțional)
    - salvare layout (persistență în browser via localStorage) + reset
    - export PNG/JPEG/PDF (din browser)

    Notă importantă:
    - În Streamlit, components.html NU poate trimite direct date înapoi în Python fără un component custom.
      De aceea, "Salvează" persistă pozițiile LOCAL în browser (localStorage) și le reîncarcă automat.
    """
    import json
    import streamlit.components.v1 as components

    # parametri UI -> valori safe (rămân exact ca la tine)
    try:
        bs = max(0.6, float(box_scale))
    except Exception:
        bs = 1.0
    try:
        fs = max(24, int(float(font_size) * 2))
    except Exception:
        fs = 28
    try:
        sp = max(0.6, float(spacing))
    except Exception:
        sp = 1.0

    # dimensiuni nod (vizibil!) + text wrap
    text_max = int(190 * bs * 8)
    node_w_min = int(220 * bs * 10)
    node_h_min = int(86 * bs * 10)
    pad = int(10 * bs * 3)
    border = int(max(2, 2 * bs * 2))

    # layout spacing (crește cu casetele)
    node_sep = int(60 * bs * sp * 6)
    rank_sep = int(165 * bs * sp * 8)

    # ===== DATA SOURCE SWITCH =====
    # 1) dacă avem payload complet (unit+pos+emp), îl folosim direct
    if elements_payload is not None:
        # folosim direct dict-ul primit, fără să atingem df_units
        payload_json = json.dumps(elements_payload, ensure_ascii=False)
    else:
        # 2) fallback legacy: doar unități din df_units (comportament identic cu cel vechi)
        df = _org_df_normalize_for_cyto(df_units)
        if df.empty:
            st.info("Nu există date pentru organigramă.")
            return

        overrides = color_overrides or {}
        nodes = []
        edges = []

        for _, r in df.iterrows():
            nid_int = int(r["id"])
            nid = str(nid_int)

            name = (str(r.get("name", "")) or "").strip()
            t = (str(r.get("type", "")) or "").strip()

            label = name
            if t:
                label = f"{label}\\n({t})"

            role, level, fill, fontc = _color_for_unit(name.upper())
            if role and role in overrides:
                fill = overrides.get(role, fill)
            if level and level in overrides:
                fill = overrides.get(level, fill)
            if "default" in overrides and (role not in overrides) and (level not in overrides):
                fill = overrides.get("default", fill)

            # IMPORTANT: adaugăm și type ca să putem trata click diferit (unit/pos/emp)
            nodes.append({"data": {"id": nid, "label": label, "bg": str(fill), "fc": str(fontc), "type": "unit"}})

            pid = r.get("parent_id", None)
            if pid is None or pd.isna(pid):
                continue
            try:
                src = str(int(pid))
            except Exception:
                continue
            edges.append({"data": {"id": f"e{src}-{nid}", "source": src, "target": nid, "type": "unit_edge"}})

        payload = {"nodes": nodes, "edges": edges}
        payload_json = json.dumps(payload, ensure_ascii=False)

    # cheie de persistare (per pagină / per DB) – suficient de stabilă
    layout_storage_key = f"SocratesHR_org_layout_{key}"



    # template HTML (FĂRĂ f-string!)
    html = r"""
<div id="__CY_ID__" style="width:100%; height:__CY_H__px; border:1px solid #e2e2e2; border-radius:10px;"></div>

<div style="display:flex; gap:8px; flex-wrap:wrap; margin-top:8px;">
  <button id="__BTN_FIT__" style="padding:6px 10px;">Fit</button>
  <button id="__BTN_EXPAND__" style="padding:6px 10px;">Expand all</button>
  <button id="__BTN_COLLAPSE__" style="padding:6px 10px;">Collapse all</button>

  <button id="__BTN_SAVE__" style="padding:6px 10px; display:__SAVE_DISPLAY__;">Salvează aranjarea</button>
  <button id="__BTN_RESET__" style="padding:6px 10px; display:__SAVE_DISPLAY__;">Resetează aranjarea</button>

  <button id="__BTN_PNG__" style="padding:6px 10px;">Export PNG</button>
  <button id="__BTN_JPG__" style="padding:6px 10px;">Export JPEG</button>
  <button id="__BTN_PDF__" style="padding:6px 10px;">Export PDF</button>
</div>

<script src="https://unpkg.com/cytoscape@3.28.1/dist/cytoscape.min.js"></script>
<script src="https://unpkg.com/dagre@0.8.5/dist/dagre.min.js"></script>
<script src="https://unpkg.com/cytoscape-dagre@2.5.0/cytoscape-dagre.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/jspdf/2.5.1/jspdf.umd.min.js"></script>

<script>
(function() {
  const data = __DATA_JSON__;
  const STORAGE_KEY = "__LAYOUT_KEY__";

  function safeParse(jsonText) {
    try { return JSON.parse(jsonText); } catch (e) { return null; }
  }

  function loadSavedPositions() {
    try {
      const txt = localStorage.getItem(STORAGE_KEY);
      if (!txt) return null;
      const obj = safeParse(txt);
      if (!obj || typeof obj !== "object") return null;
      return obj;
    } catch (e) {
      return null;
    }
  }

  function savePositions(map) {
    try {
      localStorage.setItem(STORAGE_KEY, JSON.stringify(map));
      return true;
    } catch (e) {
      return false;
    }
  }

  function clearPositions() {
    try { localStorage.removeItem(STORAGE_KEY); } catch (e) {}
  }

  let cy = cytoscape({
    container: document.getElementById("__CY_ID__"),
    elements: data,
    style: [
      {
        selector: 'node',
        style: {
          'shape': 'round-rectangle',
          'background-color': 'data(bg)',
          'border-color': '#7aa6c2',
          'border-width': __BORDER__,
          'color': 'data(fc)',
          'label': 'data(label)',
          'text-wrap': 'wrap',
          'text-max-width': __TEXT_MAX__,
          'font-size': __FONT_SIZE__,
          'text-valign': 'center',
          'text-halign': 'center',
          'padding': __PAD__,

          /* IMPORTANT: mărime vizibilă a casetei */
          'min-width': __NODE_W_MIN__,
          'width': 'label',
          'min-height': __NODE_H_MIN__,
          'height': 'auto'
        }
      },
      {
        selector: 'edge',
        style: {
          'curve-style': 'taxi',
          'taxi-direction': 'downward',
          'taxi-turn': 18,
          'width': 2,
          'line-color': '#b7b7b7',
          'target-arrow-shape': 'triangle',
          'target-arrow-color': '#b7b7b7'
        }
      },
      { selector: '.hidden', style: { 'display': 'none' } },
      { selector: '.focused-node', style: { 'border-color': '#ff4b4b', 'border-width': 6 } }
    ],
    layout: {
      name: 'dagre',
      rankDir: 'TB',
      nodeSep: __NODE_SEP__,
      edgeSep: 10,
      rankSep: __RANK_SEP__
    },
    wheelSensitivity: 0.15
  });

  // allow drag?
  cy.autoungrabify(__ENABLE_DRAG__ ? false : true);

  // map successors (tree)
  const children = {};
  data.edges.forEach(e => {
    const s = e.data.source, t = e.data.target;
    if (!children[s]) children[s] = [];
    children[s].push(t);
  });

  function descendants(root) {
    const out = [];
    const stack = (children[root] || []).slice();
    while (stack.length) {
      const n = stack.pop();
      out.push(n);
      const ch = children[n] || [];
      for (let i=0;i<ch.length;i++) stack.push(ch[i]);
    }
    return out;
  }

  function hideSubtree(rootId) {
    const desc = descendants(rootId);
    desc.forEach(id => {
      cy.getElementById(id).addClass('hidden');
      cy.edges('[target = "' + id + '"]').addClass('hidden');
    });
    cy.edges('[source = "' + rootId + '"]').addClass('hidden');
  }

  function showSubtree(rootId) {
    const desc = descendants(rootId);
    desc.forEach(id => {
      cy.getElementById(id).removeClass('hidden');
      cy.edges('[target = "' + id + '"]').removeClass('hidden');
    });
    cy.edges('[source = "' + rootId + '"]').removeClass('hidden');
  }

  const collapsed = {};

  function collapseAll() {
    cy.nodes().forEach(n => { collapsed[n.id()] = false; });

    const allTargets = new Set(data.edges.map(e => e.data.target));
    const roots = data.nodes.map(n => n.data.id).filter(id => !allTargets.has(id));

    roots.forEach(r => {
      (children[r] || []).forEach(ch => {
        collapsed[ch] = true;
        hideSubtree(ch);
      });
    });
  }

  function expandAll() {
    cy.elements().removeClass('hidden');
    Object.keys(collapsed).forEach(k => collapsed[k] = false);
  }

  function fitNow() { cy.fit(undefined, 40); }

  // 🔁 Reaplică layout salvat (dacă există)
  function applySavedIfAny() {
    if (__USE_SAVED__ !== true) return false;
    const pos = loadSavedPositions();
    if (!pos) return false;

    // aplicăm poziții
    Object.keys(pos).forEach(id => {
      const p = pos[id];
      if (!p || typeof p.x !== "number" || typeof p.y !== "number") return;
      const el = cy.getElementById(id);
      if (el && el.length) {
        el.position({ x: p.x, y: p.y });
      }
    });

    // păstrăm pozițiile (preset)
    try { cy.layout({ name: 'preset' }).run(); } catch(e) {}
    fitNow();
    return true;
  }

  // init collapse / fit
  const applied = applySavedIfAny();
  if (__START_COLLAPSED__ === true) {
    collapseAll();
  } else {
    expandAll();
  }
  if (!applied) fitNow();

  // focus din URL (dacă există ?focus=<node_id>)
  try {
    const urlAtInit = new URL(window.location.href);
    const focusId = urlAtInit.searchParams.get('focus');
    if (focusId) {
      const focusNode = cy.getElementById(focusId);
      if (focusNode && focusNode.length) {
        focusNode.addClass('focused-node');
        cy.fit(focusNode, 80);
      }
    }
  } catch (e) {}

  cy.on('tap', 'node', function(evt) {
    const n = evt.target;
    const id = n.id();
    const t = n.data('type') || 'unit';

    try {
      const url = new URL(window.location.href);
      url.searchParams.set('focus', id);

      if (t === 'emp') {
        const empId = n.data('emp_id');
        if (empId != null) {
          url.searchParams.set('emp_id', String(empId));
          url.searchParams.delete('pos_id');
          url.searchParams.delete('unit_id');
        }
      } else if (t === 'pos') {
        const empId = n.data('emp_id');
        if (empId != null) {
          url.searchParams.set('emp_id', String(empId));
          url.searchParams.delete('pos_id');
          url.searchParams.delete('unit_id');
        } else {
          const posId = n.data('pos_id');
          if (posId != null) {
            url.searchParams.set('pos_id', String(posId));
            url.searchParams.delete('emp_id');
            url.searchParams.delete('unit_id');
          }
        }
      } else {
        // unit – setăm unit_id în query params și păstrăm logica de collapse/expand
        const unitId = n.data('unit_id') != null ? n.data('unit_id') : id;
        url.searchParams.set('unit_id', String(unitId));
        url.searchParams.delete('emp_id');
        url.searchParams.delete('pos_id');

        if (children[id] && children[id].length !== 0) {
          collapsed[id] = !collapsed[id];
          if (collapsed[id]) hideSubtree(id);
          else showSubtree(id);
        }
      }

      window.location.href = url.toString();
    } catch (e) {
      // dacă URL API nu e suportat, măcar păstrăm comportamentul vechi pentru unități
      if (t === 'unit') {
        if (!children[id] || children[id].length === 0) return;
        collapsed[id] = !collapsed[id];
        if (collapsed[id]) hideSubtree(id);
        else showSubtree(id);
      }
    }
  });

  // buttons
  document.getElementById("__BTN_FIT__").onclick = function(){ fitNow(); };
  document.getElementById("__BTN_EXPAND__").onclick = function(){ expandAll(); fitNow(); };
  document.getElementById("__BTN_COLLAPSE__").onclick = function(){ collapseAll(); fitNow(); };

  // save / reset
  const btnSave = document.getElementById("__BTN_SAVE__");
  const btnReset = document.getElementById("__BTN_RESET__");

  if (btnSave) {
    btnSave.onclick = function() {
      const map = {};
      cy.nodes().forEach(n => {
        const p = n.position();
        map[n.id()] = { x: p.x, y: p.y };
      });
      const ok = savePositions(map);
      if (ok) alert("Aranjarea a fost salvată în acest browser.");
      else alert("Nu am putut salva (browser/storage blocat).");
    };
  }

  if (btnReset) {
    btnReset.onclick = function() {
      clearPositions();
      // refacem layout automat
      try {
        cy.layout({
          name: 'dagre',
          rankDir: 'TB',
          nodeSep: __NODE_SEP__,
          edgeSep: 10,
          rankSep: __RANK_SEP__
        }).run();
      } catch(e) {}
      if (__START_COLLAPSED__ === true) collapseAll();
      else expandAll();
      fitNow();
      alert("Aranjarea a fost resetată (layout automat).");
    };
  }

  function downloadURI(uri, name) {
    if (!uri || typeof uri !== "string" || uri.length < 100) {
      alert("Nu am putut genera fișierul pentru export (conținut gol). Încearcă să apeși întâi Fit / Expand all și apoi din nou Export.");
      return;
    }
    const link = document.createElement("a");
    link.download = name;
    link.href = uri;
    document.body.appendChild(link);
    link.click();
    document.body.removeChild(link);
  }

  function exportPng() {
    cy.fit(undefined, 40);
    setTimeout(function() {
      const png = cy.png({ full: true, scale: 2, bg: "#ffffff" });
      downloadURI(png, "organigrama.png");
    }, 120);
  }

  function exportJpg() {
    cy.fit(undefined, 40);
    setTimeout(function() {
      const jpg = cy.jpg({ full: true, quality: 0.95, scale: 2, bg: "#ffffff" });
      downloadURI(jpg, "organigrama.jpg");
    }, 120);
  }

  function exportPdf() {
    // Verificăm dacă jsPDF este disponibil (CDN încărcat corect)
    if (!window.jspdf || !window.jspdf.jsPDF) {
      alert("Exportul PDF nu este disponibil în acest moment (librăria jsPDF nu a putut fi încărcată). Poți folosi exportul PNG/JPEG și apoi salva ca PDF din vizualizatorul de imagini sau din funcția de imprimare a browserului.");
      return;
    }

    cy.fit(undefined, 40);
    setTimeout(function() {
      const png = cy.png({ full: true, scale: 2, bg: "#ffffff" });
      if (!png || typeof png !== "string" || png.length < 100) {
        alert("Nu am putut genera PDF-ul (imaginea de bază este goală). Încearcă să apeși întâi Fit / Expand all și apoi din nou Export PDF.");
        return;
      }
      const { jsPDF } = window.jspdf;
      const pdf = new jsPDF({ orientation: 'portrait', unit: 'pt', format: 'a4' });
      const pageW = pdf.internal.pageSize.getWidth();
      const pageH = pdf.internal.pageSize.getHeight();

      const img = new Image();
      img.onload = function() {
        const iw = img.width, ih = img.height;
        const scale = Math.min(pageW / iw, pageH / ih);
        const w = iw * scale, h = ih * scale;
        const x = (pageW - w) / 2;
        const y = 20;
        pdf.addImage(png, 'PNG', x, y, w, h);
        pdf.save("organigrama.pdf");
      };
      img.src = png;
    }, 120);
  }

  document.getElementById("__BTN_PNG__").onclick = exportPng;
  document.getElementById("__BTN_JPG__").onclick = exportJpg;
  document.getElementById("__BTN_PDF__").onclick = exportPdf;
})();
</script>
"""

    html = (
        html.replace("__CY_ID__", f"{key}_cy")
            .replace("__CY_H__", str(int(height)))
            .replace("__DATA_JSON__", payload_json)
            .replace("__BORDER__", str(border))
            .replace("__TEXT_MAX__", str(text_max))
            .replace("__FONT_SIZE__", str(fs))
            .replace("__PAD__", str(pad))
            .replace("__NODE_W_MIN__", str(node_w_min))
            .replace("__NODE_H_MIN__", str(node_h_min))
            .replace("__NODE_SEP__", str(node_sep))
            .replace("__RANK_SEP__", str(rank_sep))
            .replace("__START_COLLAPSED__", "true" if bool(start_collapsed) else "false")
            .replace("__ENABLE_DRAG__", "true" if bool(enable_drag) else "false")
            .replace("__USE_SAVED__", "true" if bool(use_saved_layout) else "false")
            .replace("__LAYOUT_KEY__", layout_storage_key)
            .replace("__BTN_FIT__", f"{key}_btn_fit")
            .replace("__BTN_EXPAND__", f"{key}_btn_expand")
            .replace("__BTN_COLLAPSE__", f"{key}_btn_collapse")
            .replace("__BTN_SAVE__", f"{key}_btn_save")
            .replace("__BTN_RESET__", f"{key}_btn_reset")
            .replace("__BTN_PNG__", f"{key}_btn_png")
            .replace("__BTN_JPG__", f"{key}_btn_jpg")
            .replace("__BTN_PDF__", f"{key}_btn_pdf")
            .replace("__SAVE_DISPLAY__", "inline-block" if bool(enable_save_layout) else "none")
    )

    components.html(html, height=height + 90, scrolling=False)


def load_stat_functii_with_employees(conn) -> pd.DataFrame:
    """
    Încarcă statul de funcții + angajații aferenți, grupați pe ID_UNITATE.

    Nu mai presupunem nume fixe de coloane.
    - Pentru STAT_FUNCTII detectăm automat:
        * ID_STAT_FUNCTIE  -> prima col. cu 'ID' și 'STAT', altfel prima coloană
        * ID_UNITATE       -> col. care conține 'UNITATE' sau 'ORG' sau 'DEPART'
        * DENUMIRE_FUNCTIE -> col. care conține 'FUNCT' sau 'POST'
        * ID_ANGAJAT       -> col. care conține 'ANGAJAT' sau 'PERS' sau 'SALARIAT'
    - Pentru tabela de angajați (DATE_ANGAJATI / ANGAJATI / EMPLOYEES):
        * ID_ANGAJAT       -> col. cu 'ID' și 'ANGAJAT' sau doar 'ID'
        * NUME             -> col. care conține 'NUME'
        * PRENUME          -> col. care conține 'PRENUME'
    """

    cur = conn.cursor()

    # Helper intern pentru detectat coloane după bucăți de nume
    def find_col(cols: list[str], must_contain: list[str], fallback: str | None = None) -> str | None:
        must_contain = [s.upper() for s in must_contain]
        for c in cols:
            up = c.upper()
            if all(s in up for s in must_contain):
                return c
        return fallback

    def safe_series(df: pd.DataFrame, col: str | None) -> pd.Series:
        if col and col in df.columns:
            return df[col]
        return pd.Series([None] * len(df), index=df.index)

    # 1. Încărcăm STAT_FUNCTII (toate coloanele)
    try:
        cur.execute("SELECT * FROM STAT_FUNCTII")
        rows = cur.fetchall()
        cols_db = [desc[0].strip().upper() for desc in cur.description]
    except Exception as e:
        st.error(f"Eroare la citirea STAT_FUNCTII: {e}")
        return pd.DataFrame(columns=[
            "ID_STAT_FUNCTIE", "ID_UNITATE", "DENUMIRE_FUNCTIE",
            "ID_ANGAJAT", "NUME", "PRENUME"
        ])

    if not rows:
        return pd.DataFrame(columns=[
            "ID_STAT_FUNCTIE", "ID_UNITATE", "DENUMIRE_FUNCTIE",
            "ID_ANGAJAT", "NUME", "PRENUME"
        ])

    df_raw = pd.DataFrame(rows, columns=cols_db)

    # Detectăm coloanele din STAT_FUNCTII
    stat_id_col = find_col(cols_db, ["ID", "STAT"])
    if not stat_id_col and cols_db:
        stat_id_col = cols_db[0]  # cădem pe prima coloană

    unit_col = find_col(cols_db, ["UNITATE"])
    if not unit_col:
        unit_col = find_col(cols_db, ["ORG"]) or find_col(cols_db, ["DEPART"])

    func_col = find_col(cols_db, ["FUNCT"])
    if not func_col:
        func_col = find_col(cols_db, ["POST"])

    ang_col = find_col(cols_db, ["ANGAJAT"])
    if not ang_col:
        ang_col = find_col(cols_db, ["PERS"])
    if not ang_col:
        ang_col = find_col(cols_db, ["SALARIAT"])

    df_stat = pd.DataFrame()
    df_stat["ID_STAT_FUNCTIE"] = safe_series(df_raw, stat_id_col)
    df_stat["ID_UNITATE"] = safe_series(df_raw, unit_col)
    df_stat["DENUMIRE_FUNCTIE"] = safe_series(df_raw, func_col)
    df_stat["ID_ANGAJAT"] = safe_series(df_raw, ang_col)

    # 2. Încercăm să găsim o tabelă de angajați
    employee_tables = ["DATE_ANGAJATI", "ANGAJATI", "EMPLOYEES"]
    df_emp = None
    found_table = None

    for tname in employee_tables:
        try:
            cur.execute(f"SELECT * FROM {tname}")
            rows_emp = cur.fetchall()
            cols_emp_db = [desc[0].strip().upper() for desc in cur.description]
            if not rows_emp:
                continue

            df_emp_raw = pd.DataFrame(rows_emp, columns=cols_emp_db)

            # detectăm coloanele pentru angajați
            emp_id_col = find_col(cols_emp_db, ["ID", "ANGAJAT"])
            if not emp_id_col:
                emp_id_col = find_col(cols_emp_db, ["ID"])

            nume_col = find_col(cols_emp_db, ["NUME"])
            prenume_col = find_col(cols_emp_db, ["PRENUME"])

            df_emp = pd.DataFrame()
            df_emp["ID_ANGAJAT"] = safe_series(df_emp_raw, emp_id_col)
            df_emp["NUME"] = safe_series(df_emp_raw, nume_col)
            df_emp["PRENUME"] = safe_series(df_emp_raw, prenume_col)

            found_table = tname
            break
        except Exception:
            continue

    # 3. Dacă nu am găsit tabelă de angajați -> întoarcem doar statul, fără nume/prenume
    if df_emp is None:
        st.warning(
            "Nu am găsit nicio tabelă de angajați dintre: DATE_ANGAJATI, ANGAJATI, EMPLOYEES. "
            "Afișez doar posturile (fără nume/prenume)."
        )
        df_stat["NUME"] = None
        df_stat["PRENUME"] = None
        return df_stat

    # 🔴 AICI era problema: tipuri diferite pe ID_ANGAJAT
    # 🔁 REZOLVARE: convertim la STRING în ambele DataFrame-uri
    df_stat["ID_ANGAJAT"] = df_stat["ID_ANGAJAT"].astype(str).str.strip()
    df_emp["ID_ANGAJAT"] = df_emp["ID_ANGAJAT"].astype(str).str.strip()

    # 4. Avem și STAT_FUNCTII și tabela de angajați -> facem merge în pandas
    try:
        df_merged = pd.merge(
            df_stat,
            df_emp,
            on="ID_ANGAJAT",
            how="left"
        )
        st.info(f"Am legat statul de funcții de tabela de angajați: {found_table}.")
        return df_merged
    except Exception as e:
        st.error(f"Eroare la combinarea statului de funcții cu angajații ({found_table}): {e}")
        df_stat["NUME"] = None
        df_stat["PRENUME"] = None
        return df_stat


    # 4. Avem și STAT_FUNCTII și tabela de angajați -> facem merge în pandas
        try:
            df_merged = pd.merge(
                df_stat,
                df_emp,
                on="ID_ANGAJAT",
                how="left"
            )
            st.info(f"Am legat statul de funcții de tabela de angajați: {found_table}.")
            return df_merged
        except Exception as e:
            st.error(f"Eroare la combinarea statului de funcții cu angajații ({found_table}): {e}")
            df_stat["NUME"] = None
            df_stat["PRENUME"] = None
            return df_stat


        # 4. Dacă am găsit o tabelă de angajați, refacem interogarea cu JOIN ca să aducem numele
        try:
            cur.execute(f"""
                SELECT
                    sf.ID_STAT_FUNCTIE,
                    sf.ID_UNITATE,
                    sf.DENUMIRE_FUNCTIE,
                    sf.ID_ANGAJAT,
                    a.NUME,
                    a.PRENUME
                FROM STAT_FUNCTII sf
                LEFT JOIN {found_table} a ON a.ID_ANGAJAT = sf.ID_ANGAJAT
            """)
            rows = cur.fetchall()
            cols = [desc[0].strip().upper() for desc in cur.description]
            df = pd.DataFrame(rows, columns=cols)

            # normalizăm coloanele, să fim siguri că există toate
            for col in ["ID_STAT_FUNCTIE", "ID_UNITATE", "DENUMIRE_FUNCTIE", "ID_ANGAJAT", "NUME", "PRENUME"]:
                if col not in df.columns:
                    df[col] = None

            st.info(f"Am legat statul de funcții de tabela de angajați: {found_table}.")
            return df

        except Exception as e:
            st.error(f"Eroare la citirea statului de funcții cu angajați din {found_table}: {e}")
            # măcar întoarcem statul simplu
            if "NUME" not in df_stat.columns:
                df_stat["NUME"] = None
            if "PRENUME" not in df_stat.columns:
                df_stat["PRENUME"] = None
            return df_stat

buf = BytesIO()  # fallback global for optional exports

nume_fisier = "export.xlsx"  # fallback global pentru nume fisier



def sync_stat_functii_from_employees(conn: sqlite3.Connection, *, overwrite: bool = True) -> int:
    """Sincronizează STATUL DE FUNCȚII din tabela employees (pentru cazul în care datele sunt deja în angajați).

    - Leagă pe cheia: Marca (dacă există), altfel CNP.
    - Creează tabela stat_functii dacă nu există.
    - Dacă tabela există, face update/insert; păstrează coloanele extra existente.
    - Dacă overwrite=True, suprascrie valorile existente în stat_functii pentru coloanele standard.

    Coloane standard populate (dacă există în employees):
      - Marca
      - Nume
      - Prenume
      - CNP
      - Functie
      - Departament
      - Loc de munca
    """
    cur = conn.cursor()
    cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='employees'")
    if cur.fetchone() is None:
        return 0

    # Citim employees (doar coloanele care există)
    emp_df = pd.read_sql_query("SELECT * FROM employees", conn)
    if emp_df is None or emp_df.empty:
        return 0

    # Identificăm coloane posibile
    emp_cols = {str(c).strip().lower(): c for c in emp_df.columns}

    col_marca = emp_cols.get("marca")
    col_cnp = emp_cols.get("cnp")
    col_ln = emp_cols.get("last_name") or emp_cols.get("nume")
    col_fn = emp_cols.get("first_name") or emp_cols.get("prenume")
    col_functie = emp_cols.get("functie") or emp_cols.get("functie_contract")
    col_depart = emp_cols.get("departament")
    col_loc = emp_cols.get("loc_munca") or emp_cols.get("loc de munca") or emp_cols.get("loc_munca".lower())

    # Cheie preferată
    key_col = col_marca or col_cnp
    if not key_col:
        return 0

    # Construim dataframe standard pentru stat_functii
    out = pd.DataFrame()
    out["Marca"] = emp_df[col_marca].astype(str) if col_marca else ""
    out["CNP"] = emp_df[col_cnp].astype(str) if col_cnp else ""
    out["Nume"] = emp_df[col_ln].astype(str) if col_ln else ""
    out["Prenume"] = emp_df[col_fn].astype(str) if col_fn else ""
    out["Functie"] = emp_df[col_functie].astype(str) if col_functie else ""
    out["Departament"] = emp_df[col_depart].astype(str) if col_depart else ""
    out["Loc de munca"] = emp_df[col_loc].astype(str) if col_loc else ""

    # Normalizări simple
    for c in ["Marca", "CNP"]:
        out[c] = out[c].fillna("").astype(str).str.strip()

    # Păstrăm doar rânduri cu cheie
    if key_col == col_marca:
        out = out[out["Marca"] != ""]
        key_name = "Marca"
    else:
        out = out[out["CNP"] != ""]
        key_name = "CNP"

    if out.empty:
        return 0

    # Verificăm dacă există stat_functii
    cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='stat_functii'")
    stat_exists = cur.fetchone() is not None

    if not stat_exists:
        out.to_sql("stat_functii", conn, if_exists="replace", index=False)
        return int(len(out))

    stat_df = pd.read_sql_query("SELECT * FROM stat_functii", conn)
    if stat_df is None:
        stat_df = pd.DataFrame()

    # Asigurăm cheia în stat
    if key_name not in stat_df.columns:
        # încercăm să mapăm dacă există altă denumire
        # dacă nu, adăugăm coloană
        stat_df[key_name] = ""

    # Facem index după cheie
    stat_df[key_name] = stat_df[key_name].fillna("").astype(str).str.strip()

    # Mergem pe update/insert: păstrăm coloanele extra
    std_cols = ["Marca", "CNP", "Nume", "Prenume", "Functie", "Departament", "Loc de munca"]
    for c in std_cols:
        if c not in stat_df.columns:
            stat_df[c] = ""

    stat_df_idx = stat_df.set_index(key_name, drop=False)
    out_idx = out.set_index(key_name, drop=False)

    updated = 0
    for k, r in out_idx.iterrows():
        if str(k).strip() == "":
            continue
        if k in stat_df_idx.index:
            # update
            for c in std_cols:
                if overwrite:
                    stat_df_idx.at[k, c] = r.get(c, "")
                else:
                    # completează doar dacă e gol
                    curv = str(stat_df_idx.at[k, c]) if k in stat_df_idx.index else ""
                    if (curv is None) or (str(curv).strip() == ""):
                        stat_df_idx.at[k, c] = r.get(c, "")
            updated += 1
        else:
            # insert new row, păstrăm și coloanele extra (cu valori goale)
            new_row = {col: "" for col in stat_df_idx.columns}
            for c in std_cols:
                new_row[c] = r.get(c, "")
            new_row[key_name] = r.get(key_name, "")
            stat_df_idx.loc[k] = new_row
            updated += 1

    # Salvăm înapoi
    stat_df_out = stat_df_idx.reset_index(drop=True)

    # Ordonare frumoasă: cheia + standard + rest
    rest_cols = [c for c in stat_df_out.columns if c not in std_cols]
    ordered = []
    if key_name in std_cols:
        # deja în std_cols
        ordered = std_cols + [c for c in rest_cols if c not in std_cols]
    else:
        ordered = [key_name] + [c for c in std_cols if c != key_name] + [c for c in rest_cols if c not in std_cols and c != key_name]
    ordered = [c for c in ordered if c in stat_df_out.columns]
    stat_df_out = stat_df_out[ordered]

    stat_df_out.to_sql("stat_functii", conn, if_exists="replace", index=False)
    return int(updated)

def render_org_tree_with_people(df_units: pd.DataFrame, df_stat: pd.DataFrame,
                                parent_id: int | None = None, level: int = 0):
    """
    Desenează organigrama ca arbore și, sub fiecare unitate, afișează posturile + oamenii
    conform statului de funcții.
    """
    if "ID_UNITATE" not in df_units.columns:
        st.error("În df_units nu există coloana 'ID_UNITATE'. Verifică schema.")
        return

    # stabilim coloana de nume unitate
    name_col = org_get_name_column(df_units)
    type_col = org_get_type_column(df_units)

    # pregătim PARENT_ID cu NaN -> -1 ca să putem filtra ușor rădăcina
    dfu = df_units.copy()
    if "PARENT_ID" in dfu.columns:
        dfu["PARENT_ID_FILL"] = dfu["PARENT_ID"].fillna(-1)
    else:
        # dacă nu ai PARENT_ID, default toate sunt de rădăcină
        dfu["PARENT_ID_FILL"] = -1

    current_parent = -1 if parent_id is None else parent_id

    children = dfu[dfu["PARENT_ID_FILL"] == current_parent] \
        .sort_values(by=name_col, ascending=True)

    for _, row in children.iterrows():
        id_unitate = int(row["ID_UNITATE"])
        nume_unitate = str(row.get(name_col, "")).strip()

        tip_text = ""
        if type_col and type_col in row and pd.notna(row[type_col]):
            tip_text = f" ({row[type_col]})"

        indent = "&nbsp;" * 4 * level
        icon = "📁" if level > 0 else "🏢"

        st.markdown(
            f"{indent}{icon} **{nume_unitate}**{tip_text}  "
        )

        # --- oamenii / posturile din statul de funcții pentru această unitate ---
        if not df_stat.empty:
            df_unit_stat = df_stat[df_stat["ID_UNITATE"] == id_unitate]

            if not df_unit_stat.empty:
                for _, rstat in df_unit_stat.iterrows():
                    functie = str(rstat.get("DENUMIRE_FUNCTIE", "")).strip()
                    nume = str(rstat.get("NUME", "") or "").strip()
                    prenume = str(rstat.get("PRENUME", "") or "").strip()

                    indent_pers = "&nbsp;" * 4 * (level + 1)
                    if nume or prenume:
                        st.markdown(
                            f"{indent_pers}• 👤 {nume} {prenume} – *{functie}*"
                        )
                    else:
                        # post prevăzut dar neocupat / fără angajat
                        st.markdown(
                            f"{indent_pers}• ⭕ Post neocupat – *{functie}*"
                        )

        # recursiv: copiii acestei unități
        render_org_tree_with_people(df_units, df_stat, parent_id=id_unitate, level=level + 1)
        
def render_org_tree_expandable(df_units: pd.DataFrame, df_stat: pd.DataFrame,
                               parent_id: int | None = None, level: int = 0):
    """
    Afișează organigrama ca listă ierarhică, dar cu posibilitatea
    de a deschide/închide fiecare unitate folosind st.expander.

    - df_units: tabelul cu unitățile (are ID_UNITATE și PARENT_ID)
    - df_stat: statul de funcții + angajați (ID_UNITATE, DENUMIRE_FUNCTIE, NUME, PRENUME)
    """

    if "ID_UNITATE" not in df_units.columns:
        st.error("În df_units nu există coloana 'ID_UNITATE'. Verifică schema.")
        return

    name_col = org_get_name_column(df_units)
    type_col = org_get_type_column(df_units)

    dfu = df_units.copy()
    if "PARENT_ID" in dfu.columns:
        dfu["PARENT_ID_FILL"] = dfu["PARENT_ID"].fillna(-1)
    else:
        dfu["PARENT_ID_FILL"] = -1

    current_parent = -1 if parent_id is None else parent_id

    children = dfu[dfu["PARENT_ID_FILL"] == current_parent] \
        .sort_values(by=name_col, ascending=True)

    for _, row in children.iterrows():
        id_unitate = int(row["ID_UNITATE"])
        nume_unitate = str(row.get(name_col, "")).strip()

        tip_text = ""
        if type_col and type_col in row and pd.notna(row[type_col]):
            tip_text = f" ({row[type_col]})"

        # un mic prefix pentru nivel (ca să vezi vizual adâncimea)
        prefix = "— " * level
        label = f"{prefix}{nume_unitate}{tip_text}"

        with st.expander(label, expanded=False):
            # 1) Posturi + oameni pentru unitatea curentă
            if df_stat is not None and not df_stat.empty:
                df_unit_stat = df_stat[df_stat["ID_UNITATE"] == id_unitate]

                if not df_unit_stat.empty:
                    st.markdown("**Posturi și ocupanți (stat de funcții):**")
                    for _, rstat in df_unit_stat.iterrows():
                        functie = str(rstat.get("DENUMIRE_FUNCTIE", "")).strip()
                        nume = str(rstat.get("NUME", "") or "").strip()
                        prenume = str(rstat.get("PRENUME", "") or "").strip()

                        if nume or prenume:
                            st.markdown(f"- 👤 **{nume} {prenume}** – _{functie}_")
                        else:
                            st.markdown(f"- ⭕ Post neocupat – _{functie}_")
                else:
                    st.caption("Nu există posturi configurate pentru această unitate în statul de funcții.")
            else:
                st.caption("Statul de funcții nu este disponibil sau nu conține date.")

            # 2) Copiii acestei unități (subunități) – recursiv
            render_org_tree_expandable(df_units, df_stat, parent_id=id_unitate, level=level + 1)


def render_org_tree_compact(df_units: pd.DataFrame, df_stat: pd.DataFrame,
                            parent_id: int | None = None, level: int = 0):
    """
    Afișează organigrama ca listă ierarhică text (fără expandere),
    cu posturi + oameni afișați direct sub fiecare unitate.
    """

    if "ID_UNITATE" not in df_units.columns:
        st.error("În df_units nu există coloana 'ID_UNITATE'. Verifică schema.")
        return

    name_col = org_get_name_column(df_units)
    type_col = org_get_type_column(df_units)

    dfu = df_units.copy()
    if "PARENT_ID" in dfu.columns:
        dfu["PARENT_ID_FILL"] = dfu["PARENT_ID"].fillna(-1)
    else:
        dfu["PARENT_ID_FILL"] = -1

    current_parent = -1 if parent_id is None else parent_id

    children = dfu[dfu["PARENT_ID_FILL"] == current_parent] \
        .sort_values(by=name_col, ascending=True)

    for _, row in children.iterrows():
        id_unitate = int(row["ID_UNITATE"])
        nume_unitate = str(row.get(name_col, "")).strip()

        tip_text = ""
        if type_col and type_col in row and pd.notna(row[type_col]):
            tip_text = f" ({row[type_col]})"

        indent = "&nbsp;" * 4 * level
        st.markdown(f"{indent}• **{nume_unitate}**{tip_text}", unsafe_allow_html=True)

        # Posturi + oameni direct sub unitate
        if df_stat is not None and not df_stat.empty:
            df_unit_stat = df_stat[df_stat["ID_UNITATE"] == id_unitate]

            for _, rstat in df_unit_stat.iterrows():
                functie = str(rstat.get("DENUMIRE_FUNCTIE", "")).strip()
                nume = str(rstat.get("NUME", "") or "").strip()
                prenume = str(rstat.get("PRENUME", "") or "").strip()

                indent_pers = "&nbsp;" * 4 * (level + 1)
                if nume or prenume:
                    st.markdown(
                        f"{indent_pers}- 👤 {nume} {prenume} – _{functie}_",
                        unsafe_allow_html=True,
                    )
                elif functie:
                    st.markdown(
                        f"{indent_pers}- ⭕ Post neocupat – _{functie}_",
                        unsafe_allow_html=True,
                    )

        # Subunități (recursiv)
        render_org_tree_compact(df_units, df_stat, parent_id=id_unitate, level=level + 1)


# -------------------------------------------------------------
# CONFIG APLICAȚIE
# -------------------------------------------------------------
APP_TITLE = ""
#DB_PATH_DEFAULT = str(ANCP_DB)
CONFIG_FILE = "socrates_config.json"

# -------------------------------------------------------------
# UTILITARE CONFIG (pentru calea DB)
# -------------------------------------------------------------
def load_config() -> Dict[str, Any]:
    try:
        with open(CONFIG_FILE, "r", encoding="utf-8") as f:
            cfg = json.load(f)
    except FileNotFoundError:
        cfg = {"db_path": DB_PATH_DEFAULT}
        save_config(cfg)
    return cfg


def save_config(cfg: Dict[str, Any]) -> None:
    with open(CONFIG_FILE, "w", encoding="utf-8") as f:
        json.dump(cfg, f, ensure_ascii=False, indent=2)


def get_db_path() -> str:
    # În cloud/VPS e mai simplu să controlăm DB path din ENV,
    # ca să nu depindem de fișierul local socrates_config.json.
    env_path = (os.getenv("SOCRATES_DB_PATH") or "").strip()
    if env_path:
        return env_path

    root = Path(__file__).resolve().parent.parent
    return str(root / "data" / "ANCPI.db")

# -------------------------------------------------------------
# ASIGURARE SCHEMA DB
# -------------------------------------------------------------
def ensure_schema(conn: sqlite3.Connection) -> None:
    cur = conn.cursor()

    # Tabela employees
    cur.execute("""
        CREATE TABLE IF NOT EXISTS employees (
            id INTEGER PRIMARY KEY AUTOINCREMENT
        );
    """)

    # Tabela COR (coduri ocupații)
    cur.execute("""
        CREATE TABLE IF NOT EXISTS cor_coduri (
            cod TEXT PRIMARY KEY,
            denumire TEXT NOT NULL,
            activ INTEGER NOT NULL DEFAULT 1
        );
    """)

    # Verificăm coloane existente
    cur.execute("PRAGMA table_info(employees);")
    existing_cols = {row[1].lower() for row in cur.fetchall()}

    # Coloană cod_cor / cod_153 / den_153
    if "cod_cor" not in existing_cols:
        cur.execute("ALTER TABLE employees ADD COLUMN cod_cor TEXT;")
        existing_cols.add("cod_cor")
    if "cod_153" not in existing_cols:
        cur.execute("ALTER TABLE employees ADD COLUMN cod_153 TEXT;")
        existing_cols.add("cod_153")
    if "den_153" not in existing_cols:
        cur.execute("ALTER TABLE employees ADD COLUMN den_153 TEXT;")
        existing_cols.add("den_153")

    # Tabela pentru poze angajați (un rând per employee_id)
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_photos (
            employee_id INTEGER PRIMARY KEY,
            mime TEXT,
            img BLOB,
            updated_at TEXT
        );
        """
    )

    # ---------------------------------------------------------
    # Documente generate (CIM, Acte adiționale, etc.) – ISTORIC
    # ---------------------------------------------------------
    # Tabela pentru ISTORIC CONTRACTE (CIM) per angajat
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_contracts (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER NOT NULL,
            cim_nr TEXT,
            cim_data TEXT,
            title TEXT,
            payload_json TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT,
            is_active INTEGER DEFAULT 0
        );
        """
    )
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_contracts_emp ON employee_contracts(employee_id)")

    # Tabel pentru documente generate per angajat (istoric + document activ)
    # IMPORTANT: acest tabel este diferit de atașamentele PDF (care sunt în employee_attachments).
    # Migrare: dacă există un vechi employee_documents cu schemă de atașamente, îl redenumim și îl migrăm.
    try:
        cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='employee_documents'")
        has_emp_docs = cur.fetchone() is not None
        if has_emp_docs:
            cur.execute("PRAGMA table_info(employee_documents)")
            _ed_cols = {r[1].lower() for r in (cur.fetchall() or [])}
            looks_like_old_attachments = (
                "file_blob" in _ed_cols and "uploaded_at" in _ed_cols and "doc_type" in _ed_cols
                and "data" not in _ed_cols and "filename" not in _ed_cols
            )
            if looks_like_old_attachments:
                # redenumim vechiul tabel ca să putem crea schema nouă
                cur.execute("ALTER TABLE employee_documents RENAME TO employee_documents_legacy_attachments")
    except Exception:
        pass

    # Schema corectă pentru employee_documents (documente generate)
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER NOT NULL,
            doc_type TEXT NOT NULL,
            doc_no TEXT,
            doc_date TEXT,
            filename TEXT,
            data BLOB NOT NULL,
            meta_json TEXT,
            created_at TEXT,
            is_active INTEGER DEFAULT 1
        );
        """
    )
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_documents_emp ON employee_documents(employee_id);")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_documents_type ON employee_documents(doc_type);")

    # ---------------------------------------------------------
    # Atașamente (PDF/scan) – CI, Studii, Medical etc.
    # ---------------------------------------------------------
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_attachments (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER NOT NULL,
            doc_type TEXT NOT NULL,
            display_name TEXT,
            file_name TEXT,
            mime_type TEXT,
            file_blob BLOB,
            uploaded_at TEXT NOT NULL,
            is_active INTEGER DEFAULT 1
        );
        """
    )
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_attachments_emp ON employee_attachments(employee_id)")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_attachments_type ON employee_attachments(employee_id, doc_type)")

    # ---------------------------------------------------------
    # L153 (anexe salarizare) – istoric per angajat
    # ---------------------------------------------------------
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_l153 (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER NOT NULL,
            anexa TEXT,
            tabel TEXT,
            functie TEXT,
            studii TEXT,
            grad TEXT,
            treapta TEXT,
            salariu REAL,
            coeficient REAL,
            data_aplicare TEXT,
            is_active INTEGER DEFAULT 1,
            created_at TEXT
        );
        """
    )
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_l153_emp ON employee_l153(employee_id)")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_l153_active ON employee_l153(employee_id, is_active)")


    # Migrare: dacă există vechiul tabel redenumit (employee_documents_legacy_attachments), copiem atașamentele.
    try:
        cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='employee_documents_legacy_attachments'")
        if cur.fetchone():
            cur.execute(
                """
                INSERT INTO employee_attachments(employee_id, doc_type, display_name, file_name, mime_type, file_blob, uploaded_at, is_active)
                SELECT employee_id, doc_type, display_name, file_name, mime_type, file_blob, uploaded_at, 1
                FROM employee_documents_legacy_attachments
                WHERE employee_id IS NOT NULL AND doc_type IS NOT NULL
                """
            )
    except Exception:
        pass

    conn.commit()

    needed_cols = {
        # bază
        "marca": "TEXT",
        "last_name": "TEXT",
        "first_name": "TEXT",
        "cnp": "TEXT",
        "cod_cor": "TEXT",
        # adresă
        "strada": "TEXT",
        "numar": "TEXT",
        "bloc": "TEXT",
        "scara": "TEXT",
        "apartament": "TEXT",
        "cod_postal": "TEXT",
        "localitate": "TEXT",
        "judet": "TEXT",
        # contact
        "telefon_fix": "TEXT",
        "mobil": "TEXT",
        "email": "TEXT",
        # CI
        "ci_tip_act": "TEXT",
        "ci_serie": "TEXT",
        "ci_numar": "TEXT",
        "ci_eliberat_de": "TEXT",
        "ci_data_eliberarii": "TEXT",
        # familie
        "stare_civila": "TEXT",
        "nr_copii": "INTEGER",
        # contract / muncă
        "functie": "TEXT",
        "departament": "TEXT",
        "data_angajare": "TEXT",
        "tip_contract": "TEXT",
        "loc_munca": "TEXT",
        "departament_organizatoric": "TEXT",
        "functie_contract": "TEXT",
        "tip_norma": "TEXT",
        "program_munca": "TEXT",
        "salariu_baza": "REAL",
        # studii & diverse
        "studii": "TEXT",
        "profesie": "TEXT",
        "calificare": "TEXT",
        "observatii": "TEXT",
        # copertă DOSAR PROFESIONAL
        "dosar_nr": "TEXT",
        "dosar_functionar_public": "TEXT",
        "dosar_data_intocmire": "TEXT",
        "dosar_autoritate": "TEXT",
        "dosar_intocmit_nume": "TEXT",
        "dosar_intocmit_functie": "TEXT",
        "dosar_intocmit_semnatura": "TEXT",
        "dosar_modificari_nume": "TEXT",
        "dosar_modificari_functie": "TEXT",
        "dosar_modificari_semnatura": "TEXT",
        "dosar_certificare_nume": "TEXT",
        "dosar_certificare_functie": "TEXT",
        "dosar_certificare_semnatura": "TEXT",
        # secțiuni suplimentare
        "activitate_in_afara_functiei": "TEXT",
        "activitate_in_cadru_institutie": "TEXT",
        "situatia_drepturi_salariale": "TEXT",
        "situatia_concedii": "TEXT",
        "situatia_disciplinara": "TEXT",
        "registru_numar": "TEXT",
        "registru_data": "TEXT",
        "registru_observatii": "TEXT",
        # stare angajat
        "activ": "INTEGER DEFAULT 1",
        "vechime_inst_ani": "INTEGER",
        "vechime_inst_luni": "INTEGER",
        "vechime_inst_fract_num": "INTEGER",
        "vechime_inst_fract_den": "INTEGER",
        "data_angajarii": "TEXT",
        "data_plecarii": "TEXT",
        "vechime_anterioara_ani": "INTEGER",
        "vechime_anterioara_luni": "INTEGER",
        "vechime_munca_ani": "INTEGER",
        "vechime_munca_luni": "INTEGER",
        "vechime_functie_ani": "INTEGER",
        "vechime_functie_luni": "INTEGER",
    
        # contract de muncă (CIM)
        "cim_numar": "TEXT",
        "cim_data": "TEXT",
        "cim_template_name": "TEXT",
        "cim_loc_munca": "TEXT",
        "cim_program": "TEXT",
        "cim_tip": "TEXT",
        "cim_durata": "TEXT",
        "cim_perioada_proba": "TEXT",
        "cim_clauze": "TEXT",
            "cim_suspendat": "INTEGER DEFAULT 0",
        "cim_suspend_start": "TEXT",
        "cim_suspend_end": "TEXT",
        "cim_suspend_motiv": "TEXT",
        "functie_cod_contract": "TEXT",
        "pensionar": "INTEGER",
        "reges_cod_grad_invaliditate": "TEXT",
        "reges_nume_grad_invaliditate": "TEXT",
        "reges_cod_grad_handicap": "TEXT",
        "reges_nume_grad_handicap": "TEXT",
        "reges_cod_tip_handicap": "TEXT",
        "reges_nume_tip_handicap": "TEXT",
        "reges_nr_cert_handicap": "TEXT",
        "reges_data_cert_handicap": "TEXT",
        "reges_termen_valabilitate": "TEXT",
    }

    for col, coltype in needed_cols.items():
        if col.lower() not in existing_cols:
            cur.execute(f"ALTER TABLE employees ADD COLUMN {col} {coltype};")

    # Tabela pentru ACCESUL LA DOSARUL PROFESIONAL
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS dosar_acces (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER NOT NULL,
            nr_crt INTEGER NOT NULL,
            nume TEXT,
            functie TEXT,
            semnatura TEXT,
            motivul TEXT,
            acces_autorizat_de TEXT,
            luat_la_cunostinta TEXT
        );
        """
    )


    # Tabela pentru PERSOANE ÎN ÎNTREȚINERE (copii / părinți)
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_dependents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER NOT NULL,
            tip TEXT NOT NULL,            -- COPIL / PARINTE
            nume TEXT,
            prenume TEXT,
            cnp TEXT,
            data_nasterii TEXT,
            grad_rudenie TEXT,
            observatii TEXT,
            activ INTEGER DEFAULT 1
        );
        """
    )

    
    # Tabel pentru șabloane globale (DOCX)
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS global_templates (
            name TEXT PRIMARY KEY,
            filename TEXT,
            data BLOB,
            updated_at TEXT
        );
        """
    )


# -------------------------------------------------------------
# COR (Cod ocupații) – helperi tabel separat `cor_coduri`
# -------------------------------------------------------------
def _cor_ensure_table(conn: sqlite3.Connection) -> None:
    """Asigură existența tabelei cor_coduri."""
    try:
        cur = conn.cursor()
        cur.execute(
            """
            CREATE TABLE IF NOT EXISTS cor_coduri (
                cod TEXT PRIMARY KEY,
                denumire TEXT NOT NULL,
                activ INTEGER NOT NULL DEFAULT 1
            );
            """
        )
        conn.commit()
    except Exception:
        # nu blocăm aplicația dacă DB e read-only / altă eroare
        pass


def _cor_norm_code(cod: str) -> str:
    s = str(cod or "").strip()
    # păstrăm doar cifre
    s = re.sub(r"\D+", "", s)
    return s


def _cor_upsert(conn: sqlite3.Connection, cod: str, denumire: str) -> None:
    """Inserează/actualizează un cod COR (activează automat)."""
    _cor_ensure_table(conn)
    c = _cor_norm_code(cod)
    d = str(denumire or "").strip()
    if not c or not d:
        return
    cur = conn.cursor()
    # SQLite: INSERT OR REPLACE păstrează PK, dar 'replace' șterge+inserează.
    # Preferăm UPSERT (disponibil pe SQLite 3.24+). Dacă nu merge, fallback la REPLACE.
    try:
        cur.execute(
            """
            INSERT INTO cor_coduri(cod, denumire, activ)
            VALUES(?, ?, 1)
            ON CONFLICT(cod) DO UPDATE SET
                denumire=excluded.denumire,
                activ=1
            """,
            (c, d),
        )
    except Exception:
        cur.execute(
            """INSERT OR REPLACE INTO cor_coduri(cod, denumire, activ) VALUES(?, ?, 1)""",
            (c, d),
        )
    conn.commit()


def _cor_soft_delete(conn: sqlite3.Connection, cod: str) -> None:
    """Dezactivează (soft delete) un cod COR."""
    _cor_ensure_table(conn)
    c = _cor_norm_code(cod)
    if not c:
        return
    cur = conn.cursor()
    cur.execute("UPDATE cor_coduri SET activ=0 WHERE cod=?", (c,))
    conn.commit()


def _cor_get(conn: sqlite3.Connection, cod: str) -> str:
    """Returnează denumirea pentru un cod COR (doar activ=1)."""
    _cor_ensure_table(conn)
    c = _cor_norm_code(cod)
    if not c:
        return ""
    cur = conn.cursor()
    cur.execute("SELECT denumire FROM cor_coduri WHERE cod=? AND activ=1", (c,))
    row = cur.fetchone()
    return str(row[0]) if row and row[0] is not None else ""


def _cor_list(conn: sqlite3.Connection, q: str = "", limit: int = 500) -> list[tuple[str, str]]:
    """Listează coduri COR (active), filtrabile după cod/denumire."""
    _cor_ensure_table(conn)
    qn = str(q or "").strip().lower()
    cur = conn.cursor()
    if qn:
        cur.execute(
            """
            SELECT cod, denumire
            FROM cor_coduri
            WHERE activ=1 AND (LOWER(cod) LIKE ? OR LOWER(denumire) LIKE ?)
            ORDER BY cod
            LIMIT ?
            """,
            (f"%{qn}%", f"%{qn}%", int(limit) if limit else 500),
        )
    else:
        cur.execute(
            """
            SELECT cod, denumire
            FROM cor_coduri
            WHERE activ=1
            ORDER BY cod
            LIMIT ?
            """,
            (int(limit) if limit else 500,),
        )
    return [(str(r[0]), str(r[1])) for r in (cur.fetchall() or [])]

# -------------------------------------------------------------
# Legea 153 – nomenclator cod funcție (tabel separat `lege153_coduri`)
# -------------------------------------------------------------
def _l153_ensure_table(conn: sqlite3.Connection) -> None:
    """Asigură existența tabelei lege153_coduri."""
    try:
        cur = conn.cursor()
        cur.execute(
            """
            CREATE TABLE IF NOT EXISTS lege153_coduri (
                cod TEXT PRIMARY KEY,
                denumire TEXT NOT NULL,
                activ INTEGER NOT NULL DEFAULT 1
            );
            """
        )
        conn.commit()
    except Exception:
        pass


def _l153_norm_code(cod: str) -> str:
    return str(cod or "").strip()


def _l153_upsert(conn: sqlite3.Connection, cod: str, denumire: str) -> None:
    """Inserează/actualizează un cod (activează automat)."""
    _l153_ensure_table(conn)
    c = _l153_norm_code(cod)
    d = str(denumire or "").strip()
    if not c or not d:
        return
    cur = conn.cursor()
    try:
        cur.execute(
            """
            INSERT INTO lege153_coduri(cod, denumire, activ)
            VALUES(?, ?, 1)
            ON CONFLICT(cod) DO UPDATE SET
                denumire=excluded.denumire,
                activ=1
            """,
            (c, d),
        )
    except Exception:
        cur.execute(
            """INSERT OR REPLACE INTO lege153_coduri(cod, denumire, activ) VALUES(?, ?, 1)""",
            (c, d),
        )
    conn.commit()


def _l153_soft_delete(conn: sqlite3.Connection, cod: str) -> None:
    _l153_ensure_table(conn)
    c = _l153_norm_code(cod)
    if not c:
        return
    cur = conn.cursor()
    cur.execute("UPDATE lege153_coduri SET activ=0 WHERE cod=?", (c,))
    conn.commit()


def _l153_get(conn: sqlite3.Connection, cod: str) -> str:
    _l153_ensure_table(conn)
    c = _l153_norm_code(cod)
    if not c:
        return ""
    cur = conn.cursor()
    cur.execute("SELECT denumire FROM lege153_coduri WHERE cod=? AND activ=1", (c,))
    row = cur.fetchone()
    return str(row[0]) if row and row[0] is not None else ""


def _l153_list(conn: sqlite3.Connection, q: str = "", limit: int = 500) -> list[tuple[str, str]]:
    _l153_ensure_table(conn)
    qn = str(q or "").strip().lower()
    cur = conn.cursor()
    if qn:
        cur.execute(
            """
            SELECT cod, denumire
            FROM lege153_coduri
            WHERE activ=1 AND (LOWER(cod) LIKE ? OR LOWER(denumire) LIKE ?)
            ORDER BY cod
            LIMIT ?
            """,
            (f"%{qn}%", f"%{qn}%", int(limit)),
        )
    else:
        cur.execute(
            """
            SELECT cod, denumire
            FROM lege153_coduri
            WHERE activ=1
            ORDER BY cod
            LIMIT ?
            """,
            (int(limit),),
        )
    return [(str(r[0]), str(r[1])) for r in cur.fetchall()]

def _l153_ensure_table(conn: sqlite3.Connection) -> None:
    """Asigură existența tabelei lege153_coduri."""
    try:
        cur = conn.cursor()
        cur.execute(
            """
            CREATE TABLE IF NOT EXISTS lege153_coduri (
                cod TEXT PRIMARY KEY,
                denumire TEXT NOT NULL,
                activ INTEGER NOT NULL DEFAULT 1
            );
            """
        )
        conn.commit()
    except Exception:
        pass


def _l153_norm_code(cod: str) -> str:
    return str(cod or "").strip()


def _l153_upsert(conn: sqlite3.Connection, cod: str, denumire: str) -> None:
    """Inserează/actualizează un cod (activează automat)."""
    _l153_ensure_table(conn)
    c = _l153_norm_code(cod)
    d = str(denumire or "").strip()
    if not c or not d:
        return
    cur = conn.cursor()
    try:
        cur.execute(
            """
            INSERT INTO lege153_coduri(cod, denumire, activ)
            VALUES(?, ?, 1)
            ON CONFLICT(cod) DO UPDATE SET
                denumire=excluded.denumire,
                activ=1
            """,
            (c, d),
        )
    except Exception:
        cur.execute(
            """INSERT OR REPLACE INTO lege153_coduri(cod, denumire, activ) VALUES(?, ?, 1)""",
            (c, d),
        )
    conn.commit()


def _l153_soft_delete(conn: sqlite3.Connection, cod: str) -> None:
    _l153_ensure_table(conn)
    c = _l153_norm_code(cod)
    if not c:
        return
    cur = conn.cursor()
    cur.execute("UPDATE lege153_coduri SET activ=0 WHERE cod=?", (c,))
    conn.commit()


def _l153_get(conn: sqlite3.Connection, cod: str) -> str:
    _l153_ensure_table(conn)
    c = _l153_norm_code(cod)
    if not c:
        return ""
    cur = conn.cursor()
    cur.execute("SELECT denumire FROM lege153_coduri WHERE cod=? AND activ=1", (c,))
    row = cur.fetchone()
    return str(row[0]) if row and row[0] is not None else ""


def _l153_list(conn: sqlite3.Connection, q: str = "", limit: int = 500) -> list[tuple[str, str]]:
    _l153_ensure_table(conn)
    qn = str(q or "").strip().lower()
    cur = conn.cursor()
    if qn:
        cur.execute(
            """
            SELECT cod, denumire
            FROM lege153_coduri
            WHERE activ=1 AND (LOWER(cod) LIKE ? OR LOWER(denumire) LIKE ?)
            ORDER BY cod
            LIMIT ?
            """,
            (f"%{qn}%", f"%{qn}%", int(limit)),
        )
    else:
        cur.execute(
            """
            SELECT cod, denumire
            FROM lege153_coduri
            WHERE activ=1
            ORDER BY cod
            LIMIT ?
            """,
            (int(limit),),
        )
    return [(str(r[0]), str(r[1])) for r in cur.fetchall()]
def get_connection() -> sqlite3.Connection:
    conn = sqlite3.connect(get_db_path(), check_same_thread=False)
    conn.row_factory = sqlite3.Row
    ensure_schema(conn)
    return conn
# -------------------------------------------------------------
# AUTH (VPS/CLOUD) – utilizatori multipli cu parole hash-uíte
# -------------------------------------------------------------
def ensure_auth_tables(conn: sqlite3.Connection) -> None:
    """Creează tabela de utilizatori (dacă nu există)."""
    cur = conn.cursor()
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS app_users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT NOT NULL UNIQUE,
            password_hash BLOB NOT NULL,
            role TEXT DEFAULT 'user',
            is_active INTEGER NOT NULL DEFAULT 1,
            created_at TEXT
        );
        """
    )
    cur.execute("CREATE INDEX IF NOT EXISTS idx_app_users_active ON app_users(is_active)")
    conn.commit()


def _hash_password(plain: str) -> bytes:
    return bcrypt.hashpw(str(plain).encode("utf-8"), bcrypt.gensalt())


def _verify_password(plain: str, pw_hash: bytes) -> bool:
    try:
        return bcrypt.checkpw(str(plain).encode("utf-8"), bytes(pw_hash))
    except Exception:
        return False


def auth_bootstrap_admin(conn: sqlite3.Connection, cfg: dict) -> None:
    """Dacă nu există niciun user, creează unul din ENV sau din Configurare (compatibil)."""
    ensure_auth_tables(conn)
    cur = conn.cursor()
    cur.execute("SELECT COUNT(1) FROM app_users")
    n = (cur.fetchone() or [0])[0] or 0
    if int(n) > 0:
        return

    # 1) preferăm ENV (pentru VPS)
    env_user = (os.getenv("SOCRATES_ADMIN_USER") or "").strip()
    env_pass = (os.getenv("SOCRATES_ADMIN_PASS") or "").strip()

    # 2) fallback: credențiale vechi din Configurare (single-user)
    cfg_user = (cfg.get("app_user") or "").strip()
    cfg_pass = (cfg.get("app_pass") or "").strip()

    username = env_user or cfg_user or "admin"
    password = env_pass or cfg_pass or "admin"

    cur.execute(
        """
        INSERT OR REPLACE INTO app_users(username, password_hash, role, is_active, created_at)
        VALUES(?, ?, 'admin', 1, ?)
        """,
        (username, sqlite3.Binary(_hash_password(password)), datetime.datetime.now().isoformat(timespec="seconds")),
    )
    conn.commit()


def auth_check(conn: sqlite3.Connection, username: str, password: str) -> tuple[bool, str]:
    ensure_auth_tables(conn)
    u = (username or "").strip()
    p = password or ""
    if not u or not p:
        return False, "Completează utilizator și parolă."

    cur = conn.cursor()
    cur.execute(
        "SELECT password_hash, role, is_active FROM app_users WHERE username = ?",
        (u,),
    )
    row = cur.fetchone()
    if not row:
        return False, "Utilizator inexistent."
    if int(row[2] or 0) != 1:
        return False, "Utilizator inactiv."

    ok = _verify_password(p, row[0])
    if not ok:
        return False, "Parolă incorectă."
    return True, str(row[1] or "user")


def require_login(conn: sqlite3.Connection, cfg: dict) -> None:
    auth_bootstrap_admin(conn, cfg)

    if str(os.getenv("SOCRATES_DISABLE_AUTH") or "").strip() in ("1", "true", "TRUE", "yes", "YES"):
        return

    if "logged_in" not in st.session_state:
        st.session_state.logged_in = False
        st.session_state.user_role = "user"
        st.session_state.username = ""

    if st.session_state.logged_in:
        return

    # Pe pagina de login ascundem header-ul Streamlit (brand din bara neagră).
    st.markdown(
        """
        <style>
        header[data-testid="stHeader"] { visibility: hidden !important; height: 0 !important; }
        </style>
        """,
        unsafe_allow_html=True,
    )

    st.markdown('<div class="bg-glow"></div>', unsafe_allow_html=True)
    # spațiu vertical ușor redus pentru a aduce hero-ul mai sus
    st.markdown("<div style='height:4vh'></div>", unsafe_allow_html=True)

    left, mid, right = st.columns([1.1, 1.8, 1.1])
    with mid:
        st.markdown("<div class='login-wrap'>", unsafe_allow_html=True)

        st.markdown(
            """
            <div class="login-hero">
              <div class="brand">Socrates@HR</div>
              <div class="login-hero-desc">Platformă instituțională pentru administrarea resurselor umane</div>
            </div>
            """,
            unsafe_allow_html=True,
        )

        st.markdown("<div class='login-card'>", unsafe_allow_html=True)
        st.markdown("<div class='login-header-bar'>AUTENTIFICARE</div>", unsafe_allow_html=True)

        # Câmpuri login într-un singur bloc logic (markup stabil pentru styling dedicat)
        def _login_field(label: str, key: str, *, is_password: bool = False, placeholder: str = "") -> str:
            st.markdown("<div class='login-field'>", unsafe_allow_html=True)
            value = st.text_input(
                label,
                key=key,
                type="password" if is_password else "default",
                placeholder=placeholder,
                label_visibility="visible",
            )
            st.markdown("</div>", unsafe_allow_html=True)
            return value

        u = _login_field("Utilizator", "login_user", is_password=False, placeholder="ex: ioana.popescu")
        p = _login_field("Parolă", "login_pass", is_password=True, placeholder="••••••••")

        # IMPORTANT: injectăm fix-ul de login după ce widget-urile există în DOM
        apply_login_fix()

        st.markdown("<div class='primary-btn'>", unsafe_allow_html=True)
        login_clicked = st.button("Intră în cont", use_container_width=True, key="login_btn")
        st.markdown("</div>", unsafe_allow_html=True)
        # Ajutor mic, sub buton (secondary)
        st.markdown("<div class='help-row'>", unsafe_allow_html=True)
        help_clicked = st.button("Ai nevoie de ajutor?", key="help_btn_link")
        st.markdown("</div>", unsafe_allow_html=True)

        if help_clicked:
            st.info("Contact IT: adaugă aici email/telefon sau instrucțiuni de resetare parolă.")


        if login_clicked:
            if not (u or "").strip() or not (p or "").strip():
                st.error("Completează utilizator și parolă.")
            else:
                ok, role_or_msg = auth_check(conn, u, p)
                if ok:
                    st.session_state.logged_in = True
                    st.session_state.user_role = role_or_msg
                    st.session_state.username = (u or "").strip()
                    # după autentificare, aterizează întotdeauna pe „🏠 Acasă”
                    st.session_state["main_choice"] = "🏠 Acasă"
                    st.rerun()
                else:
                    st.error(role_or_msg)

        st.caption("Dacă ai probleme de acces, contactează administratorul IT.")
        st.markdown("</div>", unsafe_allow_html=True)   # login-card
        st.markdown("</div>", unsafe_allow_html=True)   # login-wrap

    st.stop()

# -------------------------------------------------------------
# MODEL LOGIC (opțional)
# -------------------------------------------------------------
@dataclass
class Employee:
    id: Optional[int]
    marca: str
    last_name: str
    first_name: str
    cnp: str
    strada: str
    numar: str
    bloc: str
    scara: str
    apartament: str
    cod_postal: str
    localitate: str
    judet: str
    telefon_fix: str
    mobil: str
    email: str
    functie: str
    departament: str
    data_angajare: str
    tip_contract: str
    loc_munca: str
    departament_organizatoric: str
    functie_contract: str
    tip_norma: str
    program_munca: str
    salariu_baza: float
    ci_tip_act: str
    ci_serie: str
    ci_numar: str
    ci_eliberat_de: str
    ci_data_eliberarii: str
    stare_civila: str
    nr_copii: int
    studii: str
    profesie: str
    calificare: str
    observatii: str
    # copertă
    dosar_nr: str
    dosar_functionar_public: str
    dosar_data_intocmire: str
    dosar_autoritate: str
    dosar_intocmit_nume: str
    dosar_intocmit_functie: str
    dosar_intocmit_semnatura: str
    dosar_modificari_nume: str
    dosar_modificari_functie: str
    dosar_modificari_semnatura: str
    dosar_certificare_nume: str
    dosar_certificare_functie: str
    dosar_certificare_semnatura: str
    activitate_in_afara_functiei: str
    activitate_in_cadru_institutie: str
    situatia_drepturi_salariale: str
    situatia_concedii: str
    situatia_disciplinara: str
    registru_numar: str
    registru_data: str
    registru_observatii: str
    activ: int


# -------------------------------------------------------------
# HELPERI PENTRU VALORI
# -------------------------------------------------------------
def get_val(d: Dict[str, Any], *keys: str) -> str:
    for k in keys:
        if k in d and d[k] is not None and str(d[k]).strip() != "":
            return str(d[k])
    return ""


def safe_int(value, default=0):
    """Transformă în int și tratează corect NaN / None / ''. """
    try:
        if value is None:
            return default
        if isinstance(value, float) and math.isnan(value):
            return default
        if value == "":
            return default
        return int(value)
    except Exception:
        return default


def build_display_name(emp: Dict[str, Any]) -> str:
    ln = get_val(emp, "last_name", "LAST_NAME", "nume", "NUME")
    fn = get_val(emp, "first_name", "FIRST_NAME", "prenume", "PRENUME")
    cnp = get_val(emp, "cnp", "CNP")
    marca = get_val(emp, "marca", "MARCA")
    nume = f"{ln.strip()} {fn.strip()}".strip()
    return f"{nume} | {cnp} | {marca}"


# ---------- DOCX pentru coperta DOSAR PROFESIONAL ----------
def generate_dosar_profesional_docx(
    emp: Dict[str, Any],
    dosar_nr: str,
    functionar_public: str,
    dosar_data_intocmire: str,
    dosar_autoritate: str,
    intocmit_nume_list: List[str],
    intocmit_functie_list: List[str],
    intocmit_semnatura_list: List[str],
    modif_nume_list: List[str],
    modif_functie_list: List[str],
    modif_semnatura_list: List[str],
    cert_nume_list: List[str],
    cert_functie_list: List[str],
    cert_semnatura_list: List[str],
) -> bytes:
    # Config unitar: date instituție / conducere
    try:
        cfg = load_config()
    except Exception:
        cfg = {}
    if not dosar_autoritate or not str(dosar_autoritate).strip():
        # dacă nu este completat în fișa angajatului, folosim denumirea unității din Configurare
        dosar_autoritate = cfg.get("denumire_unitate", dosar_autoritate)
    # putem folosi în viitor și alte câmpuri din cfg (conducător, responsabil HR)

    def line(value: str, length: int = 60) -> str:
        value = (value or "").strip()
        if value:
            return value
        return "_" * length

    def normalize_group(nume_l, functie_l, semn_l):
        max_len = max(len(nume_l), len(functie_l), len(semn_l), 1)

        def pad(lst):
            lst = lst or []
            lst = lst + [""] * (max_len - len(lst))
            return lst[:max_len]

        return pad(nume_l), pad(functie_l), pad(semn_l)

    intocmit_nume_list, intocmit_functie_list, intocmit_semnatura_list = normalize_group(
        intocmit_nume_list, intocmit_functie_list, intocmit_semnatura_list
    )
    modif_nume_list, modif_functie_list, modif_semnatura_list = normalize_group(
        modif_nume_list, modif_functie_list, modif_semnatura_list
    )
    cert_nume_list, cert_functie_list, cert_semnatura_list = normalize_group(
        cert_nume_list, cert_functie_list, cert_semnatura_list
    )

    doc = Document()

    _docx_apply_antet_si_semnaturi(doc)

    p = doc.add_paragraph()
    p.alignment = 1
    run = p.add_run(f"DOSAR PROFESIONAL Nr.{line(dosar_nr, 30)}")
    run.bold = True

    p = doc.add_paragraph()
    p.alignment = 1
    run = p.add_run("FUNCTIONAR PUBLIC ")
    run.bold = True
    p.add_run(line(functionar_public, 50))

    doc.add_paragraph("")

    p = doc.add_paragraph()
    p.alignment = 1
    run = p.add_run("DATE REFERITOARE LA DOSARUL PROFESIONAL")
    run.bold = True

    doc.add_paragraph("")

    cnp_v = get_val(emp, "cnp", "CNP")

    p = doc.add_paragraph()
    p.add_run("FUNCTIONAR PUBLIC")
    p = doc.add_paragraph()
    p.add_run(line(functionar_public, 70))

    doc.add_paragraph("")

    p = doc.add_paragraph()
    p.add_run("COD NUMERIC")
    p = doc.add_paragraph()
    p.add_run("PERSONAL " + line(cnp_v, 60))

    doc.add_paragraph("")

    p = doc.add_paragraph()
    p.add_run("DATA ÎNTOCMIRII DOSARULUI")
    p = doc.add_paragraph()
    p.add_run("PROFESIONAL " + line(dosar_data_intocmire, 55))

    doc.add_paragraph("")

    p = doc.add_paragraph()
    p.add_run("AUTORITATEA SAU INSTITUȚIA PUBLICĂ CARE A ÎNTOCMIT DOSARUL")
    p = doc.add_paragraph()
    p.add_run("PROFESIONAL " + line(dosar_autoritate, 50))

    doc.add_paragraph("")

    p = doc.add_paragraph()
    run = p.add_run("PERSOANA CARE A ÎNTOCMIT DOSARUL PROFESIONAL:")
    run.bold = True

    table1 = doc.add_table(rows=1, cols=3)
    table1.style = "Table Grid"
    hdr = table1.rows[0].cells
    hdr[0].text = "NUMELE SI PRENUMELE"
    hdr[1].text = "FUNCTIA"
    hdr[2].text = "SEMNATURA"

    for i in range(len(intocmit_nume_list)):
        row = table1.add_row().cells
        row[0].text = intocmit_nume_list[i] or ""
        row[1].text = intocmit_functie_list[i] or ""
        row[2].text = intocmit_semnatura_list[i] or ""

    doc.add_paragraph("")

    p = doc.add_paragraph()
    run = p.add_run("PERSOANA AUTORIZATA SA OPEREZE MODIFICARI, COMPLETARI, RECTIFICARI:")
    run.bold = True

    table2 = doc.add_table(rows=1, cols=3)
    table2.style = "Table Grid"
    hdr = table2.rows[0].cells
    hdr[0].text = "NUMELE SI PRENUMELE"
    hdr[1].text = "FUNCTIA"
    hdr[2].text = "SEMNATURA"

    for i in range(len(modif_nume_list)):
        row = table2.add_row().cells
        row[0].text = modif_nume_list[i] or ""
        row[1].text = modif_functie_list[i] or ""
        row[2].text = modif_semnatura_list[i] or ""

    doc.add_paragraph("")

    p = doc.add_paragraph()
    run = p.add_run("PERSOANA AUTORIZATA SA CERTIFICE DATELE CUPRINSE IN DOSARUL PROFESIONAL:")
    run.bold = True

    table3 = doc.add_table(rows=1, cols=3)
    table3.style = "Table Grid"
    hdr = table3.rows[0].cells
    hdr[0].text = "NUMELE SI PRENUMELE"
    hdr[1].text = "FUNCTIA"
    hdr[2].text = "SEMNATURA"

    for i in range(len(cert_nume_list)):
        row = table3.add_row().cells
        row[0].text = cert_nume_list[i] or ""
        row[1].text = cert_functie_list[i] or ""
        row[2].text = cert_semnatura_list[i] or ""

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


# -------------------------------------------------------------
# HELPER: DOCX -> PDF (folosește docx2pdf)
# -------------------------------------------------------------
def convert_docx_to_pdf_bytes(docx_bytes: bytes) -> Optional[bytes]:
    try:
        from docx2pdf import convert
    except Exception:
        return None

    with tempfile.TemporaryDirectory() as tmpdir:
        in_path = os.path.join(tmpdir, "temp_in.docx")
        out_path = os.path.join(tmpdir, "temp_out.pdf")
        with open(in_path, "wb") as f:
            f.write(docx_bytes)
        try:
            convert(in_path, out_path)
        except Exception:
            return None
        if not os.path.exists(out_path):
            return None
        with open(out_path, "rb") as f:
            pdf_bytes = f.read()
    return pdf_bytes


# -------------------------------------------------------------
# HELPERI: GENERARE DOCX PE SECȚIUNI & DOSAR COMPLET
# -------------------------------------------------------------
def append_header_employee(doc: Document, emp: Dict[str, Any]) -> None:
    ln = get_val(emp, "last_name", "LAST_NAME", "nume", "NUME")
    fn = get_val(emp, "first_name", "FIRST_NAME", "prenume", "PRENUME")
    nume_complet = f"{ln} {fn}".strip()
    cnp = get_val(emp, "cnp", "CNP")
    marca = get_val(emp, "marca", "MARCA")

    p = doc.add_paragraph()
    p.add_run("Nume și prenume: ").bold = True
    p.add_run(nume_complet)
    p = doc.add_paragraph()
    p.add_run("CNP: ").bold = True
    p.add_run(cnp)
    p = doc.add_paragraph()
    p.add_run("Marcă: ").bold = True
    p.add_run(marca)


def append_date_personale_to_doc(doc: Document, emp: Dict[str, Any]) -> None:
    doc.add_heading("Date cu caracter personal", level=1)
    append_header_employee(doc, emp)

    doc.add_paragraph("")
    p = doc.add_paragraph()
    p.add_run("Stare civilă: ").bold = True
    p.add_run(get_val(emp, "stare_civila"))
    p = doc.add_paragraph()
    p.add_run("Număr copii: ").bold = True
    p.add_run(str(emp.get("nr_copii", "")))

    doc.add_paragraph("")
    p_title = doc.add_paragraph("Adresă:")
    p_title.runs[0].bold = True
    doc.add_paragraph(
        f"Strada {get_val(emp, 'strada')} nr. {get_val(emp, 'numar')}, "
        f"bloc {get_val(emp, 'bloc')} scara {get_val(emp, 'scara')}, ap. {get_val(emp, 'apartament')}"
    )
    doc.add_paragraph(
        f"Cod poștal {get_val(emp, 'cod_postal')}, "
        f"{get_val(emp, 'localitate')}, jud. {get_val(emp, 'judet')}"
    )

    doc.add_paragraph("")
    p_title = doc.add_paragraph("Date de contact:")
    p_title.runs[0].bold = True
    doc.add_paragraph(f"Telefon fix: {get_val(emp, 'telefon_fix')}")
    doc.add_paragraph(f"Telefon mobil: {get_val(emp, 'mobil')}")
    doc.add_paragraph(f"Email: {get_val(emp, 'email')}")

    doc.add_paragraph("")
    p_title = doc.add_paragraph("Act de identitate:")
    p_title.runs[0].bold = True
    doc.add_paragraph(
        f"Tip act: {get_val(emp, 'ci_tip_act')}  Serie: {get_val(emp, 'ci_serie')}  Număr: {get_val(emp, 'ci_numar')}"
    )
    doc.add_paragraph(
        f"Eliberat de: {get_val(emp, 'ci_eliberat_de')}  Data eliberării: {get_val(emp, 'ci_data_eliberarii')}"
    )


def create_date_personale_docx(emp: Dict[str, Any]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    append_date_personale_to_doc(doc, emp)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def append_studii_to_doc(doc: Document, emp: Dict[str, Any]) -> None:
    doc.add_heading("Studii și pregătire profesională", level=1)
    append_header_employee(doc, emp)

    doc.add_paragraph("")
    p = doc.add_paragraph()
    p.add_run("Studii: ").bold = True
    p.add_run(get_val(emp, "studii"))

    p = doc.add_paragraph()
    p.add_run("Profesia de bază: ").bold = True
    p.add_run(get_val(emp, "profesie"))

    p = doc.add_paragraph()
    p.add_run("Calificări / atestări: ").bold = True
    p.add_run(get_val(emp, "calificare"))

    p = doc.add_paragraph()
    p.add_run("Observații: ").bold = True
    p.add_run(get_val(emp, "observatii"))


def create_studii_docx(emp: Dict[str, Any]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    append_studii_to_doc(doc, emp)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def append_activitate_externa_to_doc(doc: Document, emp: Dict[str, Any]) -> None:
    doc.add_heading("Activitate în afara funcției publice", level=1)
    append_header_employee(doc, emp)

    doc.add_paragraph("")
    doc.add_paragraph(
        get_val(emp, "activitate_in_afara_functiei") or "Nu este înregistrată activitate în afara funcției publice."
    )


def create_activitate_externa_docx(emp: Dict[str, Any]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    append_activitate_externa_to_doc(doc, emp)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def append_activitate_institutie_to_doc(doc: Document, emp: Dict[str, Any]) -> None:
    doc.add_heading("Activitate în cadrul instituției", level=1)
    append_header_employee(doc, emp)

    doc.add_paragraph("")
    p = doc.add_paragraph()
    p.add_run("Funcție actuală: ").bold = True
    p.add_run(get_val(emp, "functie"))
    p = doc.add_paragraph()
    p.add_run("Compartiment / Departament: ").bold = True
    p.add_run(get_val(emp, "departament"))
    p = doc.add_paragraph()
    p.add_run("Data angajării: ").bold = True
    p.add_run(get_val(emp, "data_angajare"))
    p = doc.add_paragraph()
    p.add_run("Tip contract: ").bold = True
    p.add_run(get_val(emp, "tip_contract"))

    doc.add_paragraph("")
    p = doc.add_paragraph()
    p.add_run("Loc de muncă / Compartiment organizatoric: ").bold = True
    p.add_run(get_val(emp, "loc_munca"))
    p = doc.add_paragraph()
    p.add_run("Structura organizatorică: ").bold = True
    p.add_run(get_val(emp, "departament_organizatoric"))
    p = doc.add_paragraph()
    p.add_run("Funcție conform contractului: ").bold = True
    p.add_run(get_val(emp, "functie_contract"))
    p = doc.add_paragraph()
    p.add_run("Tip normă: ").bold = True
    p.add_run(get_val(emp, "tip_norma"))
    p = doc.add_paragraph()
    p.add_run("Program de muncă: ").bold = True
    p.add_run(get_val(emp, "program_munca"))

    doc.add_paragraph("")
    p_title = doc.add_paragraph("Descriere activitate:")
    p_title.runs[0].bold = True
    doc.add_paragraph(get_val(emp, "activitate_in_cadru_institutie"))


def create_activitate_institutie_docx(emp: Dict[str, Any]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    append_activitate_institutie_to_doc(doc, emp)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def append_situatia_salariala_to_doc(doc: Document, emp: Dict[str, Any]) -> None:
    doc.add_heading("Situația drepturilor salariale", level=1)
    append_header_employee(doc, emp)

    doc.add_paragraph("")
    p = doc.add_paragraph()
    p.add_run("Salariul de bază actual: ").bold = True
    try:
        sal = float(emp.get("salariu_baza", 0) or 0)
        p.add_run(f"{sal:.2f} lei")
    except Exception:
        p.add_run(str(emp.get("salariu_baza", "")))

    doc.add_paragraph("")
    p_title = doc.add_paragraph("Evoluția drepturilor salariale:")
    p_title.runs[0].bold = True
    doc.add_paragraph(get_val(emp, "situatia_drepturi_salariale"))


def create_situatia_salariala_docx(emp: Dict[str, Any]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    append_situatia_salariala_to_doc(doc, emp)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def append_concedii_to_doc(doc: Document, emp: Dict[str, Any]) -> None:
    doc.add_heading("Situația concediilor", level=1)
    append_header_employee(doc, emp)

    doc.add_paragraph("")
    doc.add_paragraph(get_val(emp, "situatia_concedii"))


def create_concedii_docx(emp: Dict[str, Any]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    append_concedii_to_doc(doc, emp)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def append_disciplinar_to_doc(doc: Document, emp: Dict[str, Any]) -> None:
    doc.add_heading("Situația disciplinară", level=1)
    append_header_employee(doc, emp)

    doc.add_paragraph("")
    doc.add_paragraph(get_val(emp, "situatia_disciplinara"))


def create_disciplinar_docx(emp: Dict[str, Any]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    append_disciplinar_to_doc(doc, emp)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def append_acces_to_doc(doc: Document, emp: Dict[str, Any], rows: List[sqlite3.Row]) -> None:
    doc.add_heading("Accesul la dosarul profesional", level=1)
    append_header_employee(doc, emp)

    doc.add_paragraph("")
    if not rows:
        doc.add_paragraph("Nu există înregistrări de acces la dosarul profesional.")
        return

    table = doc.add_table(rows=1, cols=7)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "Nr. crt."
    hdr[1].text = "Nume"
    hdr[2].text = "Funcție"
    hdr[3].text = "Semnătura"
    hdr[4].text = "Motivul"
    hdr[5].text = "Acces autorizat de"
    hdr[6].text = "Luat la cunoștință"

    for r in rows:
        row = table.add_row().cells
        row[0].text = str(r["nr_crt"])
        row[1].text = r["nume"] or ""
        row[2].text = r["functie"] or ""
        row[3].text = r["semnatura"] or ""
        row[4].text = r["motivul"] or ""
        row[5].text = r["acces_autorizat_de"] or ""
        row[6].text = r["luat_la_cunostinta"] or ""


def create_acces_docx(emp: Dict[str, Any], rows: List[sqlite3.Row]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    append_acces_to_doc(doc, emp, rows)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def append_registru_to_doc(doc: Document, emp: Dict[str, Any]) -> None:
    doc.add_heading("Registru evidență funcționari publici", level=1)
    append_header_employee(doc, emp)

    doc.add_paragraph("")
    p = doc.add_paragraph()
    p.add_run("Număr în registru: ").bold = True
    p.add_run(get_val(emp, "registru_numar"))
    p = doc.add_paragraph()
    p.add_run("Data înscrierii: ").bold = True
    p.add_run(get_val(emp, "registru_data"))

    doc.add_paragraph("")
    p_title = doc.add_paragraph("Observații:")
    p_title.runs[0].bold = True
    doc.add_paragraph(get_val(emp, "registru_observatii"))


def create_registru_docx(emp: Dict[str, Any]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    append_registru_to_doc(doc, emp)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def generate_dosar_complet_docx(emp: Dict[str, Any], acces_rows: List[sqlite3.Row]) -> bytes:
    doc = Document()
    _docx_apply_antet_si_semnaturi(doc)

    # Date unitate din configurare, afișate în antet
    try:
        cfg = load_config()
    except Exception:
        cfg = {}
    unit_name = cfg.get("denumire_unitate", "")
    unit_addr = cfg.get("adresa", "")

    ln = get_val(emp, "last_name", "LAST_NAME", "nume", "NUME")
    fn = get_val(emp, "first_name", "FIRST_NAME", "prenume", "PRENUME")
    nume_complet = f"{ln} {fn}".strip()
    cnp = get_val(emp, "cnp", "CNP")

    doc.add_heading("DOSAR PROFESIONAL COMPLET", level=0)
    p = doc.add_paragraph()
    p.add_run("Funcționar public: ").bold = True
    p.add_run(nume_complet)
    p = doc.add_paragraph()
    p.add_run("CNP: ").bold = True
    p.add_run(cnp)

    doc.add_page_break()
    append_date_personale_to_doc(doc, emp)

    doc.add_page_break()
    append_studii_to_doc(doc, emp)

    doc.add_page_break()
    append_activitate_externa_to_doc(doc, emp)

    doc.add_page_break()
    append_activitate_institutie_to_doc(doc, emp)

    doc.add_page_break()
    append_situatia_salariala_to_doc(doc, emp)

    doc.add_page_break()
    append_concedii_to_doc(doc, emp)

    doc.add_page_break()
    append_disciplinar_to_doc(doc, emp)

    doc.add_page_break()
    append_acces_to_doc(doc, emp, acces_rows)

    doc.add_page_break()
    append_registru_to_doc(doc, emp)

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


# -------------------------------------------------------------
# CRUD EMPLOYEES + LISTĂ CU FILTRE
# -------------------------------------------------------------
def list_employees(
    conn: sqlite3.Connection,
    active_only: bool = True,
    search: str = "",
    nume: str = "",
    prenume: str = "",
    cnp: str = "",
    marca: str = "",
) -> pd.DataFrame:
    """
    Listează angajații cu suport pentru diferite denumiri de coloane (last_name/nume etc.).

    - Face PRAGMA table_info(employees) imediat, ca să avem schema corectă indiferent de alte operații.
    - Construiește un col_map: nume_lower -> nume_real, astfel încât filtrele să folosească numele reale în SQL.
    - Funcționează atât cu sqlite3.Row (row_factory) cât și cu tuple simple.
    """
    cur = conn.cursor()

    # 1) Schema employees: citește cols_info imediat după PRAGMA, înainte de orice CREATE TABLE/INDEX
    cur.execute("PRAGMA table_info(employees);")
    cols_info = cur.fetchall() or []
    if not cols_info:
        return pd.DataFrame()

    # Compatibilitate sqlite3.Row / tuple (PRAGMA table_info: cid, name, type, ... → row[1] = name)
    def _get_col_name(row: Any) -> str:
        if hasattr(row, "keys"):
            return row["name"]
        return row[1]

    col_map: dict[str, str] = {}
    for r in cols_info:
        name = _get_col_name(r)
        col_map[name.lower()] = name

    col_names = set(col_map.keys())

    # 2) Abia acum: tabel pentru documente generate per angajat (istoric + document activ)
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER NOT NULL,
            doc_type TEXT NOT NULL,
            doc_no TEXT,
            doc_date TEXT,
            filename TEXT,
            data BLOB NOT NULL,
            meta_json TEXT,
            created_at TEXT,
            is_active INTEGER DEFAULT 1
        );
        """
    )
    cur.execute(
        "CREATE INDEX IF NOT EXISTS idx_employee_documents_emp ON employee_documents(employee_id);"
    )
    cur.execute(
        "CREATE INDEX IF NOT EXISTS idx_employee_documents_type ON employee_documents(doc_type);"
    )

    # 3) Determinăm coloanele logice (nume, prenume, cnp, marcă) folosind col_map
    def _pick_col(candidates: list[str]) -> str | None:
        for cand in candidates:
            key = cand.lower()
            if key in col_map:
                return col_map[key]
        return None

    col_last = _pick_col(["last_name", "nume"])
    col_first = _pick_col(["first_name", "prenume"])
    col_cnp = _pick_col(["cnp"])
    col_marca = _pick_col(["marca"])

    order_last = col_last or col_map.get("id", "id")
    order_first = col_first or col_map.get("id", "id")

    def _like_ci(s: str) -> str:
        """Parametru LIKE case-insensitive: %{strip+lower}%."""
        return f"%{(s or '').strip().lower()}%"

    base_sql = "SELECT * FROM employees"
    where_clauses: list[str] = []
    params: list[Any] = []

    # Filtru activi
    if active_only and "activ" in col_names:
        where_clauses.append(f"{col_map['activ']} = 1")

    # Filtre directe (case-insensitive pentru nume/prenume/marcă, CNP doar TRIM)
    if nume and col_last:
        where_clauses.append(f"LOWER(TRIM({col_last})) LIKE ?")
        params.append(_like_ci(nume))

    if prenume and col_first:
        where_clauses.append(f"LOWER(TRIM({col_first})) LIKE ?")
        params.append(_like_ci(prenume))

    if cnp and col_cnp:
        where_clauses.append(f"TRIM({col_cnp}) LIKE ?")
        params.append(f"%{(cnp or '').strip()}%")

    if marca and col_marca:
        where_clauses.append(f"LOWER(TRIM({col_marca})) LIKE ?")
        params.append(_like_ci(marca))

    # Căutare globală: dacă inputul curățat are 13 cifre → match exact CNP; altfel LIKE (CI pentru nume/prenume/marcă)
    if search:
        digits = cnp_clean(search)
        if len(digits) == 13 and col_cnp:
            where_clauses.append(f"TRIM({col_cnp}) = ?")
            params.append(digits)
        else:
            search_clauses: list[str] = []
            like_ci = _like_ci(search)
            cnp_like = f"%{(search or '').strip()}%"
            for logical_name in ["last_name", "first_name", "nume", "prenume", "marca"]:
                key = logical_name.lower()
                if key in col_map:
                    real_col = col_map[key]
                    search_clauses.append(f"LOWER(TRIM({real_col})) LIKE ?")
                    params.append(like_ci)
            if col_cnp and "cnp" in col_map:
                search_clauses.append(f"TRIM({col_cnp}) LIKE ?")
                params.append(cnp_like)
            # Căutare după ID când textul este numeric
            if "id" in col_map and (search or "").strip().isdigit():
                search_clauses.append(f"{col_map['id']} = ?")
                params.append(int((search or "").strip()))
            if search_clauses:
                where_clauses.append("(" + " OR ".join(search_clauses) + ")")

    if where_clauses:
        base_sql += " WHERE " + " AND ".join(where_clauses)

    base_sql += f" ORDER BY {order_last}, {order_first}"

    cur.execute(base_sql, params)
    col_names_result = [d[0] for d in cur.description] if cur.description else []
    rows = cur.fetchall()
    if not rows:
        return pd.DataFrame()

    # Compatibilitate Row / tuple
    if hasattr(rows[0], "keys"):
        return pd.DataFrame(rows, columns=rows[0].keys())
    return pd.DataFrame(rows, columns=col_names_result or None)



# -------------------------------------------------------------
# PAGINA: ANGAJAȚI (LISTĂ MINIMALĂ + FILTRARE)
# -------------------------------------------------------------
def _employees_columns(conn: sqlite3.Connection) -> list[str]:
    """Returnează coloanele din tabela employees (în ordinea din DB)."""
    cur = conn.cursor()
    cur.execute("PRAGMA table_info(employees);")
    cols = [r[1] for r in cur.fetchall()]  # (cid, name, type, notnull, dflt_value, pk)
    return cols


def set_current_employee(emp_id: int | None) -> None:
    """Setează angajatul curent în session (pentru Dosar/Documente/Organigramă)."""
    st.session_state["current_employee_id"] = emp_id


def get_current_employee_id() -> int | None:
    """Returnează id-ul angajatului curent din session."""
    return st.session_state.get("current_employee_id")


def _employee_fetch_min_list(conn: sqlite3.Connection, status: str = "active") -> pd.DataFrame:
    """Lista minimă pentru modulul Angajați.

    status:
      - "active"   -> activ = 1
      - "inactive" -> activ = 0
      - "all"      -> fără filtrare
    """
    status_u = (status or "active").strip().lower()
    if status_u == "inactive":
        where = "WHERE activ = 0"
    elif status_u == "all":
        where = ""
    else:
        where = "WHERE activ = 1"
    df = pd.read_sql(
        f"""
        SELECT
            id AS ID,
            COALESCE(TRIM(last_name), '')  AS NUME,
            COALESCE(TRIM(first_name), '') AS PRENUME,
            COALESCE(TRIM(cnp), '')        AS CNP,
            COALESCE(TRIM(marca), '')      AS MARCA,
            COALESCE(activ, 1)             AS ACTIV
        FROM employees
        {where}
        ORDER BY last_name, first_name
        """,
        conn,
    )
    for c in ["NUME", "PRENUME", "CNP", "MARCA"]:
        df[c] = df[c].fillna("").astype(str)
    return df


def _employee_fetch_by_id(conn: sqlite3.Connection, employee_id: int) -> dict:
    cur = conn.cursor()
    cols = _employees_columns(conn)
    cur.execute(f"SELECT {', '.join(cols)} FROM employees WHERE id = ?", (int(employee_id),))
    row = cur.fetchone()
    if not row:
        return {}
    # row poate fi sqlite3.Row sau tuple
    if hasattr(row, "keys"):
        return {k: row[k] for k in row.keys()}
    return {cols[i]: row[i] for i in range(len(cols))}


def _load_employee_bundle_impl(conn: sqlite3.Connection, emp_id: int) -> dict:
    """Încarcă employee + contracts + documents + dependents + attachments + positions (fără cache)."""
    emp = _employee_fetch_by_id(conn, int(emp_id))
    if not emp:
        return {"employee": {}, "employee_contracts": [], "employee_documents": [], "employee_dependents": [], "employee_attachments": [], "employee_positions": []}

    cur = conn.cursor()

    # employee_contracts
    cur.execute(
        """
        SELECT id, cim_nr, cim_data, title, created_at, updated_at, is_active
        FROM employee_contracts WHERE employee_id = ? ORDER BY created_at DESC, id DESC
        """,
        (int(emp_id),),
    )
    rows = cur.fetchall() or []
    contract_cols = ["id", "cim_nr", "cim_data", "title", "created_at", "updated_at", "is_active"]
    employee_contracts = [dict(zip(contract_cols, r)) for r in rows]

    # employee_documents (toate tipurile, fără BLOB data, cu meta_json)
    cur.execute(
        """
        SELECT id, doc_type, doc_no, doc_date, filename, created_at, is_active, meta_json
        FROM employee_documents WHERE employee_id = ? ORDER BY created_at DESC, id DESC
        """,
        (int(emp_id),),
    )
    rows = cur.fetchall() or []
    doc_cols = ["id", "doc_type", "doc_no", "doc_date", "filename", "created_at", "is_active", "meta_json"]
    employee_documents = [dict(zip(doc_cols, r)) for r in rows]

    # employee_dependents (listă dict)
    cur.execute(
        """
        SELECT id, employee_id, tip, nume, prenume, cnp, data_nasterii, grad_rudenie, observatii, activ
        FROM employee_dependents WHERE employee_id = ? ORDER BY nume, prenume
        """,
        (int(emp_id),),
    )
    rows = cur.fetchall() or []
    dep_cols = ["id", "employee_id", "tip", "nume", "prenume", "cnp", "data_nasterii", "grad_rudenie", "observatii", "activ"]
    employee_dependents = [dict(zip(dep_cols, r)) for r in rows]

    # employee_attachments (fără BLOB)
    cur.execute(
        """
        SELECT id, doc_type, display_name, file_name, mime_type, uploaded_at
        FROM employee_attachments WHERE employee_id = ? ORDER BY uploaded_at DESC, id DESC
        """,
        (int(emp_id),),
    )
    rows = cur.fetchall() or []
    att_cols = ["id", "doc_type", "display_name", "file_name", "mime_type", "uploaded_at"]
    employee_attachments = [dict(zip(att_cols, r)) for r in rows]

    # employee_positions (dacă tabela există)
    employee_positions = []
    cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='employee_positions'")
    if cur.fetchone():
        cur.execute(
            """
            SELECT id, employee_id, position_id, start_date, end_date
            FROM employee_positions WHERE employee_id = ? ORDER BY start_date DESC
            """,
            (int(emp_id),),
        )
        rows = cur.fetchall() or []
        pos_cols = ["id", "employee_id", "position_id", "start_date", "end_date"]
        employee_positions = [dict(zip(pos_cols, r)) for r in rows]

    return {
        "employee": emp,
        "employee_contracts": employee_contracts,
        "employee_documents": employee_documents,
        "employee_dependents": employee_dependents,
        "employee_attachments": employee_attachments,
        "employee_positions": employee_positions,
    }


@st.cache_data(ttl=60)
def load_employee_bundle(_db_path: str, emp_id: int) -> dict:
    """
    Încarcă bundle complet pentru un angajat (employees + contracts + documents + dependents + attachments + positions).
    Folosește _db_path (hashable) pentru cache. Pentru conn direct folosiți _load_employee_bundle_impl(conn, emp_id).
    """
    conn = sqlite3.connect(_db_path, check_same_thread=False)
    try:
        return _load_employee_bundle_impl(conn, int(emp_id))
    finally:
        conn.close()


def _employee_soft_delete(conn: sqlite3.Connection, employee_id: int) -> None:
    cur = conn.cursor()
    cur.execute("UPDATE employees SET activ = 0 WHERE id = ?", (int(employee_id),))
    conn.commit()
    # Actualizare imediată în UI (fără întârziere din cache).
    try:
        load_employee_bundle.clear()
    except Exception:
        pass


def normalize_employee_payload(conn: sqlite3.Connection, data: dict) -> dict:
    """
    Normalizează payload-ul pentru employees: CNP curat, validare, completare data nașterii/sex din CNP.
    Ridică ValueError cu mesaj clar dacă CNP este prezent dar invalid.
    """
    if not data:
        return {}
    result = dict(data)
    cnp_raw = (result.get("cnp") or "").strip()
    if not cnp_raw:
        return result

    result["cnp"] = cnp_clean(cnp_raw)
    valid, msg, info = cnp_validate(result["cnp"])
    if not valid:
        raise ValueError(msg)

    cols = _employees_columns(conn)
    # Data nașterii din CNP (ISO YYYY-MM-DD) dacă lipsește sau e goală
    for dob_col in ("data_nasterii", "dob", "data_nastere"):
        if dob_col not in cols:
            continue
        current = (result.get(dob_col) or "").strip()
        if not current and info.get("dob"):
            result[dob_col] = info["dob"].strftime("%Y-%m-%d")
        break
    # Sex din CNP dacă există coloană
    if "sex" in cols and info.get("sex"):
        result["sex"] = info["sex"]

    return result


def _employee_upsert(conn: sqlite3.Connection, data: dict, employee_id: int | None = None) -> int:
    """Inserează sau actualizează employees. Returnează id-ul."""
    data = normalize_employee_payload(conn, data or {})
    cols = _employees_columns(conn)
    # nu permitem modificarea directă a id-ului
    data = {k: v for k, v in data.items() if k in cols and k != "id"}

    cur = conn.cursor()
    if employee_id is None:
        if "activ" not in data:
            data["activ"] = 1
        keys = list(data.keys())
        vals = [data[k] for k in keys]
        if keys:
            cur.execute(
                f"INSERT INTO employees ({', '.join(keys)}) VALUES ({', '.join(['?']*len(keys))})",
                vals,
            )
        else:
            cur.execute("INSERT INTO employees DEFAULT VALUES")
        conn.commit()
        # Actualizare imediată în UI (fără întârziere din cache).
        try:
            load_employee_bundle.clear()
        except Exception:
            pass
        return int(cur.lastrowid)
    else:
        keys = list(data.keys())
        if keys:
            set_sql = ", ".join([f"{k} = ?" for k in keys])
            vals = [data[k] for k in keys] + [int(employee_id)]
            cur.execute(f"UPDATE employees SET {set_sql} WHERE id = ?", vals)
            conn.commit()
            # Actualizare imediată în UI (fără întârziere din cache).
            try:
                load_employee_bundle.clear()
            except Exception:
                pass
        return int(employee_id)


def _employee_field_groups(cols: list[str]) -> list[tuple[str, list[str]]]:
    """Grupează câmpurile pentru afișare/editarare (doar dacă există în schema curentă)."""
    def keep(group_cols: list[str]) -> list[str]:
        return [c for c in group_cols if c in cols]

    groups: list[tuple[str, list[str]]] = []

    groups.append(("Date de bază", keep(["marca", "last_name", "first_name", "cnp", "activ"])))
    groups.append(("Adresă", keep(["strada", "numar", "bloc", "scara", "apartament", "cod_postal", "localitate", "judet"])))
    groups.append(("Contact", keep(["telefon_fix", "mobil", "email"])))
    groups.append(("Act identitate", keep(["ci_tip_act", "ci_serie", "ci_numar", "ci_eliberat_de", "ci_data_eliberarii"])))
    groups.append(("Familie", keep(["stare_civila", "nr_copii"])))
    groups.append(("Contract / muncă", keep(["functie", "departament", "data_angajare", "tip_contract", "loc_munca",
                                           "departament_organizatoric", "functie_contract", "tip_norma",
                                           "program_munca", "salariu_baza"])))
    groups.append(("Studii & diverse", keep(["studii", "profesie", "calificare", "observatii"])))
    # Copertă dosar profesional (dacă există)
    groups.append(("Dosar profesional – copertă", keep([
        "dosar_nr", "dosar_functionar_public", "dosar_data_intocmire", "dosar_autoritate",
        "dosar_intocmit_nume", "dosar_intocmit_functie", "dosar_intocmit_semnatura",
        "dosar_modificari_nume", "dosar_modificari_functie", "dosar_modificari_semnatura",
        "dosar_certificare_nume", "dosar_certificare_functie", "dosar_certificare_semnatura",
    ])))

    # orice coloane noi (care nu sunt în grupuri) le punem la final, ca să nu pierdem nimic
    used = set()
    for _, gcols in groups:
        used.update(gcols)
    extras = [c for c in cols if c not in used and c not in ("id",)]
    if extras:
        groups.append(("Altele", extras))

    # scoatem grupurile goale
    return [(name, gcols) for name, gcols in groups if gcols]


def _employee_field_labels() -> dict[str, str]:
    """Etichete în română pentru câmpurile angajatului (afișare și formulare)."""
    return {
        "marca": "Marcă",
        "last_name": "Nume",
        "first_name": "Prenume",
        "cnp": "CNP",
        "activ": "Activ",
        "strada": "Strada",
        "numar": "Număr",
        "bloc": "Bloc",
        "scara": "Scara",
        "apartament": "Apartament",
        "cod_postal": "Cod poștal",
        "localitate": "Localitate",
        "judet": "Județ",
        "telefon_fix": "Telefon fix",
        "mobil": "Mobil",
        "email": "Email",
        "ci_tip_act": "Tip act",
        "ci_serie": "Serie",
        "ci_numar": "Număr act",
        "ci_eliberat_de": "Eliberat de",
        "ci_data_eliberarii": "Data eliberării",
        "stare_civila": "Stare civilă",
        "nr_copii": "Nr. copii",
        "functie": "Funcție",
        "departament": "Departament",
        "data_angajare": "Data angajării",
        "tip_contract": "Tip contract",
        "loc_munca": "Loc de muncă",
        "departament_organizatoric": "Departament organizatoric",
        "functie_contract": "Funcție în contract",
        "tip_norma": "Tip normă",
        "program_munca": "Program muncă",
        "salariu_baza": "Salariu de bază",
        "studii": "Studii",
        "profesie": "Profesie",
        "calificare": "Calificare",
        "observatii": "Observații",
        "dosar_nr": "Dosar nr.",
        "dosar_functionar_public": "Funcționar public",
        "dosar_data_intocmire": "Data întocmirii dosarului",
        "dosar_autoritate": "Autoritate emitentă",
        "dosar_intocmit_nume": "Persoană care a întocmit – nume",
        "dosar_intocmit_functie": "Persoană care a întocmit – funcție",
        "dosar_intocmit_semnatura": "Persoană care a întocmit – semnătură",
        "dosar_modificari_nume": "Persoană modificări – nume",
        "dosar_modificari_functie": "Persoană modificări – funcție",
        "dosar_modificari_semnatura": "Persoană modificări – semnătură",
        "dosar_certificare_nume": "Persoană certificare – nume",
        "dosar_certificare_functie": "Persoană certificare – funcție",
        "dosar_certificare_semnatura": "Persoană certificare – semnătură",
        "activitate_in_afara_functiei": "Activitate în afara funcției",
        "activitate_in_cadru_institutie": "Activitate în cadrul instituției",
        "situatia_drepturi_salariale": "Situația drepturilor salariale",
        "situatia_concedii": "Situația concediilor",
        "situatia_disciplinara": "Situația disciplinară",
        "registru_numar": "Număr în registru",
        "registru_data": "Data înscrierii în registru",
        "registru_observatii": "Observații registru",
    }


def _render_info_section(title: str, rows: list[tuple[str, str]]) -> None:
    """Randare uniformă label/value pentru secțiuni informative din fișa angajatului."""
    def _esc(text: str) -> str:
        return str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    html_rows: list[str] = []
    for label, value in rows:
        raw_val = "" if value is None else str(value)
        shown_val = raw_val.strip() if raw_val is not None else ""
        is_empty = shown_val in ("", "—")
        safe_val = "—" if is_empty else shown_val
        row_html = (
            "<div class='emp-info-row'>"
            f"<div class='emp-info-label'>{_esc(label)}</div>"
            f"<div class='emp-info-value{' is-empty' if is_empty else ''}'>{_esc(safe_val)}</div>"
            "</div>"
        )
        html_rows.append(row_html)

    st.markdown(
        (
            "<div class='emp-info-section'>"
            f"<div class='emp-info-title'>{_esc(title)}</div>"
            "<div class='emp-info-grid'>"
            + "".join(html_rows)
            + "</div></div>"
        ),
        unsafe_allow_html=True,
    )


def _emp_v(emp: dict, key: str, default: str = "—") -> str:
    """Returnează valoarea din employee ca text sigur pentru UI."""
    raw = emp.get(key, "")
    text = "" if raw is None else str(raw).strip()
    return text if text else default


def _render_cnp_field(label: str, value: str | None, *, key: str) -> str:
    """
    Secțiune dedicată pentru CNP în formularul de angajat.
    - Arată clar condițiile (13 cifre, structură, verificări).
    - Validează în timp real și afișează detalii (data nașterii, sex, județ) dacă CNP-ul e valid.
    """
    raw = st.text_input(
        label,
        value="" if value is None else str(value),
        key=key,
        help="CNP-ul trebuie să aibă 13 cifre (structură S YYMMDD JJ NNN C). "
             "Se validează automat data nașterii, sexul, județul și cifra de control.",
    )

    cleaned = cnp_clean(raw)

    if not raw.strip():
        st.caption("CNP gol (opțional). Poți lăsa necompletat dacă nu este obligatoriu.")
        return raw

    info = decode_cnp(raw, strict_county=False, allow_s_9=False)

    if info.valid:
        st.success("CNP valid.")
        detalii = []
        if info.birth_date:
            detalii.append(f"Data nașterii: {info.birth_date.strftime('%d.%m.%Y')}")
        if info.sex:
            detalii.append(f"Sex: {'Masculin' if info.sex == 'M' else 'Feminin'}")
        if info.county_name:
            detalii.append(f"Județ: {info.county_name} ({info.county_code})")
        if detalii:
            for d in detalii:
                st.caption(d)

        # Casete vizibile cu datele extrase din CNP (în secțiunea „Date de bază”)
        st.text_input(
            "Data nașterii (din CNP)",
            value=info.birth_date.strftime("%d.%m.%Y") if info.birth_date else "",
            key=f"{key}_dob_from_cnp",
            disabled=True,
        )
        st.text_input(
            "Sex (din CNP)",
            value=("Masculin" if info.sex == "M" else "Feminin") if info.sex else "",
            key=f"{key}_sex_from_cnp",
            disabled=True,
        )
        st.text_input(
            "Județ (din CNP)",
            value=f"{info.county_name} ({info.county_code})" if info.county_name else "",
            key=f"{key}_county_from_cnp",
            disabled=True,
        )
    else:
        st.error(info.error or "CNP invalid.")
        extra = []
        if info.birth_date:
            extra.append(f"Data nașterii (din cifre): {info.birth_date.strftime('%d.%m.%Y')}")
        if info.sex:
            extra.append(f"Sex (din cifre): {'M' if info.sex == 'M' else 'F'}")
        if info.county_code:
            extra.append(f"Cod județ: {info.county_code}")
        if extra:
            st.caption(" | ".join(extra))

    # Lăsăm normalizarea finală către normalize_employee_payload (care folosește cnp_clean)
    return raw


def _render_employee_form_fields(
    employee: dict,
    cols: list[str],
    *,
    prefix: str,
    single_column: bool = False,
    fine_section_titles: bool = False,
    parallel_groups: bool = False,
    section_title_class: str | None = None,
) -> dict:
    """Randare câmpuri + colectare valori (dinamic, după schema DB)."""
    out = {}
    groups = _employee_field_groups(cols)

    # helper: decide widget
    def is_number(col: str) -> bool:
        return col in ("nr_copii",)

    def is_real(col: str) -> bool:
        return col in ("salariu_baza",)

    def is_boolish(col: str) -> bool:
        return col in ("activ",)

    labels = _employee_field_labels()
    def _render_group(gname: str, gcols: list[str]) -> None:
        with st.container():
            if fine_section_titles:
                cls = "emp-block-title" if not section_title_class else f"emp-block-title {section_title_class}"
                st.markdown(f'<div class="{cls}">{gname}</div>', unsafe_allow_html=True)
            else:
                st.markdown(f"### {gname}")

            # layout: 2 coloane sau 1 coloană (pentru editare aliniată stânga)
            if single_column:
                colL, colR = st.columns([1, 0.0001])
            else:
                colL, colR = st.columns(2)
            toggle = 0
            for c in gcols:
                val = employee.get(c, None)
                # Etichetă: întâi din map-ul dedicat, altfel formatăm automat snake_case -> "Titlu Frumos"
                if c in labels:
                    label = labels[c]
                else:
                    parts = [p for p in str(c).split("_") if p]
                    label = " ".join(p.capitalize() for p in parts) if parts else str(c)
                key = f"{prefix}_{c}"

                if single_column:
                    target = colL
                else:
                    target = colL if toggle % 2 == 0 else colR
                    toggle += 1

                with target:
                    if is_boolish(c):
                        if single_column and (prefix.startswith("ang_edit_") or prefix.startswith("ang_add")):
                            st.markdown(f"**{label}**")
                            out[c] = st.checkbox(
                                "",
                                value=bool(val) if val is not None else True,
                                key=key,
                                label_visibility="collapsed",
                            )
                        else:
                            out[c] = st.checkbox(label, value=bool(val) if val is not None else True, key=key)
                    elif is_number(c):
                        if prefix.startswith("ang_add"):
                            out[c] = st.text_input(label, value="" if val in (None, "") else str(val), key=key)
                            continue
                        try:
                            v0 = int(val) if val not in (None, "") else 0
                        except Exception:
                            v0 = 0
                        out[c] = st.number_input(label, min_value=0, step=1, value=v0, key=key)
                    elif is_real(c):
                        if prefix.startswith("ang_add"):
                            out[c] = st.text_input(label, value="" if val in (None, "") else str(val), key=key)
                            continue
                        try:
                            v0 = float(val) if val not in (None, "") else 0.0
                        except Exception:
                            v0 = 0.0
                        out[c] = st.number_input(label, value=v0, key=key)
                    elif c == "cnp":
                        # CNP are o secțiune specială, cu reguli explicate și validare vizibilă
                        out[c] = _render_cnp_field(label, val, key=key)
                    else:
                        out[c] = st.text_input(label, value="" if val is None else str(val), key=key)

            st.divider()

    if parallel_groups:
        # Randăm în perechi pe rânduri separate: stânga + dreapta aliniate pe același nivel.
        for i in range(0, len(groups), 2):
            col_left, _, col_right = st.columns([1, 0.14, 1])
            gname_l, gcols_l = groups[i]
            with col_left:
                _render_group(gname_l, gcols_l)
            if i + 1 < len(groups):
                gname_r, gcols_r = groups[i + 1]
                with col_right:
                    _render_group(gname_r, gcols_r)
    else:
        for gname, gcols in groups:
            _render_group(gname, gcols)

    return out

# -------------------------------------------------------------
# PERSOANE ÎN ÎNTREȚINERE (copii / părinți / soț-soție) - CRUD
# -------------------------------------------------------------

DEPENDENT_TIP_OPTS = ["", "SOT_SOTIE", "COPIL", "PARINTE"]  # ordine UI
DEPENDENT_TIP_LABELS = {
    "": "Selectează tipul",
    "SOT_SOTIE": "Soț / Soție",
    "COPIL": "Copil",
    "PARINTE": "Părinte",
}

def _dependent_get(conn: sqlite3.Connection, dep_id: int) -> dict | None:
    cur = conn.cursor()
    cur.execute(
        """
        SELECT id, employee_id, tip, nume, prenume, cnp, data_nasterii, grad_rudenie, observatii, activ
        FROM employee_dependents
        WHERE id = ?
        """,
        (int(dep_id),),
    )
    r = cur.fetchone()
    if not r:
        return None
    return {
        "id": r[0],
        "employee_id": r[1],
        "tip": (r[2] or "").strip(),
        "nume": (r[3] or "").strip(),
        "prenume": (r[4] or "").strip(),
        "cnp": (r[5] or "").strip(),
        "data_nasterii": (r[6] or "").strip(),
        "grad_rudenie": (r[7] or "").strip(),
        "observatii": (r[8] or "").strip(),
        "activ": int(r[9] or 1),
    }


def _dependents_fetch(conn: sqlite3.Connection, employee_id: int, include_inactive: bool = False) -> pd.DataFrame:
    where = "" if include_inactive else "AND COALESCE(activ, 1) = 1"

    # ordonare “umană”: Soț/Soție -> Copil -> Părinte -> rest
    df = pd.read_sql(
        f"""
        SELECT
            id AS ID,
            employee_id AS EMPLOYEE_ID,
            COALESCE(TRIM(tip), '') AS TIP,
            COALESCE(TRIM(nume), '') AS NUME,
            COALESCE(TRIM(prenume), '') AS PRENUME,
            COALESCE(TRIM(cnp), '') AS CNP,
            COALESCE(TRIM(data_nasterii), '') AS DATA_NASTERII,
            COALESCE(TRIM(grad_rudenie), '') AS GRAD_RUDENIE,
            COALESCE(TRIM(observatii), '') AS OBSERVATII,
            COALESCE(activ, 1) AS ACTIV
        FROM employee_dependents
        WHERE employee_id = ? {where}
        ORDER BY
            CASE
                WHEN UPPER(TRIM(tip)) IN ('SOT_SOTIE','SOT','SOTIE','SOȚ','SOȚIE','SOT/SOTIE','SOT / SOTIE') THEN 0
                WHEN UPPER(TRIM(tip)) IN ('COPIL') THEN 1
                WHEN UPPER(TRIM(tip)) IN ('PARINTE','PĂRINTE') THEN 2
                ELSE 9
            END,
            NUME, PRENUME
        """,
        conn,
        params=(int(employee_id),),
    )
    return df


def _dependent_upsert(conn: sqlite3.Connection, employee_id: int, values: dict, dep_id: int | None = None) -> int:
    cur = conn.cursor()

    cols = ["employee_id", "tip", "nume", "prenume", "cnp", "data_nasterii", "grad_rudenie", "observatii", "activ"]
    data = {k: values.get(k) for k in cols}
    data["employee_id"] = int(employee_id)

    # Normalize tip
    tip = (data.get("tip") or "").strip()
    if not tip:
        tip = "COPIL"
    tip_upper = tip.upper().replace(" ", "_")

    # acceptăm și variante vechi introduse manual
    if tip_upper in ("SOT", "SOTIE", "SOȚ", "SOȚIE", "SOT/SOTIE", "SOT_/_SOTIE"):
        tip_upper = "SOT_SOTIE"
    if tip_upper in ("PĂRINTE",):
        tip_upper = "PARINTE"
    if tip_upper not in ("SOT_SOTIE", "COPIL", "PARINTE"):
        # fallback sigur
        tip_upper = "COPIL"

    data["tip"] = tip_upper

    if data.get("activ") is None:
        data["activ"] = 1

    # Trim string fields
    for k in ("nume", "prenume", "cnp", "data_nasterii", "grad_rudenie", "observatii"):
        if data.get(k) is not None:
            data[k] = str(data[k]).strip()

    if dep_id:
        sets = ", ".join([f"{c}=?" for c in cols])
        params = [data[c] for c in cols] + [int(dep_id)]
        cur.execute(f"UPDATE employee_dependents SET {sets} WHERE id=?", params)
        conn.commit()
        return int(dep_id)

    placeholders = ", ".join(["?"] * len(cols))
    params = [data[c] for c in cols]
    cur.execute(f"INSERT INTO employee_dependents ({', '.join(cols)}) VALUES ({placeholders})", params)
    conn.commit()
    return int(cur.lastrowid)


def _dependent_soft_delete(conn: sqlite3.Connection, dep_id: int) -> None:
    cur = conn.cursor()
    cur.execute("UPDATE employee_dependents SET activ = 0 WHERE id = ?", (int(dep_id),))
    conn.commit()

# ============================================================
# CNP – verificări în DB (conflict nume / duplicat)
# ============================================================
def _cnp_conflicts_in_dependents(conn, cnp: str, nume: str, prenume: str, exclude_dep_id: int | None = None):
    """
    Caută același CNP în employee_dependents și semnalează dacă numele diferă.
    """
    cnp = cnp_clean(cnp)
    if not cnp:
        return []

    sql = """
        SELECT id, employee_id, tip, nume, prenume, cnp, activ
        FROM employee_dependents
        WHERE cnp = ?
    """
    params = [cnp]
    if exclude_dep_id:
        sql += " AND id <> ?"
        params.append(int(exclude_dep_id))

    rows = conn.execute(sql, params).fetchall()
    conflicts = []
    n1 = (nume or "").strip().upper()
    p1 = (prenume or "").strip().upper()

    for r in rows:
        rid, emp_id, tip, rn, rp, rc, activ = r
        rn_u = (rn or "").strip().upper()
        rp_u = (rp or "").strip().upper()
        # conflict dacă există nume/prenume și diferă
        if (rn_u and rp_u) and ((rn_u != n1) or (rp_u != p1)):
            conflicts.append({
                "id": rid,
                "employee_id": emp_id,
                "tip": tip,
                "nume": rn,
                "prenume": rp,
                "cnp": rc,
                "activ": activ,
            })
    return conflicts


def _cnp_conflicts_in_employees(conn, cnp: str, nume: str, prenume: str, exclude_emp_id: int | None = None):
    """
    Caută același CNP în employees (dacă ai coloană cnp acolo).
    Dacă tabela employees NU are cnp, funcția va returna [] fără crash.
    """
    cnp = cnp_clean(cnp)
    if not cnp:
        return []

    # verifică dacă există coloana cnp în employees
    try:
        cols = [x[1] for x in conn.execute("PRAGMA table_info(employees)").fetchall()]
        if "cnp" not in cols:
            return []
    except Exception:
        return []

    sql = "SELECT id, nume, prenume, cnp, activ FROM employees WHERE cnp = ?"
    params = [cnp]
    if exclude_emp_id:
        sql += " AND id <> ?"
        params.append(int(exclude_emp_id))

    rows = conn.execute(sql, params).fetchall()
    conflicts = []
    n1 = (nume or "").strip().upper()
    p1 = (prenume or "").strip().upper()

    for r in rows:
        rid, rn, rp, rc, activ = r
        rn_u = (rn or "").strip().upper()
        rp_u = (rp or "").strip().upper()
        if (rn_u and rp_u) and ((rn_u != n1) or (rp_u != p1)):
            conflicts.append({
                "id": rid,
                "nume": rn,
                "prenume": rp,
                "cnp": rc,
                "activ": activ,
            })
    return conflicts

# -------------------------------------------------------------
# Șablon global DOCX (Contract de muncă)
# -------------------------------------------------------------
def _get_global_template(conn: sqlite3.Connection, name: str = "cim") -> tuple[str | None, bytes | None]:
    """Returnează (filename, data) pentru șablonul global; sau (None, None)."""
    cur = conn.cursor()
    cur.execute("SELECT filename, data FROM global_templates WHERE name = ?", (name,))
    row = cur.fetchone()
    if not row:
        return None, None
    return (row[0], row[1])


def _set_global_template(conn: sqlite3.Connection, *, name: str = "cim", filename: str, data: bytes) -> None:
    cur = conn.cursor()
    cur.execute(
        """
        INSERT INTO global_templates(name, filename, data, updated_at)
        VALUES (?, ?, ?, ?)
        ON CONFLICT(name) DO UPDATE SET
            filename=excluded.filename,
            data=excluded.data,
            updated_at=excluded.updated_at
        """,
        (name, filename, sqlite3.Binary(data), datetime.datetime.now().isoformat(timespec="seconds")),
    )
    conn.commit()

# -------------------------------------------------------------
# DB: documente angajat (CIM / adeverințe etc.) stocate ca BLOB
# -------------------------------------------------------------

# -------------------------------------------------------------
# CIM: Salvare + Istoric contracte
# -------------------------------------------------------------
def _cim_contract_save(conn, employee_id: int, payload: dict, *, cim_nr: str = "", cim_data: str = "", title: str = "") -> int:
    cur = conn.cursor()
    now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    js = json.dumps(payload or {}, ensure_ascii=False)
    cur.execute(
        """
        INSERT INTO employee_contracts(employee_id, cim_nr, cim_data, title, payload_json, created_at, updated_at, is_active)
        VALUES(?, ?, ?, ?, ?, ?, ?, 0)
        """
        , (int(employee_id), cim_nr or None, cim_data or None, title or None, js, now, now)
    )
    conn.commit()
    return int(cur.lastrowid)

def _cim_contract_list(conn, employee_id: int):
    cur = conn.cursor()
    cur.execute(
        """
        SELECT id, cim_nr, cim_data, title, created_at, updated_at, is_active
        FROM employee_contracts
        WHERE employee_id = ?
        ORDER BY created_at DESC, id DESC
        """
        , (int(employee_id),)
    )
    rows = cur.fetchall() or []
    cols = ["id","cim_nr","cim_data","title","created_at","updated_at","is_active"]
    return [dict(zip(cols, r)) for r in rows]

def _cim_contract_load(conn, contract_id: int) -> dict:
    cur = conn.cursor()
    cur.execute("SELECT payload_json FROM employee_contracts WHERE id = ?", (int(contract_id),))
    r = cur.fetchone()
    if not r:
        return {}
    try:
        return json.loads(r[0] or "{}") or {}
    except Exception:
        return {}

def _cim_contract_set_active(conn, employee_id: int, contract_id: int):
    cur = conn.cursor()
    cur.execute("UPDATE employee_contracts SET is_active = 0 WHERE employee_id = ?", (int(employee_id),))
    cur.execute("UPDATE employee_contracts SET is_active = 1 WHERE id = ?", (int(contract_id),))
    conn.commit()

def _cim_contract_delete(conn, contract_id: int):
    cur = conn.cursor()
    cur.execute("DELETE FROM employee_contracts WHERE id = ?", (int(contract_id),))
    conn.commit()


# -------------------------------------------------------------
# Documente PDF per angajat (upload + download + delete)
# -------------------------------------------------------------
def _emp_doc_save(conn, employee_id: int, doc_type: str, file_name: str, mime_type: str, file_bytes: bytes, *, display_name: str = ""):
    cur = conn.cursor()
    now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    cur.execute(
        """
        INSERT INTO employee_attachments(employee_id, doc_type, display_name, file_name, mime_type, file_blob, uploaded_at)
        VALUES(?, ?, ?, ?, ?, ?, ?)
        """
        , (int(employee_id), str(doc_type), str(display_name or doc_type), str(file_name or ""), str(mime_type or "application/pdf"), file_bytes, now)
    )
    conn.commit()

def _emp_doc_list(conn, employee_id: int, doc_type: str | None = None):
    cur = conn.cursor()
    if doc_type:
        cur.execute(
            """
            SELECT id, doc_type, display_name, file_name, mime_type, uploaded_at
            FROM employee_attachments
            WHERE employee_id = ? AND doc_type = ?
            ORDER BY uploaded_at DESC, id DESC
            """
            , (int(employee_id), str(doc_type))
        )
    else:
        cur.execute(
            """
            SELECT id, doc_type, display_name, file_name, mime_type, uploaded_at
            FROM employee_attachments
            WHERE employee_id = ?
            ORDER BY uploaded_at DESC, id DESC
            """
            , (int(employee_id),)
        )
    rows = cur.fetchall() or []
    cols = ["id","doc_type","display_name","file_name","mime_type","uploaded_at"]
    return [dict(zip(cols, r)) for r in rows]

def _emp_doc_load(conn, doc_id: int):
    cur = conn.cursor()
    cur.execute("SELECT file_name, mime_type, file_blob FROM employee_attachments WHERE id = ?", (int(doc_id),))
    r = cur.fetchone()
    if not r:
        return None, None, None
    return r[0], r[1], r[2]

def _emp_doc_delete(conn, doc_id: int):
    cur = conn.cursor()
    cur.execute("DELETE FROM employee_attachments WHERE id = ?", (int(doc_id),))
    conn.commit()

def _save_employee_document(
    conn: sqlite3.Connection,
    *,
    employee_id: int,
    doc_type: str,
    filename: str,
    data: bytes,
    doc_no: str | None = None,
    doc_date: str | None = None,
    meta: dict | None = None,
    set_active: bool = True,
) -> int:
    """Salvează un document (BLOB) și, opțional, îl marchează ca activ pentru tipul respectiv."""
    cur = conn.cursor()
    if set_active:
        cur.execute(
            "UPDATE employee_documents SET is_active = 0 WHERE employee_id = ? AND doc_type = ?",
            (int(employee_id), str(doc_type)),
        )
    cur.execute(
        """
        INSERT INTO employee_documents(employee_id, doc_type, doc_no, doc_date, filename, data, meta_json, created_at, is_active)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            int(employee_id),
            str(doc_type),
            (str(doc_no).strip() if doc_no is not None else None),
            (str(doc_date).strip() if doc_date is not None else None),
            str(filename),
            sqlite3.Binary(data),
            json.dumps(meta or {}, ensure_ascii=False),
            datetime.datetime.now().isoformat(timespec="seconds"),
            1 if bool(set_active) else 0,
        ),
    )
    conn.commit()
    return int(cur.lastrowid)


def _list_employee_documents(conn: sqlite3.Connection, *, employee_id: int, doc_type: str) -> list[dict]:
    cur = conn.cursor()
    cur.execute(
        """
        SELECT id, doc_type, doc_no, doc_date, filename, created_at, is_active
        FROM employee_documents
        WHERE employee_id = ? AND doc_type = ?
        ORDER BY created_at DESC, id DESC
        """,
        (int(employee_id), str(doc_type)),
    )
    rows = cur.fetchall() or []
    out = []
    for r in rows:
        out.append(
            {
                "id": int(r[0]),
                "doc_type": r[1],
                "doc_no": r[2],
                "doc_date": r[3],
                "filename": r[4],
                "created_at": r[5],
                "is_active": int(r[6] or 0),
            }
        )
    return out


def _get_employee_document_bytes(conn: sqlite3.Connection, *, doc_id: int) -> tuple[str | None, bytes | None]:
    cur = conn.cursor()
    cur.execute("SELECT filename, data FROM employee_documents WHERE id = ?", (int(doc_id),))
    row = cur.fetchone()
    if not row:
        return None, None
    return (row[0], row[1])
# -------------------------------------------------------------
# DOCUMENTE ANGAJAT: operații + audit
# -------------------------------------------------------------
def _employee_document_set_active(conn: sqlite3.Connection, *, doc_id: int) -> None:
    cur = conn.cursor()
    cur.execute("SELECT employee_id, doc_type FROM employee_documents WHERE id = ?", (int(doc_id),))
    row = cur.fetchone()
    if not row:
        return
    emp_id, doc_type = int(row[0]), str(row[1])
    cur.execute(
        "UPDATE employee_documents SET is_active = 0 WHERE employee_id = ? AND doc_type = ?",
        (emp_id, doc_type),
    )
    cur.execute("UPDATE employee_documents SET is_active = 1 WHERE id = ?", (int(doc_id),))
    conn.commit()

def _employee_document_delete(conn: sqlite3.Connection, *, doc_id: int) -> None:
    cur = conn.cursor()
    cur.execute("DELETE FROM employee_documents WHERE id = ?", (int(doc_id),))
    conn.commit()

def _ensure_employee_audit_table(conn: sqlite3.Connection) -> None:
    cur = conn.cursor()
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_audit (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER,
            entity TEXT,
            entity_id INTEGER,
            action TEXT,
            details_json TEXT,
            actor TEXT,
            created_at TEXT
        );
        """
    )
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_audit_emp ON employee_audit(employee_id);")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_employee_audit_entity ON employee_audit(entity, entity_id);")
    conn.commit()

def _audit_log(
    conn: sqlite3.Connection,
    *,
    employee_id: int | None,
    entity: str,
    entity_id: int | None,
    action: str,
    details: dict | None = None,
    actor: str | None = None,
) -> None:
    try:
        _ensure_employee_audit_table(conn)
    except Exception:
        return
    cur = conn.cursor()
    cur.execute(
        """
        INSERT INTO employee_audit(employee_id, entity, entity_id, action, details_json, actor, created_at)
        VALUES (?, ?, ?, ?, ?, ?, ?)
        """,
        (
            int(employee_id) if employee_id is not None else None,
            str(entity),
            int(entity_id) if entity_id is not None else None,
            str(action),
            json.dumps(details or {}, ensure_ascii=False),
            (str(actor).strip() if actor else None),
            datetime.datetime.now().isoformat(timespec="seconds"),
        ),
    )
    conn.commit()

def _audit_list(conn: sqlite3.Connection, *, employee_id: int, limit: int = 200) -> list[dict]:
    try:
        _ensure_employee_audit_table(conn)
    except Exception:
        return []
    cur = conn.cursor()
    cur.execute(
        """
        SELECT id, employee_id, entity, entity_id, action, details_json, actor, created_at
        FROM employee_audit
        WHERE employee_id = ?
        ORDER BY id DESC
        LIMIT ?
        """,
        (int(employee_id), int(limit)),
    )
    rows = cur.fetchall() or []
    out = []
    for r in rows:
        try:
            det = json.loads(r[5] or "{}")
        except Exception:
            det = {}
        out.append(
            {
                "id": int(r[0]),
                "employee_id": int(r[1]) if r[1] is not None else None,
                "entity": r[2],
                "entity_id": r[3],
                "action": r[4],
                "details": det,
                "actor": r[6],
                "created_at": r[7],
            }
        )
    return out

def _render_employee_docs_manager(
    conn: sqlite3.Connection,
    *,
    emp_id: int,
    doc_type: str,
    label: str,
    key_prefix: str,
    allow_edit: bool = True,
    bundle: dict | None = None,
    stacked: bool = False,
) -> None:
    """UI generic pentru gestionarea documentelor (upload/list/download/activ/delete). Dacă bundle e dat, folosește bundle['employee_documents']."""
    st.markdown(f"#### {label}")
    st.markdown("<div class='emp-thin-guide-line'></div>", unsafe_allow_html=True)

    actor = None
    try:
        actor = st.session_state.get("username") or st.session_state.get("user") or st.session_state.get("utilizator")
    except Exception:
        actor = None

    if allow_edit:
        if stacked:
            doc_no = st.text_input("Număr / ID document", value="", key=f"{key_prefix}docno_{doc_type}")
            doc_date = st.text_input("Data (dd.mm.yyyy)", value="", key=f"{key_prefix}docdate_{doc_type}")
            note = st.text_input("Observații (opțional)", value="", key=f"{key_prefix}note_{doc_type}")
        else:
            c1, c2, c3 = st.columns([1.4, 1.2, 2.4])
            with c1:
                doc_no = st.text_input("Număr / ID document", value="", key=f"{key_prefix}docno_{doc_type}")
            with c2:
                doc_date = st.text_input("Data (dd.mm.yyyy)", value="", key=f"{key_prefix}docdate_{doc_type}")
            with c3:
                note = st.text_input("Observații (opțional)", value="", key=f"{key_prefix}note_{doc_type}")

        up = st.file_uploader(
            "Încarcă document (PDF / DOCX / JPG / PNG)",
            type=["pdf", "docx", "jpg", "jpeg", "png"],
            key=f"{key_prefix}upl_{doc_type}",
        )
        if up is not None:
            if st.button("💾 Salvează document", key=f"{key_prefix}btn_save_{doc_type}"):
                try:
                    _bytes = up.getvalue()
                    meta = {"note": note}
                    new_id = _save_employee_document(
                        conn,
                        employee_id=int(emp_id),
                        doc_type=str(doc_type),
                        filename=str(up.name),
                        data=_bytes,
                        doc_no=(doc_no or None),
                        doc_date=(doc_date or None),
                        meta=meta,
                        set_active=True,
                    )
                    _audit_log(
                        conn,
                        employee_id=int(emp_id),
                        entity="employee_documents",
                        entity_id=int(new_id),
                        action="upload",
                        details={"doc_type": doc_type, "filename": up.name, "doc_no": doc_no, "doc_date": doc_date},
                        actor=actor,
                    )
                    st.success("Document salvat (marcat ca activ pentru tipul lui).")
                    st.rerun()
                except Exception as e:
                    st.error(f"Eroare la salvare: {e}")

    st.divider()

    if bundle is not None:
        docs = [d for d in bundle.get("employee_documents", []) if d.get("doc_type") == doc_type]
    else:
        docs = _list_employee_documents(conn, employee_id=int(emp_id), doc_type=str(doc_type))
    if not docs:
        st.info("Nu există documente încă.")
        return

    df = pd.DataFrame(docs).rename(
        columns={
            "id": "ID",
            "doc_no": "Nr",
            "doc_date": "Data",
            "filename": "Fișier",
            "created_at": "Creat la",
            "is_active": "Activ",
        }
    )
    show_cols = ["ID", "Nr", "Data", "Fișier", "Creat la", "Activ"]
    for c in show_cols:
        if c not in df.columns:
            df[c] = None
    st.dataframe(df[show_cols], use_container_width=True, height=220)

    ids = [int(x["id"]) for x in docs]
    active_id = None
    for x in docs:
        if int(x.get("is_active") or 0) == 1:
            active_id = int(x["id"])
            break
    if active_id is None:
        active_id = ids[0]

    sel = st.selectbox(
        "Alege document (pentru descărcare / activare / ștergere)",
        options=ids,
        index=ids.index(active_id) if active_id in ids else 0,
        key=f"{key_prefix}sel_{doc_type}",
    )

    fn, blob = _get_employee_document_bytes(conn, doc_id=int(sel))
    if stacked:
        if blob is not None:
            st.download_button(
                "⬇️ Descarcă",
                data=blob,
                file_name=fn or f"{doc_type}.bin",
                mime="application/octet-stream",
                key=f"{key_prefix}dl_{doc_type}",
            )
        else:
            st.button("⬇️ Descarcă", disabled=True, key=f"{key_prefix}dl_dis_{doc_type}")
        if allow_edit:
            if st.button("⭐ Setează activ", key=f"{key_prefix}setact_{doc_type}"):
                try:
                    _employee_document_set_active(conn, doc_id=int(sel))
                    _audit_log(
                        conn,
                        employee_id=int(emp_id),
                        entity="employee_documents",
                        entity_id=int(sel),
                        action="set_active",
                        details={"doc_type": doc_type, "doc_id": int(sel)},
                        actor=actor,
                    )
                    st.success("Setat ca activ.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Eroare: {e}")
        else:
            st.button("⭐ Setează activ", disabled=True, key=f"{key_prefix}setact_dis_{doc_type}")
        if allow_edit:
            if st.button("🗑 Șterge", key=f"{key_prefix}del_{doc_type}"):
                try:
                    _employee_document_delete(conn, doc_id=int(sel))
                    _audit_log(
                        conn,
                        employee_id=int(emp_id),
                        entity="employee_documents",
                        entity_id=int(sel),
                        action="delete",
                        details={"doc_type": doc_type, "doc_id": int(sel), "filename": fn},
                        actor=actor,
                    )
                    st.success("Șters.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Eroare: {e}")
        else:
            st.button("🗑 Șterge", disabled=True, key=f"{key_prefix}del_dis_{doc_type}")
    else:
        b1, b2, b3 = st.columns([1.2, 1.2, 1.2])
        with b1:
            if blob is not None:
                st.download_button(
                    "⬇️ Descarcă",
                    data=blob,
                    file_name=fn or f"{doc_type}.bin",
                    mime="application/octet-stream",
                    key=f"{key_prefix}dl_{doc_type}",
                )
            else:
                st.button("⬇️ Descarcă", disabled=True, key=f"{key_prefix}dl_dis_{doc_type}")

        with b2:
            if allow_edit:
                if st.button("⭐ Setează activ", key=f"{key_prefix}setact_{doc_type}"):
                    try:
                        _employee_document_set_active(conn, doc_id=int(sel))
                        _audit_log(
                            conn,
                            employee_id=int(emp_id),
                            entity="employee_documents",
                            entity_id=int(sel),
                            action="set_active",
                            details={"doc_type": doc_type, "doc_id": int(sel)},
                            actor=actor,
                        )
                        st.success("Setat ca activ.")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Eroare: {e}")
            else:
                st.button("⭐ Setează activ", disabled=True, key=f"{key_prefix}setact_dis_{doc_type}")

        with b3:
            if allow_edit:
                if st.button("🗑 Șterge", key=f"{key_prefix}del_{doc_type}"):
                    try:
                        _employee_document_delete(conn, doc_id=int(sel))
                        _audit_log(
                            conn,
                            employee_id=int(emp_id),
                            entity="employee_documents",
                            entity_id=int(sel),
                            action="delete",
                            details={"doc_type": doc_type, "doc_id": int(sel), "filename": fn},
                            actor=actor,
                        )
                        st.success("Șters.")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Eroare: {e}")
            else:
                st.button("🗑 Șterge", disabled=True, key=f"{key_prefix}del_dis_{doc_type}")




# -------------------------------------------------------------
# Detașări (Art.45 Codul muncii) – tabel EDITABIL + DOCX generat din cod
# -------------------------------------------------------------
def _safe_json_load(s: str | None) -> dict:
    try:
        if not s:
            return {}
        return json.loads(s)
    except Exception:
        return {}

def _safe_json_dump(d: dict) -> str:
    try:
        return json.dumps(d or {}, ensure_ascii=False)
    except Exception:
        return "{}"

def build_detasare_docx_bytes(payload: dict) -> bytes:
    """Generează un DOCX simplu pentru Detașare (fără șablon)."""
    import io
    from docx import Document

    doc = Document()

    _docx_apply_antet_si_semnaturi(doc)

    title = (payload.get("titlu") or "DECIZIE DE DETAȘARE").strip()
    doc.add_heading(title, level=1)

    # date angajator
    doc.add_paragraph(f"Angajator (expeditor): {payload.get('angajator_expeditor','')}")
    doc.add_paragraph(f"Angajator (primitor): {payload.get('angajator_primitor','')}")
    doc.add_paragraph(f"Salariat: {payload.get('salariat','')}")
    doc.add_paragraph(f"CNP: {payload.get('cnp','')}")

    doc.add_paragraph("")
    doc.add_paragraph(f"Nr.: {payload.get('nr','')} / Data: {payload.get('data_decizie','')}")
    doc.add_paragraph("")

    # temei legal + text minim
    doc.add_paragraph("Temei legal: art. 45 din Codul muncii (detașarea).")
    if payload.get("temei_legal"):
        doc.add_paragraph(f"Temei suplimentar: {payload.get('temei_legal')}")

    doc.add_paragraph("")
    doc.add_paragraph("Se dispune detașarea salariatului în interesul angajatorului primitor, astfel:")
    doc.add_paragraph(f"- Perioada: {payload.get('data_inceput','')} – {payload.get('data_sfarsit','')}")
    if payload.get("loc_munca"):
        doc.add_paragraph(f"- Loc de muncă: {payload.get('loc_munca')}")
    if payload.get("functie"):
        doc.add_paragraph(f"- Funcția: {payload.get('functie')}")
    if payload.get("motiv"):
        doc.add_paragraph(f"- Motiv: {payload.get('motiv')}")
    if payload.get("salariu_baza") not in (None, ""):
        doc.add_paragraph(f"- Salariu de bază (informativ): {payload.get('salariu_baza')}")

    if payload.get("schimbare_fel_munca") is True:
        doc.add_paragraph("")
        doc.add_paragraph("Notă: Se modifică și felul muncii pe durata detașării, cu consimțământul scris al salariatului.")

    doc.add_paragraph("")
    doc.add_paragraph("Prezenta decizie produce efecte începând cu data menționată mai sus.")
    doc.add_paragraph("")
    doc.add_paragraph("Angajator,                                Salariat,")
    doc.add_paragraph("_____________________                     _____________________")

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _detasari_list(conn: sqlite3.Connection, employee_id: int, bundle: dict | None = None) -> list[dict]:
    if bundle is not None:
        rows = [d for d in bundle.get("employee_documents", []) if d.get("doc_type") == "DETASARE"]
    else:
        rows = _list_employee_documents(conn, employee_id=int(employee_id), doc_type="DETASARE") or []
    out = []
    for r in rows:
        meta = _safe_json_load(r.get("meta_json"))
        out.append({
            "ID": int(r.get("id")),
            "Data acțiune": meta.get("data_actiune") or (r.get("created_at") or ""),
            "Data început": meta.get("data_inceput") or "",
            "Data sfârșit": meta.get("data_sfarsit") or "",
            "Angajator primitor": meta.get("angajator_primitor") or "",
            "Funcție": meta.get("functie") or "",
            "Loc muncă": meta.get("loc_munca") or "",
            "Nr. decizie": r.get("doc_no") or meta.get("nr") or "",
            "Data decizie": r.get("doc_date") or meta.get("data_decizie") or "",
            "Motiv": meta.get("motiv") or "",
            "Temei legal": meta.get("temei_legal") or "art. 45 Codul muncii",
            "Schimbare fel muncă": bool(meta.get("schimbare_fel_munca") or False),
            "Salariu de bază": meta.get("salariu_baza") or "",
            "Activ": bool(int(r.get("is_active") or 0) == 1),
            "Șterge": False,
        })
    return out


def _detasari_save_updates(conn: sqlite3.Connection, employee_id: int, edited_rows: list[dict], actor: str | None = None):
    # Update rows based on ID; delete those marked
    for row in edited_rows:
        try:
            doc_id = int(row.get("ID"))
        except Exception:
            continue

        if bool(row.get("Șterge")):
            try:
                _employee_document_delete(conn, doc_id=doc_id)
                _audit_log(
                    conn,
                    employee_id=int(employee_id),
                    entity="employee_documents",
                    entity_id=int(doc_id),
                    action="delete",
                    details={"doc_type": "DETASARE", "doc_id": doc_id},
                    actor=actor,
                )
            except Exception:
                pass
            continue

        # set active?
        if bool(row.get("Activ")):
            try:
                _employee_document_set_active(conn, doc_id=doc_id)
                _audit_log(
                    conn,
                    employee_id=int(employee_id),
                    entity="employee_documents",
                    entity_id=int(doc_id),
                    action="set_active",
                    details={"doc_type": "DETASARE", "doc_id": doc_id},
                    actor=actor,
                )
            except Exception:
                pass

        # update meta_json + doc_no/doc_date
        meta = {
            "data_actiune": (row.get("Data acțiune") or ""),
            "data_inceput": (row.get("Data început") or ""),
            "data_sfarsit": (row.get("Data sfârșit") or ""),
            "angajator_primitor": (row.get("Angajator primitor") or ""),
            "functie": (row.get("Funcție") or ""),
            "loc_munca": (row.get("Loc muncă") or ""),
            "motiv": (row.get("Motiv") or ""),
            "temei_legal": (row.get("Temei legal") or ""),
            "schimbare_fel_munca": bool(row.get("Schimbare fel muncă") or False),
            "salariu_baza": (row.get("Salariu de bază") or ""),
            "nr": (row.get("Nr. decizie") or ""),
            "data_decizie": (row.get("Data decizie") or ""),
        }

        try:
            cur = conn.cursor()
            cur.execute(
                "UPDATE employee_documents SET doc_no=?, doc_date=?, meta_json=? WHERE id=? AND employee_id=? AND doc_type='DETASARE'",
                (
                    (row.get("Nr. decizie") or None),
                    (row.get("Data decizie") or None),
                    _safe_json_dump(meta),
                    int(doc_id),
                    int(employee_id),
                ),
            )
            conn.commit()
            _audit_log(
                conn,
                employee_id=int(employee_id),
                entity="employee_documents",
                entity_id=int(doc_id),
                action="update",
                details={"doc_type": "DETASARE", "fields": {k: meta.get(k) for k in ["data_inceput","data_sfarsit","angajator_primitor","functie","loc_munca","motiv","temei_legal","salariu_baza"]}},
                actor=actor,
            )
        except Exception:
            pass


def render_detasari_editabile(conn: sqlite3.Connection, emp_id: int, emp: dict, *, allow_edit: bool = True, key_prefix: str = "", bundle: dict | None = None):
    st.markdown("#### Detașări / delegări – registru editabil (art. 45 Codul muncii)")

    actor = None
    try:
        actor = st.session_state.get("username") or st.session_state.get("user") or st.session_state.get("utilizator")
    except Exception:
        actor = None

    # ---- Adaugă detașare (fără șablon) -> generăm DOCX minim și îl salvăm la employee_documents (doc_type=DETASARE)
    if allow_edit:
        st.markdown('<div class="emp-subsection-title">➕ Adaugă detașare (generează DOCX automat)</div>', unsafe_allow_html=True)
        st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)
        st.markdown("<div class='emp-thin-guide-line'></div>", unsafe_allow_html=True)
        st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)
        nr = st.text_input("Nr. decizie", key=f"{key_prefix}det_add_nr_{emp_id}")
        data_decizie = st.text_input("Data decizie (dd.mm.yyyy)", key=f"{key_prefix}det_add_datadec_{emp_id}")
        data_actiune = st.text_input("Data acțiune (dd.mm.yyyy) (opțional)", key=f"{key_prefix}det_add_dataact_{emp_id}")
        data_inceput = st.text_input("Data început (dd.mm.yyyy)", key=f"{key_prefix}det_add_start_{emp_id}")
        data_sfarsit = st.text_input("Data sfârșit (dd.mm.yyyy)", key=f"{key_prefix}det_add_end_{emp_id}")
        salariu_baza = st.text_input("Salariu de bază (opțional)", key=f"{key_prefix}det_add_sal_{emp_id}")
        ang_prim = st.text_input("Angajator primitor", key=f"{key_prefix}det_add_prim_{emp_id}")
        functie = st.text_input("Funcție (opțional)", key=f"{key_prefix}det_add_func_{emp_id}")
        loc_munca = st.text_input("Loc de muncă (opțional)", key=f"{key_prefix}det_add_loc_{emp_id}")

        temei_legal = st.text_input("Temei legal (opțional)", value="art. 45 Codul muncii", key=f"{key_prefix}det_add_temei_{emp_id}")
        motiv = st.text_area("Motiv / interesul serviciului", key=f"{key_prefix}det_add_motiv_{emp_id}", height=90)
        schimbare_fel = st.checkbox("Se modifică și felul muncii (necesită consimțământ scris)", key=f"{key_prefix}det_add_fel_{emp_id}")

        if st.button("💾 Salvează detașare (DOCX)", key=f"{key_prefix}det_add_save_{emp_id}"):
            if not ang_prim.strip():
                st.warning("Completează «Angajator primitor».")
            elif not data_inceput.strip():
                st.warning("Completează «Data început».")
            else:
                salariat = f"{emp.get('nume','')} {emp.get('prenume','')}"
                payload = {
                    "titlu": "DECIZIE DE DETAȘARE",
                    "nr": nr,
                    "data_decizie": data_decizie,
                    "data_actiune": data_actiune,
                    "data_inceput": data_inceput,
                    "data_sfarsit": data_sfarsit,
                    "angajator_expeditor": "",
                    "angajator_primitor": ang_prim,
                    "salariat": salariat.strip(),
                    "cnp": emp.get("cnp",""),
                    "functie": functie,
                    "loc_munca": loc_munca,
                    "motiv": motiv,
                    "temei_legal": temei_legal,
                    "schimbare_fel_munca": bool(schimbare_fel),
                    "salariu_baza": salariu_baza,
                }
                docx_bytes = build_detasare_docx_bytes(payload)
                new_id = _save_employee_document(
                    conn,
                    employee_id=int(emp_id),
                    doc_type="DETASARE",
                    filename=f"Detasare_{nr or 'nou'}.docx",
                    data=docx_bytes,
                    doc_no=(nr or None),
                    doc_date=(data_decizie or None),
                    meta=payload,
                    set_active=True,
                )
                _audit_log(
                    conn,
                    employee_id=int(emp_id),
                    entity="employee_documents",
                    entity_id=int(new_id),
                    action="upload",
                    details={"doc_type": "DETASARE", "nr": nr, "data_decizie": data_decizie, "angajator_primitor": ang_prim},
                    actor=actor,
                )
                st.success("Detașare salvată.")
                st.rerun()

    st.divider()

    rows = _detasari_list(conn, employee_id=int(emp_id), bundle=bundle)
    if not rows:
        st.info("Nu există detașări încă.")
        return

    df = pd.DataFrame(rows)

    # ordonare coloane
    cols_order = [
        "ID",
        "Data acțiune",
        "Data decizie",
        "Nr. decizie",
        "Angajator primitor",
        "Data început",
        "Data sfârșit",
        "Funcție",
        "Loc muncă",
        "Motiv",
        "Temei legal",
        "Schimbare fel muncă",
        "Salariu de bază",
        "Activ",
        "Șterge",
    ]
    for c in cols_order:
        if c not in df.columns:
            df[c] = ""

    if allow_edit:
        edited = st.data_editor(
            df[cols_order],
            use_container_width=True,
            height=360,
            hide_index=True,
            disabled=["ID"],
            column_config={
                "Activ": st.column_config.CheckboxColumn("Activ"),
                "Șterge": st.column_config.CheckboxColumn("Șterge"),
                "Schimbare fel muncă": st.column_config.CheckboxColumn("Schimbare fel muncă"),
            },
            key=f"{key_prefix}det_grid_{emp_id}",
        )

        if st.button("💾 Salvează modificările", key=f"{key_prefix}det_save_{emp_id}"):
            _detasari_save_updates(conn, employee_id=int(emp_id), edited_rows=edited.to_dict(orient="records"), actor=actor)
            st.success("Actualizat.")
            st.rerun()

    else:
        st.dataframe(df[cols_order], use_container_width=True, height=360)

    # download select
    ids = [int(r["ID"]) for r in rows]
    sel = st.selectbox("Descarcă documentul de detașare", options=ids, index=0, key=f"{key_prefix}det_sel_dl_{emp_id}")
    fn, blob = _get_employee_document_bytes(conn, doc_id=int(sel))
    if blob:
        st.download_button("⬇️ Descarcă DOCX", data=blob, file_name=fn or "detasare.docx", key=f"{key_prefix}det_dl_{emp_id}")


def render_documente_informatii_salariat(conn: sqlite3.Connection, emp_id: int, emp: dict, bundle: dict | None = None):
    st.subheader("📎 Documente / Informații salariat")

    # Permisiuni (fallback): dacă ai implementat login/roluri, poți seta st.session_state["role"]
    role = None
    try:
        role = st.session_state.get("role")
    except Exception:
        role = None
    allow_edit_docs = True if role is None else (str(role).lower() in ("admin", "hr", "manager"))

    subt = st.tabs(["Decizii", "Detașări", "Pregătire profesională", "Alte documente", "Jurnal modificări"])

    with subt[0]:
        with st.container(key=f"emp_decizii_stack_{emp_id}"):
            # --- Generator decizie (DOCX) ---
            st.markdown("### 🧾 Generează decizie (DOCX)")
            st.caption("Completezi elementele necesare, iar aplicația generează documentul și îl salvează automat la «Decizii».")

            preset = st.selectbox(
                "Tip decizie (preset)",
                [
                    "Detașare (art.45)",
                    "Încetare CIM - concediere individuală (art.65)",
                    "Încetare CIM - acordul părților (art.55 lit.b)",
                    "Încetare CIM - demisie (art.81)",
                    "Suspendare de drept (art.51 lit.a)",
                    "Întrerupere activitate (art.52 lit.c)",
                    "Desfacere disciplinară (art.61 lit.a)",
                    "Inapt fizic/psihic (art.61 lit.c)",
                    "Pensionare (art.56 lit.c)",
                    "Necorespundere profesională (art.61 lit.d)",
                    "Reintegrare (art.56 lit.e)",
                    "Mandat executare (art.56 lit.f)",
                    "Altă decizie (generic)",
                ],
                key=f"dec_gen_preset_{emp_id}",
            )

            dec_nr = st.text_input("Nr. decizie", key=f"dec_gen_nr_{emp_id}")
            dec_data = st.text_input("Data (ex: 08.01.2026)", key=f"dec_gen_data_{emp_id}")
            admin_name = st.text_input("Administrator / reprezentant", key=f"dec_gen_admin_{emp_id}")
            executor = st.text_input("Persoană responsabilă (opțional)", key=f"dec_gen_exec_{emp_id}")
            tribunal = st.text_input("Instanță/Tribunal (opțional)", key=f"dec_gen_trib_{emp_id}")

            # Date angajator (din Config dacă există; fallback: câmpuri editabile)
            st.markdown('<div class="emp-subsection-title">Date angajator (pentru antet)</div>', unsafe_allow_html=True)
            company_name = st.text_input(
                "Denumire angajator",
                value=str(st.session_state.get("CFG_ANGAJATOR_DEN", "") or ""),
                key=f"dec_gen_cname_{emp_id}",
            )
            company_sediu = st.text_input(
                "Sediu",
                value=str(st.session_state.get("CFG_ANGAJATOR_SEDIU", "") or ""),
                key=f"dec_gen_csed_{emp_id}",
            )
            company_reg = st.text_input(
                "Reg.Com.",
                value=str(st.session_state.get("CFG_ANGAJATOR_REGCOM", "") or ""),
                key=f"dec_gen_creg_{emp_id}",
            )

            # Salariat (din fișa curentă)
            emp_name = f"{emp.get('nume','') or ''} {emp.get('prenume','') or ''}".strip()
            emp_functie = str(emp.get("functie","") or emp.get("functia","") or emp.get("functie_cor","") or "").strip()

            st.markdown("### Conținut decizie")
            titlu = st.text_input("Titlu document", value=preset.upper(), key=f"dec_gen_title_{emp_id}")

            preset_considerente = ""
            preset_art1 = ""
            preset_art2 = "Cu ducerea la îndeplinire a prezentei se însărcinează persoanele desemnate și se comunică celui în cauză."

            if preset.startswith("Detașare"):
                unitate = st.text_input("Unitate primitoare (SC ...)", key=f"dec_gen_unit_{emp_id}")
                perioada = st.text_input("Perioadă (ex: 30 zile / 01.02.2026-01.03.2026)", key=f"dec_gen_per_{emp_id}")
                data_start = st.text_input("Data începerii detașării", key=f"dec_gen_start_{emp_id}")
                preset_considerente = "În baza prevederilor art.45 din Codul Muncii;\nÎn temeiul prevederilor Legii nr.31/1990 și a prerogativelor stabilite prin actul constitutiv al societății,"
                preset_art1 = f"Începând cu data de {data_start} salariatul/ salariata {emp_name} având funcția de {emp_functie} se detașează la {unitate} pe o perioadă de {perioada}."
            elif "concediere individuală" in preset:
                motiv = st.text_area("Motive (pe scurt)", key=f"dec_gen_motiv_{emp_id}", height=80)
                data_inc = st.text_input("Data încetării", key=f"dec_gen_inc_{emp_id}")
                preset_considerente = f"În baza prevederilor art.65(1) din Codul Muncii;\nAvând în vedere desființarea locului de muncă ocupat de {emp_name} ca urmare a: {motiv}"
                preset_art1 = f"Începând cu data de {data_inc} încetează contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.65(1) din Codul Muncii."
            elif "acordul părților" in preset:
                cerere = st.text_input("Nr./data cerere (opțional)", key=f"dec_gen_cer_{emp_id}")
                data_inc = st.text_input("Data încetării", key=f"dec_gen_inc2_{emp_id}")
                preset_considerente = f"Având în vedere cererea {cerere} prin care {emp_name} solicită încetarea CIM prin acordul părților;\nÎn baza art.55 lit.b din Codul Muncii;"
                preset_art1 = f"Începând cu data de {data_inc} încetează contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.55 lit.b din Codul Muncii."
            elif "demisie" in preset:
                nr_dem = st.text_input("Nr./data demisie (opțional)", key=f"dec_gen_dem_{emp_id}")
                data_inc = st.text_input("Data încetării", key=f"dec_gen_inc3_{emp_id}")
                preset_considerente = f"Având în vedere demisia {nr_dem} prin care {emp_name} denunță unilateral CIM;\nÎn baza art.81 din Codul Muncii;"
                preset_art1 = f"Începând cu data de {data_inc} încetează contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.81 din Codul Muncii."
            elif "Suspendare de drept" in preset:
                data_susp = st.text_input("Data începerii suspendării", key=f"dec_gen_susp_{emp_id}")
                motiv = st.text_area("Motiv (ex: concediu creștere copil etc.)", key=f"dec_gen_suspm_{emp_id}", height=80)
                preset_considerente = "În baza prevederilor art.51 lit.a din Codul Muncii;"
                preset_art1 = f"Începând cu data de {data_susp} se suspendă contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.51 lit.a din Codul Muncii, pentru: {motiv}."
            elif "Întrerupere activitate" in preset:
                data_susp = st.text_input("Data începerii suspendării", key=f"dec_gen_intr_{emp_id}")
                motiv = st.text_area("Motive întrerupere activitate", key=f"dec_gen_intrm_{emp_id}", height=80)
                proc = st.text_input("Indemnizație % (opțional)", key=f"dec_gen_intrp_{emp_id}")
                suma = st.text_input("Suma indemnizație (opțional)", key=f"dec_gen_intrs_{emp_id}")
                preset_considerente = "În baza prevederilor art.52 lit.c și art.53 din Codul Muncii;"
                preset_art1 = f"Începând cu data de {data_susp} se suspendă contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.52 lit.c din Codul Muncii, ca urmare a: {motiv}."
                if proc or suma:
                    preset_art2 = f"Salariatul/ salariata beneficiază de indemnizație de {proc}% în sumă de {suma} lei (după caz).\n" + preset_art2
            elif "Desfacere disciplinară" in preset:
                data_inc = st.text_input("Data încetării", key=f"dec_gen_disc_{emp_id}")
                fapta = st.text_area("Abaterea disciplinară (descriere)", key=f"dec_gen_fapta_{emp_id}", height=90)
                preset_considerente = "În baza art.61 lit.a, art.251-252 din Codul Muncii;"
                preset_art1 = f"Începând cu data de {data_inc} se desface disciplinar contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.61 lit.a din Codul Muncii, pentru: {fapta}."
            elif "Inapt fizic/psihic" in preset:
                data_inc = st.text_input("Data încetării", key=f"dec_gen_inapt_{emp_id}")
                det = st.text_area("Detalii (decizie medicală etc.)", key=f"dec_gen_inaptd_{emp_id}", height=90)
                preset_considerente = "În baza art.61 lit.c, art.64 și art.76 din Codul Muncii;"
                preset_art1 = f"Începând cu data de {data_inc} încetează contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.61 lit.c din Codul Muncii. {det}"
            elif "Pensionare" in preset:
                data_inc = st.text_input("Data încetării", key=f"dec_gen_pens_{emp_id}")
                det = st.text_input("Nr./data decizie pensionare (opțional)", key=f"dec_gen_pensd_{emp_id}")
                preset_considerente = "În baza prevederilor art.56 lit.c din Codul Muncii;"
                preset_art1 = f"Începând cu data de {data_inc} încetează contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.56 lit.c din Codul Muncii. {det}"
            elif "Necorespundere profesională" in preset:
                data_inc = st.text_input("Data încetării", key=f"dec_gen_nec_{emp_id}")
                det = st.text_area("Detalii constatări/evaluare", key=f"dec_gen_necd_{emp_id}", height=90)
                preset_considerente = "În baza art.61 lit.d, art.64 și art.76 din Codul Muncii;"
                preset_art1 = f"Începând cu data de {data_inc} încetează contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.61 lit.d din Codul Muncii. {det}"
            elif "Reintegrare" in preset:
                data_inc = st.text_input("Data efect reintegrare", key=f"dec_gen_reint_{emp_id}")
                det = st.text_area("Detalii hotărâre judecătorească", key=f"dec_gen_reintd_{emp_id}", height=90)
                preset_considerente = "În baza prevederilor art.56 lit.e din Codul Muncii;"
                preset_art1 = f"Începând cu data de {data_inc} se dispune reintegrarea salariatului/salariatei, conform hotărârii judecătorești. {det}"
            elif "Mandat executare" in preset:
                data_inc = st.text_input("Data încetării", key=f"dec_gen_mand_{emp_id}")
                det = st.text_input("Nr./data mandat (opțional)", key=f"dec_gen_mandd_{emp_id}")
                preset_considerente = "În baza art.56 lit.f din Codul Muncii și a prevederilor Codului de procedură penală;"
                preset_art1 = f"Începând cu data de {data_inc} încetează contractul individual de muncă al salariatului/salariatei {emp_name}, având funcția de {emp_functie}, conform art.56 lit.f din Codul Muncii. {det}"

            considerente = st.text_area(
                "Considerente (poți edita)",
                value=preset_considerente,
                key=f"dec_gen_cons_{emp_id}",
                height=120,
            )

            a1 = st.text_area("Art.1 (poți edita)", value=preset_art1, key=f"dec_gen_a1_{emp_id}", height=90)
            a2 = st.text_area("Art.2 (poți edita)", value=preset_art2, key=f"dec_gen_a2_{emp_id}", height=80)
            a3 = st.text_area("Art.3 (opțional)", value="", key=f"dec_gen_a3_{emp_id}", height=70)

            if st.button("🧾 Generează DOCX", key=f"dec_gen_btn_{emp_id}"):
                payload = {
                    "COMPANY_NAME": company_name,
                    "COMPANY_SEDIU": company_sediu,
                    "COMPANY_REGCOM": company_reg,
                    "DEC_TITLU": titlu or "DECIZIE",
                    "DEC_NR": dec_nr,
                    "DEC_DATA": dec_data,
                    "ADMIN_NAME": admin_name,
                    "EMP_NAME": emp_name,
                    "EMP_FUNCTIE": emp_functie,
                    "CONSIDERENTE": considerente,
                    "ARTICLES": [a1, a2, a3],
                    "EXECUTOR": executor,
                    "TRIBUNAL": tribunal,
                }
                docx_bytes = build_decizie_generic_docx_bytes(payload)
                fn = f"Decizie_{(dec_nr or '').strip() or 'NR'}_{emp_name.replace(' ','_')}.docx"

                try:
                    new_id = _save_employee_document(
                        conn,
                        employee_id=int(emp_id),
                        doc_type="DECIZIE",
                        filename=fn,
                        data=docx_bytes,
                        doc_no=str(dec_nr or "").strip() or None,
                        doc_date=str(dec_data or "").strip() or None,
                        meta={
                            "preset": preset,
                            "titlu": titlu,
                            "emp_name": emp_name,
                            "emp_functie": emp_functie,
                        },
                        set_active=True,
                    )
                    st.success(f"Decizie generată și salvată (ID {new_id}).")
                except Exception as e:
                    st.error(f"Eroare la salvare în DB: {e}")

                st.download_button(
                    "⬇️ Descarcă DOCX",
                    data=docx_bytes,
                    file_name=fn,
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    key=f"dec_gen_dl_{emp_id}",
                )

        st.divider()

        _render_employee_docs_manager(
            conn,
            emp_id=int(emp_id),
            doc_type="DECIZIE",
            label="Decizii (numiri, sancțiuni, încetări, etc.)",
            key_prefix=f"ang_docs_{emp_id}_dec_",
            allow_edit=allow_edit_docs,
            bundle=bundle,
            stacked=True,
        )

    with subt[1]:
        render_detasari_editabile(
            conn,
            emp_id=int(emp_id),
            emp=emp,
            allow_edit=allow_edit_docs,
            key_prefix=f"ang_docs_{emp_id}_det_edit_",
            bundle=bundle,
        )
    with subt[2]:
        _render_employee_docs_manager(
            conn,
            emp_id=int(emp_id),
            doc_type="PREGATIRE",
            label="Pregătire profesională (cursuri, certificate, atestări)",
            key_prefix=f"ang_docs_{emp_id}_preg_",
            allow_edit=allow_edit_docs,
            bundle=bundle,
            stacked=True,
        )

    with subt[3]:
        _render_employee_docs_manager(
            conn,
            emp_id=int(emp_id),
            doc_type="ALTE",
            label="Alte documente / informații salariat",
            key_prefix=f"ang_docs_{emp_id}_alte_",
            allow_edit=allow_edit_docs,
            bundle=bundle,
            stacked=True,
        )

    with subt[4]:
        st.markdown("#### Jurnal (audit) – încărcări / ștergeri / activări")
        st.markdown("<div class='emp-thin-guide-line'></div>", unsafe_allow_html=True)
        rows = _audit_list(conn, employee_id=int(emp_id), limit=300)
        if not rows:
            st.info("Nu există evenimente în jurnal încă.")
        else:
            df_a = pd.DataFrame(
                [
                    {
                        "Creat la": r.get("created_at"),
                        "Acțiune": r.get("action"),
                        "Entitate": r.get("entity"),
                        "ID entitate": r.get("entity_id"),
                        "Actor": r.get("actor"),
                        "Detalii": json.dumps(r.get("details") or {}, ensure_ascii=False),
                    }
                    for r in rows
                ]
            )
            st.dataframe(df_a, use_container_width=True, height=360)

def page_angajati(conn: sqlite3.Connection):
    """
    Modul Angajați:
    - listă minimă (Nume, Prenume, CNP, Marcă) + filtrare
    - butoane: ADAUGĂ / MODIFICĂ / ȘTERGE
    - selectare angajat -> se deschide pagina DOAR cu acel angajat (detaliu complet)
    """
    # Marker + CSS local: scoped doar pentru pagina Angajați (pentru titluri, filtre etc.)
    st.markdown('<span id="emp-scope"></span>', unsafe_allow_html=True)
    def _inject_angajati_page_spacing_fix() -> None:
        st.markdown(
            """
            <style>
            /* Angajați: fără spații artificiale înainte de conținut */
            section.main:has(#emp-scope) .block-container{
              padding-top: 0 !important;
              margin-top: 0 !important;
            }
            section.main:has(#emp-scope) .block-container > div[data-testid="stVerticalBlock"]{
              margin-top: 4px !important;
            }
            section.main:has(#emp-scope) .block-container > div:first-child{
              margin-top: 0 !important;
              padding-top: 0 !important;
            }
            .ang-page-title{
              margin-top: 0 !important;
              padding-top: 0 !important;
              margin-bottom: 0.35rem !important;
            }
            .ang-gap-xs{ height: 0.3rem; }
            .ang-gap-sm{ height: 0.5rem; }
            .ang-gap-md{ height: 0.7rem; }
            </style>
            """,
            unsafe_allow_html=True,
        )
    _inject_angajati_page_spacing_fix()
    st.markdown(
        """
        <style>
        /* ===== Scoped doar pe pagina Angajați ===== */
        section.main:has(#emp-scope) .emp-title{
          font-size: 40px;          /* baza pentru icon + container */
          font-weight: 900;
          margin: 0 0 6px 0;
        }
        /* cuvântul ANGAJAȚI – foarte mare, uppercase */
        section.main:has(#emp-scope) .emp-title-main{
          display: inline-block;
          font-size: 56px !important;
          font-weight: 900 !important;
          text-transform: uppercase;
          letter-spacing: 0.12em;
          color: #F9FAFB !important;
        }

        section.main:has(#emp-scope) .emp-subtitle{
          font-size: 26px;
          font-weight: 600;
          opacity: .95;
          margin: 10px 0 10px 0;
        }

        /* Wrapper general Angajați – card discret, colțuri ușor rotunjite (ca firstapp) */
        section.main:has(#emp-scope) .emp-wrap{
          background: rgba(15,23,42,0.14);
          border: 1px solid rgba(148,163,184,0.12);
          border-radius: 6px;
          padding: 18px 22px;
          margin-bottom: 18px;
        }

        /* Divider subțire */
        section.main:has(#emp-scope) .thin-divider{
          width: 100%;
          height: 1px;
          border-radius: 999px;
          margin: 10px 0 14px 0;
          background: rgba(255,255,255,0.10);
        }

        /* Fără bară între butonul „Adaugă personal” și „Caută rapid” */
        section.main:has(#emp-scope) .emp-wrap hr,
        section.main:has(#emp-scope) .emp-wrap [data-testid="stHorizontalRule"]{
          display: none !important;
        }
        section.main:has(#emp-scope) .emp-wrap div:has(> .page-header),
        section.main:has(#emp-scope) .emp-wrap div:has(> .page-actions){
          border-bottom: none !important;
        }
        section.main:has(#emp-scope) .emp-wrap div:has(> .filter-card),
        section.main:has(#emp-scope) .emp-wrap [data-testid="stVerticalBlock"]:has(.filter-card){
          border-top: none !important;
        }
        /* Toate bordurile dintre blocuri în zona Angajați (listă) – elimină linia sub buton */
        section.main:has(#emp-scope) .emp-wrap [data-testid="stVerticalBlock"] > div,
        section.main:has(#emp-scope) .emp-wrap [data-testid="stVerticalBlockBorderWrapper"]{
          border-top: none !important;
          border-bottom: none !important;
        }
        /* Fără bară: cardul de filtre fără fundal/bordura, doar textul rămâne */
        section.main:has(#emp-scope) .filter-card{
          padding: 16px 18px !important;
          border: none !important;
          border-top: none !important;
          margin-top: -1px !important;
          background: transparent !important;
          box-shadow: none !important;
          backdrop-filter: none !important;
        }

        /* Search row: input + buton cu aceeași înălțime – colțuri moderate */
        section.main:has(#emp-scope) .filter-card .stTextInput>div>div input{
          height: 44px !important;
          border-radius: 6px !important;
        }
        section.main:has(#emp-scope) .filter-card .stButton>button{
          height: 44px !important;
          border-radius: 6px !important;
        }

        /* Expander mai “tight” */
        section.main:has(#emp-scope) .filter-card details > summary{
          padding: 10px 14px !important;
        }
        /* Caută rapid + Filtre avansate – spațiere și aliniere echilibrată */
        section.main:has(#emp-scope) .emp-wrap .filter-h3{
          margin: 0 0 10px 0 !important;
          padding: 0 !important;
          font-size: 14px !important;
          font-weight: 800 !important;
          opacity: .9;
        }
        section.main:has(#emp-scope) .emp-wrap .emp-search-title{
          margin-bottom: 12px !important;
        }
        /* Spațiere verticală între elemente */
        section.main:has(#emp-scope) .emp-wrap .emp-vgap{
          height: 14px;
          margin: 0;
          padding: 0;
        }
        section.main:has(#emp-scope) .emp-wrap .emp-vgap-section{
          height: 22px;
        }
        section.main:has(#emp-scope) .emp-wrap .emp-vgap-sm{
          height: 10px;
          margin: 0;
          padding: 0;
        }
        /* Aliniere: titluri și form la același nivel */
        section.main:has(#emp-scope) .emp-wrap [data-testid="stVerticalBlock"] > div,
        section.main:has(#emp-scope) .emp-wrap div[data-testid="stMarkdown"],
        section.main:has(#emp-scope) .emp-wrap [data-testid="stForm"] > div{
          padding-left: 0 !important;
          margin-left: 0 !important;
        }
        section.main:has(#emp-scope) .emp-wrap .stTextInput,
        section.main:has(#emp-scope) .emp-wrap .stButton{
          padding-left: 0 !important;
        }
        /* Spațiu între câmpurile din formular */
        section.main:has(#emp-scope) .emp-wrap [data-testid="stForm"] .stTextInput{
          margin-bottom: 4px;
        }
        /* Fără bare deasupra „Tabel angajați” și sub tabel */
        section.main:has(#emp-scope) .emp-wrap .emp-card.emp-card-table{
          border-top: none !important;
          border-bottom: none !important;
        }
        section.main:has(#emp-scope) .emp-wrap .emp-card.emp-card-table::before,
        section.main:has(#emp-scope) .emp-wrap .emp-card.emp-card-table::after{
          display: none !important;
        }
        section.main:has(#emp-scope) .emp-wrap [data-testid="stVerticalBlock"]:has([data-testid="stDataFrame"]){
          border-bottom: none !important;
        }
        section.main:has(#emp-scope) .emp-wrap div[data-testid="stDataFrame"]{
          border-bottom: none !important;
        }
        /* Bare încadrare tabel – subțiri sau eliminate */
        section.main:has(#emp-scope) .emp-wrap .emp-card.emp-card-table{
          border-left: 1px solid rgba(255,255,255,0.06) !important;
          border-right: 1px solid rgba(255,255,255,0.06) !important;
        }
        section.main:has(#emp-scope) .emp-wrap div[data-testid="stDataFrame"],
        section.main:has(#emp-scope) .emp-wrap div[data-testid="stDataFrame"] > div{
          border: 1px solid rgba(255,255,255,0.08) !important;
          border-radius: 6px !important;
          box-shadow: none !important;
        }
        section.main:has(#emp-scope) .emp-wrap [data-testid="stVerticalBlock"]:has([data-testid="stDataFrame"]) > div{
          border: none !important;
        }
        /* Fișa angajatului: header (nume + Activ/Inactiv) – colțuri moderate, ca firstapp */
        section.main:has(#emp-scope) .emp-detail-header{
          background: rgba(15,23,42,0.5);
          border: 1px solid rgba(255,255,255,0.12);
          border-radius: 6px;
          padding: 14px 18px;
          margin-bottom: 16px;
        }
        .st-key-emp_header_center_box{
          max-width: min(72vw, 760px) !important;
          margin-left: auto !important;
          margin-right: auto !important;
          padding-top: 6px !important;
          padding-bottom: 6px !important;
        }
        section.main:has(#emp-scope) .emp-detail-header .emp-header-center{
          text-align: center !important;
        }
        section.main:has(#emp-scope) .emp-detail-header .emp-header-name{
          font-size: 1.56rem !important;
          line-height: 1.2 !important;
          font-weight: 900 !important;
          color: rgba(248,250,252,0.98) !important;
          letter-spacing: 0.01em;
          margin: 0 !important;
        }
        section.main:has(#emp-scope) .emp-detail-header .emp-header-marca{
          font-size: 1.06rem !important;
          font-weight: 800 !important;
          color: rgba(226,232,240,0.96) !important;
          margin-top: 8px !important;
        }
        section.main:has(#emp-scope) .emp-detail-header .emp-header-state{
          font-size: 0.95rem !important;
          font-weight: 700 !important;
          margin-top: 10px !important;
          margin-bottom: 22px !important;
          color: rgba(148,163,184,0.96) !important;
        }
        .st-key-emp_header_actions [data-testid="stVerticalBlock"]{
          align-items: center !important;
        }
        .st-key-emp_header_actions [data-testid="stButton"]{
          width: 100% !important;
          margin-top: 8px !important;
        }
        .st-key-emp_header_actions [data-testid="stButton"]:first-of-type{
          margin-top: 8px !important;
        }
        .st-key-emp_header_actions [data-testid="stButton"] > button{
          width: 260px !important;
          min-width: 260px !important;
          max-width: 260px !important;
          min-height: 38px !important;
          height: 38px !important;
          margin-left: auto !important;
          margin-right: auto !important;
          display: block !important;
        }
        section.main:has(#emp-scope) .emp-detail-header .stSubheader,
        section.main:has(#emp-scope) .emp-detail-header [data-testid="stMarkdown"]{
          font-size: 1.25rem !important;
          font-weight: 800 !important;
          color: rgba(248,250,252,0.98) !important;
        }
        /* Tabs principale din fișa angajatului: Date salariat | Documente */
        .st-key-emp_main_tabs_box div[data-testid="stTabs"] > div[role="tablist"]{
          display: flex !important;
          gap: 0 !important;
          margin: 6px 0 14px 0 !important;
          padding: 0 !important;
          border-bottom: 1px solid rgba(148,163,184,0.28) !important;
        }
        .st-key-emp_main_tabs_box div[data-testid="stTabs"] button[role="tab"]{
          min-height: 38px !important;
          height: 38px !important;
          padding: 0 14px !important;
          border: none !important;
          border-radius: 0 !important;
          border-bottom: 2px solid transparent !important;
          background: transparent !important;
          color: rgba(226,232,240,0.88) !important;
          font-size: 0.94rem !important;
          font-weight: 700 !important;
          margin: 0 !important;
          transition: color .18s ease, border-color .18s ease !important;
        }
        .st-key-emp_main_tabs_box div[data-testid="stTabs"] button[role="tab"][aria-selected="true"]{
          background: transparent !important;
          color: #ffffff !important;
          border-bottom-color: rgba(56,189,248,0.92) !important;
          box-shadow: none !important;
        }

        /* Sub-tabs din Date salariat: pills compacte, rând dedicat, scroll orizontal fin */
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"]{
          display: flex !important;
          flex-wrap: nowrap !important;
          gap: 10px !important;
          overflow-x: auto !important;
          overflow-y: hidden !important;
          margin: 12px 0 22px 0 !important;
          padding: 2px 2px 10px 2px !important;
          border-bottom: 1px solid rgba(148,163,184,0.20) !important;
          scrollbar-width: thin;
          scrollbar-color: rgba(100,116,139,0.8) transparent;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"]::-webkit-scrollbar{
          height: 8px;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"]::-webkit-scrollbar-thumb{
          background: rgba(100,116,139,0.75);
          border-radius: 999px;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"]::-webkit-scrollbar-track{
          background: transparent;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] button[role="tab"]{
          flex: 0 0 auto !important;
          white-space: nowrap !important;
          min-height: 36px !important;
          height: 36px !important;
          padding: 0 14px !important;
          margin: 0 !important;
          border-radius: 11px !important;
          border: 1px solid rgba(148,163,184,0.35) !important;
          background: rgba(15,23,42,0.50) !important;
          color: rgba(226,232,240,0.88) !important;
          font-size: 0.87rem !important;
          font-weight: 650 !important;
          letter-spacing: 0.01em;
          transition: border-color .18s ease, background-color .18s ease, color .18s ease, box-shadow .18s ease !important;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] button[role="tab"]:hover{
          border-color: rgba(148,163,184,0.58) !important;
          background: rgba(30,41,59,0.64) !important;
          color: rgba(248,250,252,0.96) !important;
        }
        /* Mini-grup acțiuni: separare clară înainte de Edit/Șterge */
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"] > button[role="tab"]:nth-of-type(8){
          margin-left: 52px !important;
          border-color: rgba(148,163,184,0.52) !important;
          background: rgba(30,41,59,0.58) !important;
          color: rgba(248,250,252,0.96) !important;
          font-weight: 700 !important;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"] > button[role="tab"]:nth-of-type(9){
          border-color: rgba(248,113,113,0.44) !important;
          background: rgba(127,29,29,0.24) !important;
          color: rgba(254,226,226,0.98) !important;
          font-weight: 700 !important;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"] > button[role="tab"]:nth-of-type(8):hover{
          border-color: rgba(148,163,184,0.72) !important;
          background: rgba(51,65,85,0.72) !important;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"] > button[role="tab"]:nth-of-type(9):hover{
          border-color: rgba(248,113,113,0.66) !important;
          background: rgba(127,29,29,0.30) !important;
          color: #fee2e2 !important;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] button[role="tab"][aria-selected="true"]{
          background: rgba(30,41,59,0.96) !important;
          border-color: rgba(56,189,248,0.82) !important;
          color: #ffffff !important;
          box-shadow: 0 0 0 2px rgba(56,189,248,0.15) !important;
          font-weight: 750 !important;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"] > button[role="tab"]:nth-of-type(8)[aria-selected="true"]{
          border-color: rgba(148,163,184,0.84) !important;
          box-shadow: 0 0 0 2px rgba(148,163,184,0.18) !important;
        }
        .st-key-emp_subtabs_box div[data-testid="stTabs"] > div[role="tablist"] > button[role="tab"]:nth-of-type(9)[aria-selected="true"]{
          border-color: rgba(248,113,113,0.78) !important;
          background: rgba(127,29,29,0.38) !important;
          color: #fee2e2 !important;
          box-shadow: 0 0 0 2px rgba(248,113,113,0.20) !important;
        }
        /* Fișa angajatului (Date personale): secțiuni label/value, uniforme */
        .st-key-emp_subtabs_box .emp-info-section{
          margin: 0 0 16px 0;
          padding: 10px 0 0 0;
        }
        .st-key-emp_subtabs_box .emp-info-title{
          font-size: 1.02rem !important;
          font-weight: 760 !important;
          color: rgba(248,250,252,0.97) !important;
          margin: 0 0 10px 0;
          padding-bottom: 6px;
          border-bottom: none !important;
          background-image: linear-gradient(90deg, rgba(148,163,184,0.36), rgba(148,163,184,0.08));
          background-repeat: no-repeat;
          background-size: min(50vw, 620px) 1px;
          background-position: left bottom;
          letter-spacing: 0.01em;
        }
        .st-key-emp_subtabs_box .emp-info-grid{
          display: grid;
          grid-template-columns: 1fr;
          gap: 6px;
        }
        .st-key-emp_subtabs_box .emp-info-row{
          display: grid;
          grid-template-columns: minmax(170px, 240px) minmax(260px, 1fr);
          align-items: baseline;
          column-gap: 14px;
          row-gap: 4px;
          padding: 2px 0;
          border-bottom: none !important;
          background-image: linear-gradient(90deg, rgba(148,163,184,0.30), rgba(148,163,184,0.06));
          background-repeat: no-repeat;
          background-size: min(50vw, 620px) 1px;
          background-position: left bottom;
        }
        .st-key-emp_subtabs_box .emp-info-label{
          color: rgba(148,163,184,0.92) !important;
          font-size: 0.87rem !important;
          font-weight: 620 !important;
        }
        .st-key-emp_subtabs_box .emp-info-value{
          color: rgba(248,250,252,0.97) !important;
          font-size: 0.92rem !important;
          font-weight: 690 !important;
          word-break: break-word;
          overflow-wrap: anywhere;
        }
        .st-key-emp_subtabs_box .emp-info-value.is-empty{
          color: rgba(148,163,184,0.70) !important;
          font-style: italic;
          font-weight: 560 !important;
        }
        .st-key-emp_subtabs_box .emp-info-section + [data-testid="stMarkdown"] hr,
        .st-key-emp_subtabs_box .emp-info-section + [data-testid="stHorizontalRule"]{
          margin-top: 6px !important;
          margin-bottom: 12px !important;
          opacity: 0.35;
        }
        /* Coerență vizuală pentru restul subtab-urilor Angajați */
        .st-key-emp_subtabs_box .emp-block-title{
          display: block !important;
          width: 100% !important;
          font-size: 1.01rem !important;
          font-weight: 740 !important;
          color: rgba(248,250,252,0.96) !important;
          margin: 12px 0 12px 0 !important;
          padding-bottom: 6px;
          border-bottom: none !important;
          background-image: linear-gradient(90deg, rgba(148,163,184,0.36), rgba(148,163,184,0.08));
          background-repeat: no-repeat;
          background-size: min(50vw, 620px) 1px;
          background-position: left bottom;
          letter-spacing: 0.01em;
        }
        .st-key-emp_subtabs_box [data-testid="stDataFrame"]{
          border: 1px solid rgba(148,163,184,0.18) !important;
          border-radius: 8px !important;
          overflow: hidden !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box .s-card{
          background: rgba(15,23,42,0.88) !important;
          border: 1px solid rgba(148,163,184,0.20) !important;
          border-radius: 8px !important;
          padding: 14px 16px !important;
          margin-bottom: 14px !important;
        }
        .st-key-emp_subtabs_box .s-title{
          font-size: 1rem !important;
          font-weight: 730 !important;
          margin: 0.2em 0 0.55em 0 !important;
          border-bottom: none !important;
          padding-bottom: 0.35em !important;
          color: rgba(248,250,252,0.96) !important;
          background-image: linear-gradient(90deg, rgba(148,163,184,0.36), rgba(148,163,184,0.08));
          background-repeat: no-repeat;
          background-size: min(50vw, 620px) 1px;
          background-position: left bottom;
        }
        .st-key-emp_subtabs_box .s-muted{
          color: rgba(148,163,184,0.90) !important;
          font-size: 0.84rem !important;
        }
        .st-key-emp_subtabs_box .emp-vechime-note{
          color: rgba(203,213,225,0.95) !important;
          font-size: 0.90rem !important;
          margin: 4px 0 8px 0 !important;
          display: block !important;
        }
        .st-key-emp_subtabs_box .emp-vechime-status{
          color: rgba(248,250,252,0.98) !important;
          font-size: 0.93rem !important;
          margin: 2px 0 12px 0 !important;
          display: block !important;
        }
        .st-key-emp_subtabs_box .emp-overview-sep{
          height: 1px;
          width: min(50vw, 620px);
          margin: 10px 0 10px 0;
          background: linear-gradient(90deg, rgba(148,163,184,0.34), rgba(148,163,184,0.08));
          border-radius: 999px;
        }
        .st-key-emp_subtabs_box .emp-overview-pills{
          display: flex;
          flex-wrap: wrap;
          gap: 8px;
          margin: 4px 0 6px 0;
        }
        .st-key-emp_subtabs_box .emp-subsection-title{
          font-size: 0.92rem !important;
          font-weight: 700 !important;
          color: rgba(226,232,240,0.96) !important;
          margin: 10px 0 6px 0 !important;
          letter-spacing: 0.01em;
        }
        .st-key-emp_subtabs_box .emp-section-title{
          font-size: 1.02rem !important;
          font-weight: 760 !important;
          color: rgba(248,250,252,0.97) !important;
          margin: 12px 0 8px 0 !important;
          padding-bottom: 6px;
          background-image: linear-gradient(90deg, rgba(148,163,184,0.36), rgba(148,163,184,0.08));
          background-repeat: no-repeat;
          background-size: min(50vw, 620px) 1px;
          background-position: left bottom;
          letter-spacing: 0.01em;
        }
        .st-key-emp_subtabs_box .emp-subsubsection-title{
          font-size: 0.84rem !important;
          font-weight: 680 !important;
          color: rgba(203,213,225,0.95) !important;
          margin: 8px 0 5px 0 !important;
          letter-spacing: 0.01em;
        }
        .st-key-ang_edit_form [data-testid="stHorizontalBlock"]{
          justify-content: flex-start !important;
          align-items: flex-start !important;
        }
        .st-key-ang_edit_panel{
          max-width: 100% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-ang_edit_panel .stTextInput,
        .st-key-ang_edit_panel .stNumberInput,
        .st-key-ang_edit_panel .stTextArea,
        .st-key-ang_edit_panel .stSelectbox,
        .st-key-ang_edit_panel .stDateInput,
        .st-key-ang_edit_panel .stCheckbox{
          max-width: 100% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-ang_edit_panel .stCheckbox [data-testid="stCheckbox"]{
          display: flex !important;
          align-items: center !important;
          gap: 8px !important;
          flex-direction: row !important;
        }
        .st-key-ang_edit_panel [data-testid="stFormSubmitButton"] > button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-ang_delete_panel{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-ang_delete_panel [data-testid="stCheckbox"]{
          justify-content: flex-start !important;
        }
        .st-key-ang_delete_panel [data-testid="stButton"] > button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] .stTextInput,
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] .stTextArea,
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] .stSelectbox,
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] .stDateInput,
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] .stCheckbox{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] [data-testid="stButton"] > button,
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] [data-testid="stDownloadButton"] > button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
          display: block !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_dec_"] .stTextInput,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_dec_"] .stTextArea,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_dec_"] .stSelectbox,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_dec_"] .stDateInput,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_dec_"] .stFileUploader{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_dec_"] [data-testid="stButton"] > button,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_dec_"] [data-testid="stDownloadButton"] > button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
          display: block !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_det_edit_"] .stTextInput,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_det_edit_"] .stTextArea,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_det_edit_"] .stSelectbox,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_det_edit_"] .stDateInput,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_det_edit_"] .stCheckbox{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_det_edit_"] [data-testid="stButton"] > button,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_det_edit_"] [data-testid="stDownloadButton"] > button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
          display: block !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_preg_"] .stTextInput,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_preg_"] .stTextArea,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_preg_"] .stSelectbox,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_preg_"] .stDateInput,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_preg_"] .stFileUploader{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_preg_"] [data-testid="stButton"] > button,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_preg_"] [data-testid="stDownloadButton"] > button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
          display: block !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_alte_"] .stTextInput,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_alte_"] .stTextArea,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_alte_"] .stSelectbox,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_alte_"] .stDateInput,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_alte_"] .stFileUploader{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_alte_"] [data-testid="stButton"] > button,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_alte_"] [data-testid="stDownloadButton"] > button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
          display: block !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_det_edit_"] .emp-subsection-title{
          font-size: 1.08rem !important;
          font-weight: 800 !important;
          line-height: 1.25 !important;
          letter-spacing: 0.01em;
          margin: 12px 0 8px 0 !important;
          position: relative;
          padding-bottom: 8px !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_det_edit_"] .emp-subsection-title::after{
          content: "";
          display: block;
          width: min(50vw, 420px);
          max-width: 100%;
          height: 1px;
          margin-top: 7px;
          background: linear-gradient(
            90deg,
            rgba(148,163,184,0.46) 0%,
            rgba(148,163,184,0.24) 55%,
            rgba(148,163,184,0.02) 100%
          );
        }
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] .emp-subsection-title,
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] h3,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_dec_"] h4{
          position: relative;
          margin-bottom: 12px !important;
          padding-bottom: 8px !important;
        }
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] .emp-subsection-title::after,
        .st-key-emp_main_tabs_box [class*="st-key-emp_decizii_stack_"] h3::after,
        .st-key-emp_main_tabs_box [class*="st-key-ang_docs_"][class*="_dec_"] h4::after{
          content: "";
          display: block;
          width: min(50vw, 420px);
          max-width: 100%;
          height: 1px;
          margin-top: 7px;
          background: linear-gradient(
            90deg,
            rgba(148,163,184,0.46) 0%,
            rgba(148,163,184,0.24) 55%,
            rgba(148,163,184,0.02) 100%
          );
        }
        .st-key-emp_main_tabs_box .emp-thin-guide-line{
          width: min(50vw, 420px);
          max-width: 100%;
          height: 1px;
          margin: 4px 0 12px 0;
          background: linear-gradient(
            90deg,
            rgba(148,163,184,0.46) 0%,
            rgba(148,163,184,0.24) 55%,
            rgba(148,163,184,0.02) 100%
          );
        }
        /* Vechime: același pattern de aliniere/dimensiune ca în Detalii */
        .st-key-emp_subtabs_box [class*="st-key-ant_work_form_"] [data-testid="stHorizontalBlock"],
        .st-key-emp_subtabs_box [class*="st-key-vechime_form_"] [data-testid="stHorizontalBlock"]{
          justify-content: flex-start !important;
          align-items: flex-start !important;
          gap: 12px !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-ant_work_form_"] [data-testid="stHorizontalBlock"] > div,
        .st-key-emp_subtabs_box [class*="st-key-vechime_form_"] [data-testid="stHorizontalBlock"] > div{
          flex: 0 1 auto !important;
          min-width: 0 !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-ant_work_form_"] .stNumberInput,
        .st-key-emp_subtabs_box [class*="st-key-vechime_form_"] .stNumberInput,
        .st-key-emp_subtabs_box [class*="st-key-vechime_form_"] .stTextInput,
        .st-key-emp_subtabs_box [class*="st-key-vechime_form_"] .stSelectbox{
          max-width: 100% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-ant_work_form_"] [data-testid="stFormSubmitButton"] > button,
        .st-key-emp_subtabs_box [class*="st-key-vechime_form_"] [data-testid="stFormSubmitButton"] > button{
          width: 220px !important;
          min-width: 220px !important;
          max-width: 220px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
          border-radius: 10px !important;
          border: 1px solid rgba(148,163,184,0.30) !important;
          background: rgba(15,23,42,0.60) !important;
          color: rgba(248,250,252,0.96) !important;
          font-weight: 700 !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-ant_work_form_"] [data-testid="stFormSubmitButton"] > button:hover,
        .st-key-emp_subtabs_box [class*="st-key-vechime_form_"] [data-testid="stFormSubmitButton"] > button:hover{
          border-color: rgba(148,163,184,0.58) !important;
          background: rgba(30,41,59,0.72) !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] [data-testid="stHorizontalBlock"]{
          justify-content: flex-start !important;
          align-items: flex-start !important;
          gap: 12px !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] [data-testid="stHorizontalBlock"] > div{
          flex: 0 1 auto !important;
          min-width: 0 !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stTextInput,
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stTextArea,
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stSelectbox,
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stDateInput,
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stCheckbox{
          max-width: 100% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        /* Data nașterii (dep details): elimină complet haloul/albul wrapper-ului */
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stDateInput,
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stDateInput > div,
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stDateInput div[data-baseweb="input"],
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stDateInput div[data-baseweb="input"] > div{
          background: transparent !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stDateInput div[data-baseweb="input"]{
          border: 1px solid rgba(148,163,184,0.30) !important;
          border-radius: 10px !important;
          overflow: hidden !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stDateInput input{
          background: rgba(15,23,42,0.60) !important;
          color: rgba(248,250,252,0.96) !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_details_box_"] .stDateInput button{
          background: rgba(15,23,42,0.60) !important;
          border: none !important;
          box-shadow: none !important;
        }
        /* Vechime - simplu: un singur border, fără dubluri/alb */
        .st-key-emp_subtabs_box [class*="st-key-data_ang_date_"],
        .st-key-emp_subtabs_box [class*="st-key-data_plec_date_"],
        .st-key-emp_subtabs_box [class*="st-key-still_active_"],
        .st-key-emp_subtabs_box [class*="st-key-auto_calc_inst_"],
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"],
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"],
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"],
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"],
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"],
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"],
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"],
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"],
        .st-key-emp_subtabs_box [class*="st-key-vechime_report_type_"]{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-data_ang_date_"] div[data-baseweb="input"],
        .st-key-emp_subtabs_box [class*="st-key-data_plec_date_"] div[data-baseweb="input"],
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"] div[data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"] div[data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"] div[data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"] div[data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"] div[data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"] div[data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"] div[data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"] div[data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-vechime_report_type_"] div[data-baseweb="select"]{
          background: rgba(15,23,42,0.60) !important;
          border: 1px solid rgba(148,163,184,0.30) !important;
          border-radius: 10px !important;
          box-shadow: none !important;
          overflow: hidden !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-data_ang_date_"] div[data-baseweb="input"] > div,
        .st-key-emp_subtabs_box [class*="st-key-data_plec_date_"] div[data-baseweb="input"] > div,
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"] div[data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"] div[data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"] div[data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"] div[data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"] div[data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"] div[data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"] div[data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"] div[data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-vechime_report_type_"] div[data-baseweb="select"] > div{
          border: none !important;
          background: transparent !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-data_ang_date_"] input,
        .st-key-emp_subtabs_box [class*="st-key-data_plec_date_"] input,
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"] input,
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"] input,
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"] input,
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"] input,
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"] input,
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"] input,
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"] input,
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"] input{
          background: transparent !important;
          color: rgba(248,250,252,0.96) !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-ant_edit_btn_"] button,
        .st-key-emp_subtabs_box [class*="st-key-ant_lock_btn_"] button,
        .st-key-emp_subtabs_box [class*="st-key-gen_report_btn_"] button{
          width: 220px !important;
          min-width: 220px !important;
          max-width: 220px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-top: 8px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
          border-radius: 10px !important;
          border: 1px solid rgba(148,163,184,0.30) !important;
          background: rgba(15,23,42,0.60) !important;
          color: rgba(248,250,252,0.96) !important;
          box-shadow: none !important;
        }
        /* Contract de muncă: casete ca în Angajați, până la jumătate pagină */
        .st-key-emp_subtabs_box [class*="st-key-cim_"]{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-cim_"] [data-testid="stTextInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-cim_"] [data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-cim_"] [data-testid="stTextAreaRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-cim_"] .stSelectbox div[data-baseweb="select"],
        .st-key-emp_subtabs_box [class*="st-key-cim_"] .stDateInput div[data-baseweb="input"]{
          background: rgba(15,23,42,0.60) !important;
          border: 1px solid rgba(148,163,184,0.30) !important;
          border-radius: 10px !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-cim_"] [data-testid="stTextInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-cim_"] [data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-cim_"] [data-testid="stTextAreaRootElement"] > div,
        .st-key-emp_subtabs_box [class*="st-key-cim_"] .stSelectbox div[data-baseweb="select"] > div,
        .st-key-emp_subtabs_box [class*="st-key-cim_"] .stDateInput div[data-baseweb="input"] > div{
          background: transparent !important;
          border: none !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-aa_code_"]{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-aa_code_gen_"] button,
        .st-key-emp_subtabs_box [class*="st-key-aa_code_gen_pdf_"] button{
          width: 220px !important;
          min-width: 220px !important;
          max-width: 220px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-cim_code_gen_"] button,
        .st-key-emp_subtabs_box [class*="st-key-cim_code_gen_pdf_"] button,
        .st-key-emp_subtabs_box [class*="st-key-cim_"][class*="_save"] button,
        .st-key-emp_subtabs_box [class*="st-key-cim_"][class*="_set_active"] button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-reg_"]{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-reg_save_"] button,
        .st-key-emp_subtabs_box [class*="st-key-reg_cancel_"] button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        /* COD COR: doar secțiunile de după tabel la 50% */
        .st-key-emp_subtabs_box [class*="st-key-cor_sel_"],
        .st-key-emp_subtabs_box [class*="st-key-cor_cod_in_"],
        .st-key-emp_subtabs_box [class*="st-key-cor_den_in_"],
        .st-key-emp_subtabs_box [class*="st-key-cor_save_tbl_"],
        .st-key-emp_subtabs_box [class*="st-key-cor_apply_emp_"],
        .st-key-emp_subtabs_box [class*="st-key-cor_del_tbl_"]{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-cor_save_tbl_"] button,
        .st-key-emp_subtabs_box [class*="st-key-cor_apply_emp_"] button,
        .st-key-emp_subtabs_box [class*="st-key-cor_del_tbl_"] button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-l153_save_tbl_"] button,
        .st-key-emp_subtabs_box [class*="st-key-l153_apply_emp_"] button,
        .st-key-emp_subtabs_box [class*="st-key-l153_del_tbl_"] button,
        .st-key-emp_subtabs_box [class*="st-key-l153_grid_"][class*="save_active"] button{
          width: 420px !important;
          min-width: 420px !important;
          max-width: 420px !important;
          min-height: 40px !important;
          height: 40px !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        /* Legea 153: secțiunile de după tabel la 50% */
        .st-key-emp_subtabs_box [class*="st-key-l153_sel_"],
        .st-key-emp_subtabs_box [class*="st-key-l153_cod_in_"],
        .st-key-emp_subtabs_box [class*="st-key-l153_den_in_"],
        .st-key-emp_subtabs_box [class*="st-key-l153_save_tbl_"],
        .st-key-emp_subtabs_box [class*="st-key-l153_apply_emp_"],
        .st-key-emp_subtabs_box [class*="st-key-l153_del_tbl_"],
        .st-key-emp_subtabs_box [class*="st-key-l153_grid_"]{
          max-width: 50% !important;
          margin-left: 0 !important;
          margin-right: auto !important;
        }
        /* EMP subtabs: elimină complet alb-ul din casete */
        .st-key-emp_subtabs_box [data-testid="stTextInputRootElement"],
        .st-key-emp_subtabs_box [data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [data-testid="stTextAreaRootElement"],
        .st-key-emp_subtabs_box .stSelectbox div[data-baseweb="select"],
        .st-key-emp_subtabs_box .stDateInput div[data-baseweb="input"]{
          background: rgba(15,23,42,0.60) !important;
          border: 1px solid rgba(148,163,184,0.30) !important;
          border-radius: 10px !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box [data-testid="stTextInputRootElement"] > div,
        .st-key-emp_subtabs_box [data-testid="stNumberInputRootElement"] > div,
        .st-key-emp_subtabs_box [data-testid="stTextAreaRootElement"] > div,
        .st-key-emp_subtabs_box .stSelectbox div[data-baseweb="select"] > div,
        .st-key-emp_subtabs_box .stDateInput div[data-baseweb="input"] > div{
          background: transparent !important;
          border: none !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box input,
        .st-key-emp_subtabs_box textarea{
          background: transparent !important;
          color: rgba(248,250,252,0.96) !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box .stDateInput button,
        .st-key-emp_subtabs_box .stSelectbox [role="button"]{
          background: rgba(15,23,42,0.60) !important;
          border: none !important;
          box-shadow: none !important;
          color: rgba(248,250,252,0.96) !important;
        }
        /* Number input steppers (+ / -): fără alb în partea dreaptă */
        .st-key-emp_subtabs_box [data-testid="stNumberInputRootElement"] button,
        .st-key-emp_subtabs_box [data-testid="stNumberInputRootElement"] [role="button"]{
          background: rgba(15,23,42,0.60) !important;
          border: none !important;
          box-shadow: none !important;
          color: rgba(248,250,252,0.96) !important;
        }
        .st-key-emp_subtabs_box [data-testid="stNumberInputRootElement"] button:hover,
        .st-key-emp_subtabs_box [data-testid="stNumberInputRootElement"] [role="button"]:hover{
          background: rgba(30,41,59,0.72) !important;
        }
        /* VECHIME only: steppers (-/+) integrate dark, fără alb */
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"] [data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"] [data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"] [data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"] [data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"] [data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"] [data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"] [data-testid="stNumberInputRootElement"],
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"] [data-testid="stNumberInputRootElement"]{
          background: rgba(15,23,42,0.60) !important;
          border: 1px solid rgba(148,163,184,0.30) !important;
          border-radius: 10px !important;
          overflow: hidden !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"] [data-testid="stNumberInputRootElement"] button,
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"] [data-testid="stNumberInputRootElement"] button,
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"] [data-testid="stNumberInputRootElement"] button,
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"] [data-testid="stNumberInputRootElement"] button,
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"] [data-testid="stNumberInputRootElement"] button,
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"] [data-testid="stNumberInputRootElement"] button,
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"] [data-testid="stNumberInputRootElement"] button,
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"] [data-testid="stNumberInputRootElement"] button{
          background: rgba(15,23,42,0.60) !important;
          color: rgba(248,250,252,0.96) !important;
          border: none !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"] [data-testid="stNumberInputRootElement"] button:hover,
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"] [data-testid="stNumberInputRootElement"] button:hover,
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"] [data-testid="stNumberInputRootElement"] button:hover,
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"] [data-testid="stNumberInputRootElement"] button:hover,
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"] [data-testid="stNumberInputRootElement"] button:hover,
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"] [data-testid="stNumberInputRootElement"] button:hover,
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"] [data-testid="stNumberInputRootElement"] button:hover,
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"] [data-testid="stNumberInputRootElement"] button:hover{
          background: rgba(30,41,59,0.72) !important;
        }
        /* Force: elimină orice alb rezidual pe steppers BaseWeb */
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"] [data-baseweb="input"] button,
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"] [data-baseweb="input"] button,
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"] [data-baseweb="input"] button,
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"] [data-baseweb="input"] button,
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"] [data-baseweb="input"] button,
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"] [data-baseweb="input"] button,
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"] [data-baseweb="input"] button,
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"] [data-baseweb="input"] button,
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"] [data-baseweb="input"] button > div,
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"] [data-baseweb="input"] button > div,
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"] [data-baseweb="input"] button > div,
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"] [data-baseweb="input"] button > div,
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"] [data-baseweb="input"] button > div,
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"] [data-baseweb="input"] button > div,
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"] [data-baseweb="input"] button > div,
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"] [data-baseweb="input"] button > div{
          background: rgba(15,23,42,0.60) !important;
          border: none !important;
          box-shadow: none !important;
          outline: none !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-ant_ani_"] [data-baseweb="input"] button svg,
        .st-key-emp_subtabs_box [class*="st-key-ant_luni_"] [data-baseweb="input"] button svg,
        .st-key-emp_subtabs_box [class*="st-key-inst_ani_"] [data-baseweb="input"] button svg,
        .st-key-emp_subtabs_box [class*="st-key-inst_luni_"] [data-baseweb="input"] button svg,
        .st-key-emp_subtabs_box [class*="st-key-inst_fn_"] [data-baseweb="input"] button svg,
        .st-key-emp_subtabs_box [class*="st-key-inst_fd_"] [data-baseweb="input"] button svg,
        .st-key-emp_subtabs_box [class*="st-key-functie_ani_"] [data-baseweb="input"] button svg,
        .st-key-emp_subtabs_box [class*="st-key-functie_luni_"] [data-baseweb="input"] button svg{
          fill: rgba(248,250,252,0.96) !important;
          color: rgba(248,250,252,0.96) !important;
        }
        /* Persoane în întreținere - acțiuni verticale, uniforme */
        .st-key-emp_subtabs_box [class*="st-key-dep_action_btn_"] [data-testid="stFormSubmitButton"] > button{
          width: 180px !important;
          min-width: 180px !important;
          max-width: 180px !important;
          min-height: 40px !important;
          height: 40px !important;
          display: block !important;
          margin-left: 0 !important;
          margin-right: auto !important;
          border-radius: 10px !important;
          border: 1px solid rgba(148,163,184,0.30) !important;
          background: rgba(15,23,42,0.60) !important;
          color: rgba(248,250,252,0.96) !important;
          font-weight: 700 !important;
          box-shadow: none !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_action_btn_save_"] [data-testid="stFormSubmitButton"] > button{
          margin-top: 8px !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_action_btn_"] [data-testid="stFormSubmitButton"] > button:hover{
          border-color: rgba(148,163,184,0.58) !important;
          background: rgba(30,41,59,0.72) !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_action_btn_delete_"] [data-testid="stFormSubmitButton"] > button{
          border-color: rgba(248,113,113,0.44) !important;
          background: rgba(127,29,29,0.24) !important;
          color: rgba(254,226,226,0.98) !important;
        }
        .st-key-emp_subtabs_box [class*="st-key-dep_action_btn_delete_"] [data-testid="stFormSubmitButton"] > button:hover{
          border-color: rgba(248,113,113,0.66) !important;
          background: rgba(127,29,29,0.30) !important;
        }
        /* Conținut tab-uri și expandere în fișa angajat – nu toate rotunjite */
        section.main:has(#emp-scope) div[data-testid="stTabs"] > div[data-testid="stVerticalBlockBorderWrapper"],
        section.main:has(#emp-scope) div[data-testid="stExpander"]{
          border-radius: 6px !important;
        }
        section.main:has(#emp-scope) div[data-testid="stExpander"] details{
          border-radius: 6px !important;
        }
        /* Toate secțiunile și casetele în partea stângă */
        section.main:has(#emp-scope) .block-container{
          text-align: left !important;
          margin-left: 0 !important;
          margin-right: auto !important;
          padding-left: 1rem !important;
          padding-top: 0.35rem !important;
          margin-top: 0 !important;
        }
        section.main:has(#emp-scope) .page-title{
          margin-top: 0 !important;
          padding-top: 0 !important;
        }
        section.main:has(#emp-scope) div[data-testid="stMainBlockContainer"]{
          padding-top: 0 !important;
          margin-top: 0 !important;
        }
        section.main:has(#emp-scope) .block-container > div:first-child{
          margin-top: 0 !important;
          padding-top: 0 !important;
        }
        section.main:has(#emp-scope) [data-testid="stVerticalBlock"],
        section.main:has(#emp-scope) [data-testid="stVerticalBlockBorderWrapper"],
        section.main:has(#emp-scope) .emp-wrap [data-testid="stVerticalBlock"],
        section.main:has(#emp-scope) .emp-wrap [data-testid="stVerticalBlockBorderWrapper"],
        section.main:has(#emp-scope) .s-card,
        section.main:has(#emp-scope) [data-testid="stVerticalBlock"] > div,
        section.main:has(#emp-scope) form,
        section.main:has(#emp-scope) div[data-testid="stForm"]{
          text-align: left !important;
        }
        section.main:has(#emp-scope) [data-testid="stHorizontalBlock"]{
          justify-content: flex-start !important;
        }
        section.main:has(#emp-scope) [data-testid="stHorizontalBlock"] > div{
          flex: 0 1 auto !important;
          min-width: 0;
        }
        section.main:has(#emp-scope) form [data-testid="stFormSubmitButton"] > button,
        section.main:has(#emp-scope) .emp-wrap form [data-testid="stFormSubmitButton"] > button{
          font-weight: 700 !important;
          min-width: 140px;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )
    def _dep_reset_form(K: str, set_new: bool = True) -> None:
        st.session_state[f"{K}edit_id"] = None
        for kk in (
            f"{K}tip", f"{K}nume", f"{K}prenume", f"{K}cnp", f"{K}dob",
            f"{K}grad", f"{K}obs", f"{K}activ",
            f"{K}dob_manual", f"{K}cnp_confirm_conflict",
        ):
            st.session_state.pop(kk, None)
        st.session_state[f"{K}tip"] = ""
        if set_new:
            st.session_state[f"{K}force_new"] = True
    # ---------------------------
    # State UI (list / detail / add / edit / delete)
    # ---------------------------
    if "ang_view" not in st.session_state:
        st.session_state["ang_view"] = "list"
    if "ang_selected_id" not in st.session_state:
        st.session_state["ang_selected_id"] = None

    # ---------------------------
    # Sub-view: ADD / EDIT / DELETE (din listă)
    # ---------------------------
    if st.session_state.get("ang_view") in ("add", "edit", "delete"):
        cols = _employees_columns(conn)
        mode = st.session_state["ang_view"]
        emp_id = st.session_state.get("ang_selected_id")

        top1, top2 = st.columns([1, 8], vertical_alignment="center")
        with top1:
            if st.button("⬅ Înapoi la listă", key="ang_back_to_list_from_mode"):
                st.session_state["ang_view"] = "list"
                st.session_state["ang_selected_id"] = None
                st.rerun()
        with top2:
            st.markdown("", unsafe_allow_html=True)

        if mode == "add":
            st.markdown(
                """
                <style>
                .st-key-ang_add_panel{
                  max-width: 100% !important;
                  margin-left: 0 !important;
                  margin-right: auto !important;
                }
                .st-key-ang_add_panel [data-testid="stHorizontalBlock"]{
                  justify-content: flex-start !important;
                }
                .st-key-ang_add_panel .stCheckbox [data-testid="stCheckbox"]{
                  display: flex !important;
                  flex-direction: row !important;
                  align-items: center !important;
                  gap: 8px !important;
                }
                .st-key-ang_add_panel [class*="st-key-ang_add_cnp"] [data-testid="stCaptionContainer"]{
                  margin-top: -6px !important;
                  margin-bottom: 4px !important;
                }
                .st-key-ang_add_panel [class*="st-key-ang_add_cnp"] [data-testid="stCaptionContainer"] p{
                  margin-top: 0 !important;
                  margin-bottom: 0 !important;
                  line-height: 1.2 !important;
                }
                .st-key-ang_add_panel .stTextInput,
                .st-key-ang_add_panel .stNumberInput,
                .st-key-ang_add_panel .stTextArea,
                .st-key-ang_add_panel .stSelectbox,
                .st-key-ang_add_panel .stDateInput,
                .st-key-ang_add_panel .stCheckbox,
                .st-key-ang_add_panel [data-testid="stTextInputRootElement"],
                .st-key-ang_add_panel [data-testid="stNumberInputRootElement"]{
                  max-width: 58% !important;
                  width: 58% !important;
                  margin-left: 0 !important;
                  margin-right: auto !important;
                }
                .st-key-ang_add_panel [class*="st-key-ang_add_nr_copii"] [data-testid="stNumberInputRootElement"],
                .st-key-ang_add_panel [class*="st-key-ang_add_salariu_baza"] [data-testid="stNumberInputRootElement"]{
                  max-width: 58% !important;
                  width: 58% !important;
                  margin-left: 0 !important;
                  margin-right: auto !important;
                }
                .st-key-ang_add_panel [class*="st-key-ang_add_nr_copii"] [data-testid="stNumberInputRootElement"] > div,
                .st-key-ang_add_panel [class*="st-key-ang_add_salariu_baza"] [data-testid="stNumberInputRootElement"] > div{
                  max-width: 100% !important;
                  width: 100% !important;
                }
                .emp-block-title-add{
                  position: relative;
                  font-weight: 800 !important;
                  color: rgba(241,245,249,0.98) !important;
                  padding-bottom: 8px !important;
                  margin-bottom: 12px !important;
                }
                .emp-block-title-add::after{
                  content: "";
                  display: block;
                  width: min(50vw, 420px);
                  max-width: 100%;
                  height: 1px;
                  margin-top: 7px;
                  background: linear-gradient(
                    90deg,
                    rgba(148,163,184,0.46) 0%,
                    rgba(148,163,184,0.24) 55%,
                    rgba(148,163,184,0.02) 100%
                  );
                }
                .st-key-ang_add_save_btn button,
                .st-key-ang_add_cancel_btn button{
                  width: 220px !important;
                  min-width: 220px !important;
                  max-width: 220px !important;
                  height: 42px !important;
                  min-height: 42px !important;
                  max-height: 42px !important;
                }
                </style>
                """,
                unsafe_allow_html=True,
            )
            st.subheader("➕ Adaugă angajat")
            st.markdown("<div class='emp-thin-guide-line'></div>", unsafe_allow_html=True)
            st.info("Câmpuri obligatorii pentru înregistrare nouă: Marcă, Nume, Prenume și CNP (valid).")
            with st.container(key="ang_add_panel"):
                with st.form("ang_add_form"):
                    values = _render_employee_form_fields(
                        {},
                        cols,
                        prefix="ang_add",
                        single_column=True,
                        fine_section_titles=True,
                        section_title_class="emp-block-title-add",
                    )

                    # Previzualizare constantă a datelor extrase din CNP în secțiunea „Date de bază”
                    cnp_state = str(st.session_state.get("ang_add_cnp", "") or "").strip()
                    st.markdown("#### Date extrase din CNP")
                    if not cnp_state:
                        st.caption("Completează CNP-ul pentru a vedea aici data nașterii, sexul și județul extrase automat.")
                    else:
                        info = decode_cnp(cnp_state, strict_county=False, allow_s_9=False)
                        if info.valid:
                            col1, col2, col3 = st.columns(3)
                            with col1:
                                st.text_input(
                                    "Data nașterii (din CNP)",
                                    value=info.birth_date.strftime("%d.%m.%Y") if info.birth_date else "",
                                    key="ang_add_cnp_birth",
                                    disabled=True,
                                )
                            with col2:
                                st.text_input(
                                    "Sex (din CNP)",
                                    value=("Masculin" if info.sex == "M" else "Feminin") if info.sex else "",
                                    key="ang_add_cnp_sex",
                                    disabled=True,
                                )
                            with col3:
                                st.text_input(
                                    "Județ (din CNP)",
                                    value=f"{info.county_name} ({info.county_code})" if info.county_name else "",
                                    key="ang_add_cnp_county",
                                    disabled=True,
                                )
                        else:
                            st.error(info.error or "CNP invalid – nu pot extrage datele.")

                    ok = st.form_submit_button("💾 Salvează", key="ang_add_save_btn")
                    cancel = st.form_submit_button("↩ Renunță", key="ang_add_cancel_btn")

            if cancel:
                st.session_state["ang_selected_id"] = None
                st.session_state["ang_view"] = "list"
                st.rerun()

            if ok:
                # Validare obligatorie: Marcă, Nume, Prenume, CNP (și CNP valid)
                marca = str(values.get("marca") or "").strip()
                nume = str(values.get("last_name") or "").strip()
                prenume = str(values.get("first_name") or "").strip()
                cnp_raw = str(values.get("cnp") or "").strip()

                missing = []
                if not marca:
                    missing.append("Marcă")
                if not nume:
                    missing.append("Nume")
                if not prenume:
                    missing.append("Prenume")
                if not cnp_raw:
                    missing.append("CNP")

                if missing:
                    st.error("Completează câmpurile obligatorii: " + ", ".join(missing) + ".")
                else:
                    ok_cnp, msg_cnp, _info = cnp_validate(cnp_raw)
                    if not ok_cnp:
                        st.error(f"CNP invalid: {msg_cnp}")
                    else:
                        # Conversie numerică pentru câmpurile rescrise ca text în Add.
                        try:
                            _nr = str(values.get("nr_copii") or "").strip()
                            if _nr != "":
                                values["nr_copii"] = int(_nr)
                        except Exception:
                            st.error("Nr. copii trebuie să fie un număr întreg.")
                            return
                        try:
                            _sal = str(values.get("salariu_baza") or "").strip()
                            if _sal != "":
                                values["salariu_baza"] = float(_sal.replace(",", "."))
                        except Exception:
                            st.error("Salariu de bază trebuie să fie numeric.")
                            return
                        try:
                            new_id = _employee_upsert(conn, values, employee_id=None)
                        except ValueError as e:
                            st.error(str(e))
                        else:
                            st.success("Angajat adăugat. Poți completa acum restul datelor personale.")
                            st.session_state["ang_selected_id"] = new_id
                            # După prima salvare mergem în modul EDIT, nu direct în tabul „Date personale” (read-only)
                            st.session_state["ang_view"] = "edit"
                            st.rerun()

        elif mode == "edit" and emp_id:
            bundle = load_employee_bundle(get_db_path(), int(emp_id))
            emp = bundle["employee"]
            if not emp:
                st.error("Angajatul selectat nu mai există.")
                st.session_state["ang_view"] = "list"
                st.session_state["ang_selected_id"] = None
                st.rerun()

            st.subheader("📋 DATE PERSONALE")
            with st.form("ang_edit_inline_form"):
                values = _render_employee_form_fields(emp, cols, prefix=f"ang_edit_inline_{emp_id}")
                c_save, c_cancel = st.columns([1, 1])
                with c_save:
                    ok = st.form_submit_button("💾 Salvează date personale")
                with c_cancel:
                    cancel = st.form_submit_button("↩ Renunță")

            if cancel:
                st.session_state["ang_view"] = "detail"
                st.rerun()

            if ok:
                try:
                    _employee_upsert(conn, values, employee_id=int(emp_id))
                except ValueError as e:
                    st.error(str(e))
                else:
                    st.success("Modificări salvate.")
                    st.session_state["ang_view"] = "detail"
                    st.rerun()

        elif mode == "delete" and emp_id:
            bundle = load_employee_bundle(get_db_path(), int(emp_id))
            emp = bundle["employee"]
            if not emp:
                st.error("Angajatul selectat nu mai există.")
                st.session_state["ang_view"] = "list"
                st.session_state["ang_selected_id"] = None
                st.rerun()

            st.subheader("🗑 Șterge (dezactivează) angajat")
            st.warning("Ștergerea este soft delete: activ = 0.")
            confirm = st.checkbox("Confirm ștergerea", key="ang_del_inline_confirm")
            c_del, c_cancel = st.columns([1, 1])
            with c_del:
                do_del = st.button("Confirm 🗑", disabled=not confirm, key="ang_del_inline_btn")
            with c_cancel:
                cancel = st.button("↩ Renunță", key="ang_del_inline_cancel")

            if cancel:
                st.session_state["ang_view"] = "detail"
                st.rerun()

            if do_del:
                _employee_soft_delete(conn, int(emp_id))
                st.success("Angajat dezactivat.")
                st.session_state["ang_view"] = "list"
                st.session_state["ang_selected_id"] = None
                st.rerun()

        return

    # ---------------------------
    # View: DETAIL
    # ---------------------------
    if st.session_state["ang_view"] == "detail" and st.session_state["ang_selected_id"]:
        emp_id = int(st.session_state["ang_selected_id"])
        # La deschiderea fișei: scroll sus – containerul scrollabil în Streamlit e de obicei .main
        if st.session_state.pop("ang_scroll_to_top", False):
            import streamlit.components.v1 as components
            scroll_js = """
            <script>
            (function() {
                function run() {
                    var d = window.parent.document;
                    ['.main', 'section.main', '[data-testid="stMain"]', '[data-testid="stAppViewContainer"]'].forEach(function(sel) {
                        var el = d.querySelector(sel);
                        if (el) { el.scrollTop = 0; el.scrollIntoView && el.scrollIntoView({ behavior: 'instant', block: 'start' }); }
                    });
                    window.parent.scrollTo(0, 0);

                    // Forțează deschiderea pe prima secțiune: main tab + subtab.
                    var mainFirst = d.querySelector('.st-key-emp_main_tabs_box button[role="tab"]');
                    if (mainFirst && mainFirst.getAttribute('aria-selected') !== 'true') {
                        mainFirst.click();
                    }
                    var subFirst = d.querySelector('.st-key-emp_subtabs_box button[role="tab"]');
                    if (subFirst && subFirst.getAttribute('aria-selected') !== 'true') {
                        subFirst.click();
                    }
                }
                run();
                setTimeout(run, 100);
                setTimeout(run, 400);
            })();
            </script>
            """
            components.html(scroll_js, height=0)
        bundle = load_employee_bundle(get_db_path(), emp_id)
        emp = bundle["employee"]
        if not emp:
            st.error("Angajatul selectat nu mai există.")
            st.session_state["ang_view"] = "list"
            st.session_state["ang_selected_id"] = None
            st.rerun()

        st.markdown("### Fișa angajatului")
        st.markdown("<div class='emp-detail-header'>", unsafe_allow_html=True)
        top1, _sp = st.columns([1, 6])
        with top1:
            if st.button("⬅ Înapoi la listă", key="ang_back_to_list"):
                st.session_state["ang_view"] = "list"
                st.rerun()
        with st.container(key="emp_header_center_box"):
            cur_state = int(emp.get("activ", 1) or 0)
            emp_name = f"{str(emp.get('last_name','') or '').strip()} {str(emp.get('first_name','') or '').strip()}".strip()
            emp_marca = str(emp.get("marca", "") or "").strip()
            state_label = "ACTIV" if cur_state == 1 else "INACTIV"
            st.markdown(
                (
                    "<div class='emp-header-center'>"
                    f"<div class='emp-header-name'><strong>👤 {emp_name or '—'}</strong></div>"
                    f"<div class='emp-header-marca'>Marcă: {emp_marca or '—'}</div>"
                    f"<div class='emp-header-state'>Stare curentă: <strong>{state_label}</strong></div>"
                    "</div>"
                ),
                unsafe_allow_html=True,
            )
            with st.container(key="emp_header_actions"):
                if st.button("✅ Activ", disabled=(cur_state == 1), key=f"ang_set_active_{emp_id}"):
                    try:
                        _employee_upsert(conn, {**emp, "activ": 1}, employee_id=emp_id)
                    except ValueError as e:
                        st.error(str(e))
                    else:
                        st.success("Angajat activat.")
                        st.rerun()
                if st.button("⛔ Inactiv", disabled=(cur_state == 0), key=f"ang_set_inactive_{emp_id}"):
                    try:
                        _employee_upsert(conn, {**emp, "activ": 0}, employee_id=emp_id)
                    except ValueError as e:
                        st.error(str(e))
                    else:
                        st.success("Angajat dezactivat.")
                        st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)

        with st.container(key="emp_main_tabs_box"):
            main_tabs_emp = st.tabs(["👤 DATE SALARIAT", "📎 DOCUMENTE / INFORMAȚII SALARIAT"]) 
        with main_tabs_emp[0]:
                with st.container(key="emp_subtabs_box"):
                    tabs = st.tabs(
                        [
                            "Date personale",
                            "Persoane în întreținere",
                            "⏳ Vechime",
                            "📄 Contract de muncă",
                            "🧾 Reges Online",
                            "🧩 COD COR",
                            "📜 Salarizare (Legea-cadru 153/2017)", #modificat din LEGEA 153
                                            "✏️ Editează",
                            "🗑 Șterge"]
                    )
                cols = _employees_columns(conn)

                # ---------------------------
                # TAB 0: Date personale (read-only) – etichete în română, secțiuni clare
                # ---------------------------
                with tabs[0]:
                    _labels = _employee_field_labels()
                    groups = _employee_field_groups(cols)
                    for gname, gcols in groups:
                        section_rows: list[tuple[str, str]] = []
                        for c in gcols:
                            label = _labels.get(c) or " ".join(p.capitalize() for p in str(c).split("_") if p) or c
                            if c == "activ":
                                val = "Da" if int(emp.get(c, 1) or 0) == 1 else "Nu"
                            else:
                                raw = emp.get(c, "")
                                val = "" if raw is None else str(raw)
                            section_rows.append((label, val or "—"))
                        _render_info_section(gname, section_rows)
                        st.divider()

                    # Buton export date (la finalul secțiunii; modelul raportului se poate adăuga mai târziu)
                    st.markdown("---")
                    export_lines = []
                    for gname, gcols in groups:
                        export_lines.append(f"\n{gname}\n" + "-" * 40)
                        for c in gcols:
                            label = _labels.get(c) or " ".join(p.capitalize() for p in str(c).split("_") if p) or c
                            if c == "activ":
                                val = "Da" if int(emp.get(c, 1) or 0) == 1 else "Nu"
                            else:
                                raw = emp.get(c, "")
                                val = "" if raw is None else str(raw)
                            export_lines.append(f"{label}: {val or '—'}")
                    export_content = "\n".join(export_lines).strip()
                    nume_export = f"date_personale_{emp.get('marca', emp_id) or emp_id}.txt"
                    st.download_button(
                        "📥 Export date",
                        data=export_content.encode("utf-8"),
                        file_name=nume_export,
                        mime="text/plain; charset=utf-8",
                        key=f"export_date_personale_{emp_id}",
                    )

                # ---------------------------
                # TAB 1: Persoane în întreținere
                # ---------------------------
                with tabs[1]:
                
                        K = f"dep_{emp_id}_"

                        # inițializăm counter-ul o singură dată
                        if f"{K}key_counter" not in st.session_state:
                            st.session_state[f"{K}key_counter"] = 0

                        if st.session_state.pop(f"{K}reset_pending", False):
                            _dep_reset_form(K, set_new=True)
                            # ștergem toate cheile vechi
                            for _k in (f"{K}tip", f"{K}nume", f"{K}prenume", f"{K}cnp",
                                    f"{K}dob", f"{K}grad", f"{K}obs", f"{K}activ"):
                                st.session_state.pop(_k, None)
                            # regenerăm widget-urile
                            st.session_state[f"{K}key_counter"] += 1
                            st.session_state[f"{K}sel_row"] = "(nou)"
                            st.rerun()

                        # --------------------------------------------------------
                        # 2.  RESTUL CODULUI (sub-header, helpers, formular etc.)
                        # --------------------------------------------------------
                        st.subheader("Persoane în întreținere")
                        st.caption("Adaugă/editează persoane aflate în întreținerea angajatului (copil, părinte, soț/soție).")
                        _render_info_section(
                            "Context angajat",
                            [
                                ("Marcă", _emp_v(emp, "marca")),
                                ("Nume complet", f"{_emp_v(emp, 'last_name', '')} {_emp_v(emp, 'first_name', '')}".strip() or "—"),
                                ("CNP", _emp_v(emp, "cnp")),
                            ],
                        )

                        # --- helpers UI (la fel ca înainte) ---
                        TIP_OPTS = ["", "SOT_SOTIE", "COPIL", "PARINTE"]
                        TIP_LABEL = {
                            "": "Selectează tipul",
                            "SOT_SOTIE": "Soț / Soție",
                            "COPIL": "Copil",
                            "PARINTE": "Părinte",
                        }
                        GRAD_OPTS = ["", "I", "II", "III"]
                        GRAD_LABEL = {"": "—", "I": "Gradul I", "II": "Gradul II", "III": "Gradul III"}

                        def _norm_tip(x: str) -> str:
                            x = (x or "").strip().upper().replace(" ", "_")
                            if x == "":
                                return ""
                            if x in ("SOT", "SOTIE", "SOȚ", "SOȚIE", "SOT/SOTIE", "SOT_/_SOTIE"):
                                return "SOT_SOTIE"
                            if x in ("PĂRINTE",):
                                return "PARINTE"
                            if x not in ("SOT_SOTIE", "COPIL", "PARINTE"):
                                return "COPIL"
                            return x

                        def _suggest_grad_for_tip(tip: str) -> str:
                            tip = _norm_tip(tip)
                            return "I" if tip in ("COPIL", "PARINTE", "SOT_SOTIE") else ""

                        def _parse_date_any(s: str):
                            s = (s or "").strip()
                            if not s:
                                return None
                            for fmt in ("%Y-%m-%d", "%d.%m.%Y", "%d/%m/%Y"):
                                try:
                                    return datetime.datetime.strptime(s, fmt).date()
                                except Exception:
                                    pass
                            return None

                        # --- state keys ---
                        if f"{K}edit_id" not in st.session_state:
                            st.session_state[f"{K}edit_id"] = None
                        if f"{K}force_new" not in st.session_state:
                            st.session_state[f"{K}force_new"] = False

                        def _dep_reset_form(K: str, set_new: bool = True):
                            st.session_state[f"{K}edit_id"] = None
                            for kk in (
                                f"{K}tip", f"{K}nume", f"{K}prenume", f"{K}cnp", f"{K}dob",
                                f"{K}grad", f"{K}obs", f"{K}activ",
                                f"{K}dob_manual", f"{K}cnp_confirm_conflict",
                            ):
                                st.session_state.pop(kk, None)
                            st.session_state[f"{K}tip"] = ""
                            if set_new:
                                st.session_state[f"{K}force_new"] = True

                        # --- checkbox / listă / select pentru editare ---
                        show_inactive_dep = st.checkbox("Include și inactivi", value=False, key=f"{K}show_inactive")
                        dep_list = bundle["employee_dependents"]
                        if not show_inactive_dep:
                            dep_list = [d for d in dep_list if int(d.get("activ", 1) or 0) == 1]
                        df_dep = pd.DataFrame(dep_list)
                        if not df_dep.empty:
                            df_dep = df_dep.rename(columns={"id": "ID", "employee_id": "EMPLOYEE_ID", "tip": "TIP", "nume": "NUME", "prenume": "PRENUME", "cnp": "CNP", "data_nasterii": "DATA_NASTERII", "grad_rudenie": "GRAD_RUDENIE", "observatii": "OBSERVATII", "activ": "ACTIV"})
                        else:
                            df_dep = pd.DataFrame(columns=["ID", "EMPLOYEE_ID", "TIP", "NUME", "PRENUME", "CNP", "DATA_NASTERII", "GRAD_RUDENIE", "OBSERVATII", "ACTIV"])

                        st.markdown('<div class="emp-block-title">Listă</div>', unsafe_allow_html=True)
                        if df_dep.empty:
                            st.info("Nu există persoane în întreținere.")
                        else:
                            # Tabel cu coloană „Select” (bifă): bifezi rândul și se completează detaliile jos
                            df_edit = df_dep[["ID", "TIP", "NUME", "PRENUME", "CNP", "DATA_NASTERII", "GRAD_RUDENIE", "ACTIV"]].copy()
                            df_edit["TIP"] = df_edit["TIP"].map(lambda x: TIP_LABEL.get(_norm_tip(x), str(x or "").strip()))
                            df_edit = df_edit.rename(columns={"DATA_NASTERII": "DATA NAȘTERII", "GRAD_RUDENIE": "GRAD"})
                            df_edit.insert(0, "Select", df_dep["ID"].values == (st.session_state.get(f"{K}edit_id") or 0))
                            # Doar bifa e editabilă; restul coloanelor doar pentru vizualizare
                            edited_dep = st.data_editor(
                                df_edit,
                                column_config={
                                    "Select": st.column_config.CheckboxColumn("Selectat", help="Bifează pentru a încărca detaliile în formular"),
                                    "ID": st.column_config.NumberColumn("ID", disabled=True),
                                    "TIP": st.column_config.TextColumn("Tip", disabled=True),
                                    "NUME": st.column_config.TextColumn("Nume", disabled=True),
                                    "PRENUME": st.column_config.TextColumn("Prenume", disabled=True),
                                    "CNP": st.column_config.TextColumn("CNP", disabled=True),
                                    "DATA NAȘTERII": st.column_config.TextColumn("Data nașterii", disabled=True),
                                    "GRAD": st.column_config.TextColumn("Grad", disabled=True),
                                    "ACTIV": st.column_config.NumberColumn("Activ", disabled=True),
                                },
                                key=f"{K}table_editor",
                                use_container_width=True,
                                height=260,
                            )
                            # Din bifă: care rând e selectat → setăm edit_id și rerun ca să se încarce detaliile
                            selected = edited_dep[edited_dep["Select"]]
                            new_edit_id = int(selected["ID"].iloc[0]) if len(selected) >= 1 else None
                            cur_dep_id_before = st.session_state.get(f"{K}edit_id")
                            if new_edit_id != cur_dep_id_before:
                                st.session_state[f"{K}edit_id"] = new_edit_id
                                st.rerun()

                        cur_dep_id = st.session_state.get(f"{K}edit_id")
                        prev = st.session_state.get(f"{K}prev_edit_id")
                        if st.session_state.get(f"{K}edit_id") != prev:
                            st.session_state[f"{K}prev_edit_id"] = st.session_state.get(f"{K}edit_id")
                            for kk in (f"{K}tip", f"{K}nume", f"{K}prenume", f"{K}cnp", f"{K}dob", f"{K}grad", f"{K}obs", f"{K}activ"):
                                st.session_state.pop(kk, None)

                        current = {
                            "tip": "", "nume": "", "prenume": "", "cnp": "",
                            "data_nasterii": "", "grad_rudenie": "", "observatii": "", "activ": 1,
                        }
                        if cur_dep_id and not df_dep.empty:
                            row = _dependent_get(conn, cur_dep_id)
                            if row:
                                current = {
                                    "tip": _norm_tip(row["tip"]),
                                    "nume": row["nume"],
                                    "prenume": row["prenume"],
                                    "cnp": row["cnp"],
                                    "data_nasterii": row["data_nasterii"],
                                    "grad_rudenie": row["grad_rudenie"],
                                    "observatii": row["observatii"],
                                    "activ": row["activ"],
                                }

                        st.divider()
                        st.markdown('<div class="emp-block-title">Detalii</div>', unsafe_allow_html=True)

                        # --- FORMULAR Detalii: restul pe 2 coloane, Grad rudenie + Observații în stânga ---
                        counter = st.session_state[f"{K}key_counter"]

                        with st.form(key=f"{K}form_{counter}"):
                            with st.container(key=f"{K}details_box_{counter}"):
                                det_col, _det_sp = st.columns([1, 1])
                                with det_col:
                                    st.text_input("ID (auto)", value=str(cur_dep_id or ""), disabled=True, key=f"{K}id_{counter}")

                                    tip_cur = _norm_tip(current.get("tip"))
                                    tip_idx = TIP_OPTS.index(tip_cur) if tip_cur in TIP_OPTS else 0
                                    tip = st.selectbox(
                                        "Tip",
                                        options=TIP_OPTS,
                                        index=tip_idx,
                                        format_func=lambda x: TIP_LABEL.get(x, x),
                                        key=f"{K}tip_{counter}",
                                    )

                                    nume = st.text_input("Nume", value=current.get("nume", ""), key=f"{K}nume_{counter}")
                                    prenume = st.text_input("Prenume", value=current.get("prenume", ""), key=f"{K}prenume_{counter}")
                                    cnp = st.text_input(
                                        "CNP",
                                        value=current.get("cnp", ""),
                                        key=f"{K}cnp_{counter}",
                                        help="CNP-ul trebuie să aibă 13 cifre (structură S YYMMDD JJ NNN C). "
                                             "Se validează automat data nașterii, sexul, județul și cifra de control.",
                                    )
                                    cnp_cleaned = cnp_clean(cnp)
                                    if str(cnp).strip():
                                        ok_dep_cnp, msg_dep_cnp, _ = cnp_validate(cnp)
                                        if ok_dep_cnp:
                                            st.caption("CNP valid.")
                                        else:
                                            st.caption(f"CNP invalid: {msg_dep_cnp}")

                                    cnp_prefill = decode_cnp(current.get("cnp", "") or "", strict_county=False)
                                    dob_from_db = _parse_date_any(current.get("data_nasterii", ""))
                                    dob_default = (
                                        cnp_prefill.birth_date if cnp_prefill.valid and cnp_prefill.birth_date
                                        else (dob_from_db or datetime.date.today())
                                    )
                                    import datetime as _dt_dob
                                    today_dob = _dt_dob.date.today()
                                    min_dob = today_dob.replace(year=today_dob.year - 100)
                                    max_dob = today_dob.replace(year=today_dob.year + 10)
                                    dob = st.text_input(
                                        "Data nașterii",
                                        value=(dob_default.isoformat() if dob_default else ""),
                                        key=f"{K}dob_{counter}",
                                        help="Format acceptat: YYYY-MM-DD (ex: 1990-05-21).",
                                    )
                                    activ = st.checkbox("Activ", value=bool(current.get("activ", 1)), key=f"{K}activ_{counter}")

                                    # Grad rudenie și Observații (aliniere stânga, coloană îngustă)
                                    grad_cur = (current.get("grad_rudenie") or "").strip().upper()
                                    if grad_cur not in GRAD_OPTS and grad_cur:
                                        st.caption("Gradul existent nu e standard (I/II/III). Îl poți păstra la Observații.")
                                        grad_cur = ""
                                    suggested = _suggest_grad_for_tip(tip)
                                    grad_default = grad_cur or suggested
                                    grad_idx = GRAD_OPTS.index(grad_default) if grad_default in GRAD_OPTS else 0
                                    grad_rudenie = st.selectbox(
                                        "Grad rudenie",
                                        options=GRAD_OPTS,
                                        index=grad_idx,
                                        format_func=lambda x: GRAD_LABEL.get(x, x),
                                        key=f"{K}grad_{counter}",
                                    )

                                    observatii = st.text_area("Observații", value=current.get("observatii", ""), height=140, key=f"{K}obs_{counter}")

                                    is_new = cur_dep_id is None
                                    st.markdown('<div class="emp-block-title">Acțiuni</div>', unsafe_allow_html=True)
                                    save = st.form_submit_button(
                                        "Adaugă în listă" if is_new else "Salvează modificările",
                                        key=f"dep_action_btn_save_{emp_id}_{counter}",
                                        use_container_width=False,
                                    )
                                    cancel = st.form_submit_button(
                                        "Anulează",
                                        key=f"dep_action_btn_cancel_{emp_id}_{counter}",
                                        use_container_width=False,
                                    )
                                    delete = st.form_submit_button(
                                        "Dezactivează persoana",
                                        disabled=is_new,
                                        key=f"dep_action_btn_delete_{emp_id}_{counter}",
                                        use_container_width=False,
                                    )

                        if cancel:
                            st.session_state[f"{K}reset_pending"] = True
                            st.rerun()

                        if save:
                            dob_date = _parse_date_any(str(dob))
                            if not dob_date:
                                st.error("Data nașterii este invalidă. Folosește formatul YYYY-MM-DD.")
                                st.stop()
                            if dob_date < min_dob or dob_date > max_dob:
                                st.error("Data nașterii este în afara intervalului permis.")
                                st.stop()

                            if str(cnp).strip():
                                ok_dep_cnp, msg_dep_cnp, _ = cnp_validate(cnp)
                                if not ok_dep_cnp:
                                    st.error(f"CNP invalid: {msg_dep_cnp}")
                                    st.stop()

                            dep_values = {
                                "employee_id": emp_id,
                                "tip": tip,
                                "nume": nume,
                                "prenume": prenume,
                                "cnp": cnp_cleaned if str(cnp).strip() else "",
                                "data_nasterii": dob_date.isoformat() if dob_date else "",
                                "grad_rudenie": grad_rudenie,
                                "observatii": observatii,
                                "activ": 1 if activ else 0,
                            }
                            _dependent_upsert(conn, emp_id, dep_values, dep_id=int(cur_dep_id) if cur_dep_id else None)
                            load_employee_bundle.clear()
                            st.success("Adăugat în listă." if cur_dep_id is None else "Salvat.")
                            st.session_state[f"{K}reset_pending"] = True
                            st.rerun()

                        if delete and cur_dep_id:
                            _dependent_soft_delete(conn, int(cur_dep_id))
                            load_employee_bundle.clear()
                            st.success("Dezactivat (soft delete).")
                            st.session_state[f"{K}reset_pending"] = True
                            st.rerun()
                # ---------------------------------------------------------------------------------#
                # TAB 2: SECTIUNEA VECHIME (instituție + anterioară + muncă + funcție)             # 
                # ---------------------------------------------------------------------------------#
    
                with tabs[2]:
                    st.markdown('<div class="emp-block-title">Vechime</div>', unsafe_allow_html=True)
                    st.caption("Gestionare vechime instituție, vechime anterioară în muncă și overview + raport vechime PDF.")
                    _render_info_section(
                        "Context vechime",
                        [
                            ("Marcă", _emp_v(emp, "marca")),
                            ("Data angajare", _emp_v(emp, "data_angajarii", _emp_v(emp, "data_angajare"))),
                            ("Data plecării", _emp_v(emp, "data_plecarii")),
                            ("Vechime anterioară", f"{_emp_v(emp, 'vechime_anterioara_ani', '0')} ani, {_emp_v(emp, 'vechime_anterioara_luni', '0')} luni"),
                        ],
                    )

                    # =========================
                    # UTILITARE
                    # =========================

                    # calc vechime calendaristică (ani, luni, zile)
                    def _ymd_diff(start: datetime.date, end: datetime.date):
                        if not start or not end or start > end:
                            return 0, 0, 0

                        ani = end.year - start.year
                        luni = end.month - start.month
                        zile = end.day - start.day

                        if zile < 0:
                            luni -= 1
                            prev_month = end.month - 1 or 12
                            prev_year = end.year if end.month != 1 else end.year - 1
                            zile += calendar.monthrange(prev_year, prev_month)[1]

                        if luni < 0:
                            ani -= 1
                            luni += 12

                        return ani, luni, zile

                    def _i(v, d=0):
                        try:
                            return int(v) if v not in (None, "", "None") else int(d)
                        except Exception:
                            return int(d)

                    def _parse_date_any(s: str):
                        s = (s or "").strip()
                        if not s:
                            return None
                        for fmt in ("%Y-%m-%d", "%d.%m.%Y", "%d/%m/%Y"):
                            try:
                                return datetime.datetime.strptime(s, fmt).date()
                            except Exception:
                                pass
                        return None

                    # diana – PDF builder (ReportLab)
                    def _build_vechime_pdf_bytes(
                        emp: dict,
                        report_type: str,
                        data_ang: str,
                        data_plec: str,
                        still_active: bool,
                        inst_y: int,
                        inst_m: int,
                        inst_d: int,
                        inst_den: int,
                        ant_y: int,
                        ant_m: int,
                        total_y: int,
                        total_m: int,
                    ) -> bytes:
                        buffer = BytesIO()
                        c = canvas.Canvas(buffer, pagesize=A4)
                        width, height = A4

                        y_cursor = height - 60
                        c.setFont("Helvetica-Bold", 14)
                        c.drawString(50, y_cursor, "Raport Vechime")
                        y_cursor -= 24

                        c.setFont("Helvetica", 10)
                        full_name = f"{emp.get('last_name','')} {emp.get('first_name','')}".strip()
                        marca = str(emp.get("marca", "") or "")
                        c.drawString(50, y_cursor, f"Nume: {full_name}")
                        y_cursor -= 14
                        c.drawString(50, y_cursor, f"Marcă: {marca}")
                        y_cursor -= 14
                        c.drawString(50, y_cursor, f"Tip raport: {report_type}")
                        y_cursor -= 22

                        c.setFont("Helvetica-Bold", 11)
                        c.drawString(50, y_cursor, "Date angajare / plecare")
                        y_cursor -= 14
                        c.setFont("Helvetica", 10)
                        c.drawString(50, y_cursor, f"Data angajării: {data_ang or '—'}")
                        y_cursor -= 14
                        if still_active:
                            c.drawString(50, y_cursor, "Data plecării: — (încă activ)")
                        else:
                            c.drawString(50, y_cursor, f"Data plecării: {data_plec or '—'}")
                        y_cursor -= 22

                        if report_type == "Vechime în instituție":
                            c.setFont("Helvetica-Bold", 11)
                            c.drawString(50, y_cursor, "Vechime în instituție")
                            y_cursor -= 14
                            c.setFont("Helvetica", 10)
                            c.drawString(50, y_cursor, f"{inst_y} ani, {inst_m} luni, {inst_d}/{inst_den} zile")
                            y_cursor -= 22

                        if report_type == "Vechime în muncă (total)":
                            c.setFont("Helvetica-Bold", 11)
                            c.drawString(50, y_cursor, "Vechime în muncă (total)")
                            y_cursor -= 14
                            c.setFont("Helvetica", 10)
                            c.drawString(50, y_cursor, f"{total_y} ani, {total_m} luni")
                            y_cursor -= 18
                            c.drawString(50, y_cursor, f"Include: anterioară ({ant_y} ani, {ant_m} luni) + instituție ({inst_y} ani, {inst_m} luni).")
                            y_cursor -= 22

                        c.setFont("Helvetica-Oblique", 8)
                        c.drawString(50, 40, "Generat din aplicația SOCRATES HR.")
                        c.showPage()
                        c.save()

                        buffer.seek(0)
                        return buffer.read()

                    # =========================
                    # DATE DIN DB
                    # =========================

                    v_inst_ani = _i(emp.get("vechime_inst_ani", 0))
                    v_inst_luni = _i(emp.get("vechime_inst_luni", 0))
                    v_inst_fn = _i(emp.get("vechime_inst_fract_num", 0))
                    v_inst_fd = _i(emp.get("vechime_inst_fract_den", 30)) or 30

                    v_ant_ani = _i(emp.get("vechime_anterioara_ani", 0))
                    v_ant_luni = _i(emp.get("vechime_anterioara_luni", 0))

                    v_functie_ani = _i(emp.get("vechime_functie_ani", 0))
                    v_functie_luni = _i(emp.get("vechime_functie_luni", 0))

                    data_angaj = (emp.get("data_angajarii") or emp.get("data_angajare") or "").strip()
                    data_plecarii = (emp.get("data_plecarii") or "").strip()

                    # =========================
                    # 1) DATE ANGAJARE / PLECARE (LIVE)
                    # =========================

                    st.markdown('<div class="s-title">Date angajare / plecare</div>', unsafe_allow_html=True)

                    _ang_default = _parse_date_any(data_angaj) or datetime.date.today()
                    _today = datetime.date.today()
                    _min_d = _today.replace(year=_today.year - 100)
                    _max_d = _today.replace(year=_today.year + 10)
                    data_ang_date = st.date_input(
                        "Data angajării",
                        value=_ang_default,
                        min_value=_min_d,
                        max_value=_max_d,
                        key=f"data_ang_date_{emp_id}",
                    )
                    still_active = st.checkbox(
                        "Încă activ",
                        value=(data_plecarii == ""),
                        key=f"still_active_{emp_id}",
                    )
                    _plec_default = _parse_date_any(data_plecarii) or datetime.date.today()
                    data_plec_date = st.date_input(
                        "Data plecării",
                        value=_plec_default,
                        min_value=_min_d,
                        max_value=_max_d,
                        disabled=still_active,
                        key=f"data_plec_date_{emp_id}",
                    )

                    data_ang = data_ang_date.isoformat() if data_ang_date else ""
                    data_plec = "" if still_active else (data_plec_date.isoformat() if data_plec_date else "")

                    # end org-panel (wrapper eliminat din UI)

                    # =========================
                    # 2) AUTO CALC (LIVE)
                    # =========================

                    # diana – auto-calc vechime instituție + blocare câmpuri
                    auto_calc_inst = st.checkbox(
                        "Calculează automat vechimea în instituție",
                        value=True,
                        key=f"auto_calc_inst_{emp_id}",
                    )
                    lock_inst_fields = auto_calc_inst

                    # =========================
                    # 3) VECHIME ÎN INSTITUȚIE (LIVE PREVIEW + EDIT IN FORM)
                    # =========================

                    # diana – compute inst values for preview/overview
                    if auto_calc_inst and data_ang_date:
                        end_calc = datetime.date.today() if still_active else data_plec_date
                        inst_y, inst_m, inst_d = _ymd_diff(data_ang_date, end_calc)
                        inst_den = calendar.monthrange(end_calc.year, end_calc.month)[1]
                    else:
                        inst_y, inst_m, inst_d = int(v_inst_ani), int(v_inst_luni), int(v_inst_fn)
                        inst_den = int(v_inst_fd) if int(v_inst_fd) != 0 else 30

                    st.markdown('<div class="emp-block-title">Vechime în instituție</div>', unsafe_allow_html=True)
                    st.markdown('<div class="emp-vechime-note">Dacă bifezi calcul automat, câmpurile devin doar afișaj.</div>', unsafe_allow_html=True)

                    if auto_calc_inst and data_ang_date:
                        st.markdown(
                            f'<div class="emp-vechime-status"><strong>Vechime calculată:</strong> {inst_y} ani, {inst_m} luni, {inst_d} zile (salvat ca {inst_d}/{inst_den}).</div>',
                            unsafe_allow_html=True,
                        )
                    else:
                        st.markdown('<div class="emp-vechime-status">Calcul manual.</div>', unsafe_allow_html=True)

                    st.markdown("", unsafe_allow_html=True)

                    # =========================
                    # 4) VECHIME ANTERIOARĂ (SAVE SEPARAT + LOCK)
                    # =========================

                    st.markdown('<div class="emp-block-title">Vechime în muncă până la angajare</div>', unsafe_allow_html=True)
                    st.markdown('<div class="emp-vechime-note">Se salvează separat. Implicit este blocată până alegi „Modifică”.</div>', unsafe_allow_html=True)

                    # diana – lock/edit toggler pentru vechime anterioară
                    ant_edit_key = f"ant_edit_{emp_id}"
                    if ant_edit_key not in st.session_state:
                        st.session_state[ant_edit_key] = False

                    if not st.session_state[ant_edit_key]:
                        if st.button("Modifică", key=f"ant_edit_btn_{emp_id}"):
                            st.session_state[ant_edit_key] = True
                            st.rerun()
                    else:
                        if st.button("Blochează", key=f"ant_lock_btn_{emp_id}"):
                            st.session_state[ant_edit_key] = False
                            st.rerun()

                    disabled_ant = not st.session_state[ant_edit_key]

                    with st.form(key=f"ant_work_form_{emp_id}"):
                        ant_ani = st.number_input(
                            "Ani (anterioară)",
                            min_value=0,
                            max_value=60,
                            value=int(v_ant_ani),
                            step=1,
                            disabled=disabled_ant,
                            key=f"ant_ani_{emp_id}",
                        )
                        ant_luni = st.number_input(
                            "Luni (anterioară)",
                            min_value=0,
                            max_value=11,
                            value=int(v_ant_luni),
                            step=1,
                            disabled=disabled_ant,
                            key=f"ant_luni_{emp_id}",
                        )

                        save_ant = st.form_submit_button(
                            "Salvează vechime anterioară",
                            disabled=disabled_ant,
                        )

                    if save_ant:
                        upd_ant = dict(emp)
                        upd_ant.update(
                            {
                                "vechime_anterioara_ani": int(ant_ani),
                                "vechime_anterioara_luni": int(ant_luni),
                            }
                        )
                        _employee_upsert(conn, upd_ant, employee_id=emp_id)
                        st.session_state[ant_edit_key] = False
                        st.success("Vechime anterioară salvată.")
                        st.rerun()

                    # =========================
                    # 5) FORM PRINCIPAL: VECHIME INSTITUȚIE + FUNCȚIE (SAVE)
                    # =========================

                    st.markdown('<div class="s-title">Salvare vechime instituție / funcție</div>', unsafe_allow_html=True)
                    st.markdown('<div class="s-muted">Aici salvezi valorile instituției (manual sau auto-calc) și vechimea în funcție.</div>', unsafe_allow_html=True)

                    with st.form(key=f"vechime_form_{emp_id}"):

                        st.markdown('<div class="emp-block-title">🏢 Vechime în instituție (salvare)</div>', unsafe_allow_html=True)
                        inst_ani_input = st.number_input(
                            "Ani (instituție)",
                            min_value=0,
                            max_value=60,
                            value=int(v_inst_ani),
                            step=1,
                            disabled=lock_inst_fields,
                            key=f"inst_ani_{emp_id}",
                        )
                        inst_luni_input = st.number_input(
                            "Luni (instituție)",
                            min_value=0,
                            max_value=11,
                            value=int(v_inst_luni),
                            step=1,
                            disabled=lock_inst_fields,
                            key=f"inst_luni_{emp_id}",
                        )
                        inst_fn_input = st.number_input(
                            "Zile lună curentă",
                            min_value=0,
                            max_value=31,
                            value=int(v_inst_fn),
                            step=1,
                            disabled=lock_inst_fields,
                            key=f"inst_fn_{emp_id}",
                        )
                        inst_fd_input = st.number_input(
                            "Bază raportare lună (zile)",
                            min_value=28,
                            max_value=31,
                            value=int(v_inst_fd) if int(v_inst_fd) != 0 else 30,
                            step=1,
                            disabled=lock_inst_fields,
                            key=f"inst_fd_{emp_id}",
                        )

                        st.markdown('<div class="emp-block-title">Vechime în funcție</div>', unsafe_allow_html=True)
                        functie_ani = st.number_input(
                            "Ani (funcție)",
                            min_value=0,
                            max_value=60,
                            value=int(v_functie_ani),
                            step=1,
                            key=f"functie_ani_{emp_id}",
                        )
                        functie_luni = st.number_input(
                            "Luni (funcție)",
                            min_value=0,
                            max_value=11,
                            value=int(v_functie_luni),
                            step=1,
                            key=f"functie_luni_{emp_id}",
                        )

                        save_main = st.form_submit_button("Salvează (instituție + funcție)")

                    if save_main:
                        # diana – validare: data plecării nu poate fi înainte de data angajării
                        if (not still_active) and data_plec_date and data_ang_date and data_plec_date < data_ang_date:
                            st.error("Data plecării nu poate fi înainte de data angajării.")
                            st.stop()

                        # diana – dacă auto-calc e ON, suprascriem valorile care se salvează
                        if auto_calc_inst and data_ang_date:
                            end_calc = datetime.date.today() if still_active else data_plec_date
                            y, m, d = _ymd_diff(data_ang_date, end_calc)
                            den = calendar.monthrange(end_calc.year, end_calc.month)[1]
                            inst_ani_to_save = int(y)
                            inst_luni_to_save = int(m)
                            inst_fn_to_save = int(d)
                            inst_fd_to_save = int(den)
                        else:
                            inst_ani_to_save = int(inst_ani_input)
                            inst_luni_to_save = int(inst_luni_input)
                            inst_fn_to_save = int(inst_fn_input)
                            inst_fd_to_save = int(inst_fd_input) if int(inst_fd_input) != 0 else 30

                        upd = dict(emp)
                        upd.update(
                            {
                                "vechime_inst_ani": inst_ani_to_save,
                                "vechime_inst_luni": inst_luni_to_save,
                                "vechime_inst_fract_num": inst_fn_to_save,
                                "vechime_inst_fract_den": inst_fd_to_save,
                                "data_angajarii": str(data_ang).strip(),
                                "data_plecarii": str(data_plec).strip(),
                                "vechime_functie_ani": int(functie_ani),
                                "vechime_functie_luni": int(functie_luni),
                            }
                        )
                        _employee_upsert(conn, upd, employee_id=emp_id)
                        st.success("Vechime instituție + funcție salvată.")
                        st.rerun()

                    # =========================
                    # 6) OVERVIEW TOTAL (VIZUAL)
                    # =========================

                    # diana – total vechime (ANI/LUNI) = anterioară + instituție (fără conversie zile->luni)
                    total_months = (int(v_ant_ani) * 12 + int(v_ant_luni)) + (int(inst_y) * 12 + int(inst_m))
                    total_y = total_months // 12
                    total_m = total_months % 12

                    st.markdown('<div class="emp-block-title">📊 Sinteză vechime</div>', unsafe_allow_html=True)
                    st.markdown("", unsafe_allow_html=True)

                    st.markdown("**1) Vechime în instituție**")
                    st.markdown(
                        (
                            '<div class="emp-overview-pills">'
                            f'<span class="s-pill">{int(inst_y)} ani</span>'
                            f'<span class="s-pill">{int(inst_m)} luni</span>'
                            f'<span class="s-pill">{int(inst_d)}/{int(inst_den)} zile</span>'
                            '</div>'
                        ),
                        unsafe_allow_html=True,
                    )
                    st.markdown(
                        f'<div class="s-muted">Din {data_ang or "—"} până la {("azi" if still_active else (data_plec or "—"))}</div>',
                        unsafe_allow_html=True,
                    )
                    st.markdown('<div class="emp-overview-sep"></div>', unsafe_allow_html=True)
                    st.markdown("**2) Vechime anterioară (muncă)**")
                    st.markdown(
                        (
                            '<div class="emp-overview-pills">'
                            f'<span class="s-pill">{int(v_ant_ani)} ani</span>'
                            f'<span class="s-pill">{int(v_ant_luni)} luni</span>'
                            '</div>'
                        ),
                        unsafe_allow_html=True,
                    )
                    st.markdown('<div class="s-muted">Completată manual (dosar profesional).</div>', unsafe_allow_html=True)
                    st.markdown('<div class="emp-overview-sep"></div>', unsafe_allow_html=True)
                    st.markdown("**3) Vechime totală (muncă)**")
                    st.markdown(
                        (
                            '<div class="emp-overview-pills">'
                            f'<span class="s-pill">{int(total_y)} ani</span>'
                            f'<span class="s-pill">{int(total_m)} luni</span>'
                            '</div>'
                        ),
                        unsafe_allow_html=True,
                    )
                    st.markdown("", unsafe_allow_html=True)

                    # =========================
                    # 7) RAPORT PDF
                    # =========================

                    st.markdown('<div class="s-title">📄 Raport</div>', unsafe_allow_html=True)
                    st.markdown('<div class="s-muted">Generează un raport PDF pentru dosar profesional / raportări.</div>', unsafe_allow_html=True)
                    st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)

                    rt = st.selectbox(
                        "Alege tip raport",
                        ["Vechime în instituție", "Vechime în muncă (total)"],
                        key=f"vechime_report_type_{emp_id}",
                    )

                    gen = st.button("🧾 Generează raport PDF", key=f"gen_report_btn_{emp_id}")

                    if gen:
                        pdf_bytes = _build_vechime_pdf_bytes(
                            emp=emp,
                            report_type=rt,
                            data_ang=data_ang,
                            data_plec=data_plec,
                            still_active=still_active,
                            inst_y=int(inst_y),
                            inst_m=int(inst_m),
                            inst_d=int(inst_d),
                            inst_den=int(inst_den),
                            ant_y=int(v_ant_ani),
                            ant_m=int(v_ant_luni),
                            total_y=int(total_y),
                            total_m=int(total_m),
                        )
                        filename = f"raport_vechime_{emp_id}_{rt.replace(' ', '_').lower()}.pdf"
                        st.download_button(
                            "⬇️ Descarcă PDF",
                            data=pdf_bytes,
                            file_name=filename,
                            mime="application/pdf",
                            key=f"dl_report_{emp_id}",
                        )

                    st.markdown("", unsafe_allow_html=True)

                # ---------------------------
                # TAB 3: Editează (formular complet)
                # ---------------------------
                # TAB 7: Editează (formular complet)
                # ---------------------------
                with tabs[7]:
                    st.info(
                        "Editezi datele principale ale angajatului. Modificările salvate se propagă automat în modulele care citesc "
                        "din tabela employees (Organigramă, Dosar profesional, Pontaj, Centralizator concedii)."
                    )
                    with st.container(key="ang_edit_panel"):
                        with st.form("ang_edit_form"):
                            values = _render_employee_form_fields(
                                emp,
                                cols,
                                prefix=f"ang_edit_{emp_id}",
                                single_column=True,
                                fine_section_titles=True,
                                parallel_groups=True,
                            )
                            submitted = st.form_submit_button("💾 Salvează modificările")
                            cancel = st.form_submit_button("↩ Renunță")

                    if cancel:
                        st.rerun()

                    if submitted:
                        new_id = _employee_upsert(conn, values, employee_id=emp_id)
                        st.success("Salvat.")
                        st.session_state["ang_selected_id"] = new_id
                        st.rerun()

                # ---------------------------
                # TAB 4: Șterge (soft delete)
                # ---------------------------

                # ---------------------------
                # TAB 3: Contract de muncă (CIM)
                # ---------------------------
                with tabs[3]:
                    st.subheader("📄 Contract de muncă (CIM)")
                    _render_info_section(
                        "Date contract curente",
                        [
                            ("Marcă", _emp_v(emp, "marca")),
                            ("Funcție", _emp_v(emp, "functie")),
                            ("Cod COR", _emp_v(emp, "cod_cor")),
                            ("Tip contract", _emp_v(emp, "tip_contract")),
                        ],
                    )

                    # =========================================================
                    # Formular CIM (structurat) – câmpurile care trebuie să apară
                    # =========================================================
                    # mapping trebuie să existe înainte de formular (fallback sigur)
                    try:
                        mapping
                    except Exception:
                        mapping = {}

                    # Config (pentru autocompletare angajator / antet)
                    cfg_local = load_config()
                    def _cfg_pick(*keys: str) -> str:
                        for k in keys:
                            v = (cfg_local.get(k) or "")
                            if str(v).strip():
                                return str(v).strip()
                        return ""


                    if "cim_form" not in st.session_state:
                        st.session_state["cim_form"] = {}
                    _cim = st.session_state["cim_form"].setdefault(str(emp_id), {})

                    def _cim_get(key, default=""):
                        v = _cim.get(key)
                        return default if (v is None or v == "") else v

                    # Data nașterii din CNP (fallback)
                    _cnp_here = str(emp.get("cnp") or emp.get("CNP") or "").strip()
                    _dob_auto = _format_ro_date(cnp_birthdate(_cnp_here))

                    st.markdown('<div class="emp-section-title">I. Datele părților (Angajator și Angajat)</div>', unsafe_allow_html=True)
                    if True:
                        st.markdown('<div class="emp-subsection-title">A) Angajator (din Configurație)</div>', unsafe_allow_html=True)
                        _cim["ANGAJATOR_DEN"] = st.text_input("Denumire firmă", value=_cim_get("ANGAJATOR_DEN", _cfg_pick("denumire_unitate")), key=f"cim_{emp_id}_ANGAJATOR_DEN")
                        _cim["ANGAJATOR_SEDIU"] = st.text_input("Sediu social", value=_cim_get("ANGAJATOR_SEDIU", _cfg_pick("adresa")), key=f"cim_{emp_id}_ANGAJATOR_SEDIU")
                        _cim["ANGAJATOR_CUI"] = st.text_input("CUI", value=_cim_get("ANGAJATOR_CUI", _cfg_pick("cui")), key=f"cim_{emp_id}_ANGAJATOR_CUI")
                        _cim["ANGAJATOR_REG"] = st.text_input("Nr. Reg. Comerțului", value=_cim_get("ANGAJATOR_REG", _cfg_pick("reg_comertului","nr_reg_comertului")), key=f"cim_{emp_id}_ANGAJATOR_REG")
                        _cim["ANGAJATOR_REPREZENTANT"] = st.text_input("Reprezentant legal", value=_cim_get("ANGAJATOR_REPREZENTANT", _cfg_pick("conducator_nume")), key=f"cim_{emp_id}_ANGAJATOR_REPREZENTANT")
                        _cim["ANGAJATOR_CALITATE"] = st.text_input("Calitate", value=_cim_get("ANGAJATOR_CALITATE", _cfg_pick("conducator_functie")), key=f"cim_{emp_id}_ANGAJATOR_CALITATE")
                        _cim["ANGAJATOR_TEL"] = st.text_input("Telefon", value=_cim_get("ANGAJATOR_TEL", _cfg_pick("telefon")), key=f"cim_{emp_id}_ANGAJATOR_TEL")
                        _cim["ANGAJATOR_PUNCT"] = st.text_input("Punct de lucru (opțional)", value=_cim_get("ANGAJATOR_PUNCT", _cfg_pick("punct_lucru")), key=f"cim_{emp_id}_ANGAJATOR_PUNCT")

                        st.markdown('<div class="emp-subsection-title">B) Angajat (din EMPLOYEES + DATE_ANGAJATI)</div>', unsafe_allow_html=True)
                        _cim["SALARIAT_NUME"] = st.text_input("Nume", value=_cim_get("SALARIAT_NUME", str(get_val(emp, "last_name","LAST_NAME","nume","NUME") or "")), key=f"cim_{emp_id}_SALARIAT_NUME")
                        _cim["SALARIAT_PRENUME"] = st.text_input("Prenume", value=_cim_get("SALARIAT_PRENUME", str(get_val(emp, "first_name","FIRST_NAME","prenume","PRENUME") or "")), key=f"cim_{emp_id}_SALARIAT_PRENUME")
                        _cim["CNP"] = st.text_input("CNP", value=_cim_get("CNP", _cnp_here), key=f"cim_{emp_id}_CNP")

                        st.markdown("**Domiciliu**")
                        _cim["SALARIAT_LOCALITATE"] = st.text_input("Localitate", value=_cim_get("SALARIAT_LOCALITATE", mapping.get("SALARIAT_LOCALITATE","")), key=f"cim_{emp_id}_SALARIAT_LOCALITATE")
                        _cim["SALARIAT_JUDET"] = st.text_input("Județ", value=_cim_get("SALARIAT_JUDET", mapping.get("SALARIAT_JUDET","")), key=f"cim_{emp_id}_SALARIAT_JUDET")
                        _cim["SALARIAT_STRADA"] = st.text_input("Strada", value=_cim_get("SALARIAT_STRADA", mapping.get("SALARIAT_STRADA","")), key=f"cim_{emp_id}_SALARIAT_STRADA")
                        _cim["SALARIAT_NR_STR"] = st.text_input("Nr.", value=_cim_get("SALARIAT_NR_STR", mapping.get("SALARIAT_NR_STR","")), key=f"cim_{emp_id}_SALARIAT_NR_STR")

                        st.markdown("**CI / Pașaport**")
                        _cim["CI_TIP"] = st.text_input("Tip act", value=_cim_get("CI_TIP", mapping.get("CI_TIP","CI")), key=f"cim_{emp_id}_CI_TIP")
                        _cim["CI_SERIE"] = st.text_input("Serie", value=_cim_get("CI_SERIE", mapping.get("CI_SERIE","")), key=f"cim_{emp_id}_CI_SERIE")
                        _cim["CI_NR"] = st.text_input("Număr", value=_cim_get("CI_NR", mapping.get("CI_NR","")), key=f"cim_{emp_id}_CI_NR")
                        _cim["CI_ELIBERAT_DE"] = st.text_input("Eliberat de", value=_cim_get("CI_ELIBERAT_DE", mapping.get("CI_ELIBERAT_DE","")), key=f"cim_{emp_id}_CI_ELIBERAT_DE")
                        _cim["CI_DATA"] = st.text_input("Data eliberării (dd.mm.yyyy)", value=_cim_get("CI_DATA", mapping.get("CI_DATA","")), key=f"cim_{emp_id}_CI_DATA")

                        _cim["DATA_NASTERE"] = st.text_input("Data nașterii (dd.mm.yyyy) – auto din CNP", value=_cim_get("DATA_NASTERE", mapping.get("DATA_NASTERE","") or _dob_auto), key=f"cim_{emp_id}_DATA_NASTERE")
                        _cim["STUDII"] = st.text_input("Studii", value=_cim_get("STUDII", mapping.get("STUDII","")), key=f"cim_{emp_id}_STUDII")
                        _cim["OCUPATIE_TEXT"] = st.text_input("Ocupație (text)", value=_cim_get("OCUPATIE_TEXT", mapping.get("OCUPATIE_TEXT","")), key=f"cim_{emp_id}_OCUPATIE_TEXT")

                    st.markdown('<div class="emp-section-title">II. Detalii contractuale (obligatorii)</div>', unsafe_allow_html=True)
                    if True:
                        st.markdown('<div class="emp-subsubsection-title">A) Identificare contract</div>', unsafe_allow_html=True)
                        _cim["CIM_NR"] = st.text_input("Număr CIM", value=_cim_get("CIM_NR", mapping.get("CIM_NR","")), key=f"cim_{emp_id}_CIM_NR")
                        _cim["CIM_DATA"] = st.text_input("Data încheierii/semnării (dd.mm.yyyy)", value=_cim_get("CIM_DATA", mapping.get("CIM_DATA","")), key=f"cim_{emp_id}_CIM_DATA")

                        st.markdown('<div class="emp-subsubsection-title">B) Obiectul contractului</div>', unsafe_allow_html=True)
                        _cim["FUNCTIE"] = st.text_input("Funcție (denumire)", value=_cim_get("FUNCTIE", mapping.get("FUNCTIE","")), key=f"cim_{emp_id}_FUNCTIE")
                        _cim["COR"] = st.text_input("Cod COR", value=_cim_get("COR", mapping.get("COR","")), key=f"cim_{emp_id}_COR")
                        _cim["SARCINI_SUMAR"] = st.text_area("Descriere succintă sarcini", value=_cim_get("SARCINI_SUMAR", mapping.get("SARCINI_SUMAR","")), key=f"cim_{emp_id}_SARCINI_SUMAR", height=120)

                        st.markdown('<div class="emp-subsubsection-title">C) Locul muncii</div>', unsafe_allow_html=True)
                        _cim["LOC_MUNCA_ADRESA"] = st.text_input("Adresa exactă", value=_cim_get("LOC_MUNCA_ADRESA", mapping.get("LOC_MUNCA_ADRESA","")), key=f"cim_{emp_id}_LOC_MUNCA_ADRESA")
                        _cim["LOC_MUNCA_FARA_FIX"] = st.checkbox("Fără loc fix", value=bool(_cim_get("LOC_MUNCA_FARA_FIX", False)), key=f"cim_{emp_id}_LOC_MUNCA_FARA_FIX")
                        _cim["LOC_MUNCA_TEXT"] = st.text_input("Detalii (dacă e fără loc fix)", value=_cim_get("LOC_MUNCA_TEXT", mapping.get("LOC_MUNCA_TEXT","")), key=f"cim_{emp_id}_LOC_MUNCA_TEXT")

                        st.markdown('<div class="emp-subsubsection-title">D) Durata</div>', unsafe_allow_html=True)
                        _cim["DURATA_TIP"] = st.selectbox("Tip", options=["Nedeterminată", "Determinată"], index=0 if _cim_get("DURATA_TIP","Nedeterminată")=="Nedeterminată" else 1, key=f"cim_{emp_id}_DURATA_TIP")
                        _cim["DATA_START_DET"] = st.text_input("Data start determinată (dd.mm.yyyy)", value=_cim_get("DATA_START_DET", mapping.get("DATA_START_DET","")), key=f"cim_{emp_id}_DATA_START_DET")
                        _cim["DATA_END_DET"] = st.text_input("Data final determinată (dd.mm.yyyy)", value=_cim_get("DATA_END_DET", mapping.get("DATA_END_DET","")), key=f"cim_{emp_id}_DATA_END_DET")
                        _cim["DURATA_LUNI"] = st.text_input("Durata luni (opțional)", value=_cim_get("DURATA_LUNI", mapping.get("DURATA_LUNI","")), key=f"cim_{emp_id}_DURATA_LUNI")

                        st.markdown('<div class="emp-subsubsection-title">E) Program / timp de lucru</div>', unsafe_allow_html=True)
                        _cim["ORE_ZI"] = st.text_input("Ore/zi", value=_cim_get("ORE_ZI", mapping.get("ORE_ZI","8")), key=f"cim_{emp_id}_ORE_ZI")
                        _cim["ORE_SAPT"] = st.text_input("Ore/săptămână", value=_cim_get("ORE_SAPT", mapping.get("ORE_SAPT","40")), key=f"cim_{emp_id}_ORE_SAPT")
                        _cim["TIP_TIMP_LUCRU"] = st.selectbox("Timp de lucru", options=["Normal", "Redus"], index=0 if _cim_get("TIP_TIMP_LUCRU","Normal")=="Normal" else 1, key=f"cim_{emp_id}_TIP_TIMP_LUCRU")
                        _cim["PROGRAM_REPARTIZARE"] = st.text_input("Repartizare program (ex: L–V 08:00–16:00)", value=_cim_get("PROGRAM_REPARTIZARE", mapping.get("PROGRAM_REPARTIZARE","")), key=f"cim_{emp_id}_PROGRAM_REPARTIZARE")

                        st.markdown('<div class="emp-subsubsection-title">F) Perioada de probă</div>', unsafe_allow_html=True)
                        _cim["PROBA_ZILE"] = st.text_input("Zile probă", value=_cim_get("PROBA_ZILE", mapping.get("PROBA_ZILE","")), key=f"cim_{emp_id}_PROBA_ZILE")

                        st.markdown('<div class="emp-subsubsection-title">G) Salariu</div>', unsafe_allow_html=True)
                        _cim["SALARIU_BAZA"] = st.text_input("Salariu de bază brut", value=_cim_get("SALARIU_BAZA", mapping.get("SALARIU_BAZA","")), key=f"cim_{emp_id}_SALARIU_BAZA")
                        _cim["SPORURI_TEXT"] = st.text_input("Sporuri (text)", value=_cim_get("SPORURI_TEXT", mapping.get("SPORURI_TEXT","")), key=f"cim_{emp_id}_SPORURI_TEXT")
                        _cim["DATA_PLATA_SALARIU"] = st.text_input("Data plății (ex: 15)", value=_cim_get("DATA_PLATA_SALARIU", mapping.get("DATA_PLATA_SALARIU","")), key=f"cim_{emp_id}_DATA_PLATA_SALARIU")
                        _cim["MODALITATE_PLATA"] = st.text_input("Modalitate plată", value=_cim_get("MODALITATE_PLATA", mapping.get("MODALITATE_PLATA","Virament bancar")), key=f"cim_{emp_id}_MODALITATE_PLATA")

                        st.markdown('<div class="emp-subsubsection-title">H) Elemente specifice (clauze)</div>', unsafe_allow_html=True)
                        _cim["CLAUZA_CONFIDENTIALITATE"] = st.text_area("Confidențialitate", value=_cim_get("CLAUZA_CONFIDENTIALITATE", mapping.get("CLAUZA_CONFIDENTIALITATE","")), key=f"cim_{emp_id}_CLAUZA_CONFIDENTIALITATE", height=100)
                        _cim["CLAUZA_NECONCURENTA"] = st.text_area("Neconcurență", value=_cim_get("CLAUZA_NECONCURENTA", mapping.get("CLAUZA_NECONCURENTA","")), key=f"cim_{emp_id}_CLAUZA_NECONCURENTA", height=100)
                        _cim["CLAUZE_SPECIALE"] = st.text_area("Alte clauze speciale", value=_cim_get("CLAUZE_SPECIALE", mapping.get("CLAUZE_SPECIALE","")), key=f"cim_{emp_id}_CLAUZE_SPECIALE", height=120)

                    st.markdown('<div class="emp-block-title">III. Acte necesare (checklist)</div>', unsafe_allow_html=True)
                    if True:
                        _cim["ACT_OBS"] = st.text_area("Observații", value=_cim_get("ACT_OBS",""), key=f"cim_{emp_id}_ACT_OBS", height=80)



                        st.markdown('### 📎 Documente necesare (checklist și import PDF)')
                        st.caption('Bifele sunt interne (dosar/verificare). PDF-urile se salvează la angajat și pot fi descărcate/șterse.')

                        DOC_ROWS = [
                            ('CI', 'CI / Pașaport salariat'),
                            ('STUDII', 'Diplomă / Certificat studii'),
                            ('MEDICAL', 'Aviz / Certificat medical'),
                            ('ALTE', 'Alte documente'),
                        ]

                        for _dtype, _label in DOC_ROWS:
                            # Checkbox (stare internă) - rând 1
                            _key_chk = f'cim_{emp_id}_chk_{_dtype}'
                            _checked = bool(_cim.get(f'ACT_{_dtype}', False))
                            _checked = st.checkbox(_label, value=_checked, key=_key_chk)
                            _cim[f'ACT_{_dtype}'] = _checked

                            # Upload/Acțiuni - rând 2 (sub checkbox), uploader mai compact
                            c_up, c_dl, c_del = st.columns([2.0, 1.3, 0.7])
                            with c_up:
                                up = st.file_uploader(
                                    "Import PDF",
                                    type=['pdf'],
                                    key=f'cim_{emp_id}_up_{_dtype}',
                                    label_visibility="collapsed",
                                )
                                if up is not None:
                                    try:
                                        _bytes = up.getvalue()
                                        _emp_doc_save(conn, emp_id, _dtype, up.name, up.type or 'application/pdf', _bytes, display_name=_label)
                                        # dacă există document, bifăm automat
                                        _cim[f'ACT_{_dtype}'] = True
                                        st.success('Document încărcat.')
                                        st.rerun()
                                    except Exception as e:
                                        st.error(f'Eroare la încărcare: {e}')

                            _docs = [a for a in bundle.get("employee_attachments", []) if a.get("doc_type") == _dtype]
                            if _docs:
                                _last = _docs[0]
                                fn, mt, blob = _emp_doc_load(conn, _last['id'])
                                with c_dl:
                                    if blob is not None:
                                        st.download_button('⬇️ Descarcă', data=blob, file_name=fn or f'{_dtype}.pdf', mime=mt or 'application/pdf', key=f'cim_{emp_id}_dl_{_dtype}')
                                    else:
                                        st.write('')
                                with c_del:
                                    if st.button('🗑️', key=f'cim_{emp_id}_del_{_dtype}'):
                                        try:
                                            _emp_doc_delete(conn, _last['id'])
                                            st.success('Șters.')
                                            st.rerun()
                                        except Exception as e:
                                            st.error(f'Eroare la ștergere: {e}')
                            else:
                                with c_dl: st.write('')
                                with c_del: st.write('')
                            st.divider()

                    # Aplică valorile din formular peste mapping (au prioritate față de auto)
                    for _k, _v in (_cim or {}).items():
                        if _v is None:
                            continue
                        mapping[_k] = _v
                        mapping["{{%s}}" % _k] = _v

                    mapping.setdefault("SALARIAT_NUME_COMPLET", f"{mapping.get('SALARIAT_NUME','')} {mapping.get('SALARIAT_PRENUME','')}".strip())
                    # TAB 4: Reges Online (câmpuri conform șablon)
                    # ---------------------------

                    st.divider()
                    st.subheader('💾 Gestionare versiuni contract')
                    cim_title = st.text_input('Titlu/observații versiune (opțional)', value=_cim.get('CIM_TITLE',''), key=f'cim_{emp_id}_CIM_TITLE')
                    _cim['CIM_TITLE'] = cim_title
                    do_save_cim = st.button('💾 Salvează versiune CIM', key=f'cim_{emp_id}_save')
                    do_set_active_btn = st.button('⭐ Setează versiunea selectată ca activă', key=f'cim_{emp_id}_set_active')

                    if do_save_cim:
                        try:
                            payload = dict(_cim or {})
                            cid = _cim_contract_save(conn, emp_id, payload, cim_nr=str(_cim.get('CIM_NR','') or ''), cim_data=str(_cim.get('CIM_DATA','') or ''), title=str(cim_title or ''))
                            st.success(f'Versiune CIM salvată (ID={cid}).')
                        except Exception as e:
                            st.error(f'Eroare la salvare CIM: {e}')


                    # =========================================================
                    # =========================================================
                    # Generare din cod (fără șabloane) – Contract CIM + Act adițional
                    # =========================================================
                    st.markdown('<div class="emp-block-title">IV. Generează Contract CIM & Act adițional</div>', unsafe_allow_html=True)
                    if True:
                        cfg = load_config()
                        ang_den = (cfg.get("denumire_unitate") or "").strip()
                        ang_sediu = (cfg.get("adresa") or "").strip()
                        ang_cui = (cfg.get("cui") or "").strip()
                        rep_legal = (cfg.get("conducator_nume") or "").strip()
                        rep_functie = (cfg.get("conducator_functie") or "").strip()

                        # -----------------------------
                        # 1) Contract CIM (DOCX)
                        # -----------------------------
                        st.markdown('<div class="emp-subsubsection-title">📄 Contract individual de muncă (din cod)</div>', unsafe_allow_html=True)
                        st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)
                        do_gen_cim_code = st.button("🧾 Generează CIM (DOCX + salvare istoric)", key=f"cim_code_gen_{emp_id}")
                        do_gen_cim_code_pdf = st.button("📄 Generează și PDF", key=f"cim_code_gen_pdf_{emp_id}")

                        if do_gen_cim_code or do_gen_cim_code_pdf:
                            try:
                                sal_nume_complet = f"{(emp.get('nume') or '').strip()} {(emp.get('prenume') or '').strip()}".strip()
                                cim_nr = str(_cim.get("CIM_NR","") or "").strip()
                                cim_data = str(_cim.get("CIM_DATA","") or "").strip()
                                data_incepere = str(_cim.get("DATA_START","") or "").strip()
                                tip_contract = str(_cim.get("TIP_CONTRACT","") or "").strip()
                                durata_luni = str(_cim.get("DURATA_LUNI","") or "").strip()
                                data_end = str(_cim.get("DATA_END_DET","") or "").strip()
                                loc_munca = str(_cim.get("LOC_MUNCA","") or "").strip()
                                functie = str(_cim.get("FUNCTIE","") or "").strip()
                                cor = str(_cim.get("COD_COR","") or "").strip()
                                ore_zi = str(_cim.get("ORE_ZI","") or "").strip()
                                ore_sapt = str(_cim.get("ORE_SAPT","") or "").strip()
                                program_rep = str(_cim.get("PROGRAM_REPARTIZARE","") or "").strip()
                                concediu = str(_cim.get("CONCEDIU_ZILE","") or "").strip()
                                salariu_baza = str(_cim.get("SALARIU_BAZA","") or "").strip()
                                sporuri = str(_cim.get("SPORURI_TEXT","") or "").strip()
                                plata_data = str(_cim.get("DATA_PLATA_SALARIU","") or "").strip()
                                plata_mod = str(_cim.get("MODALITATE_PLATA","") or "").strip()

                                # Texte model (minim, dar coerent)
                                obiect = "Prestarea muncii în funcția stabilită prin prezentul contract."
                                if tip_contract.lower().startswith("n"):
                                    durata_text = f"Contract pe durată nedeterminată, salariatul urmând să înceapă activitatea la data de {data_incepere}."
                                else:
                                    _det = []
                                    if durata_luni:
                                        _det.append(f"de {durata_luni} luni")
                                    if data_incepere and data_end:
                                        _det.append(f"pe perioada cuprinsă între {data_incepere} și {data_end}")
                                    durata_text = "Contract pe durată determinată " + (", ".join(_det) if _det else "")

                                program_text = f"Normă întreagă, durata timpului de lucru fiind de {ore_zi or '-'} ore/zi, {ore_sapt or '-'} ore/săptămână. Repartizarea programului: {program_rep or '-'}."
                                concediu_text = f"Durata concediului anual de odihnă este de {concediu or '-'} zile lucrătoare."
                                salariu_text = f"Salariul de bază lunar brut: {salariu_baza or '-'} lei. Sporuri/alte elemente: {sporuri or '-'} . Data plății: {plata_data or '-'} . Modalitate: {plata_mod or '-'}."

                                # Domiciliu (fallback simplu)
                                domiciliu = (emp.get("adresa") or emp.get("domiciliu") or emp.get("localitate") or "").strip()
                                sal_cnp = (emp.get("cnp") or "").strip()

                                docx_bytes = build_contract_cim_docx_bytes({
                                    "CIM_NR": cim_nr,
                                    "CIM_DATA": cim_data,
                                    "ANGAJATOR_DEN": ang_den,
                                    "ANGAJATOR_SEDIU": ang_sediu,
                                    "ANGAJATOR_CUI": ang_cui,
                                    "REPREZENTANT_LEGAL": rep_legal,
                                    "REPREZENTANT_FUNCTIE": rep_functie,
                                    "SALARIAT_NUME_COMPLET": sal_nume_complet,
                                    "SALARIAT_DOMICILIU": domiciliu,
                                    "SALARIAT_CNP": sal_cnp,
                                    "OBIECT": obiect,
                                    "DURATA_TEXT": durata_text,
                                    "LOC_MUNCA": loc_munca or "-",
                                    "FUNCTIE": functie or "-",
                                    "COD_COR": cor or "-",
                                    "PROGRAM_TEXT": program_text,
                                    "CONCEDIU_TEXT": concediu_text,
                                    "SALARIU_TEXT": salariu_text,
                                })

                                file_docx_name = f"contract_cim_{emp_id}_{int(datetime.datetime.now().timestamp())}.docx"
                                _save_employee_document(
                                    conn,
                                    employee_id=int(emp_id),
                                    doc_type="CIM",
                                    filename=file_docx_name,
                                    data=docx_bytes,
                                    doc_no=cim_nr,
                                    doc_date=cim_data,
                                    meta={"generator": "code", "source": "model_2021"},
                                    set_active=True,
                                )

                                st.success("Contractul (CIM) a fost generat și salvat în istoric.")
                                st.download_button(
                                    "⬇️ Descarcă DOCX (CIM)",
                                    data=docx_bytes,
                                    file_name=file_docx_name,
                                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                    key=f"cim_code_dl_{emp_id}_{datetime.datetime.now().timestamp()}",
                                )

                                if do_gen_cim_code_pdf:
                                    try:
                                        pdf_bytes = convert_docx_to_pdf_bytes(docx_bytes)
                                        file_pdf_name = file_docx_name.replace(".docx", ".pdf")
                                        # Salvare PDF ca document separat (același doc_type)
                                        _save_employee_document(
                                            conn,
                                            employee_id=int(emp_id),
                                            doc_type="CIM",
                                            filename=file_pdf_name,
                                            data=pdf_bytes,
                                            doc_no=cim_nr,
                                            doc_date=cim_data,
                                            meta={"generator": "code", "format": "pdf"},
                                            set_active=False,
                                        )
                                        st.download_button(
                                            "⬇️ Descarcă PDF (CIM)",
                                            data=pdf_bytes,
                                            file_name=file_pdf_name,
                                            mime="application/pdf",
                                            key=f"cim_code_dl_pdf_{emp_id}_{datetime.datetime.now().timestamp()}",
                                        )
                                    except Exception as e:
                                        st.warning(f"PDF nu a putut fi generat (necesită Word/convertor). Detalii: {e}")

                            except Exception as e:
                                st.error(f"Eroare la generarea CIM din cod: {e}")

                        st.divider()

                        # -----------------------------
                        # 2) Act adițional (DOCX)
                        # -----------------------------
                        st.markdown('<div class="emp-subsubsection-title">🧾 Act adițional la C.I.M. (din cod)</div>', unsafe_allow_html=True)
                        aa_nr_code = st.text_input("Nr. act adițional", value="", key=f"aa_code_nr_{emp_id}")
                        aa_data_code = st.text_input("Data act adițional", value="", key=f"aa_code_data_{emp_id}")
                        aa_data_efect_code = st.text_input("Data efect (ex: 01.01.2026)", value="", key=f"aa_code_efect_{emp_id}")
                        revisal_nr = st.text_input("Nr. CIM în Revisal", value=str(_cim.get("CIM_NR","") or ""), key=f"aa_code_revisal_nr_{emp_id}")
                        revisal_data = st.text_input("Data CIM în Revisal", value=str(_cim.get("CIM_DATA","") or ""), key=f"aa_code_revisal_data_{emp_id}")
                        aa_continut_code = st.text_area(
                            "Conținut modificări (scrii aici exact clauzele modificate)",
                            value="",
                            height=180,
                            key=f"aa_code_continut_{emp_id}",
                        )

                        do_gen_aa_code = st.button("🧾 Generează Act adițional (DOCX + salvare istoric)", key=f"aa_code_gen_{emp_id}")
                        do_gen_aa_code_pdf = st.button("📄 Generează și PDF", key=f"aa_code_gen_pdf_{emp_id}")

                        if do_gen_aa_code or do_gen_aa_code_pdf:
                            try:
                                sal_nume_complet = f"{(emp.get('nume') or '').strip()} {(emp.get('prenume') or '').strip()}".strip()

                                docx_bytes = build_act_aditional_cim_docx_bytes({
                                    "ANGAJATOR_DEN": ang_den,
                                    "ANGAJATOR_SEDIU": ang_sediu,
                                    "REPREZENTANT_LEGAL": rep_legal,
                                    "SALARIAT_NUME_COMPLET": sal_nume_complet,
                                    "REVISAL_NR": (revisal_nr or "").strip(),
                                    "REVISAL_DATA": (revisal_data or "").strip(),
                                    "AA_NR": (aa_nr_code or "").strip(),
                                    "AA_DATA": (aa_data_code or "").strip(),
                                    "AA_DATA_EFECT": (aa_data_efect_code or "").strip(),
                                    "AA_CONTINUT_MODIFICARI": (aa_continut_code or "").strip(),
                                })

                                file_docx_name = f"act_aditional_cim_{emp_id}_{int(datetime.datetime.now().timestamp())}.docx"
                                _save_employee_document(
                                    conn,
                                    employee_id=int(emp_id),
                                    doc_type="ACT_ADITIONAL_CIM",
                                    filename=file_docx_name,
                                    data=docx_bytes,
                                    doc_no=(aa_nr_code or "").strip(),
                                    doc_date=(aa_data_code or "").strip(),
                                    meta={"generator": "code", "source": "model_2021"},
                                    set_active=True,
                                )

                                st.success("Actul adițional a fost generat și salvat în istoric.")
                                st.download_button(
                                    "⬇️ Descarcă DOCX (Act adițional)",
                                    data=docx_bytes,
                                    file_name=file_docx_name,
                                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                    key=f"aa_code_dl_{emp_id}_{datetime.datetime.now().timestamp()}",
                                )

                                if do_gen_aa_code_pdf:
                                    try:
                                        pdf_bytes = convert_docx_to_pdf_bytes(docx_bytes)
                                        file_pdf_name = file_docx_name.replace(".docx", ".pdf")
                                        _save_employee_document(
                                            conn,
                                            employee_id=int(emp_id),
                                            doc_type="ACT_ADITIONAL_CIM",
                                            filename=file_pdf_name,
                                            data=pdf_bytes,
                                            doc_no=(aa_nr_code or "").strip(),
                                            doc_date=(aa_data_code or "").strip(),
                                            meta={"generator": "code", "format": "pdf"},
                                            set_active=False,
                                        )
                                        st.download_button(
                                            "⬇️ Descarcă PDF (Act adițional)",
                                            data=pdf_bytes,
                                            file_name=file_pdf_name,
                                            mime="application/pdf",
                                            key=f"aa_code_dl_pdf_{emp_id}_{datetime.datetime.now().timestamp()}",
                                        )
                                    except Exception as e:
                                        st.warning(f"PDF nu a putut fi generat (necesită Word/convertor). Detalii: {e}")

                            except Exception as e:
                                st.error(f"Eroare la generarea Actului adițional din cod: {e}")

                    # Act adițional la CIM (DOCX/PDF) – pe șablon
                    # =========================================================

                    # (opțiunea pe șablon pentru Act adițional a fost eliminată)

                with tabs[4]:
                    st.subheader("🧾 Reges Online")
                    st.caption("Câmpuri pentru export/import REGES conform șablonului de modificare salariat.")
                    _render_info_section(
                        "Status curent REGES",
                        [
                            ("Grad invaliditate", _emp_v(emp, "reges_nume_grad_invaliditate")),
                            ("Grad handicap", _emp_v(emp, "reges_nume_grad_handicap")),
                            ("Tip handicap", _emp_v(emp, "reges_nume_tip_handicap")),
                            ("Certificat", _emp_v(emp, "reges_nr_cert_handicap")),
                        ],
                    )

                    # Liste valori (din șablon)
                    grad_invaliditate_opts = [("Fara", "Fără"), ("Grad1", "Grad 1"), ("Grad2", "Grad 2"), ("Grad3", "Grad 3")]
                    grad_handicap_opts = [("Fara", "Fără"), ("Usor", "Ușor"), ("Mediu", "Mediu"), ("Accentuat", "Accentuat"), ("Grav", "Grav")]
                    tip_handicap_opts = [
                        ("Fizic", "Fizic"),
                        ("Vizual", "Vizual"),
                        ("Auditiv", "Auditiv"),
                        ("Surdocecitate", "Surdocecitate"),
                        ("Somatic", "Somatic"),
                        ("Mintal", "Mintal"),
                        ("Psihic", "Psihic"),
                        ("HIVSIDA", "HIV/SIDA"),
                        ("Asociat", "Asociat"),
                        ("BoliRare", "Boli rare"),
                    ]

                    def _parse_date_any_ro(s: str):
                        s = (s or "").strip()
                        if not s:
                            return None
                        for fmt in ("%Y-%m-%d", "%d.%m.%Y", "%d/%m/%Y"):
                            try:
                                return datetime.datetime.strptime(s, fmt).date()
                            except Exception:
                                pass
                        return None

                    def _fmt_ro(d):
                        return d.strftime("%d.%m.%Y") if d else ""

                    # valori curente din employees
                    cod_inv = str(emp.get("reges_cod_grad_invaliditate", "") or "")
                    num_inv = str(emp.get("reges_nume_grad_invaliditate", "") or "")
                    cod_hand = str(emp.get("reges_cod_grad_handicap", "") or "")
                    num_hand = str(emp.get("reges_nume_grad_handicap", "") or "")
                    cod_tip = str(emp.get("reges_cod_tip_handicap", "") or "")
                    num_tip = str(emp.get("reges_nume_tip_handicap", "") or "")
                    nr_cert = str(emp.get("reges_nr_cert_handicap", "") or "")
                    d_cert = _parse_date_any_ro(str(emp.get("reges_data_cert_handicap", "") or ""))
                    d_val = _parse_date_any_ro(str(emp.get("reges_termen_valabilitate", "") or ""))

                    # UI (vertical, aliniat stânga)
                    inv_sel = st.selectbox(
                        "Grad invaliditate (COD)",
                        options=[o[0] for o in grad_invaliditate_opts],
                        index=max(0, [o[0] for o in grad_invaliditate_opts].index(cod_inv)) if cod_inv in [o[0] for o in grad_invaliditate_opts] else 0,
                        key=f"reg_inv_{emp_id}",
                    )
                    inv_name = dict(grad_invaliditate_opts).get(inv_sel, "")
                    st.text_input("Grad invaliditate (NUME)", value=inv_name, disabled=True, key=f"reg_inv_name_{emp_id}")

                    hand_sel = st.selectbox(
                        "Grad handicap (COD)",
                        options=[o[0] for o in grad_handicap_opts],
                        index=max(0, [o[0] for o in grad_handicap_opts].index(cod_hand)) if cod_hand in [o[0] for o in grad_handicap_opts] else 0,
                        key=f"reg_hand_{emp_id}",
                    )
                    hand_name = dict(grad_handicap_opts).get(hand_sel, "")
                    st.text_input("Grad handicap (NUME)", value=hand_name, disabled=True, key=f"reg_hand_name_{emp_id}")

                    tip_sel = st.selectbox(
                        "Tip handicap (COD)",
                        options=[o[0] for o in tip_handicap_opts],
                        index=max(0, [o[0] for o in tip_handicap_opts].index(cod_tip)) if cod_tip in [o[0] for o in tip_handicap_opts] else 0,
                        key=f"reg_tip_{emp_id}",
                    )
                    tip_name = dict(tip_handicap_opts).get(tip_sel, "")
                    st.text_input("Tip handicap (NUME)", value=tip_name, disabled=True, key=f"reg_tip_name_{emp_id}")

                    nr_cert_in = st.text_input("Număr certificat handicap", value=nr_cert, key=f"reg_nr_cert_{emp_id}")
                    _today_reg = datetime.date.today()
                    _min_reg = _today_reg.replace(year=_today_reg.year - 100)
                    _max_reg = _today_reg.replace(year=_today_reg.year + 10)
                    dc = st.date_input(
                        "Data certificat handicap",
                        value=d_cert or _today_reg,
                        min_value=_min_reg,
                        max_value=_max_reg,
                        key=f"reg_d_cert_{emp_id}",
                    )
                    dv = st.date_input(
                        "Termen de valabilitate",
                        value=d_val or dc,
                        min_value=_min_reg,
                        max_value=_max_reg,
                        key=f"reg_d_val_{emp_id}",
                    )

                    if st.button("💾 Salvează Reges", key=f"reg_save_{emp_id}"):
                        upd = dict(emp)
                        upd.update({
                            "reges_cod_grad_invaliditate": inv_sel,
                            "reges_nume_grad_invaliditate": inv_name,
                            "reges_cod_grad_handicap": hand_sel,
                            "reges_nume_grad_handicap": hand_name,
                            "reges_cod_tip_handicap": tip_sel,
                            "reges_nume_tip_handicap": tip_name,
                            "reges_nr_cert_handicap": str(nr_cert_in).strip(),
                            "reges_data_cert_handicap": dc.isoformat(),
                            "reges_termen_valabilitate": dv.isoformat(),
                        })
                        _employee_upsert(conn, upd, employee_id=emp_id)
                        st.success("Salvat.")
                        st.rerun()
                    if st.button("↩️ Renunță", key=f"reg_cancel_{emp_id}"):
                        st.rerun()


                # ---------------------------
                # TAB 5: COD COR (tabel separat)
                # ---------------------------
                # ---------------------------
                # TAB 5: COD COR (tabel separat)
                # ---------------------------
                with tabs[5]:
                    st.subheader("🧩 COD COR")
                    st.caption("Gestionare tabel COR + aplicare cod COR pe angajat (completează automat denumirea funcției).")
                    _render_info_section(
                        "Situație curentă COR",
                        [
                            ("Cod COR", _emp_v(emp, "cod_cor")),
                            ("Funcție (din angajat)", _emp_v(emp, "functie")),
                        ],
                    )

                    q = st.text_input("Caută în tabel (cod sau denumire)", value="", key=f"cor_q_{emp_id}")
                    rows = _cor_list(conn, q=q, limit=1000)
                    cor_selected_code_key = f"cor_selected_code_{emp_id}"

                    if rows:
                        df_cor = pd.DataFrame(rows, columns=["COD COR", "DENUMIRE OCUPAȚIE"])
                        selected_code = str(st.session_state.get(cor_selected_code_key, "") or "").strip()
                        if not selected_code:
                            selected_code = str(emp.get("cod_cor") or "").strip()

                        df_cor_edit = df_cor.copy()
                        df_cor_edit.insert(0, "Select", df_cor_edit["COD COR"].astype(str) == selected_code)
                        edited_cor = st.data_editor(
                            df_cor_edit,
                            column_config={
                                "Select": st.column_config.CheckboxColumn("Selectat", help="Bifează pentru a prelua în «Alege din tabel»."),
                                "COD COR": st.column_config.TextColumn("COD COR", disabled=True),
                                "DENUMIRE OCUPAȚIE": st.column_config.TextColumn("DENUMIRE OCUPAȚIE", disabled=True),
                            },
                            disabled=["COD COR", "DENUMIRE OCUPAȚIE"],
                            key=f"cor_table_editor_{emp_id}",
                            use_container_width=True,
                            height=320,
                        )
                        selected_rows = edited_cor[edited_cor["Select"]]
                        if not selected_rows.empty:
                            selected_from_table = str(selected_rows.iloc[0]["COD COR"]).strip()
                            if st.session_state.get(cor_selected_code_key) != selected_from_table:
                                st.session_state[cor_selected_code_key] = selected_from_table
                                st.rerun()
                    else:
                        st.info("Nu există încă înregistrări în tabelul COR (sau nu s-a găsit nimic pentru filtrul curent).")

                    st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)
                    st.markdown("### Aplică pentru angajat")
                    current_cod = str(emp.get("cod_cor") or "").strip()
                    selected_code = str(st.session_state.get(cor_selected_code_key, "") or "").strip()
                    selected_code = selected_code or current_cod

                    opt = [""] + [f"{c} — {d}" for (c, d) in rows]
                    default_idx = 0
                    if selected_code:
                        for i, s in enumerate(opt):
                            if s.startswith(selected_code):
                                default_idx = i
                                break

                    sel = st.selectbox("Alege din tabel", options=opt, index=default_idx, key=f"cor_sel_{emp_id}")
                    sel_cod = sel.split("—")[0].strip() if sel else ""
                    sel_den = _cor_get(conn, sel_cod) if sel_cod else ""

                    cod_in = st.text_input("Cod COR", value=sel_cod or current_cod, key=f"cor_cod_in_{emp_id}")
                    den_in = st.text_input(
                        "Denumire ocupație",
                        value=sel_den or (_cor_get(conn, current_cod) or ""),
                        key=f"cor_den_in_{emp_id}",
                    )

                    if st.button("💾 Salvează în tabel COR", key=f"cor_save_tbl_{emp_id}"):
                        if str(cod_in).strip() and str(den_in).strip():
                            _cor_upsert(conn, str(cod_in).strip(), str(den_in).strip())
                            st.success("Salvat în tabelul COR.")
                            st.rerun()
                        else:
                            st.error("Completează atât Cod COR cât și Denumire ocupație.")

                    if st.button("✅ Aplică la angajat", key=f"cor_apply_emp_{emp_id}"):
                        cod_apply = str(cod_in).strip()
                        den_apply = str(den_in).strip() if str(den_in).strip() else (_cor_get(conn, cod_apply) or "")
                        upd = dict(emp)
                        upd.update({"cod_cor": cod_apply})
                        if den_apply:
                            upd["functie"] = den_apply
                        _employee_upsert(conn, upd, employee_id=emp_id)
                        st.success("Aplicat pe angajat.")
                        st.rerun()

                    if st.button("🗑 Dezactivează cod din tabel", key=f"cor_del_tbl_{emp_id}"):
                        cod_del = str(cod_in).strip()
                        if cod_del:
                            _cor_soft_delete(conn, cod_del)
                            st.success("Cod dezactivat (soft delete) din tabelul COR.")
                            st.rerun()
                        else:
                            st.error("Completează Cod COR pentru ștergere.")


                # ---------------------------
                # TAB 6: Legea 153 (tabel separat)
                # ---------------------------
                with tabs[6]:
                    st.subheader("📜 Funcția angajatului (conform Legii 153/2017)")
                    st.caption("Selectează funcția din nomenclator și aplic-o angajatului. Denumirea se completează automat.")
                    _render_info_section(
                        "Situație curentă Legea 153",
                        [
                            ("Cod 153", _emp_v(emp, "cod_153")),
                            ("Denumire funcție", _emp_v(emp, "den_153")),
                            ("Marcă", _emp_v(emp, "marca")),
                        ],
                    )

                    q153 = st.text_input("Introdu codul sau denumirea funcției", value="", key=f"l153_q_{emp_id}")
                    rows153 = _l153_list(conn, q=q153, limit=1000)
                    l153_selected_code_key = f"l153_selected_code_{emp_id}"

                    if rows153:
                        df_153 = pd.DataFrame(rows153, columns=["COD 153", "DENUMIRE FUNCȚIE (153)"])
                        selected_code_153 = str(st.session_state.get(l153_selected_code_key, "") or "").strip()
                        if not selected_code_153:
                            selected_code_153 = str(emp.get("cod_153") or "").strip()

                        df_153_edit = df_153.copy()
                        df_153_edit.insert(0, "Select", df_153_edit["COD 153"].astype(str) == selected_code_153)
                        edited_153 = st.data_editor(
                            df_153_edit,
                            column_config={
                                "Select": st.column_config.CheckboxColumn("Selectat", help="Bifează pentru a prelua în «Alege din tabel»."),
                                "COD 153": st.column_config.TextColumn("COD 153", disabled=True),
                                "DENUMIRE FUNCȚIE (153)": st.column_config.TextColumn("DENUMIRE FUNCȚIE (153)", disabled=True),
                            },
                            disabled=["COD 153", "DENUMIRE FUNCȚIE (153)"],
                            key=f"l153_table_editor_{emp_id}",
                            use_container_width=True,
                            height=320,
                        )
                        selected_rows_153 = edited_153[edited_153["Select"]]
                        if not selected_rows_153.empty:
                            selected_from_table_153 = str(selected_rows_153.iloc[0]["COD 153"]).strip()
                            if st.session_state.get(l153_selected_code_key) != selected_from_table_153:
                                st.session_state[l153_selected_code_key] = selected_from_table_153
                                st.rerun()
                    else:
                        st.info("Nu există încă înregistrări în tabelul Legea 153 (sau nu s-a găsit nimic pentru filtrul curent).")

                    st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)
                    st.markdown("### Aplică pentru angajat")
                    cur_cod_153 = str(emp.get("cod_153") or "").strip()
                    selected_code_153 = str(st.session_state.get(l153_selected_code_key, "") or "").strip()
                    selected_code_153 = selected_code_153 or cur_cod_153

                    opt153 = [""] + [f"{c} — {d}" for (c, d) in rows153]
                    default153 = 0
                    if selected_code_153:
                        for i, s in enumerate(opt153):
                            if s.startswith(selected_code_153):
                                default153 = i
                                break

                    sel153 = st.selectbox("Alege din tabel", options=opt153, index=default153, key=f"l153_sel_{emp_id}")
                    sel153_cod = sel153.split("—")[0].strip() if sel153 else ""
                    sel153_den = _l153_get(conn, sel153_cod) if sel153_cod else ""

                    # Autocompletare robustă: folosim DOAR session_state (fără value= pe aceleași chei)
                    cod_key = f"l153_cod_in_{emp_id}"
                    den_key = f"l153_den_in_{emp_id}"
                    last_key = f"l153_last_sel_{emp_id}"

                    if cod_key not in st.session_state:
                        st.session_state[cod_key] = cur_cod_153 or ""
                    if den_key not in st.session_state:
                        st.session_state[den_key] = str(emp.get("den_153") or "").strip()

                    # Dacă utilizatorul a ales un cod din tabel, completăm automat Cod + Denumire
                    if sel153_cod and st.session_state.get(last_key) != sel153_cod:
                        st.session_state[last_key] = sel153_cod
                        st.session_state[cod_key] = sel153_cod
                        st.session_state[den_key] = (sel153_den or _l153_get(conn, sel153_cod) or "").strip()

                    cod153_in = st.text_input("Cod 153", key=cod_key)
                    den153_in = st.text_input("Denumire funcție (153)", key=den_key)
                    if st.button("💾 Salvează în tabel 153", key=f"l153_save_tbl_{emp_id}"):
                        if str(cod153_in).strip() and str(den153_in).strip():
                            _l153_upsert(conn, str(cod153_in).strip(), str(den153_in).strip())
                            st.success("Salvat în tabelul Legea 153.")
                            st.rerun()
                        else:
                            st.error("Completează atât Cod 153 cât și Denumire.")

                    if st.button("✅ Aplică la angajat", key=f"l153_apply_emp_{emp_id}"):
                        cod_apply = str(cod153_in).strip()
                        den_apply = str(den153_in).strip() if str(den153_in).strip() else (_l153_get(conn, cod_apply) or "")
                        upd = dict(emp)
                        upd.update({"cod_153": cod_apply, "den_153": den_apply})
                        _employee_upsert(conn, upd, employee_id=emp_id)
                        st.success("Aplicat pe angajat.")
                        st.rerun()

                    if st.button("🗑 Dezactivează cod din tabel", key=f"l153_del_tbl_{emp_id}"):
                        cod_del = str(cod153_in).strip()
                        if cod_del:
                            _l153_soft_delete(conn, cod_del)
                            st.success("Cod dezactivat (soft delete) din tabelul Legea 153.")
                            st.rerun()
                        else:
                            st.error("Completează Cod 153 pentru ștergere.")

                    st.divider()
                    st.subheader("💰 L153/2017 – grilă salarizare (din nomenclator)")

                    # Grila se ia din DB (Setări → Nomenclatoare → Import grilă)
                    try:
                        empty_grid = _lege153_grid_is_empty(conn)
                    except Exception:
                        empty_grid = True

                    if empty_grid:
                        st.warning("Nomenclatorul L153 (grila salarială) nu este încă importat în DB. Mergi la Setări → Nomenclatoare și importă XLSX-ul (sheet 'Normalized').")
                    else:
                        df_grid = _lege153_grid_load(conn)

                        # Precompletare din activ (dacă există)
                        active = _employee_l153_get_active(conn, emp_id)
                        KGRID = f"l153_grid_{emp_id}_"

                        def _pick_default(options, val):
                            if val in options:
                                return options.index(val)
                            return 0

                        KGRID = f"l153_grid_{emp_id}_"

                        anexas = sorted(df_grid["Anexa"].dropna().astype(str).unique().tolist())
                        anexa_sel = st.selectbox(
                            "Anexa",
                            anexas,
                            index=_pick_default(anexas, str(active.get("anexa","")).strip()) if anexas else 0,
                            key=f"{KGRID}anexa",
                        )

                        df1 = df_grid[df_grid["Anexa"].astype(str).str.strip() == str(anexa_sel).strip()] if anexa_sel else df_grid
                        tabels = sorted(df1["Tabel"].dropna().astype(str).unique().tolist())
                        tabel_sel = st.selectbox(
                            "Tabel",
                            tabels,
                            index=_pick_default(tabels, str(active.get("tabel","")).strip()) if tabels else 0,
                            key=f"{KGRID}tabel",
                        )

                        df2 = df1[df1["Tabel"].astype(str).str.strip() == str(tabel_sel).strip()] if tabel_sel else df1
                        functii = sorted(df2["Funcție"].dropna().astype(str).unique().tolist())
                        functie_sel = st.selectbox(
                            "Funcție",
                            functii,
                            index=_pick_default(functii, str(active.get("functie","")).strip()) if functii else 0,
                            key=f"{KGRID}functie",
                        )

                        df3 = df2[df2["Funcție"].astype(str).str.strip() == str(functie_sel).strip()] if functie_sel else df2

                       # helper: opțiuni curate
                        def _clean_opts(series):
                            if series is None:
                                return []
                            vals = (
                                series.dropna()
                                .astype(str)
                                .map(lambda x: x.strip())
                                .tolist()
                            )
                            vals = [v for v in vals if v and v.lower() not in ("nan", "none")]
                            # unique + sort
                            return sorted(list(dict.fromkeys(vals)))

                        def _safe_sb(label, options, key, preferred=""):
                            if not options:
                                st.caption(f"ℹ️ {label}: nu se aplică pentru selecția curentă.")
                                st.session_state[key] = ""
                                return ""
                            preferred = (preferred or "").strip()
                            cur = (st.session_state.get(key) or preferred).strip()
                            if cur not in options:
                                cur = preferred if preferred in options else options[0]
                                st.session_state[key] = cur
                            return st.selectbox(label, options, index=options.index(cur), key=key)

                        # 1) STUDII
                        studii_opts = _clean_opts(df3.get("Studii"))
                        studii_sel = _safe_sb(
                            "Studii (dacă se aplică)",
                            studii_opts,
                            key=f"{KGRID}studii",
                            preferred=active.get("studii", ""),
                        )
                        df4 = df3[df3["Studii"].astype(str).str.strip() == str(studii_sel).strip()] if studii_sel else df3

                        # 2) GRAD
                        grad_opts = _clean_opts(df4.get("Grad"))
                        grad_sel = _safe_sb(
                            "Grad (dacă se aplică)",
                            grad_opts,
                            key=f"{KGRID}grad",
                            preferred=active.get("grad", ""),
                        )
                        df5 = df4[df4["Grad"].astype(str).str.strip() == str(grad_sel).strip()] if grad_sel else df4

                        # 3) TREAPTĂ / GRADAȚIE
                        treapta_opts = _clean_opts(df5.get("Treaptă/Gradație"))
                        treapta_sel = _safe_sb(
                            "Treaptă/Gradație (dacă se aplică)",
                            treapta_opts,
                            key=f"{KGRID}treapta",
                            preferred=active.get("treapta", ""),
                        )
                        df6 = df5[df5["Treaptă/Gradație"].astype(str).str.strip() == str(treapta_sel).strip()] if treapta_sel else df5

                        # Rezultat final: dacă avem mai multe rânduri, luăm primul (de regulă e unic)
                        row = df6.iloc[0].to_dict() if not df6.empty else {}

                        s_val = row.get("Salariu", "")
                        c_val = row.get("Coeficient", "")

                        st.text_input("Salariu (auto)", value=str(s_val), key=f"{KGRID}sal", disabled=True)
                        st.text_input("Coeficient (auto)", value=str(c_val), key=f"{KGRID}coef", disabled=True)

                        _today_l153 = datetime.date.today()
                        _min_l153 = _today_l153.replace(year=_today_l153.year - 100)
                        _max_l153 = _today_l153.replace(year=_today_l153.year + 10)
                        data_apl = st.date_input(
                            "Data aplicării (de la)",
                            value=_today_l153,
                            min_value=_min_l153,
                            max_value=_max_l153,
                            key=f"{KGRID}data_apl",
                        )

                        # diana – preview selecție curentă
                        with st.container():
                            st.markdown("#### Selecție curentă")
                            st.caption(f"Anexa: **{anexa_sel}** | Tabel: **{tabel_sel}** | Funcție: **{functie_sel}**")
                            if studii_sel: st.caption(f"Studii: **{studii_sel}**")
                            if grad_sel: st.caption(f"Grad: **{grad_sel}**")
                            if treapta_sel: st.caption(f"Treaptă: **{treapta_sel}**")

                            #  diana – afișare salariu și coeficient doar dacă sunt valori valide (numere) 
                            if st.button("💾 Salvează ca activ (L153)", key=f"{KGRID}save_active"):
                                payload = {
                                    "anexa": anexa_sel,
                                    "tabel": tabel_sel,
                                    "functie": functie_sel,
                                    "studii": studii_sel,
                                    "grad": grad_sel,
                                    "treapta": treapta_sel,
                                    "salariu": float(s_val) if str(s_val).strip() not in ("", "nan", "None") else None,
                                    "coeficient": float(c_val) if str(c_val).strip() not in ("", "nan", "None") else None,
                                    "data_aplicare": str(data_apl) if data_apl else "",
                                }
                                _employee_l153_insert(conn, emp_id, payload)#set_active=True
                                st.success("Salvat L153 (activ) pentru angajat.")
                                st.rerun()

                        # Istoric
                        st.markdown("#### Istoric L153 (ultimele 20)")
                        hist = _employee_l153_list(conn, emp_id, limit=20)
                        if hist:
                            st.dataframe(pd.DataFrame(hist), use_container_width=True, height=260)
                        else:
                            st.info("Nu există încă înregistrări L153 pentru acest angajat.")

                with tabs[8]:
                    with st.container(key="ang_delete_panel"):
                        st.warning("Această acțiune dezactivează angajatul (soft delete: activ=0).")
                        confirm = st.checkbox("Confirm dezactivarea", key=f"ang_del_confirm_{emp_id}")
                        do_del = st.button("🗑 Dezactivează angajat", disabled=not confirm, key=f"ang_del_btn_{emp_id}")
                        cancel = st.button("↩ Renunță", key=f"ang_del_cancel_{emp_id}")

                    if cancel:
                        st.rerun()

                    if do_del:
                        _employee_soft_delete(conn, emp_id)
                        st.success("Angajat dezactivat.")
                        st.session_state["ang_view"] = "list"
                        st.session_state["ang_selected_id"] = None
                        st.rerun()



            # ---------------------------

        with main_tabs_emp[1]:
            render_documente_informatii_salariat(conn, emp_id, emp, bundle=bundle)
    # View: LIST
    # ---------------------------
    # Dacă suntem în fișa angajatului (detail/add/edit/delete), NU mai randăm lista.
    if st.session_state.get("ang_view", "list") != "list":
        return

    def _ang_reset_filters() -> None:
        for k in ("ang_q_ui", "ang_active_only_ui"):
            st.session_state.pop(k, None)
        st.session_state.pop("ang_applied_filters", None)
        st.session_state.pop("ang_has_searched", None)
        st.session_state["ang_selected_id"] = None
        st.rerun()

    # state filtre (aplicate explicit)
    if "ang_applied_filters" not in st.session_state:
        st.session_state["ang_applied_filters"] = {
            "search": "",
            "active_only": True,
        }
    if "ang_has_searched" not in st.session_state:
        st.session_state["ang_has_searched"] = False

    if st.session_state.get("ang_reset_pending"):
        for k in ("ang_q_ui", "ang_active_only_ui"):
            st.session_state.pop(k, None)
        st.session_state["ang_applied_filters"] = {
            "search": "",
            "active_only": True,
        }
        st.session_state["ang_has_searched"] = False
        st.session_state.pop("ang_reset_pending", None)
        st.rerun()

    # 1) Titlu pagină
    st.markdown('<h2 class="page-title ang-page-title">Angajați</h2>', unsafe_allow_html=True)
    st.markdown("<div class='ang-gap-xs'></div>", unsafe_allow_html=True)

    # 2) Buton principal
    if st.button("Adaugă personal", key="emp_add_personal"):
        st.session_state["ang_view"] = "add"
        st.session_state["ang_selected_id"] = None
        st.rerun()
    st.markdown("<div class='ang-gap-sm'></div>", unsafe_allow_html=True)

    # 3) + 4) Căutare rapidă + filtre compacte
    with st.form(key="ang_search_form"):
        st.text_input(
            "",
            value=st.session_state.get("ang_q_ui", st.session_state["ang_applied_filters"].get("search", "")),
            key="ang_q_ui",
            placeholder="Nume, ID, prenume, CNP sau marcă (Enter pentru căutare)",
            label_visibility="collapsed",
        )

        c_buttons, c_active = st.columns([0.30, 0.70])
        with c_buttons:
            submitted = st.form_submit_button("Aplică filtre", use_container_width=True, key="ang_btn_apply_filters")
            reset_clicked = st.form_submit_button("Reset filtre", use_container_width=True, key="ang_btn_reset_filters")
        with c_active:
            st.checkbox(
                "Doar activi",
                value=st.session_state.get("ang_active_only_ui", st.session_state["ang_applied_filters"].get("active_only", True)),
                key="ang_active_only_ui",
            )
    st.markdown("<div class='ang-gap-sm'></div>", unsafe_allow_html=True)

    if reset_clicked:
        st.session_state["ang_reset_pending"] = True
        st.rerun()

    if submitted:
        st.session_state["ang_applied_filters"] = {
            "search": st.session_state.get("ang_q_ui", ""),
            "active_only": st.session_state.get("ang_active_only_ui", True),
        }
        st.session_state["ang_has_searched"] = True
        st.rerun()

    applied = st.session_state.get("ang_applied_filters", {})
    df = list_employees(
        conn,
        active_only=applied.get("active_only", True),
        search=applied.get("search", ""),
        nume="",
        prenume="",
        cnp="",
        marca="",
    )

    if df.empty:
        st.warning("Niciun angajat găsit pentru criteriile introduse.")
        st.session_state["ang_selected_id"] = None
        return

    # 5) Tabel angajați
    cols_to_show = [c for c in ["id", "marca", "last_name", "first_name", "cnp", "activ"] if c in df.columns]
    dfv = df[cols_to_show].copy()
    dfv = dfv.rename(columns={
        "id": "ID",
        "marca": "Marcă",
        "last_name": "Nume",
        "first_name": "Prenume",
        "cnp": "CNP",
        "activ": "Activ",
    })
    st.caption(f"Număr angajați: {len(df)}")
    st.markdown("<div class='ang-gap-xs'></div>", unsafe_allow_html=True)

    recs = df.to_dict("records")
    options = [int(r.get("id")) for r in recs if r.get("id") is not None]
    if not options:
        st.warning("Nu am găsit ID-uri valide în tabela employees.")
        return

    def _sync_angajati_selection_from_table(event_obj) -> int | None:
        """Sursa unică de adevăr: selecția din tabel -> ang_selected_id."""
        if event_obj and getattr(event_obj, "selection", None) and getattr(event_obj.selection, "rows", None):
            rows = event_obj.selection.rows or []
            if rows:
                row_idx = rows[0]
                if 0 <= row_idx < len(recs):
                    return int(recs[row_idx].get("id"))
            return None
        return None

    # randare tabel + sincronizare selecție
    table_selected_id: int | None = None
    try:
        event = st.dataframe(
            dfv,
            use_container_width=True,
            hide_index=True,
            height=370,
            on_select="rerun",
            selection_mode="single-row",
            key="ang_df_select",
        )
        table_selected_id = _sync_angajati_selection_from_table(event)
    except TypeError:
        st.dataframe(dfv, use_container_width=True, hide_index=True, height=370)
    st.markdown("<div class='ang-gap-md'></div>", unsafe_allow_html=True)

    # Sursa unică: doar selecția din tabel actualizează ang_selected_id
    if table_selected_id is not None:
        if table_selected_id != st.session_state.get("ang_selected_id"):
            st.session_state["ang_selected_id"] = table_selected_id
            st.rerun()
    elif st.session_state.get("ang_selected_id") not in options:
        st.session_state["ang_selected_id"] = None

    # 6) Angajat selectat + sumar
    sel_id = st.session_state.get("ang_selected_id")
    sel_rec = next((r for r in recs if int(r.get("id")) == int(sel_id)), None) if sel_id else None

    # 7) Bottom section reconstruit: Angajat selectat + Acțiuni
    def _inject_angajati_bottom_section_css() -> None:
        st.markdown(
            """
            <style>
            .ang-bottom-section{
              margin-top: 12px;
              width: fit-content;
              display: flex;
              flex-direction: column;
              gap: 14px;
            }
            .ang-bottom-box{
              border: none !important;
              border-radius: 8px;
              background: transparent !important;
              padding: 10px 12px;
              box-sizing: border-box;
              box-shadow: none !important;
              width: fit-content;
              min-width: 300px;
              max-width: 330px;
            }
            .ang-bottom-title{
              margin: 0 0 10px 0 !important;
              margin-top: 0 !important;
              font-weight: 800;
              font-size: 1.02rem;
              color: rgba(248,250,252,0.98);
            }
            .ang-thin-guide-line{
              width: min(50vw, 260px);
              max-width: 100%;
              height: 1px;
              margin: 2px 0 10px 0;
              background: linear-gradient(
                90deg,
                rgba(148,163,184,0.46) 0%,
                rgba(148,163,184,0.24) 55%,
                rgba(148,163,184,0.02) 100%
              );
            }
            .ang-selected-rows{
              margin-top: 12px !important;
              display: flex;
              flex-direction: column;
              gap: 0 !important;
            }
            .ang-selected-row{
              margin: 0 0 10px 0 !important;
              font-size: 0.96rem;
              color: rgba(203,213,225,0.98);
              text-align: left;
              line-height: 1.35;
            }
            .ang-selected-row:last-child{
              margin-bottom: 0 !important;
            }
            .ang-actions-stack{
              margin-top: 12px !important;
              display: flex;
              flex-direction: column;
              align-items: center;
              gap: 14px;
            }
            .ang-actions-stack .stButton{
              margin: 0 !important;
            }
            .st-key-ang_btn_open button,
            .st-key-ang_btn_edit button,
            .st-key-ang_btn_delete button{
              width: 220px !important;
              min-width: 220px !important;
              max-width: 220px !important;
              height: 42px !important;
              min-height: 42px !important;
              max-height: 42px !important;
              border-radius: 12px !important;
              display: block !important;
              margin: 0 auto !important;
            }
            .st-key-ang_btn_apply_filters button,
            .st-key-ang_btn_reset_filters button,
            .st-key-emp_add_personal button{
              width: 220px !important;
              min-width: 220px !important;
              max-width: 220px !important;
              height: 42px !important;
              min-height: 42px !important;
              max-height: 42px !important;
            }
            .st-key-emp_add_personal button{
              margin-left: 16px !important;
            }
            </style>
            """,
            unsafe_allow_html=True,
        )

    def _render_angajat_selectat_box(emp: dict | None) -> None:
        st.markdown("<div class='ang-bottom-box'>", unsafe_allow_html=True)
        st.markdown("<p class='ang-bottom-title'>Angajat selectat</p>", unsafe_allow_html=True)
        st.markdown("<div class='ang-thin-guide-line'></div>", unsafe_allow_html=True)
        if emp:
            full_name = f"{emp.get('last_name', '')} {emp.get('first_name', '')}".strip() or "—"
            emp_id_text = emp.get("id", "—")
            st.markdown(
                f"""
                <div class='ang-selected-rows'>
                  <div class='ang-selected-row'><strong>Nume:</strong> {full_name} (ID: {emp_id_text})</div>
                  <div class='ang-selected-row'><strong>CNP:</strong> {emp.get('cnp', '') or '—'}</div>
                  <div class='ang-selected-row'><strong>Marcă:</strong> {emp.get('marca', '') or '—'}</div>
                </div>
                """,
                unsafe_allow_html=True,
            )
        else:
            st.markdown(
                """
                <div class='ang-selected-rows'>
                  <div class='ang-selected-row'>Nume: —</div>
                  <div class='ang-selected-row'>CNP: —</div>
                  <div class='ang-selected-row'>Marcă: —</div>
                </div>
                """,
                unsafe_allow_html=True,
            )
        st.markdown("</div>", unsafe_allow_html=True)

    def _render_angajati_actions_box(emp_id: int | None) -> None:
        has_emp = bool(emp_id)
        st.markdown("<div class='ang-bottom-box'>", unsafe_allow_html=True)
        st.markdown("<p class='ang-bottom-title'>Acțiuni</p>", unsafe_allow_html=True)
        st.markdown("<div class='ang-actions-stack'>", unsafe_allow_html=True)
        if st.button("Deschide fișa", key="ang_btn_open", disabled=not has_emp):
            st.session_state["ang_view"] = "detail"
            st.session_state["ang_scroll_to_top"] = True
            st.rerun()
        if st.button("Modifică personal", key="ang_btn_edit", disabled=not has_emp):
            st.session_state["ang_view"] = "edit"
            st.rerun()
        if st.button("Șterge personal", key="ang_btn_delete", disabled=not has_emp):
            st.session_state["ang_view"] = "delete"
            st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)
        st.markdown("</div>", unsafe_allow_html=True)

    _inject_angajati_bottom_section_css()
    st.markdown("<div class='ang-bottom-section'>", unsafe_allow_html=True)
    _render_angajat_selectat_box(sel_rec if sel_rec else None)
    _render_angajati_actions_box(int(sel_id) if sel_id else None)
    st.markdown("</div>", unsafe_allow_html=True)

def sync_employees_from_stat_functii(conn: sqlite3.Connection) -> int:
    """Dacă în stat_functii există coloane relevante, actualizează employees pentru a reflecta schimbările.

    Sincronizări suportate (dacă există coloanele în stat_functii):
      - Marca -> employees.marca (cheie de legare)
      - Functie -> employees.functie / employees.functie_contract (dacă există)
      - Departament -> employees.departament (dacă există)
      - Loc de munca -> employees.loc_munca (dacă există)

    Returnează: număr aproximativ de rânduri actualizate (marca potrivită).
    """
    cur = conn.cursor()
    cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='stat_functii'")
    if cur.fetchone() is None:
        return 0

    df = pd.read_sql_query("SELECT * FROM stat_functii", conn)
    if df is None or df.empty:
        return 0

    # Normalizăm numele de coloane
    cols = {str(c).strip(): c for c in df.columns}
    col_marca = cols.get("Marca") or cols.get("MARCA") or cols.get("marca")
    if not col_marca:
        return 0

    col_functie = cols.get("Functie") or cols.get("FUNCTIE") or cols.get("functie")
    col_depart = cols.get("Departament") or cols.get("DEPARTAMENT") or cols.get("departament")
    col_loc = cols.get("Loc de munca") or cols.get("LOC DE MUNCA") or cols.get("loc de munca") or cols.get("Loc_de_munca")

    # Coloane existente în employees
    cur.execute("PRAGMA table_info(employees);")
    emp_cols = {row[1].lower() for row in cur.fetchall()}

    # Nu inventăm coloane: actualizăm doar dacă există
    can_functie = "functie" in emp_cols
    can_functie_contract = "functie_contract" in emp_cols
    can_depart = "departament" in emp_cols
    can_loc = "loc_munca" in emp_cols

    if not any([col_functie, col_depart, col_loc]):
        return 0

    updated = 0
    for _, r in df.iterrows():
        marca = r.get(col_marca)
        if pd.isna(marca) or str(marca).strip() == "":
            continue
        marca_s = str(marca).strip()

        sets = []
        params = []

        if col_functie and not pd.isna(r.get(col_functie)):
            fval = str(r.get(col_functie)).strip()
            if can_functie:
                sets.append("functie = ?")
                params.append(fval)
            if can_functie_contract:
                sets.append("functie_contract = ?")
                params.append(fval)

        if col_depart and not pd.isna(r.get(col_depart)) and can_depart:
            dval = str(r.get(col_depart)).strip()
            sets.append("departament = ?")
            params.append(dval)

        if col_loc and not pd.isna(r.get(col_loc)) and can_loc:
            lval = str(r.get(col_loc)).strip()
            sets.append("loc_munca = ?")
            params.append(lval)

        if not sets:
            continue

        params.append(marca_s)
        try:
            cur.execute(f"UPDATE employees SET {', '.join(sets)} WHERE TRIM(marca) = TRIM(?)", params)
            if cur.rowcount and cur.rowcount > 0:
                updated += int(cur.rowcount)
        except Exception:
            # nu blocăm aplicația dacă nu există câmpuri/erori de date
            continue

    conn.commit()
    return updated




def sync_employees_from_stat_functii_bidir(conn: sqlite3.Connection, *, overwrite: bool = False) -> int:
    """Stat de funcții -> employees, cu control overwrite.

    Similar cu sync_employees_from_stat_functii(), dar:
      - dacă overwrite=False, completează doar câmpurile goale
      - dacă overwrite=True, suprascrie valorile
    Suportă și cod_cor / cod_153 / den_153 dacă există coloanele.

    Returnează numărul de rânduri actualizate.
    """
    cur = conn.cursor()
    cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='stat_functii'")
    if cur.fetchone() is None:
        return 0

    df = pd.read_sql_query("SELECT * FROM stat_functii", conn)
    if df.empty:
        return 0

    # Mapări posibile de coloane (stat_functii poate avea diverse denumiri)
    def _col(*names):
        for n in names:
            if n in df.columns:
                return n
        return None

    c_marca = _col('Marca','MARCA','marca')
    if not c_marca:
        return 0

    c_functie = _col('Functie','Funcție','functie')
    c_depart = _col('Departament','departament')
    c_loc = _col('Loc de munca','Loc de muncă','loc_munca','loc de munca')
    c_cor = _col('COD COR','Cod COR','cod_cor','COR','CodCor')
    c_153 = _col('COD 153','Cod 153','cod_153','Cod functie 153','Cod funcție 153','cod functie 153')
    c_den153 = _col('DEN 153','Denumire 153','den_153','Denumire functie 153','Denumire funcție 153')

    # coloane existente în employees
    cur.execute("PRAGMA table_info(employees)")
    emp_cols = {r[1].lower() for r in cur.fetchall()}

    def _has(col):
        return col and col.lower() in emp_cols

    updates = 0
    for _, r in df.iterrows():
        marca = r.get(c_marca)
        if pd.isna(marca) or str(marca).strip() == '':
            continue
        marca_s = str(marca).strip()

        # găsim angajat după marcă
        cur.execute("SELECT id, functie, departament, loc_munca, cod_cor, cod_153, den_153 FROM employees WHERE TRIM(marca)=? LIMIT 1", (marca_s,))
        row = cur.fetchone()
        if not row:
            continue
        emp_id = int(row[0])
        curr = {
            'functie': row[1] if len(row) > 1 else None,
            'departament': row[2] if len(row) > 2 else None,
            'loc_munca': row[3] if len(row) > 3 else None,
            'cod_cor': row[4] if len(row) > 4 else None,
            'cod_153': row[5] if len(row) > 5 else None,
            'den_153': row[6] if len(row) > 6 else None,
        }

        data = {}
        if c_functie and _has('functie'):
            v = r.get(c_functie)
            if overwrite or not (curr.get('functie') or '').strip():
                data['functie'] = '' if pd.isna(v) else str(v).strip()
        if c_depart and _has('departament'):
            v = r.get(c_depart)
            if overwrite or not (curr.get('departament') or '').strip():
                data['departament'] = '' if pd.isna(v) else str(v).strip()
        if c_loc and _has('loc_munca'):
            v = r.get(c_loc)
            if overwrite or not (curr.get('loc_munca') or '').strip():
                data['loc_munca'] = '' if pd.isna(v) else str(v).strip()
        if c_cor and _has('cod_cor'):
            v = r.get(c_cor)
            if overwrite or not (curr.get('cod_cor') or '').strip():
                data['cod_cor'] = '' if pd.isna(v) else str(v).strip()
        if c_153 and _has('cod_153'):
            v = r.get(c_153)
            if overwrite or not (curr.get('cod_153') or '').strip():
                data['cod_153'] = '' if pd.isna(v) else str(v).strip()
        if c_den153 and _has('den_153'):
            v = r.get(c_den153)
            if overwrite or not (curr.get('den_153') or '').strip():
                data['den_153'] = '' if pd.isna(v) else str(v).strip()

        if not data:
            continue

        set_sql = ", ".join([f"{k}=?" for k in data.keys()])
        params = list(data.values()) + [emp_id]
        cur.execute(f"UPDATE employees SET {set_sql} WHERE id=?", params)
        updates += 1

    conn.commit()
    return updates
def refresh_organigrama_from_stat(conn: sqlite3.Connection) -> Tuple[int, int]:
    """Regenerează organigrama (org_units + org_positions) din stat_functii, dacă există."""
    try:
        ensure_organigrama_tables(conn)
    except Exception:
        pass
    try:
        return generate_organigrama_from_stat(conn)
    except Exception:
        return (0, 0)


def load_org_positions_with_people(conn: sqlite3.Connection) -> pd.DataFrame:
    """Întoarce un DF compatibil cu render_org_tree_*: ID_UNITATE, DENUMIRE_FUNCTIE, NUME, PRENUME.

    Sursa adevărului:
      - org_positions (posturi) + org_units (unități)
      - employee_positions (ocupare) + employees (nume/prenume)
    """
    try:
        df_pos = pd.read_sql_query(
            """
            SELECT
                op.unit_id AS ID_UNITATE,
                op.title   AS DENUMIRE_FUNCTIE,
                e.last_name AS NUME,
                e.first_name AS PRENUME
            FROM org_positions op
            LEFT JOIN employee_positions ep
                   ON ep.position_id = op.id
                  AND (ep.end_date IS NULL OR TRIM(ep.end_date) = '')
            LEFT JOIN employees e
                   ON e.id = ep.employee_id
            WHERE op.is_active = 1
            """,
            conn,
        )
        if df_pos is None or df_pos.empty:
            return pd.DataFrame(columns=["ID_UNITATE", "DENUMIRE_FUNCTIE", "NUME", "PRENUME"])

        # normalizăm
        for c in ["ID_UNITATE", "DENUMIRE_FUNCTIE", "NUME", "PRENUME"]:
            if c not in df_pos.columns:
                df_pos[c] = ""
        df_pos["ID_UNITATE"] = df_pos["ID_UNITATE"].fillna(-1)
        return df_pos
    except Exception:
        return pd.DataFrame(columns=["ID_UNITATE", "DENUMIRE_FUNCTIE", "NUME", "PRENUME"])


def insert_employee(conn: sqlite3.Connection, data: Dict[str, Any]) -> int:
    cols = [
        "marca",
        "last_name",
        "first_name",
        "cnp",
        "functie",
        "departament",
        "data_angajare",
        "tip_contract",
        "strada",
        "numar",
        "bloc",
        "scara",
        "apartament",
        "cod_postal",
        "localitate",
        "judet",
        "telefon_fix",
        "mobil",
        "email",
        "ci_tip_act",
        "ci_serie",
        "ci_numar",
        "ci_eliberat_de",
        "ci_data_eliberarii",
        "stare_civila",
        "nr_copii",
        "loc_munca",
        "departament_organizatoric",
        "functie_contract",
        "tip_norma",
        "program_munca",
        "salariu_baza",
        "studii",
        "profesie",
        "calificare",
        "observatii",
        "dosar_nr",
        "dosar_functionar_public",
        "dosar_data_intocmire",
        "dosar_autoritate",
        "dosar_intocmit_nume",
        "dosar_intocmit_functie",
        "dosar_intocmit_semnatura",
        "dosar_modificari_nume",
        "dosar_modificari_functie",
        "dosar_modificari_semnatura",
        "dosar_certificare_nume",
        "dosar_certificare_functie",
        "dosar_certificare_semnatura",
        "activitate_in_afara_functiei",
        "activitate_in_cadru_institutie",
        "situatia_drepturi_salariale",
        "situatia_concedii",
        "situatia_disciplinara",
        "registru_numar",
        "registru_data",
        "registru_observatii",
        "activ",
    ]
    placeholders = ",".join(["?"] * len(cols))
    values = [data.get(c) for c in cols]
    cur = conn.cursor()
    cur.execute(
        f"INSERT INTO employees ({','.join(cols)}) VALUES ({placeholders})",
        values,
    )
    conn.commit()
    return cur.lastrowid


def update_employee(conn: sqlite3.Connection, emp_id: int, data: Dict[str, Any]) -> None:
    cols = [
        "marca",
        "last_name",
        "first_name",
        "cnp",
        "functie",
        "departament",
        "data_angajare",
        "tip_contract",
        "strada",
        "numar",
        "bloc",
        "scara",
        "apartament",
        "cod_postal",
        "localitate",
        "judet",
        "telefon_fix",
        "mobil",
        "email",
        "ci_tip_act",
        "ci_serie",
        "ci_numar",
        "ci_eliberat_de",
        "ci_data_eliberarii",
        "stare_civila",
        "nr_copii",
        "loc_munca",
        "departament_organizatoric",
        "functie_contract",
        "tip_norma",
        "program_munca",
        "salariu_baza",
        "studii",
        "profesie",
        "calificare",
        "observatii",
        "dosar_nr",
        "dosar_functionar_public",
        "dosar_data_intocmire",
        "dosar_autoritate",
        "dosar_intocmit_nume",
        "dosar_intocmit_functie",
        "dosar_intocmit_semnatura",
        "dosar_modificari_nume",
        "dosar_modificari_functie",
        "dosar_modificari_semnatura",
        "dosar_certificare_nume",
        "dosar_certificare_functie",
        "dosar_certificare_semnatura",
        "activitate_in_afara_functiei",
        "activitate_in_cadru_institutie",
        "situatia_drepturi_salariale",
        "situatia_concedii",
        "situatia_disciplinara",
        "registru_numar",
        "registru_data",
        "registru_observatii",
        "activ",
    ]
    set_clause = ", ".join([f"{c} = ?" for c in cols])
    values = [data.get(c) for c in cols]
    values.append(emp_id)
    cur = conn.cursor()
    cur.execute(f"UPDATE employees SET {set_clause} WHERE id = ?", values)
    conn.commit()


def delete_employee(conn: sqlite3.Connection, emp_id: int) -> None:
    cur = conn.cursor()
    cur.execute("DELETE FROM employees WHERE id = ?", (emp_id,))
    conn.commit()


# -------------------------------------------------------------
# CRUD pentru ACCESUL LA DOSAR
# -------------------------------------------------------------
def list_dosar_acces(conn: sqlite3.Connection, emp_id: int) -> List[sqlite3.Row]:
    cur = conn.cursor()
    cur.execute(
        """
        SELECT * FROM dosar_acces
        WHERE employee_id = ?
        ORDER BY nr_crt, id
        """,
        (emp_id,),
    )
    return cur.fetchall()


def insert_dosar_acces(
    conn: sqlite3.Connection,
    emp_id: int,
    nume: str,
    functie: str,
    semnatura: str,
    motivul: str,
    acces_autorizat_de: str,
    luat_la_cunostinta: str,
) -> None:
    cur = conn.cursor()
    cur.execute(
        "SELECT COALESCE(MAX(nr_crt), 0) FROM dosar_acces WHERE employee_id = ?",
        (emp_id,),
    )
    max_nr = cur.fetchone()[0] or 0
    next_nr = max_nr + 1

    cur.execute(
        """
        INSERT INTO dosar_acces (
            employee_id, nr_crt, nume, functie, semnatura,
            motivul, acces_autorizat_de, luat_la_cunostinta
        )
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (emp_id, next_nr, nume, functie, semnatura, motivul, acces_autorizat_de, luat_la_cunostinta),
    )
    conn.commit()


# -------------------------------------------------------------
# PAGINA: DOSAR PROFESIONAL
# -------------------------------------------------------------
def page_dosar_profesional(conn: sqlite3.Connection, sub_menu: str):
    st.markdown('<h2 class="page-title">Dosar profesional</h2>', unsafe_allow_html=True)

    # --------- SUS: FILTRE + LISTĂ ANGAJAȚI ---------
    st.markdown("#### Filtru angajați")

    # Nume/Prenume pe aceeași coloană, CNP/Marcă pe a doua
    col_name, col_id = st.columns(2)
    with col_name:
        filter_nume = st.text_input("Nume", key="filter_nume")
        filter_prenume = st.text_input("Prenume", key="filter_prenume")
    with col_id:
        filter_cnp = st.text_input("CNP", key="filter_cnp")
        filter_marca = st.text_input("Marcă", key="filter_marca")
    # Checkbox-ul sub filtre, pe un rând separat
    active_only = st.checkbox(
        "Doar activi",
        value=True,
        key="emp_active_only",
    )

    employees_df = list_employees(
        conn,
        active_only=active_only,
        search="",
        nume=filter_nume,
        prenume=filter_prenume,
        cnp=filter_cnp,
        marca=filter_marca,
    )

    selected_emp = None
    open_key = "dosar_open_selected"
    if employees_df.empty:
        st.warning("Nu s-a găsit niciun angajat pentru criteriile curente.")
        st.markdown("#### Adaugă angajat nou")
        with st.form("dosar_add_new_emp_compact"):
            c1, c2 = st.columns(2)
            with c1:
                new_marca = st.text_input("Marcă")
                new_last = st.text_input("Nume")
            with c2:
                new_first = st.text_input("Prenume")
                new_cnp = st.text_input("CNP")
            add_ok = st.form_submit_button("Adaugă angajat nou")
        if add_ok:
            payload = {
                "marca": (new_marca or "").strip(),
                "last_name": (new_last or "").strip(),
                "first_name": (new_first or "").strip(),
                "cnp": (new_cnp or "").strip(),
                "activ": 1,
            }
            if not payload["last_name"] or not payload["first_name"]:
                st.error("Completează cel puțin Nume și Prenume.")
            else:
                emp_id_new = insert_employee(conn, payload)
                st.success(f"Angajat nou salvat (ID = {emp_id_new}).")
                st.rerun()
        return
    else:
        records = employees_df.to_dict("records")
        options = list(range(len(records)))
        default_index = 0
        if "selected_emp_idx" in st.session_state:
            if st.session_state["selected_emp_idx"] < len(options):
                default_index = st.session_state["selected_emp_idx"]

        # Alege angajat – în listă apare DOAR Nume + Prenume
        selected_idx = st.selectbox(
            "Alege angajat",
            options=options,
            format_func=lambda idx: (
                f"{get_val(records[idx], 'last_name', 'LAST_NAME', 'nume', 'NUME')} "
                f"{get_val(records[idx], 'first_name', 'FIRST_NAME', 'prenume', 'PRENUME')}"
            ).strip(),
            index=default_index if options else 0,
            key="emp_select_box",
        )
        st.session_state["selected_emp_idx"] = selected_idx
        prev_idx = st.session_state.get("dosar_prev_selected_idx")
        if prev_idx != selected_idx:
            st.session_state[open_key] = False
        st.session_state["dosar_prev_selected_idx"] = selected_idx
        selected_emp = records[selected_idx]
        # Folosim angajatul selectat ca bază; dacă există un dosar salvat, îl reîncărcăm peste.
        emp_dosar = selected_emp
        emp_id_dosar = emp_dosar.get("id")
        if emp_id_dosar:
            bundle_dosar = load_employee_bundle(get_db_path(), int(emp_id_dosar))
            set_current_employee(int(emp_id_dosar))
            emp_dosar = bundle_dosar.get("employee") or selected_emp
        else:
            bundle_dosar = None
            set_current_employee(None)
            emp_dosar = selected_emp

        # Linie cu detalii angajat selectat – Nume | CNP | Marcă
        nume_complet = (
            f"{get_val(emp_dosar, 'last_name', 'LAST_NAME', 'nume', 'NUME')} "
            f"{get_val(emp_dosar, 'first_name', 'FIRST_NAME', 'prenume', 'PRENUME')}"
        ).strip()
        cnp_val = get_val(emp_dosar, "cnp", "CNP")
        marca_val = get_val(emp_dosar, "marca", "MARCA")

        st.markdown(
            f"**Dosar angajat:** {nume_complet} | {cnp_val} | {marca_val}"
        )
    if st.button("📂 Deschide dosarul angajatului selectat", key="dosar_open_btn"):
        st.session_state[open_key] = True
        st.rerun()
    if st.session_state.get(open_key, False):
        if st.button("📁 Închide dosarul", key="dosar_close_btn"):
            st.session_state[open_key] = False
            st.rerun()
    if not st.session_state.get(open_key, False):
        return

    st.markdown("---")

    # --------- JOS: conținut secțiune dosar ---------
    st.markdown(f"### 📁 {sub_menu}")
    if not emp_dosar:
        st.markdown(
            """
            <div class="dosar-empty">
              <div class="title">Selectează un angajat</div>
              <div class="text">
                Alege un angajat din modulul <strong>Angajați</strong> pentru a-i gestiona dosarul profesional.
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        if st.button("Deschide Angajați", key="dosar_go_angajati"):
            st.session_state["main_choice"] = "Angajați"
            st.rerun()
        return

    emp_id = emp_dosar.get("id") or emp_id_dosar
    key_prefix = f"emp_{emp_id}_"

    # --------- TIPĂRIRE DOSAR COMPLET (DOCX + PDF) ---------
    st.markdown("#### Tipărește dosarul complet")

    acces_rows_full = list_dosar_acces(conn, emp_id)
    full_docx_bytes = generate_dosar_complet_docx(emp_dosar, acces_rows_full)
    file_full_docx_name = (
        f"dosar_profesional_complet_{get_val(emp_dosar, 'cnp', 'CNP') or emp_id}.docx"
    )

    c_full1, c_full2 = st.columns(2)
    with c_full1:
        st.download_button(
            "🖨️ Tipărește dosarul complet (DOCX)",
            data=full_docx_bytes,
            file_name=file_full_docx_name,
            mime=(mimetypes.guess_type(file_full_docx_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_download_dosar_complet_docx",
        )
    with c_full2:
        full_pdf_bytes = convert_docx_to_pdf_bytes(full_docx_bytes)
        if full_pdf_bytes:
            st.download_button(
                "🖨️ Tipărește dosarul complet (PDF)",
                data=full_pdf_bytes,
                file_name=file_full_docx_name.replace(".docx", ".pdf"),
                mime="application/pdf",
                key=f"{key_prefix}btn_download_dosar_complet_pdf",
            )

    st.markdown("---")

    # ------------------ COPERTĂ ------------------
    if sub_menu == "Copertă":
        sub_tab_form, sub_tab_print = st.tabs(["Formular copertă", "Tipărire copertă"])

        with sub_tab_form:
            st.markdown("#### Coperta dosarului profesional")

            dosar_nr = st.text_input(
                "DOSAR PROFESIONAL Nr.",
                value=get_val(emp_dosar, "dosar_nr"),
                key=f"{key_prefix}dosar_nr",
            )

            nume_default = (
                get_val(emp_dosar, "last_name", "LAST_NAME", "nume", "NUME")
                + " "
                + get_val(emp_dosar, "first_name", "FIRST_NAME", "prenume", "PRENUME")
            ).strip()
            functionar_public = st.text_input(
                "FUNCȚIONAR PUBLIC",
                value=get_val(emp_dosar, "dosar_functionar_public") or nume_default,
                key=f"{key_prefix}dosar_functionar_public",
            )

            cnp_v = get_val(emp_dosar, "cnp", "CNP")
            st.text_input(
                "COD NUMERIC PERSONAL",
                value=cnp_v,
                key=f"{key_prefix}dosar_cnp_afis",
                disabled=True,
            )

            dosar_data_intocmire = st.text_input(
                "DATA ÎNTOCMIRII DOSARULUI PROFESIONAL",
                value=get_val(emp_dosar, "dosar_data_intocmire"),
                key=f"{key_prefix}dosar_data_intocmire",
            )

            dosar_autoritate = st.text_area(
                "AUTORITATEA/INSTITUȚIA PUBLICĂ CE A ÎNTOCMIT DOSARUL PROFESIONAL",
                value=get_val(emp_dosar, "dosar_autoritate"),
                key=f"{key_prefix}dosar_autoritate",
            )

            # 5. Persoana care a întocmit dosarul
            st.markdown("##### 5. Persoana care a întocmit dosarul profesional")

            intocmit_nume_stored = get_val(emp_dosar, "dosar_intocmit_nume")
            intocmit_functie_stored = get_val(emp_dosar, "dosar_intocmit_functie")
            intocmit_semnatura_stored = get_val(emp_dosar, "dosar_intocmit_semnatura")

            intocmit_nume_lines = intocmit_nume_stored.splitlines() if intocmit_nume_stored else []
            intocmit_functie_lines = intocmit_functie_stored.splitlines() if intocmit_functie_stored else []
            intocmit_semnatura_lines = (
                intocmit_semnatura_stored.splitlines() if intocmit_semnatura_stored else []
            )

            nr_intocmit_default = max(
                3,
                len(intocmit_nume_lines),
                len(intocmit_functie_lines),
                len(intocmit_semnatura_lines),
            )
            nr_intocmit = st.number_input(
                "Număr persoane (5. întocmit dosarul)",
                min_value=1,
                max_value=10,
                step=1,
                value=nr_intocmit_default,
                key=f"{key_prefix}nr_intocmit",
            )

            col_n1, col_f1, col_s1 = st.columns(3)
            intocmit_nume_list, intocmit_functie_list, intocmit_semnatura_list = [], [], []

            for i in range(nr_intocmit):
                nume_val = intocmit_nume_lines[i] if i < len(intocmit_nume_lines) else ""
                functie_val = intocmit_functie_lines[i] if i < len(intocmit_functie_lines) else ""
                semn_val = intocmit_semnatura_lines[i] if i < len(intocmit_semnatura_lines) else ""

                with col_n1:
                    n = st.text_input(
                        f"Nume și prenume persoana {i + 1}",
                        value=nume_val,
                        key=f"{key_prefix}intocmit_nume_{i}",
                    )
                with col_f1:
                    f = st.text_input(
                        f"Funcția persoana {i + 1}",
                        value=functie_val,
                        key=f"{key_prefix}intocmit_functie_{i}",
                    )
                with col_s1:
                    s = st.text_input(
                        f"Semnătura persoana {i + 1}",
                        value=semn_val,
                        key=f"{key_prefix}intocmit_semnatura_{i}",
                    )

                intocmit_nume_list.append(n)
                intocmit_functie_list.append(f)
                intocmit_semnatura_list.append(s)

            # 6. Modificări
            st.markdown(
                "##### 6. Persoana autorizata sa opereze modificari/completari/rectificari"
            )

            modif_nume_stored = get_val(emp_dosar, "dosar_modificari_nume")
            modif_functie_stored = get_val(emp_dosar, "dosar_modificari_functie")
            modif_semnatura_stored = get_val(emp_dosar, "dosar_modificari_semnatura")

            modif_nume_lines = modif_nume_stored.splitlines() if modif_nume_stored else []
            modif_functie_lines = (
                modif_functie_stored.splitlines() if modif_functie_stored else []
            )
            modif_semnatura_lines = (
                modif_semnatura_stored.splitlines() if modif_semnatura_stored else []
            )

            nr_modif_default = max(
                3,
                len(modif_nume_lines),
                len(modif_functie_lines),
                len(modif_semnatura_lines),
            )
            nr_modif = st.number_input(
                "Număr persoane (6. modificari/completari/rectificari)",
                min_value=1,
                max_value=10,
                step=1,
                value=nr_modif_default,
                key=f"{key_prefix}nr_modif",
            )

            col_n2, col_f2, col_s2 = st.columns(3)
            modif_nume_list, modif_functie_list, modif_semnatura_list = [], [], []

            for i in range(nr_modif):
                nume_val = modif_nume_lines[i] if i < len(modif_nume_lines) else ""
                functie_val = modif_functie_lines[i] if i < len(modif_functie_lines) else ""
                semn_val = modif_semnatura_lines[i] if i < len(modif_semnatura_lines) else ""

                with col_n2:
                    n = st.text_input(
                        f"Nume și prenume persoana {i + 1}",
                        value=nume_val,
                        key=f"{key_prefix}modif_nume_{i}",
                    )
                with col_f2:
                    f = st.text_input(
                        f"Funcția persoana {i + 1}",
                        value=functie_val,
                        key=f"{key_prefix}modif_functie_{i}",
                    )
                with col_s2:
                    s = st.text_input(
                        f"Semnătura persoana {i + 1}",
                        value=semn_val,
                        key=f"{key_prefix}modif_semnatura_{i}",
                    )

                modif_nume_list.append(n)
                modif_functie_list.append(f)
                modif_semnatura_list.append(s)

            # 7. Certificare
            st.markdown(
                "##### 7. Persoana autorizata sa certifice datele cuprinse in dosarul profesional"
            )

            cert_nume_stored = get_val(emp_dosar, "dosar_certificare_nume")
            cert_functie_stored = get_val(emp_dosar, "dosar_certificare_functie")
            cert_semnatura_stored = get_val(emp_dosar, "dosar_certificare_semnatura")

            cert_nume_lines = cert_nume_stored.splitlines() if cert_nume_stored else []
            cert_functie_lines = (
                cert_functie_stored.splitlines() if cert_functie_stored else []
            )
            cert_semnatura_lines = (
                cert_semnatura_stored.splitlines() if cert_semnatura_stored else []
            )

            nr_cert_default = max(
                3,
                len(cert_nume_lines),
                len(cert_functie_lines),
                len(cert_semnatura_lines),
            )
            nr_cert = st.number_input(
                "Număr persoane (7. certificare)",
                min_value=1,
                max_value=10,
                step=1,
                value=nr_cert_default,
                key=f"{key_prefix}nr_cert",
            )

            col_n3, col_f3, col_s3 = st.columns(3)
            cert_nume_list, cert_functie_list, cert_semnatura_list = [], [], []

            for i in range(nr_cert):
                nume_val = cert_nume_lines[i] if i < len(cert_nume_lines) else ""
                functie_val = cert_functie_lines[i] if i < len(cert_functie_lines) else ""
                semn_val = cert_semnatura_lines[i] if i < len(cert_semnatura_lines) else ""

                with col_n3:
                    n = st.text_input(
                        f"Nume și prenume persoana {i + 1}",
                        value=nume_val,
                        key=f"{key_prefix}cert_nume_{i}",
                    )
                with col_f3:
                    f = st.text_input(
                        f"Funcția persoana {i + 1}",
                        value=functie_val,
                        key=f"{key_prefix}cert_functie_{i}",
                    )
                with col_s3:
                    s = st.text_input(
                        f"Semnătura persoana {i + 1}",
                        value=semn_val,
                        key=f"{key_prefix}cert_semnatura_{i}",
                    )

                cert_nume_list.append(n)
                cert_functie_list.append(f)
                cert_semnatura_list.append(s)

            st.markdown("---")
            if st.button(
                "💾 Salvează coperta dosarului profesional",
                key=f"{key_prefix}btn_save_dosar",
            ):
                cur = conn.cursor()
                cur.execute(
                    """
                    UPDATE employees
                    SET dosar_nr = ?,
                        dosar_functionar_public = ?,
                        dosar_data_intocmire = ?,
                        dosar_autoritate = ?,
                        dosar_intocmit_nume = ?,
                        dosar_intocmit_functie = ?,
                        dosar_intocmit_semnatura = ?,
                        dosar_modificari_nume = ?,
                        dosar_modificari_functie = ?,
                        dosar_modificari_semnatura = ?,
                        dosar_certificare_nume = ?,
                        dosar_certificare_functie = ?,
                        dosar_certificare_semnatura = ?
                    WHERE id = ?
                    """,
                    (
                        dosar_nr,
                        functionar_public,
                        dosar_data_intocmire,
                        dosar_autoritate,
                        "\n".join(intocmit_nume_list).strip("\n"),
                        "\n".join(intocmit_functie_list).strip("\n"),
                        "\n".join(intocmit_semnatura_list).strip("\n"),
                        "\n".join(modif_nume_list).strip("\n"),
                        "\n".join(modif_functie_list).strip("\n"),
                        "\n".join(modif_semnatura_list).strip("\n"),
                        "\n".join(cert_nume_list).strip("\n"),
                        "\n".join(cert_functie_list).strip("\n"),
                        "\n".join(cert_semnatura_list).strip("\n"),
                        emp_id,
                    ),
                )
                conn.commit()
                st.success("Coperta «Dosar profesional» a fost salvată.")

        with sub_tab_print:
            st.markdown("#### Tipărire – Dosar profesional (coperta)")

            dosar_nr = get_val(emp_dosar, "dosar_nr")
            functionar_public = get_val(emp_dosar, "dosar_functionar_public")
            dosar_data_intocmire = get_val(emp_dosar, "dosar_data_intocmire")
            dosar_autoritate = get_val(emp_dosar, "dosar_autoritate")

            intocmit_nume_val = get_val(emp_dosar, "dosar_intocmit_nume")
            intocmit_functie_val = get_val(emp_dosar, "dosar_intocmit_functie")
            intocmit_semnatura_val = get_val(emp_dosar, "dosar_intocmit_semnatura")

            modif_nume_val = get_val(emp_dosar, "dosar_modificari_nume")
            modif_functie_val = get_val(emp_dosar, "dosar_modificari_functie")
            modif_semnatura_val = get_val(emp_dosar, "dosar_modificari_semnatura")

            cert_nume_val = get_val(emp_dosar, "dosar_certificare_nume")
            cert_functie_val = get_val(emp_dosar, "dosar_certificare_functie")
            cert_semnatura_val = get_val(emp_dosar, "dosar_certificare_semnatura")

            intocmit_nume_list = intocmit_nume_val.splitlines() if intocmit_nume_val else []
            intocmit_functie_list = (
                intocmit_functie_val.splitlines() if intocmit_functie_val else []
            )
            intocmit_semnatura_list = (
                intocmit_semnatura_val.splitlines() if intocmit_semnatura_val else []
            )

            modif_nume_list = modif_nume_val.splitlines() if modif_nume_val else []
            modif_functie_list = (
                modif_functie_val.splitlines() if modif_functie_val else []
            )
            modif_semnatura_list = (
                modif_semnatura_val.splitlines() if modif_semnatura_val else []
            )

            cert_nume_list = cert_nume_val.splitlines() if cert_nume_val else []
            cert_functie_list = (
                cert_functie_val.splitlines() if cert_functie_val else []
            )
            cert_semnatura_list = (
                cert_semnatura_val.splitlines() if cert_semnatura_val else []
            )

            docx_bytes = generate_dosar_profesional_docx(
                emp_dosar,
                dosar_nr,
                functionar_public,
                dosar_data_intocmire,
                dosar_autoritate,
                intocmit_nume_list,
                intocmit_functie_list,
                intocmit_semnatura_list,
                modif_nume_list,
                modif_functie_list,
                modif_semnatura_list,
                cert_nume_list,
                cert_functie_list,
                cert_semnatura_list,
            )
            file_name = f"dosar_profesional_{get_val(emp_dosar, 'cnp', 'CNP') or emp_id}.docx"
            st.download_button(
                "⬇️ Descarcă formular DOSAR PROFESIONAL (DOCX)",
                data=docx_bytes,
                file_name=file_name,
                mime=(mimetypes.guess_type(file_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_download_dosar_profesional",
            )

    # ------------------ DATE CU CARACTER PERSONAL ------------------
    elif sub_menu == "Date cu caracter personal":
        st.markdown("#### Date cu caracter personal")

        # Layout pe 2 coloane, simetric: stânga = date personale, dreapta = adresă + localizare
        col_left, col_right = st.columns(2)

        with col_left:
            last_name = st.text_input(
                "Nume",
                value=get_val(emp_dosar, "last_name", "LAST_NAME", "nume", "NUME"),
                key=f"{key_prefix}dp_last_name",
            )
            first_name = st.text_input(
                "Prenume",
                value=get_val(emp_dosar, "first_name", "FIRST_NAME", "prenume", "PRENUME"),
                key=f"{key_prefix}dp_first_name",
            )
            cnp_val = st.text_input(
                "CNP",
                value=get_val(emp_dosar, "cnp", "CNP"),
                key=f"{key_prefix}dp_cnp",
            )
            stare_civila = st.text_input(
                "Stare civilă",
                value=get_val(emp_dosar, "stare_civila"),
                key=f"{key_prefix}dp_stare_civila",
            )

            default_nr_copii = safe_int(emp_dosar.get("nr_copii", 0), default=0)
            nr_copii = st.number_input(
                "Număr copii",
                min_value=0,
                max_value=20,
                step=1,
                value=default_nr_copii,
                key=f"{key_prefix}dp_nr_copii",
            )

            telefon_fix = st.text_input(
                "Telefon fix",
                value=get_val(emp_dosar, "telefon_fix"),
                key=f"{key_prefix}dp_telefon_fix",
            )
            mobil = st.text_input(
                "Telefon mobil",
                value=get_val(emp_dosar, "mobil"),
                key=f"{key_prefix}dp_mobil",
            )
            email = st.text_input(
                "Email",
                value=get_val(emp_dosar, "email"),
                key=f"{key_prefix}dp_email",
            )

        with col_right:
            strada = st.text_input(
                "Strada",
                value=get_val(emp_dosar, "strada"),
                key=f"{key_prefix}dp_strada",
            )
            numar = st.text_input(
                "Număr",
                value=get_val(emp_dosar, "numar"),
                key=f"{key_prefix}dp_numar",
            )
            bloc = st.text_input(
                "Bloc",
                value=get_val(emp_dosar, "bloc"),
                key=f"{key_prefix}dp_bloc",
            )
            scara = st.text_input(
                "Scară",
                value=get_val(emp_dosar, "scara"),
                key=f"{key_prefix}dp_scara",
            )
            apartament = st.text_input(
                "Apartament",
                value=get_val(emp_dosar, "apartament"),
                key=f"{key_prefix}dp_apartament",
            )
            cod_postal = st.text_input(
                "Cod poștal",
                value=get_val(emp_dosar, "cod_postal"),
                key=f"{key_prefix}dp_cod_postal",
            )
            localitate = st.text_input(
                "Localitate",
                value=get_val(emp_dosar, "localitate"),
                key=f"{key_prefix}dp_localitate",
            )
            judet = st.text_input(
                "Județ",
                value=get_val(emp_dosar, "judet"),
                key=f"{key_prefix}dp_judet",
            )

        st.markdown("##### Date act identitate")

        ci_col_left, ci_col_right = st.columns(2)
        with ci_col_left:
            ci_tip_act = st.text_input(
                "Tip act (CI/BI/Pașaport)",
                value=get_val(emp_dosar, "ci_tip_act"),
                key=f"{key_prefix}dp_ci_tip_act",
            )
            ci_serie = st.text_input(
                "Serie",
                value=get_val(emp_dosar, "ci_serie"),
                key=f"{key_prefix}dp_ci_serie",
            )
            ci_numar = st.text_input(
                "Număr",
                value=get_val(emp_dosar, "ci_numar"),
                key=f"{key_prefix}dp_ci_numar",
            )
        with ci_col_right:
            ci_eliberat_de = st.text_input(
                "Eliberat de",
                value=get_val(emp_dosar, "ci_eliberat_de"),
                key=f"{key_prefix}dp_ci_eliberat_de",
            )
            ci_data_eliberarii = st.text_input(
                "Data eliberării",
                value=get_val(emp_dosar, "ci_data_eliberarii"),
                key=f"{key_prefix}dp_ci_data_eliberarii",
            )

        st.markdown("---")

        # Grupare acțiuni formular: Salvează + Export secțiune pe același rând
        act_col1, act_col2 = st.columns([0.4, 0.6])

        with act_col1:
            if st.button(
                "💾 Salvează datele cu caracter personal",
                key=f"{key_prefix}btn_save_date_pers",
            ):
                cur = conn.cursor()
                cur.execute(
                    """
                    UPDATE employees
                    SET last_name = ?,
                        first_name = ?,
                        cnp = ?,
                        stare_civila = ?,
                        nr_copii = ?,
                        strada = ?,
                        numar = ?,
                        bloc = ?,
                        scara = ?,
                        apartament = ?,
                        cod_postal = ?,
                        localitate = ?,
                        judet = ?,
                        telefon_fix = ?,
                        mobil = ?,
                        email = ?,
                        ci_tip_act = ?,
                        ci_serie = ?,
                        ci_numar = ?,
                        ci_eliberat_de = ?,
                        ci_data_eliberarii = ?
                    WHERE id = ?
                    """,
                    (
                        last_name,
                        first_name,
                        cnp_val,
                        stare_civila,
                        int(nr_copii),
                        strada,
                        numar,
                        bloc,
                        scara,
                        apartament,
                        cod_postal,
                        localitate,
                        judet,
                        telefon_fix,
                        mobil,
                        email,
                        ci_tip_act,
                        ci_serie,
                        ci_numar,
                        ci_eliberat_de,
                        ci_data_eliberarii,
                        emp_id,
                    ),
                )
                conn.commit()
                st.success("Datele cu caracter personal au fost salvate.")

        # Export Word pentru secțiune (în coloana a doua)
        emp_export = dict(emp_dosar)
        emp_export.update(
            {
                "last_name": last_name,
                "first_name": first_name,
                "cnp": cnp_val,
                "stare_civila": stare_civila,
                "nr_copii": int(nr_copii),
                "strada": strada,
                "numar": numar,
                "bloc": bloc,
                "scara": scara,
                "apartament": apartament,
                "cod_postal": cod_postal,
                "localitate": localitate,
                "judet": judet,
                "telefon_fix": telefon_fix,
                "mobil": mobil,
                "email": email,
                "ci_tip_act": ci_tip_act,
                "ci_serie": ci_serie,
                "ci_numar": ci_numar,
                "ci_eliberat_de": ci_eliberat_de,
                "ci_data_eliberarii": ci_data_eliberarii,
            }
        )
        sec_docx = create_date_personale_docx(emp_export)
        sec_name = (
            f"dosar_date_personale_{get_val(emp_export, 'cnp', 'CNP') or emp_id}.docx"
        )

        with act_col2:
            st.download_button(
                "⬇️ Exportă secțiunea în Word",
                data=sec_docx,
                file_name=sec_name,
                mime=(mimetypes.guess_type(sec_name)[0] or "application/octet-stream"),
                    key=f"{key_prefix}btn_export_date_pers_docx",
            )

    # ------------------ STUDII ------------------
    elif sub_menu == "Studii și pregătire profesională":
        st.markdown("#### Studii și pregătire profesională")

        studii = st.text_area(
            "Studii (școli absolvite, specializări, diplome)",
            value=get_val(emp_dosar, "studii"),
            key=f"{key_prefix}studii",
        )
        profesie = st.text_input(
            "Profesia de bază",
            value=get_val(emp_dosar, "profesie"),
            key=f"{key_prefix}profesie",
        )
        calificare = st.text_input(
            "Calificări/atestări",
            value=get_val(emp_dosar, "calificare"),
            key=f"{key_prefix}calificare",
        )
        observatii = st.text_area(
            "Observații generale",
            value=get_val(emp_dosar, "observatii"),
            key=f"{key_prefix}observatii",
        )

        if st.button(
            "💾 Salvează studii și pregătire profesională",
            key=f"{key_prefix}btn_save_studii",
        ):
            cur = conn.cursor()
            cur.execute(
                """
                UPDATE employees
                SET studii = ?,
                    profesie = ?,
                    calificare = ?,
                    observatii = ?
                WHERE id = ?
                """,
                (studii, profesie, calificare, observatii, emp_id),
            )
            conn.commit()
            st.success("Secțiunea «Studii și pregătire profesională» a fost salvată.")

        emp_export = dict(emp_dosar)
        emp_export.update(
            {
                "studii": studii,
                "profesie": profesie,
                "calificare": calificare,
                "observatii": observatii,
            }
        )
        sec_docx = create_studii_docx(emp_export)
        sec_name = (
            f"dosar_studii_{get_val(emp_export, 'cnp', 'CNP') or emp_id}.docx"
        )

        st.download_button(
            "⬇️ Exportă secțiunea în Word",
            data=sec_docx,
            file_name=sec_name,
            mime=(mimetypes.guess_type(sec_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_export_studii_docx",
        )

    # ------------------ ACTIVITATE ÎN AFARA FUNCȚIEI ------------------
    elif sub_menu == "Activitate în afara funcției publice":
        st.markdown("#### Activitate în afara funcției publice")

        activitate_ext = st.text_area(
            "Descrierea activităților desfășurate în afara funcției publice "
            "(contracte de muncă/colaborare, activități independente, ONG etc.)",
            value=get_val(emp_dosar, "activitate_in_afara_functiei"),
            key=f"{key_prefix}activitate_ext",
            height=200,
        )

        if st.button(
            "💾 Salvează activitatea în afara funcției publice",
            key=f"{key_prefix}btn_save_activ_ext",
        ):
            cur = conn.cursor()
            cur.execute(
                "UPDATE employees SET activitate_in_afara_functiei = ? WHERE id = ?",
                (activitate_ext, emp_id),
            )
            conn.commit()
            st.success("Secțiunea «Activitate în afara funcției publice» a fost salvată.")

        emp_export = dict(emp_dosar)
        emp_export.update({"activitate_in_afara_functiei": activitate_ext})
        sec_docx = create_activitate_externa_docx(emp_export)
        sec_name = (
            f"dosar_activitate_externa_{get_val(emp_export, 'cnp', 'CNP') or emp_id}.docx"
        )

        st.download_button(
            "⬇️ Exportă secțiunea în Word",
            data=sec_docx,
            file_name=sec_name,
            mime=(mimetypes.guess_type(sec_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_export_activ_ext_docx",
        )

    # ------------------ ACTIVITATE ÎN CADRUL INSTITUȚIEI ------------------
    elif sub_menu == "Activitate în cadrul instituției":
        st.markdown("#### Activitate în cadrul instituției")

        # -------------------------------------------------
        # 1) CITIM UNITĂȚILE DIN ORGANIGRAMĂ
        # -------------------------------------------------
        try:
            cur_org = conn.cursor()
            cur_org.execute(
                "SELECT id, name, parent_id, type FROM org_units WHERE is_active = 1 ORDER BY id"
            )
            org_rows = cur_org.fetchall()
        except Exception:
            org_rows = []

        org_paths: List[str] = []

        if org_rows:
            # mapăm id -> (name, parent_id, type)
            units_map = {
                row["id"]: (row["name"], row["parent_id"], row["type"])
                for row in org_rows
            }

            def build_path(u_id: int) -> str:
                # construim calea ierarhică: Direcție / Serviciu / Compartiment
                parts = []
                seen = set()
                while u_id and u_id in units_map and u_id not in seen:
                    seen.add(u_id)
                    name, parent_id, _tip = units_map[u_id]
                    parts.append(name)
                    u_id = parent_id
                return " / ".join(reversed(parts))

            for row in org_rows:
                path = build_path(row["id"])
                if path and path not in org_paths:
                    org_paths.append(path)

        # -------------------------------------------------
        # 2) DATE PRINCIPALE CONTRACT
        # -------------------------------------------------
        st.markdown("##### Date principale contract")

        col1, col2, col3 = st.columns(3)
        with col1:
            functie = st.text_input(
                "Funcție actuală",
                value=get_val(emp_dosar, "functie"),
                key=f"{key_prefix}functie",
            )
            departament = st.text_input(
                "Compartiment/Departament",
                value=get_val(emp_dosar, "departament"),
                key=f"{key_prefix}departament",
            )
        with col2:
            data_angajare = st.text_input(
                "Data angajării în instituție (YYYY-MM-DD)",
                value=get_val(emp_dosar, "data_angajare"),
                key=f"{key_prefix}data_angajare",
            )
            tip_contract = st.text_input(
                "Tip contract (perioadă determinată/nedeterminată etc.)",
                value=get_val(emp_dosar, "tip_contract"),
                key=f"{key_prefix}tip_contract",
            )
        with col3:
            loc_munca = st.text_input(
                "Loc de muncă/Compartiment organizatoric",
                value=get_val(emp_dosar, "loc_munca"),
                key=f"{key_prefix}loc_munca",
            )
            departament_organizatoric = st.text_input(
                "Structura organizatorică",
                value=get_val(emp_dosar, "departament_organizatoric"),
                key=f"{key_prefix}departament_organizatoric",
            )

        # -------------------------------------------------
        # 3) LEGARE CU ORGANIGRAMA (OPȚIONAL)
        # -------------------------------------------------
        loc_munca_final = loc_munca
        departament_organizatoric_final = departament_organizatoric

        if org_paths:
            st.markdown("##### Legare cu organigrama")

            # încercăm să găsim o potrivire implicită pe baza valorii deja salvate
            current_struct = (
                get_val(emp_dosar, "departament_organizatoric")
                or get_val(emp_dosar, "loc_munca")
            )
            default_idx = 0
            if current_struct:
                for i, p in enumerate(org_paths):
                    if current_struct.strip().lower() in p.strip().lower():
                        default_idx = i
                        break

            use_org = st.checkbox(
                "Preia «Loc de muncă» și «Structură organizatorică» din organigramă",
                value=True,
                key=f"{key_prefix}use_org_from_struct",
            )

            choice = st.selectbox(
                "Alege structura din organigramă",
                options=org_paths,
                index=default_idx,
                key=f"{key_prefix}org_struct_choice",
            )

            if use_org and choice:
                loc_munca_final = choice.split(" / ")[-1]  # ultimul nivel
                departament_organizatoric_final = choice

                st.info(
                    "La salvare se vor folosi valorile din organigramă pentru "
                    "«Loc de muncă» și «Structură organizatorică»."
                )

        # -------------------------------------------------
        # 4) DETALII SUPLIMENTARE
        # -------------------------------------------------
        st.markdown("##### Detalii suplimentare")

        col4, col5, col6 = st.columns(3)
        with col4:
            functie_contract = st.text_input(
                "Funcție conform contractului",
                value=get_val(emp_dosar, "functie_contract"),
                key=f"{key_prefix}functie_contract",
            )
        with col5:
            tip_norma = st.text_input(
                "Tip normă (normă întreagă/fracțiune normă)",
                value=get_val(emp_dosar, "tip_norma"),
                key=f"{key_prefix}tip_norma",
            )
        with col6:
            program_munca = st.text_input(
                "Program de muncă",
                value=get_val(emp_dosar, "program_munca"),
                key=f"{key_prefix}program_munca",
            )

        activitate_int = st.text_area(
            "Descriere sumarizată a activităților, funcțiilor deținute, "
            "promovărilor sau schimbărilor în carieră în cadrul instituției",
            value=get_val(emp_dosar, "activitate_in_cadru_institutie"),
            key=f"{key_prefix}activitate_int",
            height=200,
        )

        # -------------------------------------------------
        # 5) SALVARE ÎN DB
        # -------------------------------------------------
        if st.button(
            "💾 Salvează activitatea în cadrul instituției",
            key=f"{key_prefix}btn_save_activ_int",
        ):
            cur = conn.cursor()
            cur.execute(
                """
                UPDATE employees
                SET functie = ?,
                    departament = ?,
                    data_angajare = ?,
                    tip_contract = ?,
                    loc_munca = ?,
                    departament_organizatoric = ?,
                    functie_contract = ?,
                    tip_norma = ?,
                    program_munca = ?,
                    activitate_in_cadru_institutie = ?
                WHERE id = ?
                """,
                (
                    functie,
                    departament,
                    data_angajare,
                    tip_contract,
                    loc_munca_final,
                    departament_organizatoric_final,
                    functie_contract,
                    tip_norma,
                    program_munca,
                    activitate_int,
                    emp_id,
                ),
            )
            conn.commit()
            st.success("Secțiunea «Activitate în cadrul instituției» a fost salvată.")

        # -------------------------------------------------
        # 6) EXPORT DOCX
        # -------------------------------------------------
        emp_export = dict(emp_dosar)
        emp_export.update(
            {
                "functie": functie,
                "departament": departament,
                "data_angajare": data_angajare,
                "tip_contract": tip_contract,
                "loc_munca": loc_munca_final,
                "departament_organizatoric": departament_organizatoric_final,
                "functie_contract": functie_contract,
                "tip_norma": tip_norma,
                "program_munca": program_munca,
                "activitate_in_cadru_institutie": activitate_int,
            }
        )
        sec_docx = create_activitate_institutie_docx(emp_export)
        sec_name = (
            f"dosar_activitate_institutie_{get_val(emp_export, 'cnp', 'CNP') or emp_id}.docx"
        )

        st.download_button(
            "⬇️ Exportă secțiunea în Word",
            data=sec_docx,
            file_name=sec_name,
            mime=(mimetypes.guess_type(sec_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_export_activ_int_docx",
        )


    # ------------------ SITUAȚIA DREPTURILOR SALARIALE ------------------
    elif sub_menu == "Situația drepturilor salariale":
        st.markdown("#### Situația drepturilor salariale")

        salariu_baza = st.number_input(
            "Salariul de bază (actual)",
            min_value=0.0,
            step=1.0,
            value=float(emp_dosar.get("salariu_baza", 0) or 0),
            key=f"{key_prefix}salariu_baza",
        )

        situatia_sal = st.text_area(
            "Evoluția drepturilor salariale (majorări, indexări, sporuri, prime etc.)",
            value=get_val(emp_dosar, "situatia_drepturi_salariale"),
            key=f"{key_prefix}situatia_sal",
            height=200,
        )

        if st.button(
            "💾 Salvează situația drepturilor salariale",
            key=f"{key_prefix}btn_save_situatia_sal",
        ):
            cur = conn.cursor()
            cur.execute(
                """
                UPDATE employees
                SET salariu_baza = ?,
                    situatia_drepturi_salariale = ?
                WHERE id = ?
                """,
                (salariu_baza, situatia_sal, emp_id),
            )
            conn.commit()
            st.success("Secțiunea «Situația drepturilor salariale» a fost salvată.")

        emp_export = dict(emp_dosar)
        emp_export.update(
            {
                "salariu_baza": salariu_baza,
                "situatia_drepturi_salariale": situatia_sal,
            }
        )
        sec_docx = create_situatia_salariala_docx(emp_export)
        sec_name = (
            f"dosar_situatia_salariala_{get_val(emp_export, 'cnp', 'CNP') or emp_id}.docx"
        )

        st.download_button(
            "⬇️ Exportă secțiunea în Word",
            data=sec_docx,
            file_name=sec_name,
            mime=(mimetypes.guess_type(sec_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_export_situatia_sal_docx",
        )

    # ------------------ SITUAȚIA CONCEDIILOR ------------------
    elif sub_menu == "Situația concediilor":
        st.markdown("#### Situația concediilor")

        situatia_concedii = st.text_area(
            "Situația concediilor (de odihnă, medicale, fără plată, alte tipuri) "
            "- se poate nota pe ani sau pe perioade relevante.",
            value=get_val(emp_dosar, "situatia_concedii"),
            key=f"{key_prefix}situatia_concedii",
            height=220,
        )

        if st.button(
            "💾 Salvează situația concediilor",
            key=f"{key_prefix}btn_save_situatia_concedii",
        ):
            cur = conn.cursor()
            cur.execute(
                "UPDATE employees SET situatia_concedii = ? WHERE id = ?",
                (situatia_concedii, emp_id),
            )
            conn.commit()
            st.success("Secțiunea «Situația concediilor» a fost salvată.")

        emp_export = dict(emp_dosar)
        emp_export.update({"situatia_concedii": situatia_concedii})
        sec_docx = create_concedii_docx(emp_export)
        sec_name = (
            f"dosar_situatia_concedii_{get_val(emp_export, 'cnp', 'CNP') or emp_id}.docx"
        )

        st.download_button(
            "⬇️ Exportă secțiunea în Word",
            data=sec_docx,
            file_name=sec_name,
            mime=(mimetypes.guess_type(sec_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_export_situatia_concedii_docx",
        )

    # ------------------ SITUAȚIA DISCIPLINARĂ ------------------
    elif sub_menu == "Situația disciplinară":
        st.markdown("#### Situația disciplinară")

        situatia_disc = st.text_area(
            "Abateri disciplinare, sancțiuni aplicate, avertismente, "
            "decizii de sancționare și data încetării efectelor acestora.",
            value=get_val(emp_dosar, "situatia_disciplinara"),
            key=f"{key_prefix}situatia_disciplinara",
            height=220,
        )

        if st.button(
            "💾 Salvează situația disciplinară",
            key=f"{key_prefix}btn_save_situatia_disciplinara",
        ):
            cur = conn.cursor()
            cur.execute(
                "UPDATE employees SET situatia_disciplinara = ? WHERE id = ?",
                (situatia_disc, emp_id),
            )
            conn.commit()
            st.success("Secțiunea «Situația disciplinară» a fost salvată.")

        emp_export = dict(emp_dosar)
        emp_export.update({"situatia_disciplinara": situatia_disc})
        sec_docx = create_disciplinar_docx(emp_export)
        sec_name = (
            f"dosar_situatia_disciplinara_{get_val(emp_export, 'cnp', 'CNP') or emp_id}.docx"
        )

        st.download_button(
            "⬇️ Exportă secțiunea în Word",
            data=sec_docx,
            file_name=sec_name,
            mime=(mimetypes.guess_type(sec_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_export_situatia_disciplinara_docx",
        )

    # ------------------ ACCESUL LA DOSAR ------------------
    elif sub_menu == "Accesul la dosarul profesional":
        # marker pentru styling local
        st.markdown('<span id="dosar-acces-scope"></span>', unsafe_allow_html=True)
        st.markdown(
            """
            <style>
            section.main:has(#dosar-acces-scope) .access-form{
              max-width: 420px;
            }
            section.main:has(#dosar-acces-scope) .access-form input,
            section.main:has(#dosar-acces-scope) .access-form textarea{
              max-width: 420px;
            }
            </style>
            """,
            unsafe_allow_html=True,
        )

        st.markdown("#### Accesul la dosarul profesional")

        st.markdown("##### Istoric accesări")
        rows = list_dosar_acces(conn, emp_id)
        if rows:
            df = pd.DataFrame(rows, columns=rows[0].keys())
            if "employee_id" in df.columns:
                df = df.drop(columns=["employee_id"])
            st.dataframe(df, use_container_width=True)
        else:
            st.info("Nu există încă înregistrări de acces la dosarul profesional pentru acest angajat.")

        st.markdown("##### Adaugă o nouă accesare a dosarului")

        # Toate câmpurile pe un singur flux vertical, într-un formular mai îngust
        st.markdown('<div class="access-form">', unsafe_allow_html=True)
        nume_acces = st.text_input(
            "Numele și prenumele persoanei care accesează dosarul",
            key=f"{key_prefix}acces_nume",
        )
        functie_acces = st.text_input(
            "Funcția persoanei care accesează dosarul",
            key=f"{key_prefix}acces_functie",
        )
        semnatura_acces = st.text_input(
            "Semnătura persoanei care accesează",
            key=f"{key_prefix}acces_semnatura",
        )
        motiv_acces = st.text_area(
            "Motivul accesării dosarului",
            key=f"{key_prefix}acces_motiv",
        )
        autorizat_de = st.text_input(
            "Acces autorizat de (nume și funcție)",
            key=f"{key_prefix}acces_autorizat_de",
        )
        luat_la_cunostinta = st.text_input(
            "Luat la cunoștință (semnătura funcționarului public)",
            key=f"{key_prefix}acces_luat_la_cunostinta",
        )

        # Acțiuni: una sub alta, în ordinea corectă
        if st.button(
            "💾 Înregistrează accesul la dosar",
            key=f"{key_prefix}btn_add_acces",
        ):
            if not nume_acces or not motiv_acces:
                st.error("Completează cel puțin numele persoanei și motivul accesării.")
            else:
                insert_dosar_acces(
                    conn,
                    emp_id,
                    nume_acces,
                    functie_acces,
                    semnatura_acces,
                    motiv_acces,
                    autorizat_de,
                    luat_la_cunostinta,
                )
                st.success("Accesul la dosarul profesional a fost înregistrat.")
                rows = list_dosar_acces(conn, emp_id)

        # Export Word pentru secțiune (sub butonul de înregistrare)
        rows_for_export = list_dosar_acces(conn, emp_id)
        sec_docx = create_acces_docx(emp_dosar, rows_for_export)
        sec_name = (
            f"dosar_acces_{get_val(emp_dosar, 'cnp', 'CNP') or emp_id}.docx"
        )

        st.download_button(
            "⬇️ Exportă secțiunea în Word",
            data=sec_docx,
            file_name=sec_name,
            mime=(mimetypes.guess_type(sec_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_export_acces_docx",
        )
        st.markdown("</div>", unsafe_allow_html=True)  # end .access-form

    # ------------------ REGISTRU EVIDENȚĂ ------------------
    elif sub_menu == "Registru evidență funcționari publici":
        st.markdown("#### Registru evidență funcționari publici")

        # Număr registru și data înscrierii pe același rând
        reg_col1, reg_col2 = st.columns(2)
        with reg_col1:
            registru_numar = st.text_input(
                "Număr în registrul de evidență",
                value=get_val(emp_dosar, "registru_numar"),
                key=f"{key_prefix}registru_numar",
            )
        with reg_col2:
            import datetime as _dt

            _registru_data_raw = get_val(emp_dosar, "registru_data") or ""
            _default_date = None
            try:
                if _registru_data_raw:
                    _default_date = _dt.datetime.strptime(_registru_data_raw, "%Y-%m-%d").date()
            except Exception:
                _default_date = None

            _today = _dt.date.today()
            _min_date = _today.replace(year=_today.year - 100)
            _max_date = _today.replace(year=_today.year + 10)

            registru_data_date = st.date_input(
                "Data înscrierii în registru",
                value=_default_date,
                min_value=_min_date,
                max_value=_max_date,
                key=f"{key_prefix}registru_data",
            )
            registru_data = registru_data_date.isoformat() if registru_data_date else ""
        registru_obs = st.text_area(
            "Observații (de ex. transfer, încetare raporturi de serviciu, alte mențiuni)",
            value=get_val(emp_dosar, "registru_observatii"),
            key=f"{key_prefix}registru_obs",
            height=200,
        )

        if st.button(
            "💾 Salvează datele din registrul de evidență",
            key=f"{key_prefix}btn_save_registru",
        ):
            cur = conn.cursor()
            cur.execute(
                """
                UPDATE employees
                SET registru_numar = ?,
                    registru_data = ?,
                    registru_observatii = ?
                WHERE id = ?
                """,
                (registru_numar, registru_data, registru_obs, emp_id),
            )
            conn.commit()
            st.success("Secțiunea «Registru evidență funcționari publici» a fost salvată.")

        emp_export = dict(emp_dosar)
        emp_export.update(
            {
                "registru_numar": registru_numar,
                "registru_data": registru_data,
                "registru_observatii": registru_obs,
            }
        )
        sec_docx = create_registru_docx(emp_export)
        sec_name = (
            f"dosar_registru_{get_val(emp_export, 'cnp', 'CNP') or emp_id}.docx"
        )

        st.download_button(
            "⬇️ Exportă secțiunea în Word",
            data=sec_docx,
            file_name=sec_name,
            mime=(mimetypes.guess_type(sec_name)[0] or "application/octet-stream"),
                key=f"{key_prefix}btn_export_registru_docx",
        )


# -------------------------------------------------------------
# ORGANIGRAMA - STRUCTURA ORGANIZATORICĂ
# -------------------------------------------------------------


def ensure_org_schema(conn: sqlite3.Connection) -> None:
    """
    Creează (doar dacă nu există) tabelele și indexurile pentru organigramă.
    Nu modifică tabele existente.
    """
    cur = conn.cursor()
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS org_units (
            id INTEGER PRIMARY KEY,
            name TEXT NOT NULL,
            parent_id INTEGER NULL,
            type TEXT NULL,
            is_active INTEGER NOT NULL DEFAULT 1
        );
        """
    )
    cur.execute(
        "CREATE INDEX IF NOT EXISTS idx_org_units_parent ON org_units(parent_id);"
    )
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS org_positions (
            id INTEGER PRIMARY KEY,
            unit_id INTEGER NOT NULL,
            title TEXT NOT NULL,
            is_active INTEGER NOT NULL DEFAULT 1
        );
        """
    )
    cur.execute(
        "CREATE INDEX IF NOT EXISTS idx_org_positions_unit ON org_positions(unit_id);"
    )
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_positions (
            id INTEGER PRIMARY KEY,
            position_id INTEGER NOT NULL,
            employee_id INTEGER NOT NULL,
            start_date TEXT NULL,
            end_date TEXT NULL
        );
        """
    )
    cur.execute(
        "CREATE INDEX IF NOT EXISTS idx_emp_pos_pos ON employee_positions(position_id);"
    )
    cur.execute(
        "CREATE INDEX IF NOT EXISTS idx_emp_pos_emp ON employee_positions(employee_id);"
    )
    conn.commit()


def seed_org_demo(conn: sqlite3.Connection) -> str:
    """
    Inserează date demo în organigramă doar dacă org_units este goală.
    Returnează mesaj: "Org already populated." sau descrierea inserărilor.
    """
    cur = conn.cursor()
    try:
        n = cur.execute("SELECT COUNT(*) FROM org_units").fetchone()[0]
    except Exception:
        return "Org already populated."
    if n > 0:
        return "Org already populated."

    from datetime import date
    today = date.today().isoformat()

    # Root
    cur.execute(
        """
        INSERT INTO org_units(name, parent_id, type, is_active)
        VALUES(?, NULL, ?, 1)
        """,
        ("AGENTIA NATIONALA DE CADASTRU SI PUBLICITATE IMOBILIARA", "Instituție"),
    )
    root_id = cur.lastrowid

    # Direcția IT
    cur.execute(
        """
        INSERT INTO org_units(name, parent_id, type, is_active)
        VALUES(?, ?, ?, 1)
        """,
        ("Direcția IT", root_id, "Direcție"),
    )
    directia_it_id = cur.lastrowid

    # Serviciul Informatica 4
    cur.execute(
        """
        INSERT INTO org_units(name, parent_id, type, is_active)
        VALUES(?, ?, ?, 1)
        """,
        ("Serviciul Informatica 4", directia_it_id, "Serviciu"),
    )
    serviciu_id = cur.lastrowid

    # Posturi – doar dacă nu există deja
    cur.execute(
        """
        SELECT id FROM org_positions
        WHERE unit_id = ? AND title = ?
        LIMIT 1
        """,
        (directia_it_id, "Director IT"),
    )
    row_dir = cur.fetchone()
    if row_dir is None:
        cur.execute(
            """
            INSERT INTO org_positions(unit_id, title, is_active)
            VALUES(?, ?, 1)
            """,
            (directia_it_id, "Director IT"),
        )

    cur.execute(
        """
        SELECT id FROM org_positions
        WHERE unit_id = ? AND title = ?
        LIMIT 1
        """,
        (serviciu_id, "Inginer software"),
    )
    row_ing = cur.fetchone()
    if row_ing is None:
        cur.execute(
            """
            INSERT INTO org_positions(unit_id, title, is_active)
            VALUES(?, ?, 1)
            """,
            (serviciu_id, "Inginer software"),
        )
        inginer_pos_id = cur.lastrowid
    else:
        inginer_pos_id = row_ing[0]

    # Legătură angajat -> post "Inginer software" dacă există vreun angajat
    try:
        emp_row = cur.execute(
            "SELECT id FROM employees ORDER BY id LIMIT 1"
        ).fetchone()
    except Exception:
        emp_row = None
    if emp_row is not None and inginer_pos_id is not None:
        emp_id = emp_row[0]
        cur.execute(
            """
            INSERT INTO employee_positions(position_id, employee_id, start_date, end_date)
            VALUES(?, ?, ?, NULL)
            """,
            (inginer_pos_id, emp_id, today),
        )

    conn.commit()
    msg = (
        "Organigramă demo: 1 instituție (root), Direcția IT, Serviciul Informatica 4; "
        "posturi Director IT și Inginer software."
    )
    if emp_row is not None:
        msg += " Un angajat a fost legat de postul Inginer software."
    else:
        msg += " Niciun angajat în DB – nu s-a legat nimeni de post."
    return msg


def ensure_organigrama_tables(conn: sqlite3.Connection) -> None:
    """
    Creează (dacă nu există) tabelele pentru organigramă:
      - org_units (locuri de muncă / structuri)
      - org_positions (posturi)
      - employee_positions (repartizare angajat–post)
    """
    cur = conn.cursor()

    # Tabela unități (Direcții / Servicii / Compartimente / Birouri)
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS org_units (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            parent_id INTEGER,
            type TEXT,
            is_active INTEGER DEFAULT 1
        );
        """
    )

    # Tabela posturi
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS org_positions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            unit_id INTEGER,
            reports_to INTEGER,
            is_active INTEGER DEFAULT 1
        );
        """
    )

    # Legătura angajat - post (istoric simplu)
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_positions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER,
            position_id INTEGER,
            start_date TEXT,
            end_date TEXT
        );
        """
    )

    conn.commit()


def generate_organigrama_from_stat(conn: sqlite3.Connection) -> Tuple[int, int]:
    """
    Generează / actualizează organigrama (org_units + org_positions)
    pe baza tabelului stat_functii.

    Folosește coloanele:
      - "Loc de munca"  -> devine loc de muncă în org_units (name)
      - "Functie"       -> devine post în org_positions (title)

    Returnează: (nr_locuri_noi, nr_posturi_noi)
    """
    cur = conn.cursor()

    # Verificăm dacă există tabela stat_functii
    cur.execute(
        "SELECT name FROM sqlite_master WHERE type='table' AND name='stat_functii'"
    )
    exists = cur.fetchone() is not None
    if not exists:
        return 0, 0

    # Citim tot statul de funcții
    df_stat = pd.read_sql_query("SELECT * FROM stat_functii", conn)
    if df_stat.empty:
        return 0, 0

    if "Loc de munca" not in df_stat.columns or "Functie" not in df_stat.columns:
        # Coloanele nu sunt prezente – nu putem genera organigrama
        raise ValueError(
            'În tabelul stat_functii trebuie să existe coloanele "Loc de munca" și "Functie".'
        )

    # Citim locurile de muncă existente
    df_units = pd.read_sql_query(
        "SELECT id, name FROM org_units WHERE is_active = 1",
        conn,
    )
    name_to_id: Dict[str, int] = {}
    if not df_units.empty:
        for _, row in df_units.iterrows():
            key = str(row["name"]).strip().lower()
            name_to_id[key] = int(row["id"])

    # Citim posturile existente
    df_pos = pd.read_sql_query(
        "SELECT id, title, unit_id FROM org_positions WHERE is_active = 1",
        conn,
    )
    existing_pos: set[Tuple[str, Optional[int]]] = set()
    if not df_pos.empty:
        for _, r in df_pos.iterrows():
            title = str(r["title"]).strip().lower()
            unit_id = int(r["unit_id"]) if r["unit_id"] is not None else None
            existing_pos.add((title, unit_id))

    new_units = 0
    new_positions = 0

    for _, r in df_stat.iterrows():
        loc_raw = r.get("Loc de munca")
        functie_raw = r.get("Functie")

        loc_name = str(loc_raw).strip() if not pd.isna(loc_raw) else ""
        functie_name = str(functie_raw).strip() if not pd.isna(functie_raw) else ""

        if not loc_name and not functie_name:
            continue

        # 1) Loc de muncă -> org_units.name
        unit_id: Optional[int] = None
        if loc_name:
            key = loc_name.lower()
            if key not in name_to_id:
                cur.execute(
                    "INSERT INTO org_units (name, parent_id, type) VALUES (?, ?, ?)",
                    (loc_name, None, None),
                )
                unit_id = cur.lastrowid
                name_to_id[key] = unit_id
                new_units += 1
            else:
                unit_id = name_to_id[key]

        # 2) Funcție -> org_positions.title, legat de unit_id (dacă există)
        if functie_name:
            key_pos = (functie_name.lower(), unit_id)
            if key_pos in existing_pos:
                continue

            cur.execute(
                "INSERT INTO org_positions (title, unit_id, reports_to) VALUES (?, ?, ?)",
                (functie_name, unit_id, None),
            )
            new_positions += 1
            existing_pos.add(key_pos)

    conn.commit()
    return new_units, new_positions
    
def build_graphviz_organigrama_color(df_units):
    """
    Organigramă vizuală MARE și COLORATĂ pentru ecran + PDF vectorial.
    Acceptă coloane:
      - id, name, parent_id, type
      - sau ID_UNITATE, DENUMIRE, PARENT_ID, TIP
    """
    import pandas as pd

    df = df_units.copy()

    # Normalizăm coloanele
    if "ID_UNITATE" not in df.columns and "id" in df.columns:
        df = df.rename(
            columns={
                "id": "ID_UNITATE",
                "name": "DENUMIRE",
                "parent_id": "PARENT_ID",
                "type": "TIP",
            }
        )

    def color_for_tip(tip: str) -> str:
        t = (tip or "").strip().lower()
        if "director" in t:
            return "#4F81BD"   # albastru
        if "direc" in t:
            return "#4F81BD"
        if "serviciu" in t:
            return "#9BBB59"   # verde
        if "compartiment" in t:
            return "#F79646"   # portocaliu
        if "birou" in t:
            return "#C0504D"   # roșu
        return "#D9D9D9"       # gri deschis

    lines = []
    lines.append("digraph ORG_COLOR {")
    lines.append("  rankdir=TB;")
    # layout aerisit, font mare
    lines.append('  graph [splines=ortho, bgcolor="#FFFFFF"];')
    lines.append("  ranksep=1.8;")
    lines.append("  nodesep=1.2;")
    lines.append(
        '  node [shape=box, style="rounded,filled", '
        'fontname="Calibri", fontsize=22, color="#4A4A4A", '
        'fontcolor="#1F1F1F", margin="0.7,0.5", penwidth=2.0];'
    )
    lines.append('  edge [color="#7F7F7F", penwidth=1.6];')

    # NODURI
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        name = str(row["DENUMIRE"]).strip()
        tip = str(row.get("TIP", "") or "").strip()

        label = name
        if tip:
            label = f"{name}\\n({tip})"

        fill = color_for_tip(tip)

        lines.append(
            f'  "{uid}" [label="{label}", fillcolor="{fill}", '
            f'width=4.5, height=1.6];'
        )

    # MUCHII părinte → copil
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        parent = row.get("PARENT_ID", None)
        if parent is not None and not pd.isna(parent):
            lines.append(f'  "{int(parent)}" -> "{uid}";')

    lines.append("}")
    return "\n".join(lines)
 
def build_graphviz_organigrama(df_units):
    """
    Construiește sursa Graphviz (DOT) pentru organigramă
    pe baza dataframe-ului df_units.
    Acceptă coloane:
      - id, name, parent_id, type
    sau
      - ID_UNITATE, DENUMIRE, PARENT_ID, TIP
    """
    import pandas as pd

    df = df_units.copy()

    # Normalizăm coloanele
    if "ID_UNITATE" not in df.columns and "id" in df.columns:
        df = df.rename(
            columns={
                "id": "ID_UNITATE",
                "name": "DENUMIRE",
                "parent_id": "PARENT_ID",
                "type": "TIP",
            }
        )

    lines = []
    lines.append("digraph ORG {")
    # orientare de sus în jos
    lines.append("  rankdir=TB;")
    # mai mult aer între nivele și noduri
    lines.append("  ranksep=1.12;")
    lines.append("  nodesep=1.6;")
    # linii mai fine
    lines.append('  edge [penwidth=1.2, color="#555555"];')
    # noduri mari, stil organigramă, albastru deschis
    lines.append(
        '  node [shape=box, style="rounded,filled", '
        'fillcolor="#D9E2F3", color="#2F5597", fontcolor="#1F1F1F", '
        'fontsize=32, penwidth=1.5, margin="0.7,0.5"];'
    )

    # -------------------------------------------------
    # NODURI
    # -------------------------------------------------
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        name = str(row["DENUMIRE"]).strip()
        tip = str(row.get("TIP", "") or "").strip()

        label = name
        if tip:
            label = f"{name}\\n({tip})"

        # un pic de forțare pentru box-uri late
        lines.append(
            f'  "{uid}" [label="{label}", width=3.5, height=1.0];'
        )

    # -------------------------------------------------
    # MUCHII PĂRINTE → COPIL
    # -------------------------------------------------
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        parent = row.get("PARENT_ID", None)
        if parent is not None and not pd.isna(parent):
            lines.append(f'  "{int(parent)}" -> "{uid}";')

    # -------------------------------------------------
    # NIVELURI PE TIP (Direcție / Serviciu / Compartiment / Birou)
    # -------------------------------------------------
    nivele = ["Direcție", "Serviciu", "Compartiment", "Birou"]
    for nivel in nivele:
        noduri_nivel = df[df["TIP"] == nivel]["ID_UNITATE"].tolist()
        if noduri_nivel:
            lines.append("{ rank=same; " + " ".join(f'"{int(u)}"' for u in noduri_nivel) + " }")

    lines.append("}")
    return "\n".join(lines)

def ensure_pontaj_tables(conn: sqlite3.Connection):
    """Creează / actualizează structura de bază pentru modulul de pontaj."""
    cur = conn.cursor()
    # Dacă tabela nu există, o creăm cu toate coloanele necesare
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS pontaj (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            employee_id INTEGER NOT NULL,
            data TEXT NOT NULL,
            ore_lucru REAL DEFAULT 0,
            tip_zi TEXT,
            observatii TEXT
        );
        """
    )
    # Dacă tabela există deja dar fără unele coloane (ex: pontaj vechi), le adăugăm
    cur.execute("PRAGMA table_info(pontaj);")
    existing_cols = {row[1].lower() for row in cur.fetchall()}

    if "ore_lucru" not in existing_cols:
        cur.execute("ALTER TABLE pontaj ADD COLUMN ore_lucru REAL DEFAULT 0;")
    if "tip_zi" not in existing_cols:
        cur.execute("ALTER TABLE pontaj ADD COLUMN tip_zi TEXT;")
    if "observatii" not in existing_cols:
        cur.execute("ALTER TABLE pontaj ADD COLUMN observatii TEXT;")

    # Index pentru performanță
    cur.execute(
        "CREATE INDEX IF NOT EXISTS idx_pontaj_emp_date ON pontaj(employee_id, data);"
    )
    conn.commit()

import pandas as pd

# ==========================================================
# FUNCȚII AUXILIARE ORGANIGRAMĂ
# ==========================================================

def load_organigrama_units(conn):
    """
    Încarcă toate unitățile din organigramă.
    Presupune tabel: ORGANIGRAMA(ID_UNITATE, DENUMIRE, TIP_UNITATE, PARENT_ID)
    """
    query = """
        SELECT
            ID_UNITATE,
            DENUMIRE,
            TIP_UNITATE,
            PARENT_ID
        FROM ORGANIGRAMA
        ORDER BY ID_UNITATE
    """
    df = pd.read_sql(query, conn)
    return df


def update_parent_unit(conn, child_id, parent_id):
    """
    Setează PARENT_ID pentru unitatea child_id.
    Dacă parent_id este None -> PARENT_ID = NULL (nivel 0).
    """
    import fdb

    cursor = conn.cursor()
    if parent_id is None:
        sql = "UPDATE ORGANIGRAMA SET PARENT_ID = NULL WHERE ID_UNITATE = ?"
        params = (child_id,)
    else:
        sql = "UPDATE ORGANIGRAMA SET PARENT_ID = ? WHERE ID_UNITATE = ?"
        params = (parent_id, child_id)

    cursor.execute(sql, params)
    conn.commit()

def render_tree(nodes, roots, level=0):
    """
    Afișează în Streamlit organigrama ca listă ierarhică pe baza
    structurii returnate de build_tree().
    nodes: {id_unitate: {"info": row, "children": [ids copii]}}
    roots: listă cu ID-urile nodurilor rădăcină
    """
    indent = "   " * level
    bullet = "• "

    for uid in roots:
        info = nodes[uid]["info"]

        # denumirea unității
        nume = info.get("DENUMIRE", str(uid))

        # tipul unității (dacă există coloana TIP)
        tip = info.get("TIP", "")
        try:
            import pandas as pd
            if pd.isna(tip):
                tip = ""
        except Exception:
            pass

        if tip:
            st.markdown(f"{indent}{bullet}**{nume}** *( {tip} )*")
        else:
            st.markdown(f"{indent}{bullet}**{nume}**")

        copii = nodes[uid]["children"]
        if copii:
            render_tree(nodes, copii, level + 1)

def build_tree(df_units):
    """
    Construiește o structură arborescentă din df_units.
    Returnează un dict: {id_unitate: {"info": row, "children": [ids copii]}}
    """
    nodes = {}
    for _, row in df_units.iterrows():
        uid = row["ID_UNITATE"]
        nodes[uid] = {"info": row, "children": []}

    roots = []

    for uid, node in nodes.items():
        parent_id = node["info"]["PARENT_ID"]
        if pd.isna(parent_id):
            roots.append(uid)
        else:
            parent_id = int(parent_id)
            if parent_id in nodes:
                nodes[parent_id]["children"].append(uid)
            else:
                # Dacă PARENT_ID nu mai există în listă -> considerăm root
                roots.append(uid)

    return nodes, roots

def build_org_dot(
    df_units,
    style: str = "screen",   # "screen" sau "print"
    orientation: str = "TB"  # "TB" (sus->jos) sau "LR" (stânga->dreapta)
):
    """
    Generează sursa DOT pentru organigramă, cu 2 stiluri:
    - style="screen" -> pentru ecran
    - style="print"  -> pentru print (A3/A2, casete mari)
    """

    import pandas as pd

    df = df_units.copy()

    # Normalizăm coloanele
    if "ID_UNITATE" not in df.columns and "id" in df.columns:
        df = df.rename(
            columns={
                "id": "ID_UNITATE",
                "name": "DENUMIRE",
                "parent_id": "PARENT_ID",
                "type": "TIP",
            }
        )

    lines = []
    lines.append("digraph ORG {")
    lines.append(f"  rankdir={orientation};")

    if style == "print":
        ranksep = 2.2
        nodesep = 1.5
        fontsize = 22
        margin = "0.9,0.7"
        width = 5.0
        height = 1.8
    else:  # screen
        ranksep = 1.4
        nodesep = 1.0
        fontsize = 16
        margin = "0.6,0.4"
        width = 3.5
        height = 1.2

    lines.append(f"  ranksep={ranksep};")
    lines.append(f"  nodesep={nodesep};")
    lines.append('  edge [penwidth=1.4, color="#555555"];')
    lines.append(
        f'  node [shape=box, style="rounded,filled", '
        f'fillcolor="#D9E2F3", color="#2F5597", fontcolor="#1F1F1F", '
        f'fontsize={fontsize}, penwidth=1.6, margin="{margin}"];'
    )

    # NODURI
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        name = str(row["DENUMIRE"]).strip()
        tip = str(row.get("TIP", "") or "").strip()

        label = name
        if tip:
            label = f"{name}\\n({tip})"

        lines.append(
            f'  "{uid}" [label="{label}", width={width}, height={height}];'
        )

    # MUCHII
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        parent = row.get("PARENT_ID", None)
        if parent is not None and not pd.isna(parent):
            lines.append(f'  "{int(parent)}" -> "{uid}";')

    # NIVELURI pe tip (opțional)
    nivele = ["Direcție", "Serviciu", "Compartiment", "Birou"]
    for nivel in nivele:
        noduri_nivel = df[df["TIP"] == nivel]["ID_UNITATE"].tolist()
        if noduri_nivel:
            lines.append("{ rank=same; " + " ".join(f'"{int(u)}"' for u in noduri_nivel) + " }")

    lines.append("}")
    return "\n".join(lines)

# ========================= ORGANIGRAMA – LAYOUT MANUAL (Drag&Drop) =========================

def ensure_org_layout_table(conn):
    """Creează tabela ORG_LAYOUT dacă nu există (SQLite/Firebird friendly)."""
    try:
        cur = conn.cursor()
        cur.execute(
            """
            CREATE TABLE IF NOT EXISTS ORG_LAYOUT (
                ID_UNITATE INTEGER PRIMARY KEY,
                X REAL NOT NULL,
                Y REAL NOT NULL,
                UPDATED_AT TEXT
            )
            """
        )
        conn.commit()
    except Exception:
        # fallback fără IF NOT EXISTS
        try:
            cur = conn.cursor()
            cur.execute(
                """
                CREATE TABLE ORG_LAYOUT (
                    ID_UNITATE INTEGER PRIMARY KEY,
                    X REAL NOT NULL,
                    Y REAL NOT NULL,
                    UPDATED_AT TEXT
                )
                """
            )
            conn.commit()
        except Exception:
            pass


def load_org_layout(conn) -> dict:
    """Returnează dict: {id_unitate(str): (x,y)}"""
    try:
        cur = conn.cursor()
        cur.execute("SELECT ID_UNITATE, X, Y FROM ORG_LAYOUT")
        rows = cur.fetchall()
        out = {}
        for rid, x, y in rows:
            out[str(rid)] = (float(x), float(y))
        return out
    except Exception:
        return {}


def reset_org_layout(conn):
    try:
        cur = conn.cursor()
        cur.execute("DELETE FROM ORG_LAYOUT")
        conn.commit()
    except Exception:
        pass


def save_org_layout(conn, layout_json: dict):
    """layout_json format: {id: {'x':..,'y':..}}"""
    from datetime import datetime
    cur = conn.cursor()
    now = datetime.now().isoformat(timespec="seconds")
    for sid, pos in (layout_json or {}).items():
        try:
            x = float(pos.get("x", 0))
            y = float(pos.get("y", 0))
        except Exception:
            continue
        try:
            cur.execute(
                "INSERT OR REPLACE INTO ORG_LAYOUT (ID_UNITATE, X, Y, UPDATED_AT) VALUES (?, ?, ?, ?)",
                (int(sid), x, y, now),
            )
        except Exception:
            try:
                # Firebird: UPDATE OR INSERT
                cur.execute(
                    "UPDATE OR INSERT INTO ORG_LAYOUT (ID_UNITATE, X, Y, UPDATED_AT) VALUES (?, ?, ?, ?) MATCHING (ID_UNITATE)",
                    (int(sid), x, y, now),
                )
            except Exception:
                pass
    try:
        conn.commit()
    except Exception:
        pass


def auto_layout_from_hierarchy(df_units, scale=1.0, dx=240, dy=140):
    """Auto layout după ierarhie: 1 sus, apoi împarte pe niveluri.
    Acceptă df cu coloane (ID_UNITATE/DENUMIRE/PARENT_ID) sau (id/name/parent_id).
    """
    import pandas as pd

    if df_units is None or df_units.empty:
        return {}

    df = df_units.copy()
    cols = {c.upper(): c for c in df.columns}
    rename = {}
    if "ID_UNITATE" in cols and "id" not in df.columns:
        rename[cols["ID_UNITATE"]] = "id"
    if "DENUMIRE" in cols and "name" not in df.columns:
        rename[cols["DENUMIRE"]] = "name"
    if "PARENT_ID" in cols and "parent_id" not in df.columns:
        rename[cols["PARENT_ID"]] = "parent_id"
    if rename:
        df = df.rename(columns=rename)

    df["id"] = pd.to_numeric(df["id"], errors="coerce")
    df["parent_id"] = pd.to_numeric(df["parent_id"], errors="coerce")

    # roots: parent_id null sau parent inexistent
    ids = set([int(x) for x in df["id"].dropna().tolist()])
    roots = []
    for _, r in df.iterrows():
        if pd.isna(r["id"]):
            continue
        cid = int(r["id"])
        pid = r.get("parent_id", None)
        if pid is None or pd.isna(pid) or int(pid) not in ids:
            roots.append(cid)

    # dacă avem "DIRECTOR GENERAL" îl punem primul
    def _name(i):
        row = df[df["id"] == i]
        if row.empty:
            return ""
        return str(row.iloc[0].get("name", "")).upper()

    roots = sorted(roots, key=lambda i: (0 if "DIRECTOR GENERAL" in _name(i) else 1, _name(i)))

    # adjacency
    children = {}
    for _, r in df.iterrows():
        if pd.isna(r["id"]):
            continue
        cid = int(r["id"])
        pid = r.get("parent_id", None)
        if pid is None or pd.isna(pid):
            continue
        try:
            pid = int(pid)
        except Exception:
            continue
        children.setdefault(pid, []).append(cid)

    # BFS levels
    level = {}
    queue = []
    for r in roots:
        level[r] = 0
        queue.append(r)

    while queue:
        u = queue.pop(0)
        for v in children.get(u, []):
            if v in level:
                continue
            level[v] = level[u] + 1
            queue.append(v)

    max_level = max(level.values()) if level else 0
    by_level = {k: [] for k in range(max_level + 1)}
    for nid, lv in level.items():
        by_level.setdefault(lv, []).append(nid)

    # sort nodes in each level by name
    for lv in by_level:
        by_level[lv] = sorted(by_level[lv], key=lambda i: _name(i))

    dx = float(dx) * float(scale)
    dy = float(dy) * float(scale)

    pos = {}
    x_center = 600.0 * float(scale)  # centru canvas

    for lv in range(max_level + 1):
        nodes = by_level.get(lv, [])
        if not nodes:
            continue
        n = len(nodes)
        total_w = (n - 1) * dx
        x0 = x_center - total_w / 2.0
        y = 60.0 * float(scale) + lv * dy
        for i, nid in enumerate(nodes):
            x = x0 + i * dx
            pos[str(nid)] = (x, y)

    return pos


def render_org_drag_editor(df_units, positions: dict, scale=1.0):
    """Editor drag&drop în HTML cu muchii (SVG). Returnează layout_json dict."""
    import json as _json
    import pandas as pd

    if df_units is None or df_units.empty:
        st.info("Nu există unități pentru editor.")
        return {}

    df = df_units.copy()
    cols = {c.upper(): c for c in df.columns}
    rename = {}
    if "ID_UNITATE" in cols and "id" not in df.columns:
        rename[cols["ID_UNITATE"]] = "id"
    if "DENUMIRE" in cols and "name" not in df.columns:
        rename[cols["DENUMIRE"]] = "name"
    if "PARENT_ID" in cols and "parent_id" not in df.columns:
        rename[cols["PARENT_ID"]] = "parent_id"
    if rename:
        df = df.rename(columns=rename)

    df["id"] = pd.to_numeric(df["id"], errors="coerce")
    df["parent_id"] = pd.to_numeric(df["parent_id"], errors="coerce")

    layout = {}
    edges = []
    for _, r in df.iterrows():
        if pd.isna(r["id"]):
            continue
        sid = str(int(r["id"]))
        name = str(r.get("name", ""))
        x, y = positions.get(sid, (0.0, 0.0))
        layout[sid] = {"x": float(x), "y": float(y), "label": name}
        pid = r.get("parent_id", None)
        if pid is not None and not pd.isna(pid):
            try:
                edges.append([str(int(pid)), sid])
            except Exception:
                pass

    data_json = _json.dumps({"nodes": layout, "edges": edges}, ensure_ascii=False)

    html = f"""
    <style>
      .org-canvas {{
        position: relative;
        width: 100%;
        height: 720px;
        border: 1px solid #263244;
        background: rgba(11,18,32,0.92);
        overflow: auto;
        border-radius: 14px;
      }}
      .org-svg {{
        position: absolute;
        left: 0; top: 0;
        width: 2400px;
        height: 1800px;
        pointer-events: none;
      }}
      .org-node {{
        position: absolute;
        padding: 8px 10px;
        border: 1px solid rgba(34,197,94,0.55);
        border-radius: 12px;
        background: rgba(15,23,42,0.95);
        color: #F8FAFC;
        font-family: Arial, sans-serif;
        font-size: 12px;
        cursor: grab;
        user-select: none;
        white-space: nowrap;
        box-shadow: none;
      }}
    </style>

    <div class="org-canvas" id="orgCanvas">
      <svg class="org-svg" id="orgSvg">
        <defs>
          <marker id="arrow" markerWidth="10" markerHeight="10" refX="7" refY="3" orient="auto" markerUnits="strokeWidth">
            <path d="M0,0 L0,6 L9,3 z" fill="#999"></path>
          </marker>
        </defs>
      </svg>
    </div>

    <textarea id="orgLayoutOut" style="width:100%; height:120px; margin-top:8px;">{data_json}</textarea>

    <script>
      const canvas = document.getElementById('orgCanvas');
      const svg = document.getElementById('orgSvg');
      const out = document.getElementById('orgLayoutOut');
      let payload = JSON.parse(out.value);
      let nodes = payload.nodes || {{}};
      let edges = payload.edges || [];

      function nodeCenter(el){{
        const r = el.getBoundingClientRect();
        const cr = canvas.getBoundingClientRect();
        return {{
          x: (r.left - cr.left) + r.width/2 + canvas.scrollLeft,
          y: (r.top - cr.top) + r.height/2 + canvas.scrollTop
        }};
      }}

      function nodeRect(el){{
  const r = el.getBoundingClientRect();
  const cr = canvas.getBoundingClientRect();
  return {{
    x: (r.left - cr.left) + canvas.scrollLeft,
    y: (r.top - cr.top) + canvas.scrollTop,
    w: r.width,
    h: r.height
  }};
}}

      function redrawEdges(){{
        // clear old lines
        const olds = svg.querySelectorAll('line, polyline');
        olds.forEach(l => l.remove());

        edges.forEach(([p,c]) => {{
          const pe = document.getElementById('node_'+p);
          const ce = document.getElementById('node_'+c);
          if(!pe || !ce) return;
          const pr = nodeRect(pe);
const cr = nodeRect(ce);

// ancore: jos din părinte -> sus în copil
const p0 = {{ x: pr.x + pr.w/2, y: pr.y + pr.h }};
const c0 = {{ x: cr.x + cr.w/2, y: cr.y }};

// mid-Y inițial (între părinte și copil)
let yMid = (p0.y + c0.y) / 2;

// evităm să "tăiem" prin alte casete: dacă segmentul orizontal la yMid
// intersectează o casetă, căutăm un yMid liber (în pași de 28px)
const xMin = Math.min(p0.x, c0.x), xMax = Math.max(p0.x, c0.x);
const step = 28;

function intersectsAtY(yTest) {{
  for (const nid in nodes) {{
    // ignorăm capetele
    if (nid === p || nid === c) continue;
    const elN = document.getElementById('node_' + nid);
    if (!elN) continue;
    const rr = nodeRect(elN);
    const yIn = (yTest >= rr.y && yTest <= (rr.y + rr.h));
    if (!yIn) continue;
    // overlap pe X cu segmentul orizontal
    const xOverlap = !(xMax < rr.x || xMin > (rr.x + rr.w));
    if (xOverlap) return true;
  }}
  return false;
}}

if (intersectsAtY(yMid)) {{
  for (let k = 1; k <= 20; k++) {{
    const up = yMid - step * k;
    if (up > 0 && !intersectsAtY(up)) {{ yMid = up; break; }}
    const down = yMid + step * k;
    if (down < 1750 && !intersectsAtY(down)) {{ yMid = down; break; }}
  }}
}}

const pts = [
  `${{p0.x}},${{p0.y}}`,
  `${{p0.x}},${{yMid}}`,
  `${{c0.x}},${{yMid}}`,
  `${{c0.x}},${{c0.y}}`
].join(' ');

const pl = document.createElementNS('http://www.w3.org/2000/svg','polyline');
pl.setAttribute('points', pts);
pl.setAttribute('fill', 'none');
pl.setAttribute('stroke', '#999');
pl.setAttribute('stroke-width', '2');
pl.setAttribute('marker-end', 'url(#arrow)');
svg.appendChild(pl);
        }});
      }}

      function render(){{
        // clear nodes
        const olds = canvas.querySelectorAll('.org-node');
        olds.forEach(n => n.remove());

        Object.keys(nodes).forEach(id => {{
          const n = nodes[id];
          const el = document.createElement('div');
          el.className = 'org-node';
          el.id = 'node_'+id;
          el.innerText = n.label || id;
          el.style.left = (n.x || 0) + 'px';
          el.style.top  = (n.y || 0) + 'px';
          canvas.appendChild(el);

          let dragging=false, ox=0, oy=0;

          el.addEventListener('mousedown', (e) => {{
            dragging=true;
            el.style.cursor='grabbing';
            ox = e.clientX - el.offsetLeft;
            oy = e.clientY - el.offsetTop;
          }});

          window.addEventListener('mousemove', (e) => {{
            if(!dragging) return;
            const x = e.clientX - ox;
            const y = e.clientY - oy;
            el.style.left = x + 'px';
            el.style.top  = y + 'px';
            nodes[id].x = x;
            nodes[id].y = y;
            out.value = JSON.stringify({{nodes:nodes, edges:edges}});
            redrawEdges();
          }});

          window.addEventListener('mouseup', () => {{
            if(!dragging) return;
            dragging=false;
            el.style.cursor='grab';
          }});
        }});

        redrawEdges();
      }}

      render();
      canvas.addEventListener('scroll', redrawEdges);
      window.addEventListener('resize', redrawEdges);
    </script>
    """

    st.components.v1.html(html, height=900, scrolling=True)

    # parse back JSON
    try:
        payload2 = _json.loads(st.session_state.get("_dummy_org_layout", "{}"))  # fallback (nu folosit)
    except Exception:
        payload2 = None

    # st.components nu ne dă direct valoare, deci folosim textarea ca copy/paste:
    st.caption("📌 Layout JSON se actualizează automat în caseta de mai sus. Apasă **Salvează layout-ul** ca să îl scriem în DB.")
    try:
        raw = st.text_area("Layout JSON (editabil)", value=data_json, height=140, key="org_layout_json_text")
        payload = _json.loads(raw)
        return payload.get("nodes", {})
    except Exception:
        return layout

def generate_org_files_from_dot(dot_src: str, export_scale: float = 1.0):
    """
    Dintr-un DOT (Graphviz), generează:
    - PNG (bytes)
    - PDF (bytes)
    - PPTX (bytes) cu un slide ce conține organigrama

    Returnează: (png_bytes, pdf_bytes, pptx_bytes)
    """
    import graphviz
    import tempfile
    import os
    import io
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation
    from pptx.util import Inches

    src = graphviz.Source(dot_src)

    with tempfile.TemporaryDirectory() as tmpdir:
        png_path = src.render(
            filename=os.path.join(tmpdir, "organigrama_full"),
            format="png",
            cleanup=True,
        )

        # PNG
        with open(png_path, "rb") as f:
            png_bytes = f.read()

        # Aplicăm zoom export REAL (mărire/micșorare) pe imagine, astfel încât
        # ce vezi în preview = ce ajunge în PDF/PNG.
        try:
            s = float(export_scale or 1.0)
        except Exception:
            s = 1.0
        if s <= 0:
            s = 1.0
        if abs(s - 1.0) > 1e-3:
            try:
                im0 = Image.open(io.BytesIO(png_bytes)).convert("RGBA")
                new_w = max(1, int(im0.size[0] * s))
                new_h = max(1, int(im0.size[1] * s))
                im1 = im0.resize((new_w, new_h), resample=Image.LANCZOS)
                buf_png = io.BytesIO()
                im1.save(buf_png, format="PNG")
                png_bytes = buf_png.getvalue()
            except Exception:
                pass


        # PDF
        img = Image.open(io.BytesIO(png_bytes))
        pdf_buffer = io.BytesIO()
        img.save(pdf_buffer, format="PDF", resolution=300.0)
        pdf_buffer.seek(0)
        pdf_bytes = pdf_buffer.getvalue()

        # PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        img_stream = io.BytesIO(png_bytes)
        pic = slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.5),
                                       width=Inches(11))  # adaptăm la lățime
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        pptx_bytes = pptx_buffer.getvalue()

    return png_bytes, pdf_bytes, pptx_bytes

def print_tree_text(nodes, roots, level=0):
    """
    Afișează în Streamlit organigrama ca listă ierarhică text.
    """
    indent = "   " * level
    bullet = "• "

    for uid in roots:
        info = nodes[uid]["info"]
        denumire = info["DENUMIRE"]
        tip = info.get("TIP_UNITATE", "")
        st.markdown(f"{indent}{bullet}**{denumire}** *( {tip} )*")
        children = nodes[uid]["children"]
        if children:
            print_tree_text(nodes, children, level + 1)

# ==========================================================
# PAGINA ORGANIGRAMA (VERSIUNEA COMPLETĂ ȘI CORECTĂ)
# ==========================================================

# ==========================================================
# FUNCȚII AJUTĂTOARE PENTRU ORGANIGRAMĂ
# ==========================================================

def build_tree(df_units):
    """
    Construiește structura arborescentă din dataframe-ul cu unități.
    Acceptă coloane:
      - ID_UNITATE, DENUMIRE, PARENT_ID, TIP
    sau
      - id, name, parent_id, type
    """
    import pandas as pd

    df = df_units.copy()

    # Normalizăm numele coloanelor
    if "ID_UNITATE" not in df.columns and "id" in df.columns:
        df = df.rename(
            columns={
                "id": "ID_UNITATE",
                "name": "DENUMIRE",
                "parent_id": "PARENT_ID",
                "type": "TIP",
            }
        )

    nodes = {}
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        nodes[uid] = {"info": row, "children": []}

    roots = []

    for uid, node in nodes.items():
        parent = node["info"].get("PARENT_ID", None)
        if parent is None or (isinstance(parent, float) and pd.isna(parent)):
            roots.append(uid)
        else:
            parent = int(parent)
            if parent in nodes:
                nodes[parent]["children"].append(uid)
            else:
                roots.append(uid)

    return nodes, roots


def render_tree(nodes, roots, level=0):
    """
    Afișează organigrama ca listă ierarhică în Streamlit,
    pe baza structurii (nodes, roots) întoarsă de build_tree().
    """
    indent = "   " * level
    bullet = "• "

    for uid in roots:
        info = nodes[uid]["info"]

        nume = info.get("DENUMIRE", str(uid))
        tip = info.get("TIP", "")

        try:
            import pandas as pd
            if pd.isna(tip):
                tip = ""
        except Exception:
            pass

        if tip:
            st.markdown(f"{indent}{bullet}**{nume}** *( {tip} )*")
        else:
            st.markdown(f"{indent}{bullet}**{nume}**")

        copii = nodes[uid]["children"]
        if copii:
            render_tree(nodes, copii, level + 1)

def build_graphviz_organigrama(df_units):
    """
    Construiește sursa Graphviz (DOT) pentru organigramă
    pe baza dataframe-ului df_units.
    Acceptă coloane:
      - id, name, parent_id, type
    sau
      - ID_UNITATE, DENUMIRE, PARENT_ID, TIP
    """
    import pandas as pd

    df = df_units.copy()

    # Normalizăm coloanele
    if "ID_UNITATE" not in df.columns and "id" in df.columns:
        df = df.rename(
            columns={
                "id": "ID_UNITATE",
                "name": "DENUMIRE",
                "parent_id": "PARENT_ID",
                "type": "TIP",
            }
        )

    lines = []
    lines.append("digraph ORG {")
    lines.append('  rankdir=TB;')
    lines.append('  node [shape=box, style="rounded,filled", fontsize=10];')

    # Noduri
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        name = str(row["DENUMIRE"])
        tip = str(row.get("TIP", "") or "")

        label = name
        if tip:
            label = f"{name}\\n({tip})"

        lines.append(f'  "{uid}" [label="{label}"];')

    # Muchii părinte → copil
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        parent = row.get("PARENT_ID", None)
        if parent is not None and not pd.isna(parent):
            lines.append(f'  "{int(parent)}" -> "{uid}";')

    lines.append("}")
    return "\n".join(lines)
    
def build_graphviz_organigrama_print(df_units):
    """
    Variantă specială pentru PRINT (A3/A2):
    - casete foarte mari
    - font mare
    - spațiu mare între niveluri
    Acceptă coloane:
      - id, name, parent_id, type
    sau
      - ID_UNITATE, DENUMIRE, PARENT_ID, TIP
    """
    import pandas as pd

    df = df_units.copy()

    # Normalizăm coloanele
    if "ID_UNITATE" not in df.columns and "id" in df.columns:
        df = df.rename(
            columns={
                "id": "ID_UNITATE",
                "name": "DENUMIRE",
                "parent_id": "PARENT_ID",
                "type": "TIP",
            }
        )

    lines = []
    lines.append("digraph ORG_PRINT {")
    lines.append("  rankdir=TB;")
    lines.append("  ranksep=2.2;")
    lines.append("  nodesep=1.5;")
    lines.append('  edge [penwidth=1.6, color="#555555"];')
    lines.append(
        '  node [shape=box, style="rounded,filled", '
        'fillcolor="#D9E2F3", color="#2F5597", fontcolor="#1F1F1F", '
        'fontsize=22, penwidth=1.8, margin="0.9,0.7"];'
    )

    # NODURI
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        name = str(row["DENUMIRE"]).strip()
        tip = str(row.get("TIP", "") or "").strip()

        label = name
        if tip:
            label = f"{name}\\n({tip})"

        lines.append(
            f'  "{uid}" [label="{label}", width=5.0, height=1.8];'
        )

    # MUCHII
    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        parent = row.get("PARENT_ID", None)
        if parent is not None and not pd.isna(parent):
            lines.append(f'  "{int(parent)}" -> "{uid}";')

    # NIVELURI pe tip
    nivele = ["Direcție", "Serviciu", "Compartiment", "Birou"]
    for nivel in nivele:
        noduri_nivel = df[df["TIP"] == nivel]["ID_UNITATE"].tolist()
        if noduri_nivel:
            lines.append("{ rank=same; " + " ".join(f'"{int(u)}"' for u in noduri_nivel) + " }")

    lines.append("}")
    return "\n".join(lines)
    
def build_organigrama_tree_structure(df_units):
    """
    Construiește structura arborelui de organigramă
    și întoarce:
      - nodes: dict[id] = {id, name, tip_raw, parent_id}
      - children_by_parent: dict[parent_id] = [child_id, ...]
      - roots: list de id-uri fără părinte
    Normalizăm automat coloanele (id/name/parent_id/type).
    """
    import pandas as pd

    df = df_units.copy()

    # Normalizăm coloanele
    if "ID_UNITATE" not in df.columns and "id" in df.columns:
        df = df.rename(
            columns={
                "id": "ID_UNITATE",
                "name": "DENUMIRE",
                "parent_id": "PARENT_ID",
                "type": "TIP",
            }
        )

    nodes = {}
    children_by_parent = {}

    for _, row in df.iterrows():
        uid = int(row["ID_UNITATE"])
        name = str(row["DENUMIRE"]).strip()
        tip_raw = str(row.get("TIP", "") or "").strip()
        parent = row.get("PARENT_ID", None)

        nodes[uid] = {
            "id": uid,
            "name": name,
            "tip_raw": tip_raw,
            "parent_id": int(parent) if pd.notna(parent) else None,
        }

        pid = int(parent) if pd.notna(parent) else None
        children_by_parent.setdefault(pid, []).append(uid)

    roots = children_by_parent.get(None, [])

    # Sortare alfabetică simplă pentru toți copiii
    def sort_key(uid: int):
        return str(nodes[uid]["name"]).lower()

    for pid, lst in children_by_parent.items():
        children_by_parent[pid] = sorted(lst, key=sort_key)

    return nodes, children_by_parent, roots
   
def render_organigrama_tree(df_units: pd.DataFrame, df_stat: pd.DataFrame | None = None, *, key_prefix: str = "org_tree"):
    """
    Afișează organigrama ca arbore ierarhic, FĂRĂ st.expander (Streamlit nu permite expandere imbricate).

    - Extinderea/colapsarea pe nod se face cu st.checkbox (salvat în st.session_state).
    - Opțional: dacă primește df_stat (ID_UNITATE, DENUMIRE_FUNCTIE, NUME, PRENUME),
      afișează posturile + ocupanții sub fiecare unitate (legare cu statul de funcții).

    df_units trebuie să aibă (sau să poată fi mapat la):
      - ID_UNITATE, DENUMIRE, PARENT_ID, TIP
      sau
      - id, name, parent_id, type
    """
    if df_units is None or df_units.empty:
        st.info("Nu există unități în organigramă.")
        return

    # Injectăm un mic CSS ca să arate mai compact checkbox-ul în arbore
    st.markdown(
        """
        <style>
        /* micșorăm spațiul vertical dintre checkbox-uri în arbore */
        div[data-testid="stCheckbox"] { margin-bottom: -6px; }
        </style>
        """,
        unsafe_allow_html=True,
    )

    nodes, children_by_parent, roots = build_organigrama_tree_structure(df_units)

    if not roots:
        st.info("Nu sunt definite noduri rădăcină (PARENT_ID = NULL).")
        return

    # Culori pe NIVEL
    def level_bg(level: int) -> str:
        if level <= 0:
            return "#1f77b4"   # albastru
        elif level == 1:
            return "#2ca02c"   # verde
        elif level == 2:
            return "#ff7f0e"   # portocaliu
        else:
            return "#d62728"   # roșu

    # Icoane pe NIVEL
    def level_icon(level: int) -> str:
        if level <= 0:
            return "🏛️"
        elif level == 1:
            return "🏢"
        elif level == 2:
            return "📂"
        else:
            return "📌"

    def node_label(uid: int) -> str:
        n = nodes[uid]
        if n["tip_raw"]:
            return f"{n['name']} ({n['tip_raw']})"
        return n["name"]

    def node_badge_html(uid: int, level: int, bold: bool = False) -> str:
        label = node_label(uid)
        icon = level_icon(level)
        bg = level_bg(level)
        weight = "bold" if bold else "normal"
        return (
            f"<span style='background-color:{bg}; color:white; "
            f"padding:3px 10px; border-radius:10px; "
            f"font-weight:{weight}; font-size:0.9rem;'>"
            f"{icon} {label}</span>"
        )

    def _render_positions_for_unit(id_unitate: int, level: int):
        if df_stat is None or df_stat.empty:
            return
        if "ID_UNITATE" not in df_stat.columns:
            return
        try:
            df_unit_stat = df_stat[df_stat["ID_UNITATE"] == id_unitate]
        except Exception:
            return
        if df_unit_stat is None or df_unit_stat.empty:
            return

        # Afișăm posturile sub unitate (indentat)
        for _, rstat in df_unit_stat.iterrows():
            functie = str(rstat.get("DENUMIRE_FUNCTIE", "") or "").strip()
            nume = str(rstat.get("NUME", "") or "").strip()
            prenume = str(rstat.get("PRENUME", "") or "").strip()

            indent_p = "&nbsp;" * 4 * (level + 1)
            if nume or prenume:
                st.markdown(
                    f"{indent_p}- 👤 {nume} {prenume} – _{functie}_",
                    unsafe_allow_html=True,
                )
            elif functie:
                st.markdown(
                    f"{indent_p}- ⭕ Post neocupat – _{functie}_",
                    unsafe_allow_html=True,
                )

    # Recursiv: folosim checkbox ca expand/collapse (fără expandere imbricate)
    def render_node(uid: int, level: int = 0):
        has_children = uid in children_by_parent and len(children_by_parent[uid]) > 0
        indent_html = "&nbsp;" * 4 * level
        badge = node_badge_html(uid, level, bold=has_children)

        if has_children:
            # cheie stabilă în session_state
            state_key = f"{key_prefix}_open_{uid}"
            # label simplu (fără HTML în checkbox), badge-ul îl afișăm separat
            label_txt = node_label(uid)

            # rând: checkbox + badge
            c1, c2 = st.columns([0.12, 0.88], gap="small")
            with c1:
                opened = st.checkbox(" ", value=bool(st.session_state.get(state_key, False)), key=state_key)
            with c2:
                st.markdown(f"{indent_html}• {badge}", unsafe_allow_html=True)

            # posturi + oameni sub unitate (indiferent dacă e deschis/închis)
            _render_positions_for_unit(int(uid), level)

            if opened:
                for cid in children_by_parent.get(uid, []):
                    render_node(cid, level + 1)
        else:
            st.markdown(f"{indent_html}• {badge}", unsafe_allow_html=True)
            _render_positions_for_unit(int(uid), level)

    # Afișăm toate rădăcinile
    for rid in sorted(roots, key=lambda u: nodes[u]["name"].lower()):
        render_node(rid, level=0)


def get_organigrama_subtree(df_units, root_id: int):
    """
    Întoarce un dataframe DOAR cu nodurile din subarborele
    care pornește de la root_id (root + toate subunitățile lui).
    """
    import pandas as pd

    df = df_units.copy()

    # Normalizăm coloanele
    if "ID_UNITATE" not in df.columns and "id" in df.columns:
        df = df.rename(
            columns={
                "id": "ID_UNITATE",
                "name": "DENUMIRE",
                "parent_id": "PARENT_ID",
                "type": "TIP",
            }
        )

    # Construim structura pentru parcurgere
    nodes, children_by_parent, roots = build_organigrama_tree_structure(df)

    # Colectăm toți descendenții lui root_id
    subtree_ids = set()

    def collect(u: int):
        subtree_ids.add(u)
        for c in children_by_parent.get(u, []):
            collect(c)

    collect(int(root_id))

    # Filtrăm doar nodurile din subarbore
    df_sub = df[df["ID_UNITATE"].isin(subtree_ids)].copy()
    return df_sub

def load_org_units(conn):
    """Încarcă unitățile din tabela org_units (sursa unică pentru organigramă)."""
    try:
        df_units = pd.read_sql_query(
            "SELECT id, name, parent_id, type, is_active "
            "FROM org_units WHERE is_active = 1 ORDER BY id",
            conn,
        )
    except Exception as e:
        try:
            st.error(f"Eroare la încărcarea unităților din org_units: {e}")
        except Exception:
            pass
        return pd.DataFrame(columns=["id","name","parent_id","type","is_active"])
    return df_units

def fetch_org_units(conn):
        return pd.read_sql_query("""
            SELECT id, name, parent_id, type
            FROM org_units
            WHERE is_active = 1
            ORDER BY id
        """, conn)

def fetch_positions_with_people(conn):
        return pd.read_sql_query("""
            SELECT
                op.id          AS pos_id,
                op.unit_id     AS unit_id,
                COALESCE(op.title, 'Post') AS pos_title,
                e.id           AS emp_id,
                e.first_name   AS first_name,
                e.last_name    AS last_name,
                COALESCE(e.functie, '') AS functie
            FROM org_positions op
            LEFT JOIN employee_positions ep
                ON ep.position_id = op.id
                AND (ep.end_date IS NULL OR TRIM(ep.end_date) = '')
            LEFT JOIN employees e
                ON e.id = ep.employee_id
            WHERE op.is_active = 1
            ORDER BY op.unit_id, op.id
        """, conn)

def build_cyto_elements(df_units: pd.DataFrame, df_pos: pd.DataFrame):
    elements = []

    # Units nodes
    for _, r in df_units.iterrows():
        uid = int(r["id"])
        name = str(r["name"])
        # folosim aceeași logică de culoare ca pentru organigrama clasică
        _, _, fill, fontc = _color_for_unit(name.upper())
        elements.append({
            "data": {
                "id": f"unit:{uid}",
                "label": name,
                "type": "unit",
                "unit_id": uid,
                "parent_id": int(r["parent_id"]) if pd.notna(r["parent_id"]) else None,
                "unit_type": str(r.get("type") or ""),
                "bg": str(fill),
                "fc": str(fontc),
            }
        })

    # Unit edges (parent -> child)
    for _, r in df_units.iterrows():
        if pd.notna(r["parent_id"]):
            elements.append({
                "data": {
                    "id": f"edge:unit:{int(r['parent_id'])}->{int(r['id'])}",
                    "source": f"unit:{int(r['parent_id'])}",
                    "target": f"unit:{int(r['id'])}",
                    "type": "unit_edge",
                }
            })

    # Positions + people
    for _, r in df_pos.iterrows():
        pid = int(r["pos_id"])
        uid = int(r["unit_id"])
        pos_label = str(r.get("pos_title") or "Post")

        # position node – trebuie să aibă bg/fc pentru stilul Cytoscape
        elements.append({
            "data": {
                "id": f"pos:{pid}",
                "label": pos_label,
                "type": "pos",
                "pos_id": pid,
                "unit_id": uid,
                "emp_id": int(r["emp_id"]) if pd.notna(r["emp_id"]) else None,
                "bg": "#ffffff",
                "fc": "#000000",
            }
        })
        # edge unit -> position
        elements.append({
            "data": {
                "id": f"edge:unit:{uid}->pos:{pid}",
                "source": f"unit:{uid}",
                "target": f"pos:{pid}",
                "type": "unit_pos_edge",
            }
        })

        # employee node + edge pos -> employee (if assigned)
        if pd.notna(r["emp_id"]):
            emp_id = int(r["emp_id"])
            emp_name = f"{str(r.get('last_name') or '').strip()} {str(r.get('first_name') or '').strip()}".strip()
            if not emp_name:
                emp_name = f"Angajat #{emp_id}"

            # employee node – trebuie să aibă bg/fc pentru stilul Cytoscape
            elements.append({
                "data": {
                    "id": f"emp:{emp_id}",
                    "label": emp_name,
                    "type": "emp",
                    "emp_id": emp_id,
                    "pos_id": pid,
                    "unit_id": uid,
                    "bg": "#f5f5f5",
                    "fc": "#000000",
                }
            })
            elements.append({
                "data": {
                    "id": f"edge:pos:{pid}->emp:{emp_id}",
                    "source": f"pos:{pid}",
                    "target": f"emp:{emp_id}",
                    "type": "pos_emp_edge",
                }
            })

    return elements

def cyto_elements_to_payload(elements: list[dict]) -> dict:
        nodes, edges = [], []
        for el in elements or []:
            d = (el or {}).get("data", {})
            if not isinstance(d, dict):
                continue
            if "source" in d and "target" in d:
                edges.append(el)
            else:
                nodes.append(el)
        return {"nodes": nodes, "edges": edges}

def page_organigrama(conn):
    # Marker pentru a putea controla padding-ul doar pe Organigramă
    st.markdown('<span id="org-scope"></span>', unsafe_allow_html=True)
    def _inject_organigrama_page_css() -> None:
        st.markdown(
            """
            <style>
            section.main:has(#org-scope) .org-page .org-title{
              margin: 0 0 0.5rem 0 !important;
              padding: 0 !important;
            }
            section.main:has(#org-scope) .org-page .org-toolbar,
            section.main:has(#org-scope) .org-page .org-filters,
            section.main:has(#org-scope) .org-page .org-canvas,
            section.main:has(#org-scope) .org-page .org-bottom{
              border: 1px solid rgba(148,163,184,0.22);
              border-radius: 8px;
              background: rgba(15,23,42,0.18);
              padding: 12px 14px;
              margin: 0 0 12px 0;
              text-align: left;
            }
            section.main:has(#org-scope) .org-page .org-filters{
              margin-bottom: 18px;
            }
            section.main:has(#org-scope) .org-page .org-toolbar{
              max-width: 220px;
              display: flex;
              flex-direction: column;
              gap: 10px;
            }
            section.main:has(#org-scope) .org-page .org-bottom{ max-width: 560px; }
            section.main:has(#org-scope) .org-page .org-section-title{
              margin: 0 0 6px 0;
              font-size: 1.18rem;
              font-weight: 900;
              color: rgba(248,250,252,0.96);
              letter-spacing: 0.01em;
            }
            section.main:has(#org-scope) .org-page .org-section-desc{
              margin: 0;
              max-width: 860px;
              color: rgba(203,213,225,0.90);
              font-size: 0.94rem;
              line-height: 1.45;
            }
            section.main:has(#org-scope) .org-page .org-view-radio{
              margin-top: 12px;
            }
            section.main:has(#org-scope) .org-page .org-view-radio div[data-testid="stRadio"] [role="radiogroup"]{
              gap: 12px !important;
            }
            section.main:has(#org-scope) .org-page .org-view-radio div[data-testid="stRadio"] label p,
            section.main:has(#org-scope) .org-page .org-view-radio div[data-testid="stRadio"] label span{
              font-weight: 800 !important;
              color: rgba(248,250,252,0.96) !important;
            }
            section.main:has(#org-scope) .org-page div[data-testid="stTabs"] > div[role="tablist"]{
              display: flex !important;
              flex-direction: column !important;
              align-items: flex-start !important;
              gap: 10px !important;
              max-width: 220px !important;
              margin-bottom: 12px !important;
            }
            section.main:has(#org-scope) .org-page div[data-testid="stTabs"] > div[role="tablist"] button[role="tab"]{
              width: 220px !important;
              min-width: 220px !important;
              max-width: 220px !important;
              height: 42px !important;
              min-height: 42px !important;
              max-height: 42px !important;
              border-radius: 12px !important;
              justify-content: flex-start !important;
              text-align: left !important;
              font-weight: 700 !important;
            }
            section.main:has(#org-scope) .org-page .org-selected-row{
              margin: 0 0 6px 0;
              color: rgba(203,213,225,0.98);
            }
            section.main:has(#org-scope) .org-page .org-selected-row:last-child{ margin-bottom: 0; }
            .st-key-org_btn_top_refresh button,
            .st-key-org_btn_clear_selection button,
            .st-key-org_btn_sel_refresh button,
            .st-key-org_btn_sel_clear button{
              width: 220px !important;
              min-width: 220px !important;
              max-width: 220px !important;
              height: 42px !important;
              min-height: 42px !important;
              max-height: 42px !important;
            }
            </style>
            """,
            unsafe_allow_html=True,
        )
    _inject_organigrama_page_css()
    st.markdown("<div class='org-page'>", unsafe_allow_html=True)
    st.markdown('<h2 class="page-title org-title">Organigramă</h2>', unsafe_allow_html=True)

    try:
        ensure_org_schema(conn)
    except Exception:
        pass

    cur = conn.cursor()
    # ---------------------------------------------------------
    # 1) ÎNCĂRCARE UNITĂȚI
    # ---------------------------------------------------------
    def _load_org_units_any(conn):
        """Încearcă să încarce unitățile din tabele posibile, fără să crape dacă lipsesc."""
        candidates = [
            ("org_units", "SELECT id AS ID_UNITATE, name AS DENUMIRE, parent_id AS PARENT_ID, type AS TIP, is_active AS IS_ACTIVE FROM org_units WHERE is_active = 1 ORDER BY id"),
            ("ORG_UNITS", "SELECT id AS ID_UNITATE, name AS DENUMIRE, parent_id AS PARENT_ID, type AS TIP, is_active AS IS_ACTIVE FROM ORG_UNITS WHERE is_active = 1 ORDER BY id"),
            ("ORGANIGRAMA", "SELECT id AS ID_UNITATE, denumire AS DENUMIRE, parent_id AS PARENT_ID, tip AS TIP FROM ORGANIGRAMA ORDER BY id"),
            ("ORG_STRUCTURA", "SELECT id AS ID_UNITATE, denumire AS DENUMIRE, parent_id AS PARENT_ID, tip AS TIP FROM ORG_STRUCTURA ORDER BY id"),
        ]

        tables = set()
        try:
            tt = pd.read_sql_query("SELECT name FROM sqlite_master WHERE type='table'", conn)
            if tt is not None and not tt.empty and "name" in tt.columns:
                tables = set([str(x).lower() for x in tt["name"].tolist()])
        except Exception:
            pass

        last_err = None
        for tname, q in candidates:
            try:
                if tables and tname.lower() not in tables:
                    continue
                df = pd.read_sql_query(q, conn)
                if df is not None and not df.empty:
                    for c in ["ID_UNITATE", "DENUMIRE", "PARENT_ID", "TIP"]:
                        if c not in df.columns:
                            df[c] = None
                    return df
            except Exception as e:
                last_err = e
                continue

        raise RuntimeError(f"Nu am găsit o tabelă de organigramă compatibilă. Ultima eroare: {last_err}")

    try:
        df_units = _load_org_units_any(conn)
    except Exception as e:
        st.error(f"Eroare la încărcarea unităților: {e}")
        df_units = pd.DataFrame(columns=["id", "name", "parent_id", "type", "is_active"])

    if df_units.empty:
        st.warning("Nu există unități în organigramă.")

    # ---------------------------------------------------------
    # 2) VALIDARE STRUCTURĂ
    # ---------------------------------------------------------

    def validate_org_units(df_units: pd.DataFrame) -> list[str]:
        issues = []
        if df_units is None or df_units.empty:
            return ["Nu există unități în org_units."]

        df = df_units.copy()
        df["id"] = pd.to_numeric(df["id"], errors="coerce")
        df["parent_id"] = pd.to_numeric(df["parent_id"], errors="coerce")

        ids = set(df["id"].dropna().astype(int).tolist())

        # parent lipsă
        for _, r in df.iterrows():
            pid = r["parent_id"]
            if pd.notna(pid) and int(pid) not in ids:
                issues.append(f"Unitatea {int(r['id'])} are parent_id={int(pid)} care nu există.")

        # multiple roots / no root
        roots = df[df["parent_id"].isna()]["id"].dropna().astype(int).tolist()
        if len(roots) == 0:
            issues.append("Nu există rădăcină (nicio unitate cu parent_id NULL).")
        elif len(roots) > 1:
            issues.append(f"Există {len(roots)} rădăcini (parent_id NULL). Exemplu: {roots[:10]}")

        return issues

    # ---------------------------------------------------------
    # 3) ORGANIGRAMĂ IERARHICĂ (LISTĂ) – dezactivată
    # ---------------------------------------------------------
    # Lista text clasică a fost înlocuită cu o versiune în tab-ul "Organigramă".

    # ---------------------------------------------------------
    # TABURI: 1) ORGANIGRAMĂ  2) STRUCTURĂ INTERNĂ  3) IMPORT
    # ---------------------------------------------------------

    # Wrapper pentru layout-ul de organigramă (stilizat în CSS)
    st.markdown("<div class='org-wrap'>", unsafe_allow_html=True)

    tab_org, tab_struct, tab_import = st.tabs(
        ["Organigramă", "Structură internă", "Import organigramă"]
    )

    # ===========================
    # TAB 1 – ORGANIGRAMĂ
    # ===========================
    with tab_org:
        if df_units is None or df_units.empty:
            st.info("Nu există unități active în organigramă.")
            return

        # Asigurăm existența coloanei 'type' pentru filtre, indiferent de sursa de date
        if "type" not in df_units.columns:
            df_units["type"] = None

        # ===== PANOU CONTROL (sus) – doar View + Export (fără căutare/filtre în UI) =====
        st.markdown('<div class="org-filters">', unsafe_allow_html=True)

        # Variabile interne pentru eventuale filtre viitoare (UI ascuns momentan)
        q = ""
        sel_types: list[str] = []

        st.markdown("<div class='org-section-title'>Mod de vizualizare</div>", unsafe_allow_html=True)
        st.markdown(
            "<p class='org-section-desc'>Alege cum dorești să fie afișată organigrama.</p>",
            unsafe_allow_html=True,
        )
        st.markdown("<div class='org-view-radio'>", unsafe_allow_html=True)
        view_mode = st.radio(
            "",
            ["Interactivă", "Ierarhică"],
            index=0,
            key="org_view_mode",
            label_visibility="collapsed",
        )
        st.markdown("</div>", unsafe_allow_html=True)

        # Exporturile sunt gestionate mai jos, în secțiunea unificată „Export selecție”
        st.markdown("</div>", unsafe_allow_html=True)

        # ===== Filter data =====
        df_view = df_units.copy()
        if sel_types:
            df_view = df_view[df_view["type"].astype(str).isin(sel_types)]
        if q.strip():
            qq = q.strip().lower()
            df_view = df_view[df_view["name"].astype(str).str.lower().str.contains(qq)]

        # ===== ORGANIGRAMĂ – full width sub panoul de control =====
        st.markdown('<div class="org-canvas">', unsafe_allow_html=True)
        st.markdown("<div class='org-section-title'>Organigramă</div>", unsafe_allow_html=True)
        st.markdown(
            "<p class='org-section-desc'>Vizualizează structura organizațională completă și accesează atât varianta ierarhică, cât și cea interactivă.</p>",
            unsafe_allow_html=True,
        )

        if view_mode == "Ierarhică":
            # Afișăm direct organigrama ierarhică (fără expander suplimentar)
            try:
                df_pos = load_org_positions_with_people(conn)
            except Exception:
                df_pos = pd.DataFrame()

            render_organigrama_tree(df_units, df_stat=df_pos, key_prefix="org_units_tree")

        else:
            # Modul vizual interactiv (Cytoscape) – afișat direct când este selectat
            cyto_box_scale = st.slider("📦 Mărime casete", 0.6, 3.0, 1.2, 0.1, key="cyto_box_scale_main")
            cyto_font_size = st.slider("🔤 Font", 10, 34, 14, 1, key="cyto_font_size_main")
            cyto_spacing = st.slider("↔️ Spacing", 0.6, 2.4, 1.0, 0.05, key="cyto_spacing_main")
            cyto_start_collapsed = st.checkbox("Pornește colapsat (doar nivel 1)", True, key="cyto_start_collapsed_main")
            cyto_enable_drag = st.checkbox("🖐️ Mută casete (Drag)", True, key="cyto_enable_drag_main")

            cyto_custom_colors = st.checkbox("🎨 Culori personalizate", False, key="cyto_custom_colors_main")
            cyto_colors = {}
            if cyto_custom_colors:
                cc1, cc2, cc3 = st.columns(3)
                with cc1:
                    cyto_colors["director_general"] = st.color_picker("DG", "#0B3C5D", key="cyto_col_dg")
                with cc2:
                    cyto_colors["directie"] = st.color_picker("Direcție", "#457B9D", key="cyto_col_directie")
                with cc3:
                    cyto_colors["default"] = st.color_picker("Default", "#D3E3F3", key="cyto_col_default")

            # Construim payload complet (unități + poziții + oameni)
            df_units_full = fetch_org_units(conn)
            df_pos = fetch_positions_with_people(conn)
            elements = build_cyto_elements(df_units_full, df_pos)
            payload = cyto_elements_to_payload(elements)

            render_org_cytoscape_collapsible(
                None,
                elements_payload=payload,
                height=780,
                key="org_cyto_main",
                box_scale=float(st.session_state.get("cyto_box_scale_main", 1.2)),
                font_size=int(st.session_state.get("cyto_font_size_main", 14)),
                spacing=float(st.session_state.get("cyto_spacing_main", 1.0)),
                start_collapsed=bool(st.session_state.get("cyto_start_collapsed_main", True)),
                enable_drag=bool(st.session_state.get("cyto_enable_drag_main", True)),
                enable_save_layout=True,
                use_saved_layout=True,
                color_overrides=cyto_colors if st.session_state.get("cyto_custom_colors_main") else None,
            )

            # Detalii despre selecție (unitate / persoană / poziție)
            emp_q = st.query_params.get("emp_id")
            pos_q = st.query_params.get("pos_id")
            unit_q = st.query_params.get("unit_id")

            try:
                emp_id = int(str(emp_q)) if emp_q is not None else None
            except ValueError:
                emp_id = None

            try:
                unit_id = int(str(unit_q)) if unit_q is not None else None
            except ValueError:
                unit_id = None

            # Detalii selecție – UI fără panou dedicat, detaliile se afișează direct sub organigramă

            # 1) Dacă avem emp_id în query params, afișăm direct profilul angajatului (din bundle)
            if emp_id is not None:
                bundle_org = load_employee_bundle(get_db_path(), emp_id)
                emp_row = bundle_org.get("employee") or {}
                if not emp_row:
                    st.warning("Persoana selectată nu a putut fi găsită.")
                else:
                    _eid = emp_row.get("id") or emp_id
                    full_name = f"{ (emp_row.get('last_name') or '').strip() } { (emp_row.get('first_name') or '').strip() }".strip() or "(necunoscut)"
                    functie = (emp_row.get("functie") or "").strip()
                    email = (emp_row.get("email") or "").strip()
                    telefon = (emp_row.get("telefon") or emp_row.get("mobil") or "").strip()
                    departament = (emp_row.get("departament") or "").strip()
                    unit_name = (emp_row.get("unit_name") or emp_row.get("loc_munca") or "").strip()

                    # Foto angajat (dacă există) – query separat (nu e în bundle)
                    prow = None
                    try:
                        cur = conn.cursor()
                        prow = cur.execute(
                            "SELECT mime, img FROM employee_photos WHERE employee_id = ?",
                            (_eid,),
                        ).fetchone()
                    except Exception:
                        pass

                    if prow and prow[1]:
                        mime, img_blob = prow
                        try:
                            st.image(BytesIO(img_blob), caption=full_name or "Fotografie", use_container_width=True)
                        except Exception:
                            pass

                    st.write(f"**Nume**: {full_name or '(necunoscut)'}")
                    if functie:
                        st.write(f"**Funcție**: {functie}")
                    if email:
                        st.write(f"**Email**: {email}")
                    if telefon:
                        st.write(f"**Telefon**: {telefon}")
                    if unit_name or departament:
                        st.write(
                            f"**Unitate / departament**: "
                            f"{unit_name or ''}{' – ' if unit_name and departament else ''}{departament or ''}"
                        )

                    # Upload / actualizare fotografie
                    uploaded = st.file_uploader(
                        "Încarcă / actualizează fotografia",
                        type=["png", "jpg", "jpeg"],
                        key=f"emp_photo_{_eid}",
                    )
                    if uploaded is not None:
                        try:
                            img_bytes = uploaded.read()
                            mime = uploaded.type or "image/jpeg"
                            cur.execute(
                                """
                                INSERT INTO employee_photos(employee_id, mime, img, updated_at)
                                VALUES (?, ?, ?, datetime('now'))
                                ON CONFLICT(employee_id) DO UPDATE SET
                                    mime=excluded.mime,
                                    img=excluded.img,
                                    updated_at=excluded.updated_at
                                """,
                                (_eid, mime, img_bytes),
                            )
                            conn.commit()
                            st.success("Fotografia a fost salvată.")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Eroare la salvarea fotografiei: {e}")

                    if st.button("Înapoi la unitate", key="back_to_unit_from_emp"):
                        qs = dict(st.query_params)
                        qs.pop("emp_id", None)
                        st.query_params.clear()
                        for k, v in qs.items():
                            st.query_params[k] = v
                        st.rerun()

                    # 2) Dacă avem pos_id în query params (și nu avem emp_id), afișăm detalii despre post
                    elif pos_q is not None:
                        try:
                            pos_id_int = int(str(pos_q))
                        except ValueError:
                            pos_id_int = None

                        if pos_id_int is None:
                            st.warning("Postul selectat nu este valid.")
                        else:
                            cur = conn.cursor()
                            try:
                                prow = cur.execute(
                                    """
                                    SELECT
                                        op.id,
                                        op.title,
                                        op.unit_id,
                                        ou.name AS unit_name
                                    FROM org_positions op
                                    LEFT JOIN org_units ou ON ou.id = op.unit_id
                                    WHERE op.id = ?
                                    """,
                                    (pos_id_int,),
                                ).fetchone()
                            except Exception:
                                prow = None

                            if not prow:
                                st.warning("Postul selectat nu a putut fi găsit.")
                            else:
                                _, p_title, p_unit_id, p_unit_name = prow
                                st.write(f"**Post**: {p_title or '(fără titlu)'}")
                                if p_unit_name:
                                    st.write(f"**Unitate**: {p_unit_name}")

                                # Pentru admin permitem legarea unui angajat de acest post
                                is_admin = str(st.session_state.get("user_role", "user")) == "admin"
                                if is_admin:
                                    with st.expander("Leagă angajat", expanded=False):
                                        # Construim lista de angajați (id + nume complet)
                                        try:
                                            df_emp = pd.read_sql_query(
                                                """
                                                SELECT id,
                                                       TRIM(COALESCE(last_name, '') || ' ' || COALESCE(first_name, '')) AS full_name
                                                FROM employees
                                                ORDER BY full_name
                                                """,
                                                conn,
                                            )
                                        except Exception as e:
                                            df_emp = pd.DataFrame()
                                            st.error(f"Eroare la citirea angajaților: {e}")

                                        if df_emp.empty:
                                            st.info("Nu există angajați în baza de date.")
                                        else:
                                            emp_options = [
                                                f"{int(row['id'])} – {str(row['full_name'] or '').strip() or '(fără nume)'}"
                                                for _, row in df_emp.iterrows()
                                            ]
                                            selected = st.selectbox(
                                                "Alege angajatul",
                                                options=emp_options,
                                                key=f"link_emp_to_pos_{pos_id_int}",
                                            )
                                            selected_id = None
                                            if selected:
                                                try:
                                                    selected_id = int(selected.split("–", 1)[0].strip())
                                                except Exception:
                                                    selected_id = None

                                            if st.button("Leagă", key=f"btn_link_emp_to_pos_{pos_id_int}"):
                                                if selected_id is None:
                                                    st.error("Selectează un angajat.")
                                                else:
                                                    try:
                                                        # Marcăm orice ocupare activă pe acest post ca închisă
                                                        cur.execute(
                                                            """
                                                            UPDATE employee_positions
                                                            SET end_date = CASE
                                                                WHEN end_date IS NULL OR TRIM(end_date) = '' THEN date('now')
                                                                ELSE end_date
                                                            END
                                                            WHERE position_id = ? AND (end_date IS NULL OR TRIM(end_date) = '')
                                                            """,
                                                            (pos_id_int,),
                                                        )
                                                        # Inserăm noua legătură
                                                        cur.execute(
                                                            """
                                                            INSERT INTO employee_positions(employee_id, position_id, start_date, end_date)
                                                            VALUES(?, ?, date('now'), NULL)
                                                            """,
                                                            (selected_id, pos_id_int),
                                                        )
                                                        conn.commit()
                                                        st.success("Angajatul a fost legat de post.")
                                                        st.rerun()
                                                    except Exception as e:
                                                        st.error(f"Eroare la legarea angajatului de post: {e}")

                    # 3) Dacă avem unitate selectată, afișăm info + persoane din stat_functii
                    elif unit_id is not None:
                        cur = conn.cursor()
                        try:
                            urow = cur.execute(
                                "SELECT id, name, type, parent_id FROM org_units WHERE id = ?",
                                (unit_id,),
                            ).fetchone()
                        except Exception:
                            urow = None

                        if not urow:
                            st.warning("Unitatea selectată nu a putut fi găsită.")
                        else:
                            _, u_name, u_type, _ = urow
                            st.write(f"**Unitate**: {u_name}")
                            if u_type:
                                st.write(f"**Tip**: {u_type}")

                            # Panel administrare pentru admin (adăugare posturi în unitate)
                            is_admin = str(st.session_state.get("user_role", "user")) == "admin"
                            if is_admin:
                                with st.expander("Administrare (admin)", expanded=False):
                                    st.markdown("**Adaugă post în această unitate**")
                                    new_pos_title = st.text_input(
                                        "Titlu post",
                                        value="",
                                        key=f"admin_add_pos_title_{unit_id}",
                                        placeholder="Ex: Consilier, Inspector, Referent",
                                    ).strip()
                                    if st.button("Salvează", key=f"admin_add_pos_save_{unit_id}"):
                                        if not new_pos_title:
                                            st.error("Completează titlul postului.")
                                        else:
                                            cur = conn.cursor()
                                            try:
                                                cur.execute(
                                                    """
                                                    INSERT INTO org_positions(unit_id, title, is_active)
                                                    VALUES(?, ?, 1)
                                                    """,
                                                    (int(unit_id), new_pos_title),
                                                )
                                                conn.commit()
                                                st.success("Postul a fost adăugat.")
                                                st.rerun()
                                            except Exception as e:
                                                st.error(f"Eroare la adăugarea postului: {e}")

                            # Încercăm să citim statul de funcții și să găsim persoanele din această unitate
                            try:
                                df_stat = pd.read_sql_query("SELECT * FROM stat_functii", conn)
                            except Exception as e:
                                st.error(f"Eroare la citirea tabelului stat_functii: {e}")
                                df_stat = pd.DataFrame()

                            if df_stat.empty:
                                st.info("Nu există date în tabelul stat_functii.")
                            else:
                                cols_lower = {c.lower(): c for c in df_stat.columns}

                                def _find_col(possible_names: list[str]):
                                    for name in possible_names:
                                        if name.lower() in cols_lower:
                                            return cols_lower[name.lower()]
                                    return None

                                col_loc = _find_col(["Loc de munca", "LOC_DE_MUNCA", "LOC MUNCA", "LOC_MUNCA"])
                                col_functie = _find_col(["Functie", "FUNCȚIE", "FUNCŢIE", "DEN_FUNCTIE", "DENUMIRE_FUNCTIE"])
                                col_nume = _find_col(["Nume", "NUME"])
                                col_prenume = _find_col(["Prenume", "PRENUME"])
                                col_cor = _find_col(["COD_COR", "Cod COR", "COD COR", "COR"])

                                # Încercăm să detectăm o coloană ID de unitate (ID_UNITATE / ID_LOC_DE_MUNCA / ceva cu 'unitate'+'id')
                                col_id_unit = None
                                for c in df_stat.columns:
                                    name_l = c.lower()
                                    if "id_unitate" in name_l or "id unitate" in name_l.replace("_", " "):
                                        col_id_unit = c
                                        break
                                    if "id_loc_de_munca" in name_l:
                                        col_id_unit = c
                                        break
                                    if "id" in name_l and "unitate" in name_l:
                                        col_id_unit = c
                                        break

                                if not col_loc and not col_id_unit:
                                    st.warning(
                                        "Nu am găsit în stat_functii coloanele necesare pentru a filtra persoanele din unitate "
                                        "(nici ID_UNITATE / ID_LOC_DE_MUNCA, nici „Loc de munca”)."
                                    )
                                    df_sel = pd.DataFrame()
                                else:
                                    # 1) Preferăm filtrarea directă pe ID dacă avem coloană numerică de unitate
                                    df_sel = pd.DataFrame()
                                    if col_id_unit:
                                        try:
                                            vals = pd.to_numeric(df_stat[col_id_unit], errors="coerce")
                                            mask_id = vals == int(unit_id)
                                            df_sel = df_stat[mask_id].copy()
                                        except Exception:
                                            df_sel = pd.DataFrame()

                                    # 2) Fallback: filtrare text pe „Loc de munca” care conține numele unității
                                    if df_sel.empty and col_loc and u_name:
                                        mask_txt = df_stat[col_loc].astype(str).str.lower().str.contains(str(u_name).lower())
                                        df_sel = df_stat[mask_txt].copy()

                                    if df_sel.empty:
                                        st.info("Nu există persoane în stat_functii pentru această unitate.")
                                    else:
                                        st.caption(f"{len(df_sel)} persoane găsite în stat_functii pentru această unitate.")
                                        # Încercăm să mapăm fiecare rând din stat_functii la un employee_id din employees
                                        emp_df = pd.read_sql_query("SELECT * FROM employees", conn)
                                        emp_cols = {c.lower(): c for c in emp_df.columns}

                                        col_emp_id = None
                                        for c in df_sel.columns:
                                            if "angajat" in c.lower():
                                                col_emp_id = c
                                                break

                                        col_stat_cnp = None
                                        for cand in ["CNP", "cnp"]:
                                            if cand in df_sel.columns:
                                                col_stat_cnp = cand
                                                break

                                        col_emp_cnp = emp_cols.get("cnp")

                                        col_stat_nume = col_nume
                                        col_stat_prenume = col_prenume

                                        # Preprocesăm employees pentru căutări
                                        if col_emp_cnp:
                                            emp_df["_cnp_clean"] = emp_df[col_emp_cnp].astype(str).str.replace(r"\\D+", "", regex=True)
                                        else:
                                            emp_df["_cnp_clean"] = ""

                                        emp_df["_full_name_clean"] = (
                                            emp_df.get("last_name", "")
                                            .astype(str)
                                            .str.strip()
                                            .str.lower()
                                            + " "
                                            + emp_df.get("first_name", "").astype(str).str.strip().str.lower()
                                        ).str.strip()

                                        emp_index_by_id = {int(r["id"]): r for _, r in emp_df.iterrows() if "id" in r}

                                        rows_out = []
                                        for _, r in df_sel.iterrows():
                                            linked_emp_id = None

                                            # 1) mapping direct prin ID_ANGAJAT (sau similar)
                                            if col_emp_id and pd.notna(r.get(col_emp_id)):
                                                try:
                                                    cand_id = int(str(r.get(col_emp_id)))
                                                    if cand_id in emp_index_by_id:
                                                        linked_emp_id = cand_id
                                                except ValueError:
                                                    pass

                                            # 2) fallback: CNP
                                            if linked_emp_id is None and col_stat_cnp:
                                                cnp_val = str(r.get(col_stat_cnp) or "")
                                                cnp_clean = re.sub(r"\\D+", "", cnp_val)
                                                if cnp_clean:
                                                    match = emp_df[emp_df["_cnp_clean"] == cnp_clean]
                                                    if not match.empty:
                                                        linked_emp_id = int(match.iloc[0]["id"])

                                            # 3) fallback: nume + prenume
                                            if linked_emp_id is None and col_stat_nume and col_stat_prenume:
                                                nume = str(r.get(col_stat_nume) or "").strip().lower()
                                                pren = str(r.get(col_stat_prenume) or "").strip().lower()
                                                full = f"{nume} {pren}".strip()
                                                if full:
                                                    match = emp_df[emp_df["_full_name_clean"] == full]
                                                    if not match.empty:
                                                        linked_emp_id = int(match.iloc[0]["id"])

                                            rows_out.append(
                                                {
                                                    "Nume": str(r.get(col_nume) or ""),
                                                    "Prenume": str(r.get(col_prenume) or ""),
                                                    "Funcție": str(r.get(col_functie) or ""),
                                                    "Cod COR": str(r.get(col_cor) or ""),
                                                    "_employee_id": linked_emp_id,
                                                }
                                            )

                                        df_view_sf = pd.DataFrame(rows_out)
                                        if df_view_sf.empty:
                                            st.info("Nu există persoane în stat_functii pentru această unitate.")
                                        else:
                                            # Afișăm lista de persoane (Nume Prenume – Funcție – Cod COR) + buton Profil (dacă avem employee_id)
                                            for idx, row in df_view_sf.iterrows():
                                                emp_row_id = row.get("_employee_id")
                                                cols = st.columns([7, 2])
                                                with cols[0]:
                                                    line = f"{row['Nume']} {row['Prenume']}".strip()
                                                    if row["Funcție"]:
                                                        line += f" – {row['Funcție']}"
                                                    if row["Cod COR"]:
                                                        line += f" – {row['Cod COR']}"
                                                    st.write(line or "(fără nume)")
                                                with cols[1]:
                                                    if emp_row_id is not None:
                                                        if st.button("Profil", key=f"sf_prof_{unit_id}_{idx}"):
                                                            # Păstrăm unit_id în query params și setăm emp_id
                                                            st.query_params["emp_id"] = str(int(emp_row_id))
                                                            st.rerun()

                    # 3) Dacă nu avem nici unitate, nici persoană explicită, dar avem doar poziție, păstrăm mesajul generic
                    elif emp_id is None and pos_q is None:
                        st.info("Selectează o persoană/poziție în organigramă.")

                    st.markdown("</div>", unsafe_allow_html=True)

            st.markdown("</div>", unsafe_allow_html=True)

        # ===== ELEMENT SELECTAT + ACȚIUNI (compact, similar cu Angajați) =====
        emp_q_bottom = st.query_params.get("emp_id")
        pos_q_bottom = st.query_params.get("pos_id")
        unit_q_bottom = st.query_params.get("unit_id")
        try:
            emp_id_bottom = int(str(emp_q_bottom)) if emp_q_bottom is not None else None
        except Exception:
            emp_id_bottom = None
        try:
            pos_id_bottom = int(str(pos_q_bottom)) if pos_q_bottom is not None else None
        except Exception:
            pos_id_bottom = None
        try:
            unit_id_bottom = int(str(unit_q_bottom)) if unit_q_bottom is not None else None
        except Exception:
            unit_id_bottom = None

        st.markdown("<div class='org-bottom'>", unsafe_allow_html=True)
        st.markdown("<div class='org-section-title'>Element selectat</div>", unsafe_allow_html=True)
        if emp_id_bottom is not None:
            bundle_bottom = load_employee_bundle(get_db_path(), int(emp_id_bottom))
            emp_bottom = bundle_bottom.get("employee") or {}
            emp_name_bottom = f"{(emp_bottom.get('last_name') or '').strip()} {(emp_bottom.get('first_name') or '').strip()}".strip() or f"ID {emp_id_bottom}"
            st.markdown(f"<div class='org-selected-row'><strong>Persoană:</strong> {emp_name_bottom}</div>", unsafe_allow_html=True)
            st.markdown(f"<div class='org-selected-row'><strong>CNP:</strong> {emp_bottom.get('cnp') or '—'}</div>", unsafe_allow_html=True)
            st.markdown(f"<div class='org-selected-row'><strong>Marcă:</strong> {emp_bottom.get('marca') or '—'}</div>", unsafe_allow_html=True)
        elif pos_id_bottom is not None:
            row_pos_bottom = conn.execute(
                "SELECT op.title, ou.name FROM org_positions op LEFT JOIN org_units ou ON ou.id = op.unit_id WHERE op.id = ?",
                (int(pos_id_bottom),),
            ).fetchone()
            st.markdown("<div class='org-selected-row'><strong>Tip:</strong> Post</div>", unsafe_allow_html=True)
            st.markdown(f"<div class='org-selected-row'><strong>Denumire:</strong> {(row_pos_bottom[0] if row_pos_bottom else '—') or '—'}</div>", unsafe_allow_html=True)
            st.markdown(f"<div class='org-selected-row'><strong>Unitate:</strong> {(row_pos_bottom[1] if row_pos_bottom else '—') or '—'}</div>", unsafe_allow_html=True)
        elif unit_id_bottom is not None:
            row_unit_bottom = conn.execute(
                "SELECT name, type, parent_id FROM org_units WHERE id = ?",
                (int(unit_id_bottom),),
            ).fetchone()
            st.markdown("<div class='org-selected-row'><strong>Tip:</strong> Unitate</div>", unsafe_allow_html=True)
            st.markdown(f"<div class='org-selected-row'><strong>Denumire:</strong> {(row_unit_bottom[0] if row_unit_bottom else '—') or '—'}</div>", unsafe_allow_html=True)
            st.markdown(f"<div class='org-selected-row'><strong>Părinte:</strong> {(row_unit_bottom[2] if row_unit_bottom else '—') if row_unit_bottom else '—'}</div>", unsafe_allow_html=True)
        else:
            st.markdown("<div class='org-selected-row'>Nicio selecție activă.</div>", unsafe_allow_html=True)
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown("<div class='org-bottom'>", unsafe_allow_html=True)
        st.markdown("<div class='org-section-title'>Acțiuni</div>", unsafe_allow_html=True)
        if st.button("Reîmprospătează", key="org_btn_sel_refresh"):
            st.rerun()
        if st.button("Curăță selecția", key="org_btn_sel_clear", disabled=not any(v is not None for v in (emp_id_bottom, pos_id_bottom, unit_id_bottom))):
            qs_bottom = dict(st.query_params)
            qs_bottom.pop("emp_id", None)
            qs_bottom.pop("pos_id", None)
            qs_bottom.pop("unit_id", None)
            st.query_params.clear()
            for k, v in qs_bottom.items():
                st.query_params[k] = v
            st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)

        # ===== EXPORT – ORGANIGRAMĂ IERARHICĂ COMPLETĂ (TXT / DOCX) =====
        st.markdown("### 📤 Export organigramă ierarhică")
        st.caption(
            "Descarcă structura completă a organigramei în format text sau document, pentru utilizare internă."
        )

        nodes_full, children_full, roots_full = build_organigrama_tree_structure(df_units)
        lines_full: list[str] = []

        def _build_text_node_full(uid: int, level: int = 0) -> None:
            prefix = "    " * level + "- "
            label = nodes_full[uid]["name"]
            if nodes_full[uid]["tip_raw"]:
                label = f"{nodes_full[uid]['name']} ({nodes_full[uid]['tip_raw']})"
            lines_full.append(prefix + label)
            for cid in children_full.get(uid, []):
                _build_text_node_full(cid, level + 1)

        for rid in sorted(roots_full, key=lambda u: nodes_full[u]["name"].lower()):
            _build_text_node_full(rid, 0)

        text_export_full = "\n".join(lines_full)
        txt_bytes_full = text_export_full.encode("utf-8")
        st.download_button(
            "⬇️ Descarcă organigrama completă (TXT)",
            data=txt_bytes_full,
            file_name="organigrama_completa.txt",
            mime="text/plain",
            key="org_export_full_txt",
        )

        try:
            import io

            doc_full = Document()
            _docx_apply_antet_si_semnaturi(doc_full)
            doc_full.add_heading("Organigramă ierarhică – structură completă", level=1)
            for line in lines_full:
                doc_full.add_paragraph(line)

            docx_buffer_full = io.BytesIO()
            doc_full.save(docx_buffer_full)
            docx_buffer_full.seek(0)

            st.download_button(
                "⬇️ Descarcă organigrama completă (DOCX)",
                data=docx_buffer_full,
                file_name="organigrama_completa.docx",
                mime=(mimetypes.guess_type("organigrama_completa.docx")[0] or "application/octet-stream"),
                key="org_export_full_docx",
            )
        except Exception as e:
            st.caption(
                "Pentru exportul complet în DOCX ai nevoie de pachetul `python-docx` instalat. "
                f"(Detalii tehnice: {e})"
            )

        st.markdown("---")

        # 2) ORGANIGRAMĂ IERARHICĂ – SELECȚIE CURENTĂ (secțiune ascunsă din UI)
        # Determinăm selecția curentă din query params (angajat / post / unitate)
        sel_emp_q = st.query_params.get("emp_id")
        sel_pos_q = st.query_params.get("pos_id")
        sel_unit_q = st.query_params.get("unit_id")

        export_unit_id: int | None = None
        scope_parts: list[str] = []

        # 2.1) Angajat selectat → deducem unitatea curentă din employee_positions/org_positions
        if sel_emp_q is not None:
            try:
                sel_emp_id = int(str(sel_emp_q))
            except ValueError:
                sel_emp_id = None

            if sel_emp_id is not None:
                try:
                    cur = conn.cursor()
                    row_emp = cur.execute(
                        """
                        SELECT op.unit_id,
                               TRIM(COALESCE(e.last_name, '') || ' ' || COALESCE(e.first_name, '')) AS full_name
                        FROM employee_positions ep
                        JOIN org_positions op ON op.id = ep.position_id
                        JOIN employees e ON e.id = ep.employee_id
                        WHERE ep.employee_id = ?
                        ORDER BY COALESCE(ep.end_date, '9999-12-31') DESC,
                                 COALESCE(ep.start_date, '0001-01-01') DESC
                        LIMIT 1
                        """,
                        (sel_emp_id,),
                    ).fetchone()
                except Exception:
                    row_emp = None

                if row_emp:
                    export_unit_id = int(row_emp[0])
                    full_name = (row_emp[1] or "").strip()
                    scope_parts.append(f"angajat: {full_name or sel_emp_id}")

        # 2.2) Dacă nu am putut deduce din angajat, încercăm după post selectat
        if export_unit_id is None and sel_pos_q is not None:
            try:
                sel_pos_id = int(str(sel_pos_q))
            except ValueError:
                sel_pos_id = None

            if sel_pos_id is not None:
                try:
                    cur = conn.cursor()
                    row_pos = cur.execute(
                        """
                        SELECT
                            op.unit_id,
                            op.title,
                            ou.name AS unit_name
                        FROM org_positions op
                        LEFT JOIN org_units ou ON ou.id = op.unit_id
                        WHERE op.id = ?
                        """,
                        (sel_pos_id,),
                    ).fetchone()
                except Exception:
                    row_pos = None

                if row_pos:
                    export_unit_id = int(row_pos[0])
                    p_title = (row_pos[1] or "").strip()
                    p_unit_name = (row_pos[2] or "").strip()
                    label = p_title or f"post #{sel_pos_id}"
                    if p_unit_name:
                        scope_parts.append(f"post: {label} – {p_unit_name}")
                    else:
                        scope_parts.append(f"post: {label}")

        # 2.3) Dacă avem unitate selectată explicit
        if export_unit_id is None and sel_unit_q is not None:
            try:
                export_unit_id = int(str(sel_unit_q))
            except ValueError:
                export_unit_id = None

        if export_unit_id is None:
            # Secțiunea de export pe selecție este dezactivată vizual în acest moment
            pass
        else:
            nodes_sel, children_sel, _roots_sel = build_organigrama_tree_structure(df_units)

            if export_unit_id not in nodes_sel:
                st.caption("Selecția curentă nu are o structură ierarhică disponibilă pentru export.")
            else:
                scope_text = ", ".join(scope_parts) if scope_parts else None
                if scope_text:
                    st.caption(
                        f"Se va exporta structura ierarhică aferentă selecției curente ({scope_text})."
                    )
                else:
                    st.caption("Se va exporta structura ierarhică aferentă selecției curente.")

                lines_sel: list[str] = []

                def _build_text_node_sel(uid: int, level: int = 0) -> None:
                    prefix = "    " * level + "- "
                    label = nodes_sel[uid]["name"]
                    if nodes_sel[uid]["tip_raw"]:
                        label = f"{nodes_sel[uid]['name']} ({nodes_sel[uid]['tip_raw']})"
                    lines_sel.append(prefix + label)
                    for cid in children_sel.get(uid, []):
                        _build_text_node_sel(cid, level + 1)

                _build_text_node_sel(export_unit_id, 0)

                text_export_sel = "\n".join(lines_sel)
                txt_bytes_sel = text_export_sel.encode("utf-8")
                st.download_button(
                    "⬇️ Descarcă structura selecției (TXT)",
                    data=txt_bytes_sel,
                    file_name="organigrama_selectie.txt",
                    mime="text/plain",
                    key="org_export_sel_txt",
                )

                try:
                    import io

                    doc_sel = Document()
                    _docx_apply_antet_si_semnaturi(doc_sel)
                    doc_sel.add_heading("Organigramă ierarhică – selecție curentă", level=1)
                    if scope_text:
                        doc_sel.add_paragraph(scope_text)
                    for line in lines_sel:
                        doc_sel.add_paragraph(line)

                    docx_buffer_sel = io.BytesIO()
                    doc_sel.save(docx_buffer_sel)
                    docx_buffer_sel.seek(0)

                    st.download_button(
                        "⬇️ Descarcă structura selecției (DOCX)",
                        data=docx_buffer_sel,
                        file_name="organigrama_selectie.docx",
                        mime=(mimetypes.guess_type("organigrama_selectie.docx")[0] or "application/octet-stream"),
                        key="org_export_sel_docx",
                    )
                except Exception as e:
                    st.caption(
                        "Pentru export DOCX ai nevoie de pachetul `python-docx` instalat. "
                        f"(Detalii tehnice: {e})"
                    )

        # (exportul PNG/JPEG/PDF pentru organigrama interactivă rămâne disponibil direct în viewer)

    # ===========================
    # TAB 2 – STRUCTURA INTERNĂ
    # ===========================
    with tab_struct:
        st.markdown("### STRUCTURA INTERNĂ: posturi și personal")
        # Încercăm să citim statul de funcții
        try:
            df_stat = pd.read_sql_query("SELECT * FROM stat_functii", conn)
        except Exception as e:
            st.error(f"Eroare la citirea tabelului stat_functii: {e}")
            df_stat = pd.DataFrame()

        if df_stat.empty:
            st.info("Nu există date în tabelul stat_functii.")
        else:
            # Detectăm coloanele relevante, indiferent de scriere
            cols_lower = {c.lower(): c for c in df_stat.columns}

            def find_col(possible_names):
                for name in possible_names:
                    if name.lower() in cols_lower:
                        return cols_lower[name.lower()]
                return None

            col_loc = find_col(["Loc de munca", "LOC_DE_MUNCA", "LOC MUNCA", "LOC_MUNCA"])
            col_functie = find_col(["Functie", "FUNCȚIE", "FUNCŢIE", "DEN_FUNCTIE", "DENUMIRE_FUNCTIE"])
            col_nume = find_col(["Nume", "NUME"])
            col_prenume = find_col(["Prenume", "PRENUME"])
            col_cor = find_col(["COD_COR", "Cod COR", "COD COR", "COR"])

            # Construim un view doar cu coloanele găsite
            cols_view = []
            headers = []

            if col_loc:
                cols_view.append(col_loc)
                headers.append("Loc de muncă (Direcție/Serviciu/Compartiment)")
            if col_nume:
                cols_view.append(col_nume)
                headers.append("Nume")
            if col_prenume:
                cols_view.append(col_prenume)
                headers.append("Prenume")
            if col_functie:
                cols_view.append(col_functie)
                headers.append("Funcție")
            if col_cor:
                cols_view.append(col_cor)
                headers.append("Cod COR")

            if not cols_view:
                st.warning(
                    "Nu am găsit în stat_functii coloanele necesare pentru afișarea structurii interne "
                    "(Loc de munca, Nume, Prenume, Functie, Cod COR)."
                )
            else:
                # View de lucru simplu: căutare, valori sintetice, tabel (fără filtre suplimentare / preview)
                df_view = df_stat[cols_view].copy()
                df_view.columns = headers

                # 1) Căutare / filtrare text
                search_si = st.text_input(
                    "Caută:",
                    value="",
                    key="struct_interna_search",
                    placeholder="loc de muncă, nume, funcție sau cod COR",
                )

                if search_si.strip():
                    txt = search_si.strip().lower()
                    df_view = df_view[
                        df_view.apply(
                            lambda r: any(
                                txt in str(val).lower() for val in r.values
                            ),
                            axis=1,
                        )
                    ]

                # 2) Valori sintetice (Total înregistrări / Structuri unice / Funcții unice)
                col_loc_header = "Loc de muncă (Direcție/Serviciu/Compartiment)"
                n_total = len(df_view)
                n_loc = df_view[col_loc_header].nunique() if col_loc_header in df_view.columns else 0
                if "Funcție" in df_view.columns:
                    n_func = df_view["Funcție"].nunique()
                elif "Cod COR" in df_view.columns:
                    n_func = df_view["Cod COR"].nunique()
                else:
                    n_func = 0

                st.markdown(
                    f"**Total înregistrări:** {n_total}  \n"
                    f"**Structuri unice:** {n_loc}  \n"
                    f"**Funcții unice:** {n_func}"
                )

                # 3) Tabel principal
                st.dataframe(df_view, use_container_width=True)

                # 4) Export CSV – sub tabel, aliniat la stânga, doar dacă avem rezultate
                if not df_view.empty:
                    col_btn, _ = st.columns([0.3, 0.7])
                    with col_btn:
                        csv_bytes = df_view.to_csv(index=False).encode("utf-8")
                        st.download_button(
                            "Export CSV",
                            data=csv_bytes,
                            file_name="structura_interna.csv",
                            mime="text/csv",
                            key="struct_interna_download_csv",
                            use_container_width=True,
                        )

    # ===========================
    # TAB 3 – IMPORT ORGANIGRAMĂ
    # ===========================
    with tab_import:
        # 1) TITLU SECȚIUNE
        st.subheader("Import / sincronizare organigramă")

        # 2) STATUS ORGANIGRAMĂ (KPI-uri)
        try:
            units_active = conn.execute(
                "SELECT COUNT(*) FROM org_units WHERE is_active = 1"
            ).fetchone()[0]
        except Exception:
            units_active = 0
        try:
            positions_active = conn.execute(
                "SELECT COUNT(*) FROM org_positions WHERE is_active = 1"
            ).fetchone()[0]
        except Exception:
            positions_active = 0
        try:
            roots = conn.execute(
                "SELECT COUNT(*) FROM org_units WHERE is_active = 1 AND parent_id IS NULL"
            ).fetchone()[0]
            legaturi_lipsa = max(0, roots - 1)
        except Exception:
            legaturi_lipsa = 0
        last_action = st.session_state.get("org_last_action") or "—"
        last_action_escaped = str(last_action).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        # Header: buton Refresh în colțul stânga (fără titlu suplimentar)
        _status_header_left, _status_header_right = st.columns([0.18, 1])
        with _status_header_left:
            if st.button("🔄 Refresh status", key="org_status_refresh_btn"):
                st.rerun()
        with _status_header_right:
            pass

        st.markdown(
            "<div class='org-card'>"
            "<div class=\"org-status-head\">"
            "  <div>"
            "    <div class=\"title\">Status organigramă</div>"
            "  </div>"
            "</div>"
            "<div class='org-kpis'>"
            f"<div class='org-kpi'><span class='label'>Unități active</span><span class='value'>{units_active}</span></div>"
            f"<div class='org-kpi'><span class='label'>Posturi active</span><span class='value'>{positions_active}</span></div>"
            f"<div class='org-kpi'><span class='label'>Legături lipsă</span><span class='value'>{legaturi_lipsa}</span></div>"
            "</div>"
            "</div>",
            unsafe_allow_html=True,
        )

        # 3) SINCRONIZARE CU STATUL DE FUNCȚII – bloc dedicat
        st.markdown("<div>", unsafe_allow_html=True)
        st.markdown("### Sincronizare cu statul de funcții")

        st.caption(
            "Poți genera sau actualiza automat structura organigramei pe baza tabelului "
            "`stat_functii`, folosind coloanele „Loc de munca” și „Functie”."
        )

        if st.button("Preia / sincronizează organigrama din statul de funcții", key="org_sync_stat"):
            try:
                new_units, new_pos = generate_organigrama_from_stat(conn)
                st.success(
                    f"Sincronizare reușită: {new_units} locuri de muncă noi și {new_pos} posturi noi "
                    "au fost preluate din statul de funcții."
                )
                st.rerun()
            except Exception as e:
                st.error(f"Eroare la sincronizarea cu statul de funcții: {e}")

        # -----------------------------------------------------
        # Demo seed (doar admin)
        # -----------------------------------------------------
        st.markdown("</div>", unsafe_allow_html=True)

        # 5) IMPORT EXCEL
        st.markdown("<div>", unsafe_allow_html=True)
        st.markdown("### Importă organigrama din Excel")

        fisier_org = st.file_uploader(
            "Încarcă fișierul Excel cu structura organigramei",
            type=["xlsx", "xls"],
            key="org_import_file",
        )

        if fisier_org is not None:
            try:
                df_import = pd.read_excel(fisier_org)
                st.success("Fișierul a fost încărcat cu succes.")
                st.dataframe(df_import.head(50), use_container_width=True)
            except Exception as e:
                st.error(f"Eroare la citirea fișierului Excel: {e}")

        st.markdown("</div>", unsafe_allow_html=True)

        # 6) LEGĂTURI STRUCTURĂ (copil → părinte) + VALIDARE
        st.markdown("<div>", unsafe_allow_html=True)
        st.markdown("### Legături structură (copil → părinte) și validare")

        df_units = pd.read_sql_query(
            "SELECT id, name, parent_id, type, is_active "
            "FROM org_units WHERE is_active = 1 ORDER BY id",
            conn,
        )

        # 🔗 După aceste date, formularul de legături copil → părinte rămâne la fel
        if df_units.empty:
            st.caption("Nu există elemente pentru stabilirea legăturilor.")
        else:
            cur = conn.cursor()
            df_units_sorted = df_units.sort_values("id")

            st.markdown("#### Legătură simplă copil → părinte")
            # Marker pentru CSS: text alb la label + valoare pentru cele 2 selectbox-uri
            st.markdown(
                '<div id="org-link-copil-parent-anchor" style="display:none" aria-hidden="true"></div>',
                unsafe_allow_html=True,
            )
            st.markdown(
                """
                <style>
                /* Text ALB: label „Alege unitatea COPIL/PĂRINTE” și valoarea selectată (ex. 2 – AGENTIA...) */
                #org-link-copil-parent-anchor ~ div [data-testid="stSelectbox"] p,
                #org-link-copil-parent-anchor ~ div [data-testid="stSelectbox"] label,
                #org-link-copil-parent-anchor ~ div [data-testid="stSelectbox"] div[data-baseweb="select"],
                #org-link-copil-parent-anchor ~ div [data-testid="stSelectbox"] div[data-baseweb="select"] span,
                #org-link-copil-parent-anchor ~ div [data-testid="stSelectbox"] div[data-baseweb="select"] div {
                  color: #ffffff !important;
                }
                div[role="listbox"] div[role="option"],
                div[role="listbox"] div[role="option"] span {
                  color: #ffffff !important;
                }
                </style>
                """,
                unsafe_allow_html=True,
            )
            child_labels = [
                f"{row['id']} – {row['name']} ({row['type']})"
                for _, row in df_units_sorted.iterrows()
            ]
            child_ids = df_units_sorted["id"].tolist()

            idx_child = st.selectbox(
                "Alege unitatea COPIL",
                options=list(range(len(child_labels))),
                format_func=lambda i: child_labels[i],
                key="org_link_child_single",
            )

            parent_labels = child_labels
            parent_ids = child_ids

            idx_parent = st.selectbox(
                "Alege unitatea PĂRINTE",
                options=list(range(len(parent_labels))),
                format_func=lambda i: parent_labels[i],
                key="org_link_parent_single",
            )

            child_id = child_ids[idx_child]
            parent_id = parent_ids[idx_parent]

            if child_id == parent_id:
                st.error("O unitate nu poate fi propriul ei părinte.")
            else:
                if st.button("💾 Salvează legătura copil → părinte", key="org_btn_save_link_single"):
                    try:
                        cur.execute(
                            "UPDATE org_units SET parent_id = ? WHERE id = ?",
                            (parent_id, child_id),
                        )
                        conn.commit()
                        st.success("Legătura a fost salvată.")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Eroare la salvare: {e}")

            # -----------------------------------------------------
            # 3) LEGĂTURI MULTIPLE – UN PĂRINTE → MAI MULȚI COPII
            # -----------------------------------------------------
            st.markdown("### 🔗 Legături multiple pentru același părinte")

            st.caption(
                "1. Alege PĂRINTELE.\n"
                "2. Filtrează și selectează unitățile COPIL fără părinte.\n"
                "3. Salvează toate legăturile dintr-o dată."
            )

            parent_search = st.text_input(
                "Caută în lista de posibili PĂRINȚI (după denumire sau tip)",
                value="",
                key="org_parent_search_multi",
            )

            parent_df = df_units_sorted.copy()
            if parent_search.strip():
                txt = parent_search.strip().lower()
                parent_df = parent_df[
                    parent_df["name"].astype(str).str.lower().str.contains(txt)
                    | parent_df["type"].astype(str).str.lower().str.contains(txt)
                ]

            if parent_df.empty:
                st.info("Nu există părinți care să corespundă filtrului.")
            else:
                parent_options = [
                    f"{row['id']} – {row['name']} ({row['type']})"
                    for _, row in parent_df.iterrows()
                ]
                parent_ids_multi = parent_df["id"].tolist()

                idx_parent_multi = st.selectbox(
                    "Alege nodul PĂRINTE (pentru legături multiple)",
                    options=list(range(len(parent_options))),
                    format_func=lambda i: parent_options[i],
                    key="org_multi_parent_idx",
                )
                parent_id_multi = int(parent_ids_multi[idx_parent_multi])

                df_children_all = df_units_sorted[
                    df_units_sorted["parent_id"].isna()
                    & (df_units_sorted["id"] != parent_id_multi)
                ]

                child_search = st.text_input(
                    "Caută în nodurile COPIL (după denumire sau tip)",
                    value="",
                    key="org_child_search_multi",
                )
                if child_search.strip():
                    txtc = child_search.strip().lower()
                    df_children_all = df_children_all[
                        df_children_all["name"].astype(str).str.lower().str.contains(txtc)
                        | df_children_all["type"].astype(str).str.lower().str.contains(txtc)
                    ]

                if df_children_all.empty:
                    st.info(
                        "Nu există noduri fără părinte care să corespundă filtrului și să poată fi legate sub acest părinte."
                    )

        st.markdown("</div>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# CENTRALIZATOR CONCEDII (din PONTAJ)
# -------------------------------------------------------------

def page_stat_de_functii(conn: sqlite3.Connection):
    """Pagină pentru STATUL DE FUNCȚII: salvare/editare + import din Excel cu mapare de câmpuri."""
    st.markdown('<h2 class="page-title">Stat de funcții</h2>', unsafe_allow_html=True)

    # Marker + CSS local pentru structurare pe carduri
    st.markdown('<span id="stat-scope"></span>', unsafe_allow_html=True)
    st.markdown(
        """
        <style>
        section.main:has(#stat-scope) .stat-card{
          background: #162a3d;
          border: 1px solid #2c445c;
          border-radius: 10px;
          padding: 24px 28px;
          margin-bottom: 28px;
        }
        section.main:has(#stat-scope) .stat-card h4{
          margin-top: 0;
          margin-bottom: 16px;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    tab_saved, tab_import = st.tabs(
        ["Stat de funcții salvat", "Import din fișier Excel"]
    )

    # ---------------------------------------------------------
    # TAB 1: STAT DE FUNCȚII SALVAT (EDITABIL)
    # ---------------------------------------------------------
    with tab_saved:
        st.markdown("### Stat de funcții salvat (editare directă)")

        # Card 1 – Sincronizare bidirecțională Stat ↔ Angajați
        st.markdown('<div class="stat-card">', unsafe_allow_html=True)
        with st.expander("🔁 Sincronizare (Stat ↔ Angajați)", expanded=False):
            colA, colB = st.columns(2)
            with colA:
                overwrite = st.checkbox(
                    "Suprascrie datele existente (overwrite)",
                    value=False,
                    key="sync_bidir_overwrite_stat",
                )
            with colB:
                confirm = st.checkbox(
                    "Confirm că vreau să rulez sincronizarea.",
                    value=False,
                    key="sync_bidir_confirm_stat",
                )

            c1, c2 = st.columns(2)
            with c1:
                if st.button("🔁 Angajați → Stat de funcții", key="btn_sync_emp_to_stat"):
                    if not confirm:
                        st.warning("Bifează confirmarea înainte de sincronizare.")
                    else:
                        try:
                            n_imp = sync_stat_functii_from_employees(conn, overwrite=bool(overwrite))
                            st.success(f"Stat de funcții actualizat din Angajați: {n_imp} rânduri.")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Eroare la preluarea din Angajați: {e}")

            with c2:
                if st.button("🔁 Stat de funcții → Angajați", key="btn_sync_stat_to_emp_bidir"):
                    if not confirm:
                        st.warning("Bifează confirmarea înainte de sincronizare.")
                    else:
                        try:
                            n_upd = sync_employees_from_stat_functii_bidir(conn, overwrite=bool(overwrite))
                            st.success(f"Angajați actualizați din Stat de funcții: {n_upd}.")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Eroare la preluarea din Stat de funcții: {e}")
        st.markdown("</div>", unsafe_allow_html=True)


        cur = conn.cursor()
        cur.execute(
            "SELECT name FROM sqlite_master WHERE type='table' AND name='stat_functii'"
        )
        exists = cur.fetchone() is not None

        if not exists:
            st.info(
                "Nu există încă un stat de funcții salvat în baza de date. "
                "Importă mai întâi un fișier Excel în tab-ul „Import din fișier Excel”."
            )
        else:
            df = pd.read_sql_query("SELECT * FROM stat_functii", conn)

            # Card 2 – Configurare câmpuri (coloane)
            st.markdown('<div class="stat-card">', unsafe_allow_html=True)
            st.markdown("#### Configurare câmpuri (coloane)")
            col1, col2 = st.columns(2)

            with col1:
                new_col_name = st.text_input(
                    "Nume câmp nou (coloană)",
                    key="stat_new_col",
                    help="Ex: „Observații”, „Cod intern”, etc.",
                )
                if st.button("➕ Adaugă câmp", key="btn_add_stat_col"):
                    if new_col_name:
                        new_col_name = new_col_name.strip()
                        if new_col_name and new_col_name not in df.columns:
                            df[new_col_name] = ""
                        else:
                            st.warning("Câmpul există deja sau este gol.")

            with col2:
                cols_to_drop = st.multiselect(
                    "Alege câmpurile de șters",
                    options=list(df.columns),
                    key="stat_cols_drop",
                )
                if st.button("🗑️ Șterge câmpurile selectate", key="btn_drop_stat_cols"):
                    if cols_to_drop:
                        df = df.drop(columns=list(cols_to_drop))
            st.markdown("</div>", unsafe_allow_html=True)

            # Card 3 – Editare stat de funcții (tabel mare + salvare)
            st.markdown('<div class="stat-card">', unsafe_allow_html=True)
            st.markdown("#### Editare date")
            edited_df = st.data_editor(
                df,
                num_rows="dynamic",
                use_container_width=True,
                key="stat_editor",
            )

            if st.button("💾 Salvează modificările în statul de funcții", key="btn_save_stat"):
                edited_df.to_sql("stat_functii", conn, if_exists="replace", index=False)

                # 1) Stat -> Organigramă (org_units + org_positions)
                nu, np = refresh_organigrama_from_stat(conn)

                # 2) Stat -> employees (pentru ca Dosar Profesional / Pontaj să reflecte schimbările)
                n_upd = sync_employees_from_stat_functii(conn)

                msg = "Statul de funcții a fost salvat în baza de date."
                if nu or np:
                    msg += f" | Organigramă actualizată: {nu} unități noi, {np} posturi noi."
                if n_upd:
                    msg += f" | Angajați sincronizați: {n_upd}."
                st.success(msg)

                # Reîncărcare UI ca să se vadă imediat în Organigramă / Centralizator / Dosar / Pontaj
                st.rerun()
            st.markdown("</div>", unsafe_allow_html=True)

            # Card 4 – Export & sincronizare globală
            st.markdown('<div class="stat-card">', unsafe_allow_html=True)
            st.markdown("#### Export și sincronizare globală")

            # Export rapid în Excel
            excel_bytes = _df_to_xlsx_bytes(edited_df, sheet_name="Stat_de_functii")
            if excel_bytes is None:
                # Fallback CSV (nu necesită openpyxl/xlsxwriter)
                excel_bytes = edited_df.to_csv(index=False).encode('utf-8')
                mime = 'text/csv'
            else:
                mime = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            st.download_button(
                "⬇️ Exportă statul de funcții salvat (Excel)",
                data=excel_bytes,
                file_name="stat_de_functii_salvat.xlsx",
                mime=(mimetypes.guess_type("stat_de_functii_salvat.xlsx")[0] or "application/octet-stream"),
                key="btn_export_stat_functii_salvat",
            )

            # Export în Word/PDF
            try:
                from docx import Document as _DocxDoc
                doc_stat = _DocxDoc()
                doc_stat.add_heading("Stat de funcții", level=1)
                try:
                    cfg = load_config()
                except Exception:
                    cfg = {}
                unit_name = cfg.get("denumire_unitate", "")
                if unit_name:
                    doc_stat.add_paragraph(f"Unitate: {unit_name}")
                table = doc_stat.add_table(rows=1, cols=len(edited_df.columns))
                hdr_cells = table.rows[0].cells
                for j, col_name in enumerate(edited_df.columns):
                    hdr_cells[j].text = str(col_name)
                for _, row in edited_df.iterrows():
                    row_cells = table.add_row().cells
                    for j, col_name in enumerate(edited_df.columns):
                        val = row.get(col_name, "")
                        row_cells[j].text = "" if pd.isna(val) else str(val)
                buf_docx = BytesIO()
                doc_stat.save(buf_docx)
                buf_docx.seek(0)
                st.download_button(
                    "⬇️ Exportă statul de funcții salvat (Word)",
                    data=buf_docx.getvalue(),
                    file_name="stat_de_functii_salvat.docx",
                    mime=(mimetypes.guess_type("stat_de_functii_salvat.docx")[0] or "application/octet-stream"),
                    key="btn_export_stat_functii_salvat_docx",
                )
                pdf_bytes = convert_docx_to_pdf_bytes(buf_docx.getvalue())
                if pdf_bytes:
                    st.download_button(
                        "⬇️ Exportă statul de funcții salvat (PDF)",
                        data=pdf_bytes,
                        file_name="stat_de_functii_salvat.pdf",
                        mime="application/pdf",
                        key="btn_export_stat_functii_salvat_pdf",
                    )
            except Exception:
                pass

            st.markdown("#### Sincronizare cu lista de angajați (employees)")
            if st.button("🔄 Importă / actualizează employees din statul de funcții", key="btn_sync_stat_to_emp"):
                # construim un index de angajați existenți după Marcă
                emp_df = list_employees(conn, active_only=False)
                existing_by_marca = {}
                if not emp_df.empty:
                    for _, r_emp in emp_df.iterrows():
                        m = r_emp.get("marca") or r_emp.get("MARCA") or r_emp.get("Marca")
                        if m is not None:
                            existing_by_marca[str(m)] = r_emp

                inserate = 0
                actualizate = 0

                for _, r in edited_df.iterrows():
                    marca = r.get("Marca")
                    if pd.isna(marca) or not str(marca).strip():
                        continue
                    marca_str = str(marca).strip()

                    nume = r.get("Nume") or ""
                    prenume = r.get("Prenume") or ""
                    cnp = r.get("CNP") or ""
                    functie = r.get("Functie") or ""
                    loc_munca = r.get("Loc de munca") or ""

                    if marca_str in existing_by_marca:
                        row_emp = existing_by_marca[marca_str]
                        emp_id = int(row_emp["id"])
                        data = {}
                        for col in [
                            "marca","last_name","first_name","cnp",
                            "functie","departament","activ"
                        ]:
                            data[col] = None
                        data["marca"] = marca_str
                        data["last_name"] = str(nume).strip()
                        data["first_name"] = str(prenume).strip()
                        data["cnp"] = str(cnp).strip()
                        data["functie"] = str(functie).strip()
                        data["departament"] = str(loc_munca).strip()
                        data["activ"] = 1
                        update_employee(conn, emp_id, data)
                        actualizate += 1
                    else:
                        data = {}
                        for col in [
                            "marca",
                            "last_name",
                            "first_name",
                            "cnp",
                            "functie",
                            "departament",
                            "data_angajare",
                            "tip_contract",
                            "strada",
                            "numar",
                            "bloc",
                            "scara",
                            "apartament",
                            "cod_postal",
                            "localitate",
                            "judet",
                            "telefon_fix",
                            "mobil",
                            "email",
                            "ci_tip_act",
                            "ci_serie",
                            "ci_numar",
                            "ci_eliberat_de",
                            "ci_data_eliberarii",
                            "stare_civila",
                            "nr_copii",
                            "loc_munca",
                            "departament_organizatoric",
                            "functie_contract",
                            "tip_norma",
                            "program_munca",
                            "salariu_baza",
                            "studii",
                            "profesie",
                            "calificare",
                            "observatii",
                            "dosar_nr",
                            "dosar_functionar_public",
                            "dosar_data_intocmire",
                            "dosar_autoritate",
                            "dosar_intocmit_nume",
                            "dosar_intocmit_functie",
                            "dosar_intocmit_semnatura",
                            "dosar_modificari_nume",
                            "dosar_modificari_functie",
                            "dosar_modificari_semnatura",
                            "dosar_certificare_nume",
                            "dosar_certificare_functie",
                            "dosar_certificare_semnatura",
                            "activitate_in_afara_functiei",
                            "activitate_in_cadru_institutie",
                            "situatia_drepturi_salariale",
                            "situatia_concedii",
                            "situatia_disciplinara",
                            "registru_numar",
                            "registru_data",
                            "registru_observatii",
                            "activ",
                        ]:
                            data[col] = None
                        data["marca"] = marca_str
                        data["last_name"] = str(nume).strip()
                        data["first_name"] = str(prenume).strip()
                        data["cnp"] = str(cnp).strip()
                        data["functie"] = str(functie).strip()
                        data["departament"] = str(loc_munca).strip()
                        data["activ"] = 1
                        insert_employee(conn, data)
                        inserate += 1

                st.success(f"Sincronizare finalizată: {inserate} angajați noi, {actualizate} actualizați.")
# ---------------------------------------------------------
    # TAB 2: IMPORT DIN EXCEL (MAPARE CÂMPURI)
    # ---------------------------------------------------------
    with tab_import:
        st.markdown("### Stat de funcții completat (import din Excel)")

        uploaded_file = st.file_uploader(
            "Alege fișierul Excel cu datele pentru statul de funcții",
            type=["xlsx"],
            key="stat_functii_xlsx",
        )

        if uploaded_file is not None:
            try:
                df_src = pd.read_excel(uploaded_file)
            except Exception as e:
                st.error(f"Eroare la citirea fișierului Excel: {e}")
                return

            st.markdown("#### Previzualizare fișier Excel")
            st.dataframe(df_src.head(), use_container_width=True)

            st.markdown("#### Mapare câmpuri către modelul de stat de funcții")

            campuri_tinta = [
                "Marca",
                "CNP",
                "Nume",
                "Prenume",
                "Loc de munca",
                "Stare contract",
                "Tip angajat",
                "Tip raport",
                "Cod COR",
                "Functie",
                "Alta functie",
                "Explicatie alta functie",
                "Studii",
                "Gradatie",
                "Grad",
                "Vechime(ani)",
                "Norma",
                "Salariu de incadrare",
                "Norme salariu",
                "Data angajarii",
                "Data incetarii",
            ]

            col_options = ["(nu completa)"] + list(df_src.columns)
            mapping = {}

            for target in campuri_tinta:
                mapping[target] = st.selectbox(
                    f"Coloană pentru „{target}”",
                    options=col_options,
                    key=f"map_stat_{target.replace(' ', '_').replace('(', '').replace(')', '').replace('.', '')}",
                )

            if st.button("Generează și salvează stat de funcții din Excel"):
                n = len(df_src)
                out = {}
                for target in campuri_tinta:
                    src_col = mapping.get(target)
                    if src_col and src_col != "(nu completa)" and src_col in df_src.columns:
                        out[target] = df_src[src_col]
                    else:
                        out[target] = [""] * n

                df_stat_xlsx = pd.DataFrame(out)

                # Afișăm rezultatul
                st.markdown("#### Stat de funcții completat (din Excel)")
                st.dataframe(df_stat_xlsx, use_container_width=True)

                # Salvăm în baza de date pentru utilizare ulterioară
                df_stat_xlsx.to_sql("stat_functii", conn, if_exists="replace", index=False)
                st.success(
                    "Statul de funcții a fost generat din Excel și salvat în baza de date. "
                    "Îl poți edita în tab-ul „Stat de funcții salvat”."
                )

                excel_bytes = _df_to_xlsx_bytes(df_stat_xlsx, sheet_name="Stat_de_functii")
                if excel_bytes is None:
                    # Fallback CSV (nu necesită openpyxl/xlsxwriter)
                    excel_bytes = df_stat_xlsx.to_csv(index=False).encode("utf-8")
                    mime = "text/csv"
                    fname = "stat_de_functii_completat_import.csv"
                else:
                    mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    fname = "stat_de_functii_completat_import.xlsx"

                st.download_button(
                    "⬇️ Exportă statul de funcții completat (din Excel)",
                    data=excel_bytes,
                    file_name=fname,
                    mime=(mimetypes.guess_type(fname)[0] or "application/octet-stream"),
                    key="btn_export_stat_functii_completat_xlsx",
                )



def page_centralizator_concedii(conn: sqlite3.Connection):
    """Centralizator concedii bazat pe pontaj (CO/CM/CFP/FS/ALTA)."""
    ensure_pontaj_tables(conn)
    st.markdown('<h2 class="page-title">Centralizator concedii</h2>', unsafe_allow_html=True)

    col_y, col_m = st.columns(2)
    with col_y:
        year = st.number_input(
            "An",
            min_value=2000,
            max_value=2100,
            value=pd.Timestamp.today().year,
            step=1,
            key="centr_year",
        )
    with col_m:
        opt_luna = st.selectbox(
            "Lună",
            options=["Tot anul"] + [str(m) for m in range(1, 13)],
            index=0,
            key="centr_month",
        )

    if opt_luna == "Tot anul":
        period_start = pd.Timestamp(year=int(year), month=1, day=1)
        period_end = pd.Timestamp(year=int(year), month=12, day=31)
    else:
        m_sel = int(opt_luna)
        period_start = pd.Timestamp(year=int(year), month=m_sel, day=1)
        period_end = (period_start + pd.offsets.MonthEnd(0)).normalize()

    st.markdown(f"**Perioadă analizată:** {period_start.date()} - {period_end.date()}")

    tipuri_default = ["CO", "CM", "CFP", "FS"]
    tipuri_select = st.multiselect(
        "Tipuri de zile de inclus în centralizator",
        options=["CO", "CM", "CFP", "FS", "ALTA"],
        default=tipuri_default,
        key="centr_tipuri",
    )

    if not tipuri_select:
        st.warning("Selectează cel puțin un tip de zi pentru centralizator.")
        return

    placeholders = ",".join(["?"] * len(tipuri_select))
    params = [str(period_start.date()), str(period_end.date())] + tipuri_select

    sql = f"""
    SELECT
        e.marca AS Marca,
        e.last_name AS Nume,
        e.first_name AS Prenume,
        p.tip_zi,
        COUNT(DISTINCT p.data) AS Zile,
        SUM(p.ore_lucru) AS Ore
    FROM pontaj p
    LEFT JOIN employees e ON e.id = p.employee_id
    WHERE date(p.data) BETWEEN date(?) AND date(?)
      AND p.tip_zi IN ({placeholders})
    GROUP BY e.marca, e.last_name, e.first_name, p.tip_zi
    ORDER BY e.last_name, e.first_name, p.tip_zi
    """

    try:
        df_raw = pd.read_sql_query(sql, conn, params=params)
    except Exception as e:
        st.error(f"Eroare la citirea datelor din pontaj: {e}")
        return

    if df_raw.empty:
        st.info("Nu există concedii / zile de tipurile selectate în perioada aleasă.")
        return

    pivot = df_raw.pivot_table(
        index=["Marca", "Nume", "Prenume"],
        columns="tip_zi",
        values=["Zile", "Ore"],
        aggfunc="sum",
        fill_value=0,
    )

    pivot.columns = [f"{metric}_{tip}" for (metric, tip) in pivot.columns]
    df_c = pivot.reset_index()

    tipuri_in_data = sorted(df_raw["tip_zi"].unique())
    zile_cols = [f"Zile_{t}" for t in tipuri_in_data if f"Zile_{t}" in df_c.columns]
    ore_cols = [f"Ore_{t}" for t in tipuri_in_data if f"Ore_{t}" in df_c.columns]

    if zile_cols:
        df_c["Total_zile"] = df_c[zile_cols].sum(axis=1)
    else:
        df_c["Total_zile"] = 0

    if ore_cols:
        df_c["Total_ore"] = df_c[ore_cols].sum(axis=1)
    else:
        df_c["Total_ore"] = 0

    # Îmbogățim cu Loc de munca și Funcție din stat_functii (dacă există)
    try:
        df_stat = pd.read_sql_query(
            'SELECT Marca, "Loc de munca" AS Loc_de_munca, Functie FROM stat_functii',
            conn,
        )
        if not df_stat.empty:
            df_stat = df_stat.drop_duplicates(subset=["Marca"])
            df_c = df_c.merge(df_stat, on="Marca", how="left")
    except Exception:
        pass

    front_cols = ["Marca", "Nume", "Prenume"]
    opt_cols = ["Loc_de_munca", "Functie"]
    other_cols = [c for c in df_c.columns if c not in front_cols + opt_cols]
    final_cols = front_cols + [c for c in opt_cols if c in df_c.columns] + other_cols
    df_c = df_c[final_cols]

    st.markdown("### Centralizator concedii (zile și ore pe tip de zi)")
    st.dataframe(df_c, use_container_width=True)

    excel_bytes = _df_to_xlsx_bytes(df_c, sheet_name="Centralizator")
    if excel_bytes is None:
        # Fallback CSV (nu necesită openpyxl/xlsxwriter)
        excel_bytes = df_c.to_csv(index=False).encode("utf-8")
        mime = "text/csv"
        nume_fisier = f"centralizator_concedii_{int(year)}"
        if opt_luna != "Tot anul":
            nume_fisier += f"_{int(opt_luna):02d}"
        nume_fisier += ".csv"
    else:
        mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        nume_fisier = f"centralizator_concedii_{int(year)}"
        if opt_luna != "Tot anul":
            nume_fisier += f"_{int(opt_luna):02d}"
        nume_fisier += ".xlsx"

    st.download_button(
        "⬇️ Exportă centralizator concedii",
        data=excel_bytes,
        file_name=nume_fisier,
        mime=(mimetypes.guess_type(nume_fisier)[0] or "application/octet-stream"),
        key="btn_export_centralizator_concedii",
    )

    

    # Export în Word/PDF
    try:
        from docx import Document as _DocDoc
        doc_c = _DocDoc()
        doc_c.add_heading("Centralizator concedii", level=1)
        try:
            cfg = load_config()
        except Exception:
            cfg = {}
        unit_name = cfg.get("denumire_unitate", "")
        if unit_name:
            doc_c.add_paragraph(f"Unitate: {unit_name}")
        table = doc_c.add_table(rows=1, cols=len(df_c.columns))
        hdr_cells = table.rows[0].cells
        for j, col_name in enumerate(df_c.columns):
            hdr_cells[j].text = str(col_name)
        for _, row in df_c.iterrows():
            row_cells = table.add_row().cells
            for j, col_name in enumerate(df_c.columns):
                val = row.get(col_name, "")
                row_cells[j].text = "" if pd.isna(val) else str(val)
        buf_docx = BytesIO()
        doc_c.save(buf_docx)
        buf_docx.seek(0)
        st.download_button(
            "⬇️ Exportă centralizatorul concediilor (Word)",
            data=buf_docx.getvalue(),
            file_name=f"centralizator_concedii_{int(year)}.docx",
            mime=(mimetypes.guess_type(f"centralizator_concedii_{int(year)}.docx")[0] or "application/octet-stream"),
                key="centr_export_docx",
        )
        pdf_bytes = convert_docx_to_pdf_bytes(buf_docx.getvalue())
        if pdf_bytes:
            st.download_button(
                "⬇️ Exportă centralizatorul concediilor (PDF)",
                data=pdf_bytes,
                file_name=f"centralizator_concedii_{int(year)}.pdf",
                mime="application/pdf",
                key="centr_export_pdf",
            )
    except Exception:
        pass

# -------------------------------------------------------------
# PONTAJ
# -------------------------------------------------------------
def page_pontaj(conn: sqlite3.Connection):
    """Modulul de pontaj angajați, legat de Statul de funcții și Centralizator concedii."""
    ensure_pontaj_tables(conn)
    st.markdown('<h2 class="page-title">Pontaj angajați</h2>', unsafe_allow_html=True)

    col_y, col_m = st.columns(2)
    with col_y:
        year = st.number_input(
            "An",
            min_value=2000,
            max_value=2100,
            value=pd.Timestamp.today().year,
            step=1,
            key="pontaj_year",
        )
    with col_m:
        month = st.number_input(
            "Lună",
            min_value=1,
            max_value=12,
            value=pd.Timestamp.today().month,
            step=1,
            key="pontaj_month",
        )

    period_start = pd.Timestamp(year=int(year), month=int(month), day=1)
    period_end = (period_start + pd.offsets.MonthEnd(0)).normalize()

    st.markdown(f"**Perioadă pontaj:** {period_start.date()} - {period_end.date()}")

    st.markdown("### Adaugă pontaj pentru un angajat")

    emp_df = list_employees(conn, active_only=True)
    if emp_df.empty:
        st.info("Nu există angajați activi în baza de date.")
        return

    recs = emp_df.to_dict("records")
    idx_emp = st.selectbox(
        "Alege angajatul",
        options=list(range(len(recs))),
        format_func=lambda i: build_display_name(recs[i]),
        key="pontaj_emp_idx",
    )
    emp_row = recs[idx_emp]
    emp_id = (
        emp_row.get("id")
        or emp_row.get("ID")
        or emp_row.get("id_employee")
        or emp_row.get("ID_ANGAJAT")
    )

    marca_emp = (
        emp_row.get("marca")
        or emp_row.get("Marca")
        or emp_row.get("MARCA")
    )
    if marca_emp is not None:
        try:
            cur_info = conn.cursor()
            cur_info.execute(
                'SELECT "Loc de munca", Functie FROM stat_functii WHERE Marca = ? LIMIT 1',
                (str(marca_emp),),
            )
            row_info = cur_info.fetchone()
            if row_info:
                loc_munca, functie = row_info[0], row_info[1]
                st.markdown(
                    f"**Funcție:** {functie or '-'}  •  **Loc de muncă:** {loc_munca or '-'}"
                )
        except Exception:
            pass

    # Formular pontaj
    col_d, col_ore, col_tip = st.columns(3)
    with col_d:
        import datetime as _dt_pontaj
        _today_p = _dt_pontaj.date.today()
        _min_p = _today_p.replace(year=_today_p.year - 100)
        _max_p = _today_p.replace(year=_today_p.year + 10)
        data_p = st.date_input(
            "Data",
            value=period_start.date(),
            min_value=_min_p,
            max_value=_max_p,
            key="pontaj_data",
        )
    with col_ore:
        ore_lucru = st.number_input(
            "Ore lucrate",
            min_value=0.0,
            max_value=24.0,
            step=0.5,
            value=8.0,
            key="pontaj_ore",
        )
    with col_tip:
        tip_zi = st.selectbox(
            "Tip zi",
            options=["LUCRU", "CO", "CM", "CFP", "FS", "ALTA"],
            index=0,
            key="pontaj_tip_zi",
        )

    observatii = st.text_input(
        "Observații (opțional)",
        key="pontaj_observatii",
    )

    if st.button("💾 Adaugă în pontaj", key="pontaj_btn_add"):
        cur = conn.cursor()
        cur.execute(
            """
            INSERT INTO pontaj (employee_id, data, ore_lucru, tip_zi, observatii)
            VALUES (?, ?, ?, ?, ?)
            """,
            (int(emp_id), str(data_p), float(ore_lucru), tip_zi, observatii.strip() or None),
        )
        conn.commit()
        st.success("Înregistrarea de pontaj a fost salvată.")

    st.markdown("---")
    st.markdown("### Pontajul lunii curente (editabil)")

    cur = conn.cursor()
    cur.execute(
        """
        SELECT
            p.id,
            p.data,
            p.ore_lucru,
            p.tip_zi,
            p.observatii,
            e.marca,
            e.last_name,
            e.first_name
        FROM pontaj p
        LEFT JOIN employees e ON e.id = p.employee_id
        WHERE date(p.data) BETWEEN date(?) AND date(?)
        ORDER BY p.data, e.last_name, e.first_name
        """,
        (str(period_start.date()), str(period_end.date())),
    )
    rows = cur.fetchall()
    if not rows:
        st.info("Nu există pontaj în perioada selectată.")
        return

    df_p = pd.DataFrame(rows, columns=rows[0].keys())

    edited_df = st.data_editor(
        df_p,
        use_container_width=True,
        num_rows="dynamic",
        key="pontaj_editor",
    )

    if st.button("💾 Salvează modificările de pontaj", key="pontaj_btn_save_all"):
        # rescriem tabela pontaj pentru perioada aleasă
        cur = conn.cursor()
        cur.execute(
            "DELETE FROM pontaj WHERE date(data) BETWEEN date(?) AND date(?)",
            (str(period_start.date()), str(period_end.date())),
        )
        conn.commit()

        for _, r in edited_df.iterrows():
            try:
                emp_marca = r.get("marca") or r.get("Marca") or r.get("MARCA")
                emp_id_row = None
                if emp_marca:
                    cur.execute(
                        "SELECT id FROM employees WHERE marca = ? OR MARCA = ?",
                        (str(emp_marca), str(emp_marca)),
                    )
                    found = cur.fetchone()
                    if found:
                        emp_id_row = found[0]

                if not emp_id_row:
                    # fallback: dacă nu găsim după marcă, lăsăm employee_id None
                    # (sau poți ignora rândul)
                    continue

                cur.execute(
                    """
                    INSERT INTO pontaj (employee_id, data, ore_lucru, tip_zi, observatii)
                    VALUES (?, ?, ?, ?, ?)
                    """,
                    (
                        int(emp_id_row),
                        str(r["data"]),
                        float(r.get("ore_lucru", 0) or 0),
                        str(r.get("tip_zi") or ""),
                        str(r.get("observatii") or "") or None,
                    ),
                )
            except Exception:
                continue
        conn.commit()
        st.success("Pontajul a fost rescris pentru perioada selectată.")

    st.markdown("#### Operațiuni suplimentare")
    col_a, col_b = st.columns(2)
    with col_a:
        if st.button(
            "🗑️ Șterge TOT pontajul pentru luna selectată",
            key="pontaj_btn_del_month",
        ):
            cur = conn.cursor()
            cur.execute(
                "DELETE FROM pontaj WHERE date(data) BETWEEN date(?) AND date(?)",
                (str(period_start.date()), str(period_end.date())),
            )
            conn.commit()
            st.success("Pontajul pentru luna selectată a fost șters.")
    with col_b:
        excel_bytes = _df_to_xlsx_bytes(edited_df, sheet_name="Pontaj")
        if excel_bytes is None:
            # Fallback CSV (nu necesită openpyxl/xlsxwriter)
            excel_bytes = edited_df.to_csv(index=False).encode("utf-8")
            mime = "text/csv"
            fname = f"pontaj_{int(year)}_{int(month):02d}.csv"
        else:
            mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            fname = f"pontaj_{int(year)}_{int(month):02d}.xlsx"

        st.download_button(
            "⬇️ Exportă pontajul",
            data=excel_bytes,
            file_name=fname,
            mime=(mimetypes.guess_type(fname)[0] or "application/octet-stream"),
            key="pontaj_btn_export",
        )


# -------------------------------------------------------------
# MAIN
# -------------------------------------------------------------

def page_pontaj_hub(conn, cfg: dict):
    """Hub Pontaj: include aplicația Pontaj separată + centralizator concedii (în același tab)."""
    st.markdown('<h2 class="page-title">Pontaj</h2>', unsafe_allow_html=True)
    pontaj_url_default = cfg.get("pontaj_url") or "http://localhost:8502"
    with st.expander("🔗 Deschide aplicația Pontaj (modul separat)", expanded=True):
        st.info("Pontajul rulează ca aplicație separată (modul independent) pe același server, dar pe alt port.")
        pontaj_url = st.text_input("URL Pontaj", value=pontaj_url_default, key="pontaj_url_inline")
        st.markdown(f"👉 Deschide Pontaj: {pontaj_url}")
        st.link_button("🚀 Deschide Pontaj", pontaj_url)
        st.caption("Dacă nu se deschide, verifică firewall/port (ex: 8502) și că serviciul pontaj.rulează.")
    st.markdown("---")
    sub = st.radio("Alege în Pontaj", ["📌 Centralizator concedii", "⏱️ Pontaj (în Socrates – opțional)"], index=0, key="pontaj_hub_choice")
    if sub.startswith("📌"):
        page_centralizator_concedii(conn)
    else:
        page_pontaj(conn)


def main():
    # Layout și temă globală
    # st.set_page_config(page_title=APP_TITLE, layout="wide")
    apply_premium_theme()
    apply_centered_layout(1280)
    # Spacing global: primul conținut util începe imediat sub top bar
    st.markdown(
        """
        <style>
        section[data-testid="stMain"] .block-container { padding-top: 0.5rem !important; margin-top: 0 !important; }
        div[data-testid="stAppViewBlockContainer"] { padding-top: 0.5rem !important; margin-top: 0 !important; }
        div[data-testid="stMainBlockContainer"] { padding-top: 0.5rem !important; margin-top: 0 !important; }
        section[data-testid="stMain"] div[data-testid="stVerticalBlock"] { gap: 0.35rem !important; }
        section[data-testid="stMain"] div[data-testid="stVerticalBlock"] > div:first-child {
          margin-top: 0 !important;
          padding-top: 0 !important;
          min-height: 0 !important;
        }
        section[data-testid="stMain"] h1,
        section[data-testid="stMain"] h2,
        section[data-testid="stMain"] h3 { margin-top: 0 !important; padding-top: 0 !important; }
        </style>
        """,
        unsafe_allow_html=True,
    )

    cfg = load_config()
    db_path = cfg.get("db_path", DB_PATH_DEFAULT)

    # Conectare DB + autentificare (recomandat pentru VPS/cloud)
    conn = get_connection()

    # Asigurăm schema organigramei și pe conexiunea globală HR (dacă este folosită)
    try:
        ensure_org_schema(conn_hr)
    except Exception:
        pass

    require_login(conn, cfg)

    # După login, afișăm la loc header-ul Streamlit.
    if st.session_state.get("logged_in", False):
        st.markdown(
            """
            <style>
            header[data-testid="stHeader"] { visibility: visible !important; height: auto !important; }
            </style>
            """,
            unsafe_allow_html=True,
        )

    # Sidebar: brand
    st.sidebar.markdown(
        '<div class="sb-brand">Socrates@HR</div>',
        unsafe_allow_html=True,
    )

    # Sidebar: user card (username ca nume principal)
    _user = (st.session_state.get("username") or "").strip() or "—"
    st.sidebar.markdown(
        f'<div class="sb-user">'
        f'<div class="sb-user-title">UTILIZATOR</div>'
        f'<div class="name">{_user}</div>'
        f'</div>',
        unsafe_allow_html=True,
    )

    # Navigare principală: meniu unitar, cu spacing controlat, dar păstrând aceleași valori logice.
    OPS = ["🏠 Acasă", "Organigramă", "Angajați"]
    ADM = ["Stat de funcții", "Pontaj", "Dosar profesional", "Configurare"]
    NAV = OPS + ADM

    main_choice = st.session_state.get("main_choice", "🏠 Acasă")
    if main_choice not in NAV:
        main_choice = "🏠 Acasă"
        st.session_state["main_choice"] = main_choice
    st.session_state.setdefault("main_choice", main_choice)

    nav_index = NAV.index(main_choice)

    # Un singur st.radio pentru toate butoanele principale – spacing uniform, flow vertical coerent.
    nav_choice = st.sidebar.radio(
        "",
        NAV,
        index=nav_index,
        key="sidebar_nav_pick",
        label_visibility="collapsed",
    )

    # Păstrăm comportamentul existent: main_choice rămâne sursa adevărului.
    st.session_state["main_choice"] = nav_choice

    # Sincronizăm vechile chei de stare astfel încât orice logică existentă bazată pe ele
    # să continue să funcționeze identic.
    if nav_choice in OPS:
        st.session_state["sidebar_ops_pick"] = nav_choice
        # Dacă nu e selectat nimic în ADM, păstrăm ultima alegere sau primul element.
        st.session_state.setdefault("sidebar_adm_pick", ADM[0])
    else:
        st.session_state["sidebar_adm_pick"] = nav_choice
        st.session_state.setdefault("sidebar_ops_pick", OPS[0])

    main_choice = st.session_state.get("main_choice", "🏠 Acasă")

    # Fundal per pagină: imagine doar pe Acasă, albastru solid pe restul
    # Scoatem emoji-ul pentru comparația cu apply_page_background
    if main_choice.startswith("🏠"):
        logical_page = "Acasă"
    else:
        logical_page = main_choice
    apply_page_background(logical_page, home_bg_css=HOME_BG_CSS)

    if main_choice == "Dosar profesional":
        st.session_state.setdefault(
            "sub_menu",
            "Copertă",
        )
        st.sidebar.markdown('<div class="sb-subexp sb-dosar">', unsafe_allow_html=True)
        with st.sidebar.expander("📁 Dosar profesional", expanded=True):
            st.radio(
                "Secțiuni dosar profesional",
                [
                    "Copertă",
                    "Date cu caracter personal",
                    "Studii și pregătire profesională",
                    "Activitate în afara funcției publice",
                    "Activitate în cadrul instituției",
                    "Situația drepturilor salariale",
                    "Situația concediilor",
                    "Situația disciplinară",
                    "Accesul la dosarul profesional",
                    "Registru evidență funcționari publici",
                ],
                key="sub_menu",
            )
        st.sidebar.markdown("</div>", unsafe_allow_html=True)

    # Footer sidebar: separator + butoane
    st.sidebar.markdown('<div class="sb-sep"></div>', unsafe_allow_html=True)
    st.sidebar.markdown('<div class="sb-foot">', unsafe_allow_html=True)
    _col_logout, _col_help = st.sidebar.columns(2)
    with _col_logout:
        if st.sidebar.button("🚪 Deconectare", key="sidebar_btn_logout"):
            st.session_state.logged_in = False
            st.session_state.user_role = "user"
            st.session_state.username = ""
            st.rerun()
    with _col_help:
        if st.sidebar.button("❓ Ajutor", key="sidebar_btn_help"):
            st.session_state["sidebar_show_help"] = True
            st.rerun()
    if st.session_state.get("sidebar_show_help"):
        st.sidebar.info("Contact IT pentru resetare parolă sau suport. Adaugă aici instrucțiuni sau link.")
        st.session_state["sidebar_show_help"] = False

    # Footer context: versiune + DB (basename)
    try:
        db_name = Path(get_db_path()).name
    except Exception:
        db_name = "necunoscut"
    st.sidebar.markdown(
        f"<div class='small'>v1.0 • Local<br>DB: {db_name}</div>",
        unsafe_allow_html=True,
    )

    st.sidebar.markdown('</div>', unsafe_allow_html=True)

    if main_choice == "🏠 Acasă":
        page_home(conn, cfg)

    elif main_choice == "Angajați":
        page_angajati(conn)

    elif main_choice == "Organigramă":
        page_organigrama(conn)
        
    elif main_choice == "Stat de funcții":
        page_stat_de_functii(conn)

    elif main_choice == "Dosar profesional":
        page_dosar_profesional(conn, st.session_state["sub_menu"])
    elif main_choice == "Pontaj":
        page_pontaj_hub(conn, cfg)

    elif main_choice == "Configurare":
        # Marker + CSS scoped doar pentru pagina Configurare
        st.markdown('<span id="cfg-scope"></span>', unsafe_allow_html=True)
        st.markdown(
            """
            <style>
            /* ===== Configurare – layout single column, premium ===== */
            section.main:has(#cfg-scope) .page-container{
              max-width: 880px;
              margin-left: 60px;
              margin-right: auto;
              padding-top: 40px;
              padding-bottom: 80px;
            }

            section.main:has(#cfg-scope) .section-card{
              background: #162a3d;
              border: 1px solid #2c445c;
              border-radius: 10px;
              padding: 32px;
              margin-bottom: 48px;
            }

            section.main:has(#cfg-scope) .section-title{
              font-size: 18px;
              font-weight: 600;
              margin-bottom: 24px;
              color: #e6edf5;
            }

            section.main:has(#cfg-scope) .section-title::after{
              content: "";
              display: block;
              width: 40px;
              height: 2px;
              background: #1d4ed8;
              margin-top: 8px;
              opacity: 0.6;
            }

            section.main:has(#cfg-scope) .form-group{
              margin-bottom: 20px;
            }

            section.main:has(#cfg-scope) label{
              font-size: 13px;
              color: #9fb3c8 !important;
              margin-bottom: 6px;
              display: block;
            }

            section.main:has(#cfg-scope) input,
            section.main:has(#cfg-scope) select,
            section.main:has(#cfg-scope) textarea{
              height: 40px;
              padding: 8px 12px;
            }

            section.main:has(#cfg-scope) .sub-block{
              background: #132636;
              border-radius: 8px;
              padding: 20px;
              margin-bottom: 20px;
            }

            section.main:has(#cfg-scope) .upload-zone{
              background: #132636;
              border: 1px dashed #2c445c;
              border-radius: 8px;
              padding: 20px;
              transition: border-color 0.2s ease;
            }

            section.main:has(#cfg-scope) .upload-zone:hover{
              border-color: #1d4ed8;
            }

            /* Upload Antet – elimină caseta albă internă și o face albastru închis */
            section.main:has(#cfg-scope) .upload-zone [data-testid="stFileUploader"]{
              background: transparent;
            }

            section.main:has(#cfg-scope) .upload-zone [data-testid="stFileUploader"] > div:first-child{
              background: #0f172a;
              border-radius: 10px;
            }

            section.main:has(#cfg-scope) .upload-zone [data-testid="stFileUploaderDropzone"]{
              background: #0f172a;
              border: 1px dashed #2c445c;
              border-radius: 10px;
            }

            section.main:has(#cfg-scope) .save-bar{
              position: sticky;
              bottom: 0;
              background: #0f1f2e;
              padding: 20px 0;
              border-top: 1px solid #2c445c;
              margin-top: 16px;
            }

            section.main:has(#cfg-scope) .save-bar div[data-testid="stButton"] > button{
              background: #1d4ed8 !important;
              padding: 10px 22px !important;
              border-radius: 6px !important;
              font-weight: 500 !important;
              border: none !important;
            }

            section.main:has(#cfg-scope) .save-bar div[data-testid="stButton"] > button:hover{
              background: #1e40af !important;
            }

            /* ===== Utilizatori (admin) – government clean dark table ===== */
            section.main:has(#cfg-scope) .admin-table{
              width: 100%;
              border-collapse: separate;
              border-spacing: 0;
            }

            section.main:has(#cfg-scope) .admin-table thead{
              background: #1c3146;
            }

            section.main:has(#cfg-scope) .admin-table th{
              padding: 12px 16px;
              text-align: left;
              font-weight: 600;
              color: #e6edf5;
              border-bottom: 1px solid #2c445c;
            }

            section.main:has(#cfg-scope) .admin-table td{
              padding: 12px 16px;
              background: #162a3d;
              border-bottom: 1px solid #2c445c;
              color: #e6edf5;
            }

            section.main:has(#cfg-scope) .admin-table tr:hover td{
              background: #1a2e42;
            }

            /* colțuri rotunjite tabel */
            section.main:has(#cfg-scope) .admin-table thead tr:first-child th:first-child{
              border-top-left-radius: 8px;
            }
            section.main:has(#cfg-scope) .admin-table thead tr:first-child th:last-child{
              border-top-right-radius: 8px;
            }
            section.main:has(#cfg-scope) .admin-table tbody tr:last-child td:first-child{
              border-bottom-left-radius: 8px;
            }
            section.main:has(#cfg-scope) .admin-table tbody tr:last-child td:last-child{
              border-bottom-right-radius: 8px;
            }

            /* badge status */
            section.main:has(#cfg-scope) .badge-active{
              background: rgba(22, 163, 74, 0.15);
              color: #22c55e;
              padding: 4px 10px;
              border-radius: 20px;
              font-size: 12px;
              font-weight: 500;
              display: inline-block;
            }
            section.main:has(#cfg-scope) .badge-inactive{
              background: rgba(220, 38, 38, 0.15);
              color: #ef4444;
              padding: 4px 10px;
              border-radius: 20px;
              font-size: 12px;
              font-weight: 500;
              display: inline-block;
            }

            /* zona “acțiuni periculoase” */
            section.main:has(#cfg-scope) .danger-section{
              background: rgba(220, 38, 38, 0.08);
              border: 1px solid rgba(220, 38, 38, 0.25);
              border-radius: 8px;
              padding: 20px;
              margin-top: 8px;
            }

            section.main:has(#cfg-scope) .danger-section div[data-testid="stButton"] > button{
              background: #7f1d1d !important;
              color: #ffffff !important;
            }

            section.main:has(#cfg-scope) .danger-section div[data-testid="stButton"] > button:hover{
              background: #991b1b !important;
            }

            /* ===== Formular Utilizatori (admin) – compact, 2 coloane ===== */
            section.main:has(#cfg-scope) .user-admin-card{
              background: #162a3d;
              border: 1px solid #2c445c;
              border-radius: 10px;
              padding: 24px 28px;
              margin-bottom: 24px;
              max-width: 700px;
            }

            section.main:has(#cfg-scope) .user-form{
              max-width: 700px;
            }

            section.main:has(#cfg-scope) .user-form .form-grid{
              display: grid;
              grid-template-columns: 280px 280px;
              gap: 16px 24px;
              margin-bottom: 20px;
            }

            section.main:has(#cfg-scope) .user-form .form-group{
              display: flex;
              flex-direction: column;
            }

            section.main:has(#cfg-scope) .user-form input,
            section.main:has(#cfg-scope) .user-form select{
              height: 36px;
              font-size: 14px;
              padding: 6px 10px;
              max-width: 280px;
            }

            section.main:has(#cfg-scope) .user-form .button-group{
              display: flex;
              flex-direction: column;
              gap: 12px;
              align-items: flex-start;
            }

            /* tabel admin – structură enterprise */
            section.main:has(#cfg-scope) .admin-table{
              width: 100%;
              border-collapse: separate;
              border-spacing: 0;
              font-size: 14px;
            }

            section.main:has(#cfg-scope) .admin-table th,
            section.main:has(#cfg-scope) .admin-table td{
              padding: 12px 16px;
              border-bottom: 1px solid #2c445c;
              border-right: 1px solid #2c445c;
            }

            section.main:has(#cfg-scope) .admin-table th:last-child,
            section.main:has(#cfg-scope) .admin-table td:last-child{
              border-right: none;
            }

            section.main:has(#cfg-scope) .admin-table thead{
              background: #1c3146;
            }

            section.main:has(#cfg-scope) .admin-table thead th{
              font-weight: 600;
              border-bottom: 2px solid #2c445c;
            }

            section.main:has(#cfg-scope) .admin-table tbody tr{
              background: #162a3d;
              font-weight: 400;
            }

            section.main:has(#cfg-scope) .admin-table tbody tr:hover{
              background: #1a2e42;
            }

            section.main:has(#cfg-scope) .admin-table td{
              white-space: nowrap;
            }

            section.main:has(#cfg-scope) .table-wrapper{
              overflow-x: auto;
            }

            </style>
            """,
            unsafe_allow_html=True,
        )

        st.markdown('<div class="page-container">', unsafe_allow_html=True)
        st.markdown('<h2 class="section-title">⚙ Configurare</h2>', unsafe_allow_html=True)

        # Hub navigare configurare
        cfg_view = st.session_state.get("config_view", "")
        st.markdown('<div class="config-navigation">', unsafe_allow_html=True)
        nav_pontaj = st.button("🕒 Pontaj (modul separat)", key="cfg_nav_pontaj")
        nav_users = st.button("👥 Utilizatori (admin)", key="cfg_nav_users")
        st.markdown("</div>", unsafe_allow_html=True)

        if nav_pontaj:
            st.session_state["config_view"] = "pontaj"
            st.rerun()
        if nav_users:
            st.session_state["config_view"] = "users"
            st.rerun()

        cfg_view = st.session_state.get("config_view", "")

        # --- Pontaj (modul separat) – pagină dedicată în Configurare ---
        if cfg_view == "pontaj":
            st.markdown('<div class="section-card">', unsafe_allow_html=True)
            st.markdown('<div class="section-title">🕒 Pontaj (modul separat)</div>', unsafe_allow_html=True)
            cfg["pontaj_url"] = st.text_input(
                "URL aplicație Pontaj (ex: http://SERVER:8502)",
                value=str(cfg.get("pontaj_url") or "http://localhost:8502"),
                key="cfg_pontaj_url",
            )
            if st.button("Închide", key="cfg_close_pontaj"):
                st.session_state["config_view"] = ""
                st.rerun()
            st.markdown("</div>", unsafe_allow_html=True)

        # --- Administrare utilizatori (doar pentru admin) – pagină dedicată în Configurare ---
        is_admin = str(st.session_state.get("user_role", "user")) == "admin"
        if is_admin and cfg_view == "users":
            st.markdown('<div class="admin-users">', unsafe_allow_html=True)
            # Buton de închidere pentru secțiunea Utilizatori (admin)
            if st.button("Închide", key="cfg_close_users"):
                st.session_state["config_view"] = ""
                st.rerun()

            ensure_auth_tables(conn)
            curu = conn.cursor()
            curu.execute("SELECT username, role, is_active, created_at FROM app_users ORDER BY username")
            rows_u = curu.fetchall() or []
            if rows_u:
                from html import escape as _escape

                body_rows = []
                for r in rows_u:
                    username, role, is_active, created_at = r
                    badge = (
                        '<span class="badge-active">Activ</span>'
                        if int(is_active or 0)
                        else '<span class="badge-inactive">Inactiv</span>'
                    )
                    created_txt = _escape(str(created_at or ""))
                    body_rows.append(
                        f"<tr>"
                        f"<td>{_escape(str(username or ''))}</td>"
                        f"<td>{_escape(str(role or ''))}</td>"
                        f"<td>{badge}</td>"
                        f"<td>{created_txt}</td>"
                        f"</tr>"
                    )

                table_html = """
                <div class="table-wrapper">
                  <table class="admin-table">
                    <thead>
                      <tr>
                        <th>Username</th>
                        <th>Rol</th>
                        <th>Status</th>
                        <th>Creat la</th>
                      </tr>
                    </thead>
                    <tbody>
                """ + "\n".join(body_rows) + """
                    </tbody>
                  </table>
                </div>
                """
                st.markdown(table_html, unsafe_allow_html=True)
            else:
                st.info("Nu există utilizatori în DB (se va crea automat unul la login).")

            st.markdown('<div class="user-admin-card">', unsafe_allow_html=True)
            st.markdown("##### Adaugă / resetează utilizator")

            st.markdown('<div class="user-form">', unsafe_allow_html=True)

            # Grid intern: Username, Parolă nouă, Rol
            st.markdown('<div class="form-grid">', unsafe_allow_html=True)

            st.markdown('<div class="form-group">', unsafe_allow_html=True)
            nu = st.text_input("Username", key="adm_new_user").strip()
            st.markdown('</div>', unsafe_allow_html=True)

            st.markdown('<div class="form-group">', unsafe_allow_html=True)
            npw = st.text_input("Parolă nouă", type="password", key="adm_new_pass")
            st.markdown('</div>', unsafe_allow_html=True)

            st.markdown('<div class="form-group">', unsafe_allow_html=True)
            role = st.selectbox("Rol", ["user", "admin"], key="adm_new_role")
            st.markdown('</div>', unsafe_allow_html=True)

            st.markdown('</div>', unsafe_allow_html=True)  # end .form-grid

            # Câmp dezactivare sub grid
            st.markdown('<div class="form-group">', unsafe_allow_html=True)
            du = st.text_input("Dezactivează user", key="adm_disable_user").strip()
            st.markdown('</div>', unsafe_allow_html=True)

            # Grup butoane – afișate vertical, aliniate stânga
            st.markdown('<div class="button-group">', unsafe_allow_html=True)
            if st.button("Salvează (create/update)", key="adm_save_user"):
                if not nu or not npw:
                    st.error("Completează username și parolă.")
                else:
                    curu.execute(
                        """
                        INSERT INTO app_users(username, password_hash, role, is_active, created_at)
                        VALUES(?, ?, ?, 1, ?)
                        ON CONFLICT(username) DO UPDATE SET
                            password_hash=excluded.password_hash,
                            role=excluded.role,
                            is_active=1
                        """,
                        (nu, sqlite3.Binary(_hash_password(npw)), role, datetime.datetime.now().isoformat(timespec="seconds")),
                    )
                    conn.commit()
                    st.success("Utilizator salvat.")
                    st.rerun()
            if st.button("Dezactivează", key="adm_disable_btn"):
                if not du:
                    st.error("Scrie username-ul.")
                else:
                    curu.execute("UPDATE app_users SET is_active=0 WHERE username=?", (du,))
                    conn.commit()
                    st.success("Utilizator dezactivat (dacă exista).")
                    st.rerun()
            st.markdown('</div>', unsafe_allow_html=True)  # end .button-group

            # Acțiuni periculoase rămân separate, vizual
            with st.expander("⚠️ Acțiuni periculoase", expanded=False):
                st.markdown('<div class="danger-section">', unsafe_allow_html=True)
                st.caption("Oprește procesul aplicației (local/VPS). Folosește doar când e necesar.")
                confirm = st.checkbox("Confirm oprirea aplicației", value=False, key="shutdown_confirm")

                if st.button("⛔ Oprește aplicația", disabled=not confirm, key="btn_close_app"):
                    try:
                        import os as _os
                        _os._exit(0)
                    except Exception:
                        st.warning("Aplicația a fost oprită. Poți închide această fereastră.")
                        st.stop()
                st.markdown("</div>", unsafe_allow_html=True)

                st.markdown("</div>", unsafe_allow_html=True)  # end .user-form
                st.markdown("</div>", unsafe_allow_html=True)  # end .user-admin-card

                st.markdown("</div>", unsafe_allow_html=True)  # end .admin-users

        # --- Date unitate / instituție ---
        st.markdown('<div class="section-card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Date unitate / instituție</div>', unsafe_allow_html=True)

        denumire_unitate = st.text_input(
            "Denumire unitate",
            value=cfg.get("denumire_unitate", ""),
            key="cfg_denumire",
        )
        cui = st.text_input(
            "CUI / CIF",
            value=cfg.get("cui", ""),
            key="cfg_cui",
        )
        adresa = st.text_input(
            "Adresă (stradă, număr, localitate, județ)",
            value=cfg.get("adresa", ""),
            key="cfg_adresa",
        )
        cont_bancar = st.text_input(
            "Cont bancar (IBAN)",
            value=cfg.get("cont_bancar", ""),
            key="cfg_cont_bancar",
        )
        banca = st.text_input(
            "Bancă",
            value=cfg.get("banca", ""),
            key="cfg_banca",
        )

        telefon = st.text_input(
            "Telefon",
            value=cfg.get("telefon", ""),
            key="cfg_tel",
        )
        fax = st.text_input(
            "Fax",
            value=cfg.get("fax", ""),
            key="cfg_fax",
        )
        email = st.text_input(
            "Email",
            value=cfg.get("email", ""),
            key="cfg_email",
        )
        st.markdown("</div>", unsafe_allow_html=True)

        # --- Antet ---
        st.markdown('<div class="section-card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Antet documente</div>', unsafe_allow_html=True)
        st.markdown('<div class="upload-zone">', unsafe_allow_html=True)
        sigla_file = st.file_uploader(
            "Siglă unitate (PNG/JPG) – opțional (se va afișa în stânga antetului)",
            type=["png", "jpg", "jpeg"],
            key="cfg_sigla_upl",
        )
        st.markdown("</div>", unsafe_allow_html=True)
        afiseaza_iban_in_antet = st.checkbox(
            "Afișează IBAN în antet (opțional)",
            value=bool(cfg.get("afiseaza_iban_in_antet", False)),
            key="cfg_show_iban",
        )
        st.markdown("</div>", unsafe_allow_html=True)

        # --- Semnare documente + Semnături ---
        st.markdown('<div class="section-card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Semnare documente</div>', unsafe_allow_html=True)

        semnare_metoda = st.selectbox(
            "Metodă semnare",
            options=["Nesemnate", "Olograf (scan)", "Electronic (token)"],
            index={"Nesemnate": 0, "Olograf (scan)": 1, "Electronic (token)": 2}.get(cfg.get("semnare_metoda", "Nesemnate"), 0),
            key="cfg_semnare_metoda",
        )

        st.markdown('<div class="sub-block">', unsafe_allow_html=True)
        st.markdown("**Conducător unitate**", unsafe_allow_html=True)
        conducator_nume = st.text_input(
            "Nume și prenume",
            value=cfg.get("conducator_nume", ""),
            key="cfg_conducator_nume",
        )
        conducator_functie = st.text_input(
            "Funcție",
            value=cfg.get("conducator_functie", ""),
            key="cfg_conducator_functie",
        )
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown('<div class="sub-block">', unsafe_allow_html=True)
        st.markdown("**Responsabil HR**", unsafe_allow_html=True)
        responsabil_hr_nume = st.text_input(
            "Nume și prenume",
            value=cfg.get("responsabil_hr_nume", ""),
            key="cfg_responsabil_hr_nume",
        )
        responsabil_hr_functie = st.text_input(
            "Funcție",
            value=cfg.get("responsabil_hr_functie", ""),
            key="cfg_responsabil_hr_functie",
        )
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown("</div>", unsafe_allow_html=True)

        # --- User aplicație ---
        st.markdown('<div class="section-card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">User aplicație</div>', unsafe_allow_html=True)

        app_user = st.text_input(
            "Utilizator aplicație",
            value=cfg.get("app_user", ""),
            key="cfg_app_user",
        )
        app_pass = st.text_input(
            "Parolă aplicație",
            value=cfg.get("app_pass", ""),
            type="password",
            key="cfg_app_pass",
        )
        st.markdown("</div>", unsafe_allow_html=True)

        # --- Bază de date ---
        st.markdown('<div class="section-card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Bază de date</div>', unsafe_allow_html=True)

        new_db_path = st.text_input(
            "Calea către baza de date (db_path)",
            value=db_path,
            key="cfg_db_path",
        )
        st.markdown("</div>", unsafe_allow_html=True)

        # --- Câmpuri suplimentare (custom) ---
        st.markdown('<div class="section-card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Câmpuri suplimentare</div>', unsafe_allow_html=True)

        custom_fields: Dict[str, str] = cfg.get("custom_fields", {})
        if not isinstance(custom_fields, dict):
            custom_fields = {}

        # Afișăm câmpurile existente
        if custom_fields:
            st.markdown("#### Câmpuri existente")
            for k, v in custom_fields.items():
                new_v = st.text_input(
                    k,
                    value=str(v),
                    key=f"cfg_custom_{k}",
                )
                custom_fields[k] = new_v

        st.markdown("#### Adaugă câmp nou")
        new_custom_name = st.text_input(
            "Nume câmp nou",
            value="",
            key="cfg_custom_new_name",
            help="Ex: 'Telefon instituție', 'Website', etc.",
        )
        new_custom_value = st.text_input(
            "Valoare câmp nou",
            value="",
            key="cfg_custom_new_value",
        )

        if st.button("➕ Adaugă câmp", key="cfg_btn_add_custom"):
            if not new_custom_name.strip():
                st.error("Numele câmpului este obligatoriu.")
            else:
                custom_fields[new_custom_name.strip()] = new_custom_value
                st.success(
                    f"Câmpul «{new_custom_name.strip()}» a fost adăugat. "
                    "La următoarea deschidere va apărea în lista de mai sus."
                )

        st.markdown("</div>", unsafe_allow_html=True)

        # -------------------------------------------------------------
        # 📚 NOMENCLATOARE (COR + Legea 153) – Import din Excel/CSV
        # -------------------------------------------------------------
        st.markdown('<div class="section-card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">📚 Nomenclatoare (import din Excel/CSV)</div>', unsafe_allow_html=True)

        tab_cor, tab_153 = st.tabs(["🧩 COR", "📜 Legea 153"])

        # ----------------------------
        # COR
        # ----------------------------
        with tab_cor:
            st.caption("Importă/actualizează codurile COR în baza aplicației (tabela: cor_coduri).")
            upl = st.file_uploader(
                "Încarcă Excel/CSV (coloane: cod, denumire) – COR",
                type=["xlsx", "xls", "csv"],
                key="cfg_upl_cor",
            )

            def _read_table_file(uploaded):
                if uploaded is None:
                    return None
                name = (uploaded.name or "").lower()
                try:
                    if name.endswith(".csv"):
                        return pd.read_csv(uploaded)
                    else:
                        return pd.read_excel(uploaded)
                except Exception as e:
                    st.error(f"Nu am putut citi fișierul: {e}")
                    return None

            df_in = _read_table_file(upl)
            if df_in is not None and not df_in.empty:
                # normalize columns
                cols = {str(c).strip().lower(): c for c in df_in.columns}
                c_cod = cols.get("cod") or cols.get("cod cor") or cols.get("cod_cor") or df_in.columns[0]
                c_den = cols.get("denumire") or cols.get("denumire ocupatie") or cols.get("denumire ocupație") or cols.get("ocupatie") or cols.get("ocupație") or (df_in.columns[1] if len(df_in.columns) > 1 else df_in.columns[0])

                df_norm = pd.DataFrame()
                df_norm["cod"] = df_in[c_cod].astype(str).str.strip()
                df_norm["denumire"] = df_in[c_den].astype(str).str.strip()

                # cod = doar cifre, 6 caractere (păstrăm doar cele valide)
                df_norm["cod"] = df_norm["cod"].str.replace(r"\D+", "", regex=True).str.zfill(6)
                df_norm = df_norm[(df_norm["cod"].str.len() == 6) & (df_norm["denumire"] != "")]
                df_norm = df_norm.drop_duplicates(subset=["cod"]).reset_index(drop=True)

                st.write("Previzualizare (primele 50 rânduri):")
                st.dataframe(df_norm.head(50), use_container_width=True, height=260)

                col1, col2 = st.columns([1, 1])
                with col1:
                    if st.button("📥 Importă/Actualizează COR", key="cfg_btn_import_cor"):
                        try:
                            _cor_ensure_table(conn)
                            cur = conn.cursor()
                            cur.executemany(
                                """
                                INSERT INTO cor_coduri (cod, denumire, activ)
                                VALUES (?, ?, 1)
                                ON CONFLICT(cod) DO UPDATE SET
                                    denumire=excluded.denumire,
                                    activ=1
                                """,
                                list(df_norm[["cod", "denumire"]].itertuples(index=False, name=None)),
                            )
                            conn.commit()
                            st.success(f"Import COR reușit: {len(df_norm)} coduri.")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Eroare la import COR: {e}")
                with col2:
                    if st.button("🧹 Șterge toate codurile COR (dezactivează)", key="cfg_btn_disable_cor"):
                        try:
                            _cor_ensure_table(conn)
                            cur = conn.cursor()
                            cur.execute("UPDATE cor_coduri SET activ=0;")
                            conn.commit()
                            st.warning("Toate codurile COR au fost dezactivate (activ=0).")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Eroare: {e}")
            else:
                st.info("Încarcă un fișier Excel/CSV pentru import COR.")

            # Status curent
            try:
                _cor_ensure_table(conn)
                cur = conn.cursor()
                cur.execute("SELECT COUNT(*) FROM cor_coduri WHERE activ=1;")
                n_active = int(cur.fetchone()[0] or 0)
                st.caption(f"Coduri COR active în aplicație: **{n_active}**")
            except Exception:
                pass

        # ----------------------------
        # Legea 153
        # ----------------------------
        with tab_153:
            st.caption("Importă/actualizează codurile de funcții (Legea 153) în baza aplicației (tabela: lege153_coduri).")
            upl2 = st.file_uploader(
                "Încarcă Excel/CSV (coloane: cod, denumire) – Legea 153",
                type=["xlsx", "xls", "csv"],
                key="cfg_upl_153",
            )
            df_in2 = _read_table_file(upl2)
            if df_in2 is not None and not df_in2.empty:
                cols2 = {str(c).strip().lower(): c for c in df_in2.columns}
                c_cod2 = cols2.get("cod") or cols2.get("cod functie") or cols2.get("cod funcție") or cols2.get("cod_153") or cols2.get("cod 153") or df_in2.columns[0]
                c_den2 = cols2.get("denumire") or cols2.get("denumire functie") or cols2.get("denumire funcție") or cols2.get("functie") or cols2.get("funcție") or (df_in2.columns[1] if len(df_in2.columns) > 1 else df_in2.columns[0])

                df_norm2 = pd.DataFrame()
                df_norm2["cod"] = df_in2[c_cod2].astype(str).str.strip()
                df_norm2["denumire"] = df_in2[c_den2].astype(str).str.strip()
                df_norm2 = df_norm2[(df_norm2["cod"] != "") & (df_norm2["denumire"] != "")]
                df_norm2 = df_norm2.drop_duplicates(subset=["cod"]).reset_index(drop=True)

                st.write("Previzualizare (primele 50 rânduri):")
                st.dataframe(df_norm2.head(50), use_container_width=True, height=260)

                col1, col2 = st.columns([1, 1])
                with col1:
                    if st.button("📥 Importă/Actualizează Legea 153", key="cfg_btn_import_153"):
                        try:
                            _l153_ensure_table(conn)
                            cur = conn.cursor()
                            cur.executemany(
                                """
                                INSERT INTO lege153_coduri (cod, denumire, activ)
                                VALUES (?, ?, 1)
                                ON CONFLICT(cod) DO UPDATE SET
                                    denumire=excluded.denumire,
                                    activ=1
                                """,
                                list(df_norm2[["cod", "denumire"]].itertuples(index=False, name=None)),
                            )
                            conn.commit()
                            st.success(f"Import Legea 153 reușit: {len(df_norm2)} coduri.")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Eroare la import Legea 153: {e}")
                with col2:
                    if st.button("🧹 Șterge toate codurile 153 (dezactivează)", key="cfg_btn_disable_153"):
                        try:
                            _l153_ensure_table(conn)
                            cur = conn.cursor()
                            cur.execute("UPDATE lege153_coduri SET activ=0;")
                            conn.commit()
                            st.warning("Toate codurile Legea 153 au fost dezactivate (activ=0).")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Eroare: {e}")
            else:
                st.info("Încarcă un fișier Excel/CSV pentru import Legea 153.")

            # Status curent
            try:
                _l153_ensure_table(conn)
                cur = conn.cursor()
                cur.execute("SELECT COUNT(*) FROM lege153_coduri WHERE activ=1;")
                n_active2 = int(cur.fetchone()[0] or 0)
                st.caption(f"Coduri Legea 153 active în aplicație: **{n_active2}**")
            except Exception:
                pass

            st.divider()
            st.subheader("💰 Legea 153 – grilă salarizare (anexe)")

            st.caption(
                "Importă nomenclatorul complet (Anexa/Tabel/Funcție/Studii/Grad/Treaptă/Salariu/Coeficient) "
                "în baza aplicației (tabela: lege153_grid). Sursa recomandată: sheet-ul 'Normalized'."
            )

            upl_grid = st.file_uploader(
                "Încarcă XLSX – Grilă L153 (sheet: Normalized)",
                type=["xlsx", "xls"],
                key="cfg_upl_153_grid",
            )

            df_grid = None
            if upl_grid is not None:
                try:
                    df_grid_raw = pd.read_excel(upl_grid, sheet_name="Normalized")
                    df_grid = _l153_norm_df(df_grid_raw)
                    st.write("Previzualizare grilă (primele 50 rânduri):")
                    st.dataframe(df_grid.head(50), use_container_width=True, height=260)
                except Exception as e:
                    st.error(f"Nu am putut citi XLSX-ul (sheet 'Normalized'): {e}")

            cgi1, cgi2 = st.columns([1, 2])
            with cgi1:
                if st.button("📥 Importă grilă L153 în DB", key="cfg_btn_import_153_grid"):
                    if df_grid is None or df_grid.empty:
                        st.error("Încarcă mai întâi fișierul XLSX cu sheet 'Normalized'.")
                    else:
                        try:
                            n = _lege153_grid_replace(conn, df_grid)
                            try:
                                _lege153_grid_df.clear()
                            except Exception:
                                pass
                            st.success(f"Import grilă L153 reușit: {n} rânduri.")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Eroare la import grilă L153: {e}")

            with cgi2:
                try:
                    _lege153_grid_ensure_table(conn)
                    cur = conn.cursor()
                    cur.execute("SELECT COUNT(1) FROM lege153_grid")
                    cnt = (cur.fetchone() or [0])[0]
                    st.info(f"În DB: {cnt} rânduri în lege153_grid.")
                except Exception:
                    pass


        # Sticky save bar la finalul paginii de configurare (unicul buton de salvare)
        st.markdown('<div class="save-bar">', unsafe_allow_html=True)
        save_clicked = st.button("💾 Salvează configurarea", key="cfg_btn_save_all")
        st.markdown("</div>", unsafe_allow_html=True)

        if save_clicked:
            cfg["db_path"] = new_db_path
            cfg["denumire_unitate"] = denumire_unitate
            cfg["cui"] = cui
            cfg["adresa"] = adresa
            cfg["cont_bancar"] = cont_bancar
            cfg["banca"] = banca
            cfg["telefon"] = telefon
            cfg["fax"] = fax
            cfg["email"] = email
            cfg["afiseaza_iban_in_antet"] = bool(afiseaza_iban_in_antet)
            cfg["semnare_metoda"] = semnare_metoda

            # Salvează sigla (dacă a fost încărcată)
            try:
                if sigla_file is not None:
                    base_dir = os.path.dirname(CONFIG_FILE)
                    assets_dir = os.path.join(base_dir, "assets")
                    os.makedirs(assets_dir, exist_ok=True)
                    ext = os.path.splitext(sigla_file.name or "")[1].lower()
                    if ext not in [".png", ".jpg", ".jpeg"]:
                        ext = ".png"
                    sigla_path = os.path.join(assets_dir, "sigla_unitate" + ext)
                    with open(sigla_path, "wb") as f:
                        f.write(sigla_file.getvalue())
                    cfg["sigla_path"] = sigla_path
            except Exception:
                pass

            cfg["conducator_nume"] = conducator_nume
            cfg["conducator_functie"] = conducator_functie
            cfg["responsabil_hr_nume"] = responsabil_hr_nume
            cfg["responsabil_hr_functie"] = responsabil_hr_functie
            cfg["app_user"] = app_user
            cfg["app_pass"] = app_pass
            cfg["custom_fields"] = custom_fields

            save_config(cfg)
            st.success(
                "Configurarea a fost salvată. "
                "La următoarea rulare se va folosi noua bază de date și noile date de unitate."
            )

        st.markdown("</div>", unsafe_allow_html=True)  # end .section-card Nomenclatoare
        st.markdown("</div>", unsafe_allow_html=True)  # end .page-container

    # IMPORTANT: ca la login – după ce pagina este randată
    if st.session_state.get("logged_in", False):
        apply_app_post_render_fix()
        apply_toolbar_runtime_fix()


if __name__ == "__main__":
    main()
